"""
TDL Equity Analyzer - Backend API with PostgreSQL
Cross-device sync for portfolio analyses and overviews
"""

from flask import Flask, request, jsonify, Response
from flask_cors import CORS
import requests
import re
import os
import json
import base64
import threading
import queue
import psycopg2
from psycopg2.extras import RealDictCursor
from psycopg2.pool import ThreadedConnectionPool
from datetime import datetime
import anthropic
import openai
from google import genai
from google.genai import types as genai_types

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB max upload size
CORS(app, origins=[
    "https://equity-analyzer.tonydlee.workers.dev",
    "https://charlie-deployment.tonydlee.workers.dev",
    "http://localhost:3000",
    "http://localhost:5000",
    "http://127.0.0.1:3000",
    "http://127.0.0.1:5000",
])


# ============================================
# TELEGRAM NOTIFICATIONS
# ============================================
_TELEGRAM_BOT_TOKEN = os.environ.get("TELEGRAM_BOT_TOKEN", "")
_TELEGRAM_CHAT_ID = os.environ.get("TELEGRAM_CHAT_ID", "")

def _notify_telegram(message: str) -> None:
    """Send a Telegram notification. Silently fails if not configured."""
    if not _TELEGRAM_BOT_TOKEN or not _TELEGRAM_CHAT_ID:
        return
    try:
        requests.post(
            f"https://api.telegram.org/bot{_TELEGRAM_BOT_TOKEN}/sendMessage",
            json={"chat_id": _TELEGRAM_CHAT_ID, "text": message, "parse_mode": "Markdown"},
            timeout=10,
        )
    except Exception:
        pass

# ============================================
# IN-MEMORY CACHE
# ============================================
import time
import threading

class SimpleCache:
    """Thread-safe in-memory cache with TTL and manual invalidation."""
    def __init__(self):
        self._data = {}
        self._lock = threading.Lock()

    def get(self, key):
        with self._lock:
            entry = self._data.get(key)
            if entry and time.time() < entry['expires']:
                return entry['value']
            if entry:
                del self._data[key]
            return None

    def set(self, key, value, ttl=300):
        with self._lock:
            self._data[key] = {'value': value, 'expires': time.time() + ttl}

    def invalidate(self, *keys):
        with self._lock:
            for key in keys:
                self._data.pop(key, None)

cache = SimpleCache()

# ============================================
# AUTHENTICATION
# ============================================

import hashlib, hmac, uuid
CHARLIE_PASSWORD = os.environ.get('CHARLIE_PASSWORD', '')
CHARLIE_API_KEY = os.environ.get('CHARLIE_API_KEY', '')
_AUTH_SECRET = os.environ.get('CHARLIE_AUTH_SECRET', CHARLIE_API_KEY or 'charlie-default-secret')

# Public paths that don't require auth
AUTH_EXEMPT_PATHS = {'/health', '/api/agent/health', '/api/auth/config', '/api/auth/login'}

def _make_session_token(password):
    """Create an HMAC session token from the password."""
    return hmac.new(_AUTH_SECRET.encode(), password.encode(), hashlib.sha256).hexdigest()

def _verify_session_token(token):
    """Verify a session token matches the expected password hash."""
    if not CHARLIE_PASSWORD:
        return True
    expected = _make_session_token(CHARLIE_PASSWORD)
    return hmac.compare_digest(token, expected)

@app.before_request
def require_auth():
    """Global auth gate — runs before every request."""
    if request.method == 'OPTIONS':
        return None
    if request.path in AUTH_EXEMPT_PATHS:
        return None

    # No auth configured = dev mode, allow all
    if not CHARLIE_PASSWORD and not CHARLIE_API_KEY:
        return None

    auth_header = request.headers.get('Authorization', '')

    # API key auth (for local agent)
    if auth_header.startswith('ApiKey ') and CHARLIE_API_KEY:
        key = auth_header[7:]
        if key == CHARLIE_API_KEY:
            return None
        return jsonify({'error': 'Invalid API key'}), 401

    # Session token auth (password-based)
    if auth_header.startswith('Bearer '):
        token = auth_header[7:]
        if _verify_session_token(token):
            return None
        return jsonify({'error': 'Invalid or expired session'}), 401

    # No auth header
    if CHARLIE_PASSWORD:
        return jsonify({'error': 'Authentication required'}), 401
    return None


@app.route('/api/auth/config', methods=['GET'])
def auth_config():
    """Public endpoint: returns whether auth is required."""
    return jsonify({'authRequired': bool(CHARLIE_PASSWORD)})


@app.route('/api/auth/login', methods=['POST'])
def auth_login():
    """Public endpoint: verify password and return session token."""
    data = request.get_json() or {}
    password = data.get('password', '')
    if not CHARLIE_PASSWORD:
        return jsonify({'token': 'no-auth', 'message': 'Auth not configured'})
    if password == CHARLIE_PASSWORD:
        token = _make_session_token(password)
        return jsonify({'token': token})
    return jsonify({'error': 'Incorrect password'}), 401


# ============================================
# AGENT ALERTS
# ============================================

@app.route('/api/alerts', methods=['GET'])
def get_alerts():
    """Get all alerts, optionally filtered by status."""
    status = request.args.get('status', '')
    limit = int(request.args.get('limit', 50))
    try:
        with get_db() as (_, cur):
            if status:
                cur.execute('SELECT * FROM agent_alerts WHERE status = %s ORDER BY created_at DESC LIMIT %s', (status, limit))
            else:
                cur.execute('SELECT * FROM agent_alerts ORDER BY created_at DESC LIMIT %s', (limit,))
            rows = cur.fetchall()
        return jsonify({'alerts': [{
            'id': r['id'],
            'alertType': r['alert_type'],
            'ticker': r['ticker'],
            'title': r['title'],
            'detail': r['detail'] if isinstance(r['detail'], dict) else json.loads(r['detail'] or '{}'),
            'status': r['status'],
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
        } for r in rows], 'total': len(rows)})
    except Exception as e:
        print(f"Error fetching alerts: {e}")
        return jsonify({'alerts': [], 'total': 0})


@app.route('/api/alerts', methods=['POST'])
def create_alert():
    """Create a new alert (called by local agent or backend jobs)."""
    data = request.get_json()
    alert_type = data.get('alertType', 'general')
    ticker = data.get('ticker', '')
    title = data.get('title', '')
    detail = data.get('detail', {})
    if not title:
        return jsonify({'error': 'Title required'}), 400

    alert_id = str(uuid.uuid4())
    try:
        with get_db(commit=True) as (conn, cur):
            # Deduplicate: don't create if same type+ticker+title exists and is still 'new'
            cur.execute("SELECT id FROM agent_alerts WHERE alert_type = %s AND ticker = %s AND title = %s AND status = 'new'",
                       (alert_type, ticker, title))
            if cur.fetchone():
                return jsonify({'id': None, 'message': 'Duplicate alert exists'})
            cur.execute('''
                INSERT INTO agent_alerts (id, alert_type, ticker, title, detail, status, created_at)
                VALUES (%s, %s, %s, %s, %s, 'new', NOW())
            ''', (alert_id, alert_type, ticker, title, json.dumps(detail)))
        return jsonify({'id': alert_id})
    except Exception as e:
        print(f"Error creating alert: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/alerts/<alert_id>/dismiss', methods=['POST'])
def dismiss_alert(alert_id):
    """Dismiss an alert."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE agent_alerts SET status = 'dismissed' WHERE id = %s", (alert_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/alerts/<alert_id>/action', methods=['POST'])
def action_alert(alert_id):
    """Mark alert as actioned (e.g., synthesis triggered)."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE agent_alerts SET status = 'actioned' WHERE id = %s", (alert_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/alerts/dismiss-all', methods=['POST'])
def dismiss_all_alerts():
    """Dismiss all new alerts."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE agent_alerts SET status = 'dismissed' WHERE status = 'new'")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/alerts/count', methods=['GET'])
def alert_count():
    """Get count of new alerts (for badge)."""
    try:
        with get_db() as (_, cur):
            cur.execute("SELECT COUNT(*) as cnt FROM agent_alerts WHERE status = 'new'")
            cnt = cur.fetchone()['cnt']
        return jsonify({'count': cnt})
    except Exception:
        return jsonify({'count': 0})


# ============================================
# MEDIA TRACKER FEEDS — CRUD
# ============================================

def _row_to_feed(r):
    return {
        'id': r['id'],
        'sourceType': r['source_type'],
        'name': r['name'],
        'feedUrl': r['feed_url'],
        'sectorTags': r['sector_tags'] or [],
        'muted': r['muted'],
        'lastPolledAt': r['last_polled_at'].isoformat() if r['last_polled_at'] else None,
        'lastEpisodeAt': r['last_episode_at'].isoformat() if r['last_episode_at'] else None,
        'pollIntervalMin': r['poll_interval_min'],
        'errorCount': r['error_count'],
        'lastError': r['last_error'],
        'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
    }


@app.route('/api/media/feeds', methods=['GET'])
def media_feeds_list():
    with get_db() as (_c, cur):
        cur.execute("SELECT * FROM media_feeds ORDER BY name ASC")
        rows = cur.fetchall()
    feeds = [_row_to_feed(r) for r in rows]
    return jsonify({'feeds': feeds, 'total': len(feeds)})


@app.route('/api/media/feeds', methods=['POST'])
def media_feeds_create():
    data = request.get_json() or {}
    name = (data.get('name') or '').strip()
    feed_url = (data.get('feedUrl') or '').strip()
    source_type = (data.get('sourceType') or 'podcast').strip()
    sector_tags = data.get('sectorTags') or []
    poll_interval = int(data.get('pollIntervalMin') or 30)
    if not name or not feed_url:
        return jsonify({'error': 'name and feedUrl required'}), 400
    feed_id = str(uuid.uuid4())
    with get_db(commit=True) as (_c, cur):
        cur.execute('''
            INSERT INTO media_feeds (id, source_type, name, feed_url, sector_tags, poll_interval_min)
            VALUES (%s, %s, %s, %s, %s, %s)
            RETURNING *
        ''', (feed_id, source_type, name, feed_url, sector_tags, poll_interval))
        row = cur.fetchone()
    return jsonify({'feed': _row_to_feed(row)})


@app.route('/api/media/feeds/<feed_id>', methods=['PATCH'])
def media_feeds_update(feed_id):
    data = request.get_json() or {}
    fields, values = [], []
    for k, col in [('muted', 'muted'), ('name', 'name'), ('feedUrl', 'feed_url'),
                   ('sectorTags', 'sector_tags'), ('pollIntervalMin', 'poll_interval_min')]:
        if k in data:
            fields.append(f"{col} = %s")
            values.append(data[k])
    if not fields:
        return jsonify({'error': 'no updatable fields'}), 400
    values.append(feed_id)
    with get_db(commit=True) as (_c, cur):
        cur.execute(f"UPDATE media_feeds SET {', '.join(fields)} WHERE id = %s RETURNING *", values)
        row = cur.fetchone()
    if not row:
        return jsonify({'error': 'not found'}), 404
    return jsonify({'feed': _row_to_feed(row)})


@app.route('/api/media/feeds/<feed_id>', methods=['DELETE'])
def media_feeds_delete(feed_id):
    with get_db(commit=True) as (_c, cur):
        cur.execute("DELETE FROM media_feeds WHERE id = %s", (feed_id,))
    return jsonify({'success': True})


# ============================================
# MEDIA TRACKER — SIGNALS WATCHLIST
# ============================================

def _row_to_signal(r):
    return {
        'id': r['id'],
        'kind': r['kind'],
        'value': r['value'],
        'associatedTicker': r['associated_ticker'],
        'muted': r['muted'],
        'note': r['note'],
        'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
    }


@app.route('/api/media/watchlist', methods=['GET'])
def media_watchlist_list():
    with get_db() as (_c, cur):
        cur.execute("SELECT * FROM signals_watchlist ORDER BY kind, value")
        rows = cur.fetchall()
    return jsonify({'signals': [_row_to_signal(r) for r in rows], 'total': len(rows)})


@app.route('/api/media/watchlist', methods=['POST'])
def media_watchlist_create():
    data = request.get_json() or {}
    kind = (data.get('kind') or '').strip()
    value = (data.get('value') or '').strip()
    if kind not in ('ticker', 'keyword', 'exec') or not value:
        return jsonify({'error': 'kind must be ticker|keyword|exec; value required'}), 400
    sid = str(uuid.uuid4())
    try:
        with get_db(commit=True) as (_c, cur):
            cur.execute('''
                INSERT INTO signals_watchlist (id, kind, value, associated_ticker, note)
                VALUES (%s, %s, %s, %s, %s)
                RETURNING *
            ''', (sid, kind, value, data.get('associatedTicker'), data.get('note')))
            row = cur.fetchone()
    except Exception as e:
        if 'duplicate key' in str(e).lower() or 'unique' in str(e).lower():
            return jsonify({'error': 'already exists'}), 409
        raise
    return jsonify({'signal': _row_to_signal(row)})


@app.route('/api/media/watchlist/<signal_id>', methods=['PATCH'])
def media_watchlist_update(signal_id):
    data = request.get_json() or {}
    fields, values = [], []
    for k, col in [('muted', 'muted'), ('note', 'note'), ('associatedTicker', 'associated_ticker')]:
        if k in data:
            fields.append(f"{col} = %s")
            values.append(data[k])
    if not fields:
        return jsonify({'error': 'no updatable fields'}), 400
    values.append(signal_id)
    with get_db(commit=True) as (_c, cur):
        cur.execute(f"UPDATE signals_watchlist SET {', '.join(fields)} WHERE id = %s RETURNING *", values)
        row = cur.fetchone()
    if not row:
        return jsonify({'error': 'not found'}), 404
    return jsonify({'signal': _row_to_signal(row)})


@app.route('/api/media/watchlist/<signal_id>', methods=['DELETE'])
def media_watchlist_delete(signal_id):
    with get_db(commit=True) as (_c, cur):
        cur.execute("DELETE FROM signals_watchlist WHERE id = %s", (signal_id,))
    return jsonify({'success': True})


# ============================================
# MEDIA TRACKER — FEED (firehose) + SCANNER
# ============================================

_scanner_runs = {}  # in-memory scan status; OK to lose across deploys


@app.route('/api/media/feed', methods=['GET'])
def media_feed_read():
    source  = request.args.get('source')
    ticker  = request.args.get('ticker')
    sector  = request.args.get('sector')
    days    = int(request.args.get('days', 7))
    q       = request.args.get('q', '').strip()
    material_only = request.args.get('material') == 'true'
    limit   = int(request.args.get('limit', 100))

    where = ["e.created_at > NOW() - %s::interval"]
    params = [f'{days} days']
    if source:
        where.append("f.source_type = %s"); params.append(source)
    if material_only:
        where.append("EXISTS (SELECT 1 FROM media_digest_points p2 WHERE p2.episode_id=e.id AND p2.material)")
    if q:
        where.append("(e.title ILIKE %s OR e.show_notes ILIKE %s)")
        params.extend([f'%{q}%', f'%{q}%'])

    with get_db() as (_c, cur):
        cur.execute(f'''
            SELECT e.*, f.name AS feed_name, f.sector_tags AS feed_sector_tags
              FROM media_episodes e
              JOIN media_feeds f ON f.id = e.feed_id
             WHERE {' AND '.join(where)}
               AND e.status = 'done'
             ORDER BY e.published_at DESC NULLS LAST
             LIMIT %s
        ''', params + [limit])
        episodes = cur.fetchall()

        episode_ids = [e['id'] for e in episodes]
        points = []
        if episode_ids:
            pq = "SELECT * FROM media_digest_points WHERE episode_id = ANY(%s)"
            pparams = [episode_ids]
            if ticker:
                pq += " AND %s = ANY(tickers)"; pparams.append(ticker)
            if sector:
                pq += " AND %s = ANY(sector_tags)"; pparams.append(sector)
            pq += " ORDER BY episode_id, point_order"
            cur.execute(pq, pparams)
            points = cur.fetchall()

    by_ep = {}
    for p in points:
        by_ep.setdefault(p['episode_id'], []).append({
            'id': p['id'], 'text': p['text'], 'tickers': p['tickers'] or [],
            'sectorTags': p['sector_tags'] or [], 'themeTags': p['theme_tags'] or [],
            'material': p['material'], 'timestampSec': p['timestamp_sec'],
        })

    result = []
    for e in episodes:
        pts = by_ep.get(e['id'], [])
        if (ticker or sector) and not pts:
            continue  # filtered all points out — hide the episode
        result.append({
            'id': e['id'], 'feedId': e['feed_id'], 'feedName': e['feed_name'],
            'title': e['title'],
            'publishedAt': e['published_at'].isoformat() if e['published_at'] else None,
            'sourceUrl': e['source_url'], 'points': pts,
        })
    return jsonify({'episodes': result, 'total': len(result)})


@app.route('/api/media/run-scanner', methods=['POST'])
def media_run_scanner():
    scan_id = str(uuid.uuid4())
    _scanner_runs[scan_id] = {'status': 'running', 'started_at': datetime.utcnow().isoformat()}

    def _run():
        try:
            from media_trackers import poller
            poller.poll_all_feeds()
            _scanner_runs[scan_id]['status'] = 'done'
        except Exception as e:
            _scanner_runs[scan_id]['status'] = 'failed'
            _scanner_runs[scan_id]['error'] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return jsonify({'scanId': scan_id}), 202


@app.route('/api/media/scan/<scan_id>', methods=['GET'])
def media_scan_status(scan_id):
    rec = _scanner_runs.get(scan_id)
    if not rec:
        return jsonify({'error': 'unknown scan id'}), 404
    return jsonify(rec)


# ============================================
# MULTI-MODEL LLM FALLBACK
# ============================================

class LLMError(Exception):
    """Raised when all LLM providers fail."""
    def __init__(self, errors):
        self.errors = errors  # list of (provider, model, exception)
        messages = [f"{p}/{m}: {e}" for p, m, e in errors]
        super().__init__(f"All LLM providers failed: {'; '.join(messages)}")

MODEL_TIERS = {
    "fast": [
        ("anthropic", "claude-haiku-4-5-20251001"),
        ("gemini",    "gemini-2.0-flash"),
        ("openai",    "gpt-4o-mini"),
    ],
    "standard": [
        ("anthropic", "claude-sonnet-4-5-20250929"),
        ("gemini",    "gemini-2.0-flash"),
        ("openai",    "gpt-4o"),
    ],
    "advanced": [
        ("anthropic", "claude-opus-4-6"),
        ("gemini",    "gemini-2.0-pro"),
        ("openai",    "gpt-4o"),
    ],
}

def _extract_json(text):
    """Extract and parse JSON from LLM response text.

    Handles markdown fences, surrounding prose, truncated JSON (auto-closes
    brackets), and trailing commas.  Mirrors the robust parse_mp_json() logic.
    """
    text = text.strip()
    # Strip markdown code fences
    if text.startswith("```"):
        lines = text.split("\n")
        lines = lines[1:]  # skip opening ```json line
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        text = "\n".join(lines).strip()
    # Try direct parse first
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # Try to find outermost JSON object/array
    for sc, ec in [("{", "}"), ("[", "]")]:
        s = text.find(sc)
        e = text.rfind(ec)
        if s != -1 and e != -1 and e > s:
            try:
                return json.loads(text[s:e + 1])
            except json.JSONDecodeError:
                continue
    # Repair truncated JSON by closing open brackets/braces
    import re
    for sc, ec in [("{", "}"), ("[", "]")]:
        s = text.find(sc)
        if s != -1:
            fragment = text[s:]
            for _attempt in range(8):
                try:
                    return json.loads(fragment)
                except json.JSONDecodeError:
                    # Strip trailing partial string value (e.g. `"key": "truncated te`)
                    fragment = re.sub(r',\s*"[^"]*":\s*"[^"]*$', '', fragment)
                    # Strip trailing partial key (e.g. `"key": `)
                    fragment = re.sub(r',\s*"[^"]*"\s*:\s*$', '', fragment)
                    # Strip trailing comma
                    fragment = re.sub(r',\s*$', '', fragment)
                    # Close open brackets/braces
                    open_braces = fragment.count('{') - fragment.count('}')
                    open_brackets = fragment.count('[') - fragment.count(']')
                    fragment = fragment.rstrip()
                    # Remove trailing comma before closing
                    if fragment.endswith(','):
                        fragment = fragment[:-1]
                    fragment += ']' * max(0, open_brackets) + '}' * max(0, open_braces)
                    try:
                        return json.loads(fragment)
                    except json.JSONDecodeError:
                        # Strip back further — remove last complete entry and try again
                        last_brace = fragment.rfind('},')
                        last_bracket = fragment.rfind('],')
                        cut_point = max(last_brace, last_bracket)
                        if cut_point > 0:
                            fragment = fragment[:cut_point + 1]
                        else:
                            break
    raise json.JSONDecodeError("Could not parse JSON from LLM response", text[:500], 0)


def _get_api_keys(anthropic_api_key="", gemini_api_key="", openai_api_key=""):
    """Resolve API keys from explicit params or environment variables."""
    return {
        "anthropic": anthropic_api_key or os.environ.get("ANTHROPIC_API_KEY", ""),
        "gemini":    gemini_api_key or os.environ.get("GEMINI_API_KEY", "") or os.environ.get("GOOGLE_API_KEY", ""),
        "openai":    openai_api_key or os.environ.get("OPENAI_API_KEY", ""),
    }

def _is_retryable(provider, error):
    """Determine if an error should trigger fallback to next provider."""
    if provider == "anthropic":
        if isinstance(error, anthropic.AuthenticationError):
            return False
        if isinstance(error, (anthropic.RateLimitError, anthropic.APIConnectionError)):
            return True
        if isinstance(error, anthropic.APIStatusError):
            return error.status_code in (429, 500, 502, 503, 529)
        if isinstance(error, requests.Timeout):
            return True
        return True  # Network errors etc.
    elif provider == "gemini":
        err_str = str(error)
        return "429" in err_str or "RESOURCE_EXHAUSTED" in err_str or "500" in err_str or "503" in err_str or isinstance(error, (ConnectionError, TimeoutError))
    elif provider == "openai":
        if isinstance(error, openai.AuthenticationError):
            return False
        if isinstance(error, (openai.RateLimitError, openai.APIConnectionError)):
            return True
        if isinstance(error, openai.APIStatusError):
            return error.status_code in (429, 500, 502, 503)
        return True
    return True

def _call_anthropic(*, messages, system, model, max_tokens, timeout, api_key):
    """Call Anthropic API using the SDK. Returns normalized response dict."""
    client = anthropic.Anthropic(api_key=api_key, timeout=timeout)
    kwargs = {
        "model": model,
        "max_tokens": max_tokens,
        "messages": messages,
    }
    if system:
        kwargs["system"] = system
    response = client.messages.create(**kwargs)
    text = ""
    for block in response.content:
        if hasattr(block, "text"):
            text += block.text
    return {
        "text": text,
        "usage": {
            "input_tokens": response.usage.input_tokens,
            "output_tokens": response.usage.output_tokens,
        },
        "provider": "anthropic",
        "model": model,
    }

def _call_gemini(*, messages, system, model, max_tokens, timeout, api_key):
    """Call Google Gemini API. Returns normalized response dict."""
    client = genai.Client(api_key=api_key)
    contents = []
    for msg in messages:
        role = "user" if msg["role"] == "user" else "model"
        if isinstance(msg["content"], list):
            parts = []
            for block in msg["content"]:
                if block.get("type") == "text":
                    parts.append(genai_types.Part.from_text(text=block["text"]))
                elif block.get("type") == "image":
                    source = block["source"]
                    parts.append(genai_types.Part.from_bytes(
                        data=base64.b64decode(source["data"]),
                        mime_type=source["media_type"],
                    ))
                elif block.get("type") == "document":
                    source = block["source"]
                    parts.append(genai_types.Part.from_bytes(
                        data=base64.b64decode(source["data"]),
                        mime_type=source.get("media_type", "application/pdf"),
                    ))
            contents.append(genai_types.Content(role=role, parts=parts))
        else:
            contents.append(genai_types.Content(
                role=role,
                parts=[genai_types.Part.from_text(text=msg["content"])]
            ))
    config = genai_types.GenerateContentConfig(max_output_tokens=max_tokens)
    if system:
        config.system_instruction = system
    response = client.models.generate_content(
        model=model,
        contents=contents,
        config=config,
    )
    text = response.text or ""
    usage = {"input_tokens": 0, "output_tokens": 0}
    if hasattr(response, "usage_metadata") and response.usage_metadata:
        usage["input_tokens"] = getattr(response.usage_metadata, "prompt_token_count", 0) or 0
        usage["output_tokens"] = getattr(response.usage_metadata, "candidates_token_count", 0) or 0
    return {
        "text": text,
        "usage": usage,
        "provider": "gemini",
        "model": model,
    }

def _call_openai(*, messages, system, model, max_tokens, timeout, api_key):
    """Call OpenAI API. Returns normalized response dict."""
    client = openai.OpenAI(api_key=api_key, timeout=timeout)
    oai_messages = []
    if system:
        oai_messages.append({"role": "system", "content": system})
    for msg in messages:
        role = msg["role"]
        if isinstance(msg["content"], list):
            oai_content = []
            for block in msg["content"]:
                if block.get("type") == "text":
                    oai_content.append({"type": "text", "text": block["text"]})
                elif block.get("type") == "image":
                    source = block["source"]
                    data_uri = f"data:{source['media_type']};base64,{source['data']}"
                    oai_content.append({
                        "type": "image_url",
                        "image_url": {"url": data_uri}
                    })
                elif block.get("type") == "document":
                    # OpenAI chat completions don't support inline PDFs — skip
                    oai_content.append({"type": "text", "text": "[PDF document — content not available for this provider]"})
            oai_messages.append({"role": role, "content": oai_content})
        else:
            oai_messages.append({"role": role, "content": msg["content"]})
    response = client.chat.completions.create(
        model=model,
        messages=oai_messages,
        max_tokens=max_tokens,
    )
    text = response.choices[0].message.content or ""
    return {
        "text": text,
        "usage": {
            "input_tokens": response.usage.prompt_tokens if response.usage else 0,
            "output_tokens": response.usage.completion_tokens if response.usage else 0,
        },
        "provider": "openai",
        "model": model,
    }

_LLM_ADAPTERS = {
    "anthropic": _call_anthropic,
    "gemini":    _call_gemini,
    "openai":    _call_openai,
}

def call_llm(*, messages, system="", tier="standard", max_tokens=4096,
             timeout=120, anthropic_api_key="", gemini_api_key="", openai_api_key=""):
    """Call LLM with automatic multi-provider fallback.

    Returns dict with keys: text, usage, provider, model.
    Raises LLMError if all providers fail.
    """
    api_keys = _get_api_keys(anthropic_api_key, gemini_api_key, openai_api_key)
    chain = MODEL_TIERS.get(tier, MODEL_TIERS["standard"])
    errors = []
    for provider, model in chain:
        key = api_keys.get(provider, "")
        if not key:
            continue
        try:
            print(f"[LLM Fallback] Trying {provider}/{model}...")
            result = _LLM_ADAPTERS[provider](
                messages=messages, system=system, model=model,
                max_tokens=max_tokens, timeout=timeout, api_key=key,
            )
            print(f"[LLM Fallback] Success with {provider}/{model}")
            return result
        except Exception as e:
            print(f"[LLM Fallback] {provider}/{model} failed: {type(e).__name__}: {e}")
            errors.append((provider, model, e))
            if not _is_retryable(provider, e):
                break
    raise LLMError(errors)

def call_llm_stream(*, messages, system="", tier="standard", max_tokens=16384,
                    anthropic_api_key="", gemini_api_key="", openai_api_key=""):
    """Generator: yields keep-alive spaces, then final result dict.

    Usage:
        for chunk in call_llm_stream(...):
            if isinstance(chunk, dict):
                llm_result = chunk  # final result
            else:
                yield chunk  # keep-alive space
    """
    api_keys = _get_api_keys(anthropic_api_key, gemini_api_key, openai_api_key)
    chain = MODEL_TIERS.get(tier, MODEL_TIERS["standard"])
    errors = []
    for provider, model in chain:
        key = api_keys.get(provider, "")
        if not key:
            continue
        try:
            print(f"[LLM Stream Fallback] Trying {provider}/{model}...")
            if provider == "anthropic":
                client = anthropic.Anthropic(api_key=key, timeout=300)
                result_text = ""
                kwargs = {"model": model, "max_tokens": max_tokens, "messages": messages}
                if system:
                    kwargs["system"] = system
                with client.messages.stream(**kwargs) as stream:
                    for text in stream.text_stream:
                        result_text += text
                        yield " "
                    response = stream.get_final_message()
                print(f"[LLM Stream Fallback] Success with {provider}/{model}")
                yield {
                    "text": result_text,
                    "usage": {
                        "input_tokens": response.usage.input_tokens,
                        "output_tokens": response.usage.output_tokens,
                    },
                    "provider": provider,
                    "model": model,
                }
                return
            else:
                # Gemini/OpenAI: run non-streaming in background thread, yield keep-alive
                future_result = {}
                future_error = {}
                def _run_call(p=provider, m=model, k=key):
                    try:
                        future_result["data"] = _LLM_ADAPTERS[p](
                            messages=messages, system=system, model=m,
                            max_tokens=max_tokens, timeout=300, api_key=k,
                        )
                    except Exception as exc:
                        future_error["err"] = exc
                thread = threading.Thread(target=_run_call)
                thread.start()
                while thread.is_alive():
                    yield " "
                    thread.join(timeout=2.0)
                if "err" in future_error:
                    raise future_error["err"]
                print(f"[LLM Stream Fallback] Success with {provider}/{model}")
                yield future_result["data"]
                return
        except Exception as e:
            print(f"[LLM Stream Fallback] {provider}/{model} failed: {type(e).__name__}: {e}")
            errors.append((provider, model, e))
            if not _is_retryable(provider, e):
                break
    raise LLMError(errors)

# ============================================
# DATABASE CONNECTION
# ============================================

from contextlib import contextmanager

# Connection pool — initialized lazily, one per Gunicorn worker process
_pool = None

def _get_database_url():
    database_url = os.environ.get('DATABASE_URL')
    if not database_url:
        raise Exception('DATABASE_URL environment variable not set')
    # Render uses postgres:// but psycopg2 needs postgresql://
    if database_url.startswith('postgres://'):
        database_url = database_url.replace('postgres://', 'postgresql://', 1)
    return database_url

def _get_pool():
    """Get or create the connection pool (thread-safe, lazy init)."""
    global _pool
    if _pool is None or _pool.closed:
        # 3 workers × 2 threads = 6 handlers; pool up to 10 per worker
        _pool = ThreadedConnectionPool(
            minconn=2, maxconn=10,
            dsn=_get_database_url(),
            cursor_factory=RealDictCursor
        )
        print(f"DB connection pool created (min=2, max=10)")
    return _pool

def get_db_connection():
    """Get a connection from the pool."""
    return _get_pool().getconn()

@contextmanager
def get_db(commit=False):
    """Context manager for pooled database connections. Returns connection to pool on exit."""
    pool = _get_pool()
    conn = pool.getconn()
    cur = conn.cursor()
    try:
        yield conn, cur
        if commit:
            conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        cur.close()
        # Roll back any implicit transaction from reads before returning to pool
        try:
            conn.rollback()
        except Exception:
            pass
        try:
            if conn.closed:
                pool.putconn(conn, close=True)
            else:
                pool.putconn(conn)
        except Exception:
            pass

def init_db():
    """Initialize database tables"""
    try:
        with get_db(commit=True) as (_, cur):
            # Portfolio Analyses table
            cur.execute('''
                CREATE TABLE IF NOT EXISTS portfolio_analyses (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) UNIQUE NOT NULL,
                    company VARCHAR(255),
                    analysis JSONB,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Stock Overviews table
            cur.execute('''
                CREATE TABLE IF NOT EXISTS stock_overviews (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) UNIQUE NOT NULL,
                    company_name VARCHAR(255),
                    company_overview TEXT,
                    business_model TEXT,
                    business_mix TEXT,
                    opportunities TEXT,
                    risks TEXT,
                    conclusion TEXT,
                    raw_content TEXT,
                    history JSONB DEFAULT '[]',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Add business_mix column if it doesn't exist (migration)
            cur.execute('''
                DO $$
                BEGIN
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='stock_overviews' AND column_name='business_mix') THEN
                        ALTER TABLE stock_overviews ADD COLUMN business_mix TEXT;
                    END IF;
                END $$;
            ''')

            # Chat Histories table
            cur.execute('''
                CREATE TABLE IF NOT EXISTS chat_histories (
                    id VARCHAR(100) PRIMARY KEY,
                    title VARCHAR(255),
                    messages JSONB DEFAULT '[]',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Meeting Summaries table
            cur.execute('''
                CREATE TABLE IF NOT EXISTS meeting_summaries (
                    id VARCHAR(100) PRIMARY KEY,
                    title VARCHAR(255),
                    raw_notes TEXT,
                    summary TEXT,
                    questions TEXT,
                    topic VARCHAR(100) DEFAULT 'General',
                    topic_type VARCHAR(20) DEFAULT 'other',
                    source_type VARCHAR(20) DEFAULT 'paste',
                    source_files JSONB DEFAULT '[]',
                    doc_type VARCHAR(50) DEFAULT 'other',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Add columns if they don't exist (migration)
            cur.execute('''
                DO $$
                BEGIN
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='topic') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN topic VARCHAR(100) DEFAULT 'General';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='topic_type') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN topic_type VARCHAR(20) DEFAULT 'other';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='source_type') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN source_type VARCHAR(20) DEFAULT 'paste';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='source_files') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN source_files JSONB DEFAULT '[]';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='doc_type') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN doc_type VARCHAR(50) DEFAULT 'other';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='assessment') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN assessment TEXT DEFAULT '';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='categories') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN categories JSONB DEFAULT '[]';
                    END IF;
                END $$;
            ''')

            # Document Files table - stores actual document content for re-analysis
            cur.execute('''
                CREATE TABLE IF NOT EXISTS document_files (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    filename VARCHAR(255) NOT NULL,
                    file_data TEXT NOT NULL,
                    file_type VARCHAR(50),
                    mime_type VARCHAR(100),
                    metadata JSONB DEFAULT '{}',
                    file_size INTEGER,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    UNIQUE(ticker, filename)
                )
            ''')

            # Research Categories (tickers + topics)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_categories (
                    id VARCHAR(100) PRIMARY KEY,
                    name VARCHAR(255) NOT NULL,
                    type VARCHAR(20) DEFAULT 'ticker',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Research Documents (files/text under categories)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_documents (
                    id VARCHAR(100) PRIMARY KEY,
                    category_id VARCHAR(100) REFERENCES research_categories(id) ON DELETE CASCADE,
                    name VARCHAR(255) NOT NULL,
                    content TEXT,
                    file_names JSONB DEFAULT '[]',
                    smart_name VARCHAR(500),
                    original_filename VARCHAR(500),
                    published_date VARCHAR(100),
                    doc_type VARCHAR(50) DEFAULT 'other',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Migration: Add new columns if they don't exist
            cur.execute('''
                DO $$
                BEGIN
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='research_documents' AND column_name='smart_name') THEN
                        ALTER TABLE research_documents ADD COLUMN smart_name VARCHAR(500);
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='research_documents' AND column_name='original_filename') THEN
                        ALTER TABLE research_documents ADD COLUMN original_filename VARCHAR(500);
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='research_documents' AND column_name='published_date') THEN
                        ALTER TABLE research_documents ADD COLUMN published_date VARCHAR(100);
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='research_documents' AND column_name='has_stored_files') THEN
                        ALTER TABLE research_documents ADD COLUMN has_stored_files BOOLEAN DEFAULT FALSE;
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='research_documents' AND column_name='doc_type') THEN
                        ALTER TABLE research_documents ADD COLUMN doc_type VARCHAR(50) DEFAULT 'other';
                    END IF;
                END $$;
            ''')

            # Research Document Files (stored PDFs/files for re-analysis)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_document_files (
                    id SERIAL PRIMARY KEY,
                    document_id VARCHAR(100) REFERENCES research_documents(id) ON DELETE CASCADE,
                    filename VARCHAR(500) NOT NULL,
                    file_type VARCHAR(100),
                    file_data TEXT,
                    file_size INTEGER,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Summary Files (stored PDFs/files for summaries)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS summary_files (
                    id SERIAL PRIMARY KEY,
                    summary_id VARCHAR(100) REFERENCES meeting_summaries(id) ON DELETE CASCADE,
                    filename VARCHAR(500) NOT NULL,
                    file_type VARCHAR(100),
                    file_data TEXT,
                    file_size INTEGER,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Migration: Add has_stored_files to meeting_summaries if not exists
            cur.execute('''
                DO $$
                BEGIN
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='meeting_summaries' AND column_name='has_stored_files') THEN
                        ALTER TABLE meeting_summaries ADD COLUMN has_stored_files BOOLEAN DEFAULT FALSE;
                    END IF;
                END $$;
            ''')

            # Research Analyses (framework results under documents)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_analyses (
                    id VARCHAR(100) PRIMARY KEY,
                    document_id VARCHAR(100) REFERENCES research_documents(id) ON DELETE CASCADE,
                    prompt_id VARCHAR(100),
                    prompt_name VARCHAR(255),
                    prompt_icon VARCHAR(10),
                    result TEXT,
                    usage JSONB DEFAULT '{}',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Create index for faster lookups
            cur.execute('''
                CREATE INDEX IF NOT EXISTS idx_document_files_ticker
                ON document_files(ticker)
            ''')

            # ============================================
            # MEETING PREP TABLES
            # ============================================

            cur.execute('''
                CREATE TABLE IF NOT EXISTS mp_companies (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) UNIQUE NOT NULL,
                    name VARCHAR(255),
                    sector VARCHAR(100),
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS mp_meetings (
                    id SERIAL PRIMARY KEY,
                    company_id INTEGER REFERENCES mp_companies(id),
                    meeting_date DATE,
                    meeting_type VARCHAR(50) DEFAULT 'other',
                    status VARCHAR(20) DEFAULT 'draft',
                    notes TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_mp_meetings_company ON mp_meetings(company_id)')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS mp_documents (
                    id SERIAL PRIMARY KEY,
                    meeting_id INTEGER REFERENCES mp_meetings(id) ON DELETE CASCADE,
                    filename VARCHAR(500) NOT NULL,
                    file_data TEXT,
                    doc_type VARCHAR(50) DEFAULT 'other',
                    doc_date VARCHAR(20),
                    page_count INTEGER,
                    token_estimate INTEGER,
                    extracted_text TEXT,
                    upload_order INTEGER,
                    file_size INTEGER,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_mp_documents_meeting ON mp_documents(meeting_id)')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS mp_question_sets (
                    id SERIAL PRIMARY KEY,
                    meeting_id INTEGER REFERENCES mp_meetings(id) ON DELETE CASCADE,
                    version INTEGER DEFAULT 1,
                    status VARCHAR(20) DEFAULT 'ready',
                    topics_json TEXT,
                    synthesis_json TEXT,
                    generation_model VARCHAR(100),
                    generation_tokens INTEGER,
                    error_message TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_mp_question_sets_meeting ON mp_question_sets(meeting_id)')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS mp_past_questions (
                    id SERIAL PRIMARY KEY,
                    company_id INTEGER REFERENCES mp_companies(id),
                    meeting_id INTEGER REFERENCES mp_meetings(id) ON DELETE SET NULL,
                    question TEXT NOT NULL,
                    topic VARCHAR(255),
                    response_notes TEXT,
                    status VARCHAR(20) DEFAULT 'asked',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_mp_past_questions_company ON mp_past_questions(company_id)')

            # ============================================
            # SLIDE GENERATOR TABLES
            # ============================================

            cur.execute('''
                CREATE TABLE IF NOT EXISTS slide_projects (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20),
                    title VARCHAR(255) NOT NULL,
                    theme VARCHAR(50) DEFAULT 'sketchnote',
                    status VARCHAR(20) DEFAULT 'draft',
                    total_slides INTEGER DEFAULT 0,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS slide_items (
                    id SERIAL PRIMARY KEY,
                    project_id INTEGER REFERENCES slide_projects(id) ON DELETE CASCADE,
                    slide_number INTEGER NOT NULL,
                    title VARCHAR(255) NOT NULL,
                    type VARCHAR(30) DEFAULT 'content',
                    content TEXT,
                    illustration_hints JSONB DEFAULT '[]',
                    no_header BOOLEAN DEFAULT FALSE,
                    image_data TEXT,
                    content_hash VARCHAR(64),
                    status VARCHAR(20) DEFAULT 'new',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_slide_items_project ON slide_items(project_id)')
            cur.execute('CREATE UNIQUE INDEX IF NOT EXISTS idx_slide_items_project_num ON slide_items(project_id, slide_number)')

            # ============================================
            # STUDIO TABLES
            # ============================================

            cur.execute('''
                CREATE TABLE IF NOT EXISTS studio_design_themes (
                    id SERIAL PRIMARY KEY,
                    name VARCHAR(255) NOT NULL,
                    description TEXT,
                    style_prompt TEXT NOT NULL,
                    preview_image TEXT,
                    is_default BOOLEAN DEFAULT FALSE,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS studio_outputs (
                    id SERIAL PRIMARY KEY,
                    title VARCHAR(500) NOT NULL,
                    type VARCHAR(30) NOT NULL,
                    status VARCHAR(50) DEFAULT 'pending',
                    theme_id INTEGER REFERENCES studio_design_themes(id) ON DELETE SET NULL,
                    source_config JSONB DEFAULT '{}',
                    settings JSONB DEFAULT '{}',
                    content JSONB DEFAULT '{}',
                    image_data TEXT,
                    progress_current INTEGER DEFAULT 0,
                    progress_total INTEGER DEFAULT 0,
                    error_message TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS studio_slide_images (
                    id SERIAL PRIMARY KEY,
                    output_id INTEGER REFERENCES studio_outputs(id) ON DELETE CASCADE,
                    slide_number INTEGER NOT NULL,
                    image_data TEXT,
                    content_hash VARCHAR(64),
                    status VARCHAR(20) DEFAULT 'new',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_studio_slide_images_output ON studio_slide_images(output_id)')

            # Seed default studio themes from SLIDE_THEMES if table is empty
            cur.execute('SELECT COUNT(*) as cnt FROM studio_design_themes')
            if cur.fetchone()['cnt'] == 0:
                for key, theme in SLIDE_THEMES.items():
                    cur.execute('''
                        INSERT INTO studio_design_themes (name, description, style_prompt, is_default)
                        VALUES (%s, %s, %s, %s)
                    ''', (theme['name'], theme.get('illustration_guidance', ''), theme['style_prefix'], True))

            # App Settings table - persists API keys across PWA reinstalls
            cur.execute('''
                CREATE TABLE IF NOT EXISTS app_settings (
                    key VARCHAR(100) PRIMARY KEY,
                    value TEXT NOT NULL,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Analysis Jobs table - background analysis processing
            cur.execute('''
                CREATE TABLE IF NOT EXISTS analysis_jobs (
                    id VARCHAR(100) PRIMARY KEY,
                    ticker VARCHAR(20),
                    status VARCHAR(20) DEFAULT 'pending',
                    progress TEXT DEFAULT '',
                    batch_num INTEGER DEFAULT 0,
                    total_batches INTEGER DEFAULT 0,
                    chars_received INTEGER DEFAULT 0,
                    api_key TEXT NOT NULL,
                    request_payload JSONB NOT NULL,
                    result JSONB,
                    error TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Widen analysis_jobs.id to support UUID (was VARCHAR(20), needs VARCHAR(100))
            cur.execute('''
                DO $$ BEGIN
                    IF EXISTS (
                        SELECT 1 FROM information_schema.columns
                        WHERE table_name='analysis_jobs' AND column_name='id' AND character_maximum_length < 100
                    ) THEN
                        ALTER TABLE analysis_jobs ALTER COLUMN id TYPE VARCHAR(100);
                    END IF;
                END $$;
            ''')

            # Thesis scorecard data (traffic-light thresholds for formatted exports)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_scorecard_data (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) UNIQUE NOT NULL,
                    scorecard_data JSONB,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Thesis infographic history (versioned image storage)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_infographic_history (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    mode VARCHAR(5) NOT NULL DEFAULT '1',
                    detail VARCHAR(10) NOT NULL DEFAULT 'full',
                    style VARCHAR(50) NOT NULL DEFAULT 'professional',
                    slide_images JSONB NOT NULL DEFAULT '[]',
                    edit_prompt TEXT,
                    parent_id INTEGER REFERENCES thesis_infographic_history(id) ON DELETE SET NULL,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_infographic_history_ticker ON thesis_infographic_history(ticker)')
            cur.execute("ALTER TABLE thesis_infographic_history ADD COLUMN IF NOT EXISTS color_scheme VARCHAR(50) DEFAULT 'default'")

            # Infographic templates (saved favorites for reuse across tickers)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_infographic_templates (
                    id SERIAL PRIMARY KEY,
                    name VARCHAR(255) NOT NULL,
                    mode VARCHAR(5) NOT NULL DEFAULT '1',
                    detail VARCHAR(10) NOT NULL DEFAULT 'full',
                    style VARCHAR(50) NOT NULL DEFAULT 'professional',
                    color_scheme VARCHAR(50) DEFAULT 'default',
                    show_risk_detail BOOLEAN DEFAULT false,
                    include_company BOOLEAN DEFAULT false,
                    reference_image TEXT NOT NULL,
                    source_ticker VARCHAR(20),
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Thesis condensed versions (shorter input source for formats)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_condensed (
                    ticker VARCHAR(20) PRIMARY KEY,
                    condensed_analysis JSONB NOT NULL,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Thesis full versions (curated selection from detailed)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_full (
                    ticker VARCHAR(20) PRIMARY KEY,
                    full_analysis JSONB NOT NULL,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')

            # Add history columns to thesis tables (migration)
            cur.execute('''
                DO $$
                BEGIN
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='thesis_full' AND column_name='history') THEN
                        ALTER TABLE thesis_full ADD COLUMN history JSONB DEFAULT '[]';
                    END IF;
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns
                                  WHERE table_name='thesis_condensed' AND column_name='history') THEN
                        ALTER TABLE thesis_condensed ADD COLUMN history JSONB DEFAULT '[]';
                    END IF;
                END $$;
            ''')

            # Thesis format history (versioned PDF/DOCX storage)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_format_history (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    format_type VARCHAR(50) NOT NULL DEFAULT 'executive',
                    output_type VARCHAR(10) NOT NULL DEFAULT 'pdf',
                    file_data TEXT NOT NULL,
                    filename VARCHAR(255),
                    file_size INTEGER,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_format_history_ticker ON thesis_format_history(ticker)')

            # Thesis evolution snapshots (track thesis/scorecard changes over time)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS thesis_snapshots (
                    id SERIAL PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    snapshot_type VARCHAR(20) NOT NULL DEFAULT 'analysis',
                    thesis_summary TEXT,
                    pillar_count INTEGER DEFAULT 0,
                    signpost_statuses JSONB DEFAULT '[]',
                    risk_statuses JSONB DEFAULT '[]',
                    green_count INTEGER DEFAULT 0,
                    yellow_count INTEGER DEFAULT 0,
                    red_count INTEGER DEFAULT 0,
                    total_items INTEGER DEFAULT 0,
                    conviction TEXT,
                    raw_snapshot JSONB,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_thesis_snapshots_ticker ON thesis_snapshots(ticker)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_thesis_snapshots_ticker_date ON thesis_snapshots(ticker, created_at DESC)')

            # Research Pipeline Jobs table - batch processing queue
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_pipeline_jobs (
                    id TEXT PRIMARY KEY,
                    batch_id TEXT,
                    ticker TEXT NOT NULL,
                    job_type TEXT NOT NULL DEFAULT 'process',
                    status TEXT NOT NULL DEFAULT 'queued',
                    progress INTEGER DEFAULT 0,
                    current_step TEXT DEFAULT '',
                    total_steps INTEGER DEFAULT 7,
                    steps_detail JSONB DEFAULT '[]',
                    result JSONB,
                    error TEXT,
                    created_at TIMESTAMP DEFAULT NOW(),
                    updated_at TIMESTAMP DEFAULT NOW(),
                    completed_at TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_pipeline_jobs_batch ON research_pipeline_jobs(batch_id)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_pipeline_jobs_ticker ON research_pipeline_jobs(ticker)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_pipeline_jobs_status ON research_pipeline_jobs(status)')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_notes (
                    id TEXT PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    version VARCHAR(10) DEFAULT '1.0',
                    note_markdown TEXT,
                    sources_markdown TEXT,
                    changelog_markdown TEXT,
                    note_docx TEXT,
                    charts JSONB DEFAULT '[]',
                    metadata JSONB DEFAULT '{}',
                    created_at TIMESTAMP DEFAULT NOW(),
                    updated_at TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_research_notes_ticker ON research_notes(ticker)')

            # TradingAgents runs
            cur.execute('''
                CREATE TABLE IF NOT EXISTS agent_runs (
                    id TEXT PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    analysis_date DATE NOT NULL,
                    llm_provider VARCHAR(50) NOT NULL,
                    llm_model VARCHAR(100) NOT NULL,
                    status VARCHAR(20) DEFAULT 'running',
                    logs JSONB DEFAULT '[]',
                    decision TEXT,
                    report TEXT,
                    created_at TIMESTAMP DEFAULT NOW(),
                    completed_at TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_agent_runs_ticker ON agent_runs(ticker)')

            # Add cost tracking columns to agent_runs
            cur.execute('''
                DO $$ BEGIN
                    IF NOT EXISTS (SELECT 1 FROM information_schema.columns WHERE table_name='agent_runs' AND column_name='estimated_cost') THEN
                        ALTER TABLE agent_runs ADD COLUMN estimated_cost NUMERIC(10,4) DEFAULT 0;
                        ALTER TABLE agent_runs ADD COLUMN input_tokens INTEGER DEFAULT 0;
                        ALTER TABLE agent_runs ADD COLUMN output_tokens INTEGER DEFAULT 0;
                    END IF;
                END $$;
            ''')

            # Validation history
            cur.execute('''
                CREATE TABLE IF NOT EXISTS validation_runs (
                    id TEXT PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    content_type VARCHAR(20) NOT NULL,
                    overall_score INTEGER DEFAULT 0,
                    total_issues INTEGER DEFAULT 0,
                    high_confidence_issues INTEGER DEFAULT 0,
                    consensus_issues JSONB DEFAULT '[]',
                    reviewers JSONB DEFAULT '{}',
                    created_at TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_validation_runs_ticker ON validation_runs(ticker)')

            # Contextual chat
            cur.execute('''
                CREATE TABLE IF NOT EXISTS content_chats (
                    id TEXT PRIMARY KEY,
                    ticker VARCHAR(20) NOT NULL,
                    content_type VARCHAR(20) NOT NULL,
                    messages JSONB DEFAULT '[]',
                    created_at TIMESTAMP DEFAULT NOW(),
                    updated_at TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_content_chats_ticker ON content_chats(ticker)')

            # Mark stale processing jobs as failed (server restart recovery)
            cur.execute('''
                UPDATE analysis_jobs
                SET status = 'failed', error = 'Server restarted during processing', api_key = ''
                WHERE status IN ('pending', 'processing') AND updated_at < NOW() - INTERVAL '5 minutes'
            ''')

            # Clean up old completed/failed jobs
            cur.execute('''
                DELETE FROM analysis_jobs
                WHERE (status IN ('complete', 'failed') AND created_at < NOW() - INTERVAL '24 hours')
            ''')

            # Clean up stale pipeline jobs (stuck running for > 30 min)
            cur.execute("""
                UPDATE research_pipeline_jobs SET status = 'failed', error = 'Job timed out', updated_at = NOW()
                WHERE status IN ('queued', 'running') AND created_at < NOW() - INTERVAL '30 minutes'
            """)

            # Clean up old completed/failed pipeline jobs (older than 7 days)
            cur.execute("""
                DELETE FROM research_pipeline_jobs
                WHERE status IN ('complete', 'failed') AND created_at < NOW() - INTERVAL '7 days'
            """)

            # Ticker Settings table (custom sector assignments, etc.)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS ticker_settings (
                    ticker VARCHAR(20) PRIMARY KEY,
                    sector VARCHAR(100),
                    custom_company VARCHAR(255),
                    updated_at TIMESTAMP DEFAULT NOW()
                )
            ''')

            # Deep Research runs
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_runs (
                    id TEXT PRIMARY KEY,
                    query TEXT NOT NULL,
                    mode VARCHAR(20) DEFAULT 'topic',
                    status VARCHAR(20) DEFAULT 'running',
                    progress INTEGER DEFAULT 0,
                    current_step TEXT DEFAULT '',
                    agents JSONB DEFAULT '{}',
                    synthesis TEXT,
                    created_at TIMESTAMP DEFAULT NOW(),
                    completed_at TIMESTAMP
                )
            ''')

            # Catalyst folders (local agent reports subfolder contents)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS catalyst_folders (
                    ticker VARCHAR(20) PRIMARY KEY,
                    folders JSONB DEFAULT '[]',
                    updated_at TIMESTAMP DEFAULT NOW()
                )
            ''')

            # Research exports (research -> slides/infographics)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS research_exports (
                    id VARCHAR(100) PRIMARY KEY,
                    research_run_id TEXT NOT NULL,
                    slide_project_id INTEGER REFERENCES slide_projects(id) ON DELETE SET NULL,
                    format VARCHAR(50) NOT NULL,
                    type VARCHAR(20) NOT NULL,
                    slide_count INTEGER,
                    status VARCHAR(20) DEFAULT 'queued',
                    progress INTEGER DEFAULT 0,
                    error_message TEXT,
                    created_at TIMESTAMP DEFAULT NOW(),
                    completed_at TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_research_exports_run ON research_exports(research_run_id)')

            # Agent alerts
            cur.execute('''
                CREATE TABLE IF NOT EXISTS agent_alerts (
                    id VARCHAR(100) PRIMARY KEY,
                    alert_type VARCHAR(50) NOT NULL,
                    ticker VARCHAR(20),
                    title TEXT NOT NULL,
                    detail JSONB DEFAULT '{}',
                    status VARCHAR(20) DEFAULT 'new',
                    created_at TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_agent_alerts_status ON agent_alerts(status)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_agent_alerts_created ON agent_alerts(created_at DESC)')

            cur.execute('''
                CREATE TABLE IF NOT EXISTS media_feeds (
                    id                VARCHAR(100) PRIMARY KEY,
                    source_type       VARCHAR(20) NOT NULL,
                    name              TEXT NOT NULL,
                    feed_url          TEXT NOT NULL,
                    sector_tags       TEXT[] DEFAULT '{}',
                    muted             BOOLEAN DEFAULT FALSE,
                    last_polled_at    TIMESTAMP,
                    last_episode_at   TIMESTAMP,
                    poll_interval_min INT DEFAULT 30,
                    error_count       INT DEFAULT 0,
                    last_error        TEXT,
                    created_at        TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS media_episodes (
                    id                VARCHAR(100) PRIMARY KEY,
                    feed_id           VARCHAR(100) REFERENCES media_feeds(id) ON DELETE CASCADE,
                    guid              TEXT NOT NULL,
                    title             TEXT NOT NULL,
                    published_at      TIMESTAMP,
                    audio_url         TEXT,
                    source_url        TEXT,
                    show_notes        TEXT,
                    duration_sec      INT,
                    transcript        TEXT,
                    transcript_source VARCHAR(20),
                    status            VARCHAR(20) DEFAULT 'new',
                    error_message     TEXT,
                    cost_usd          NUMERIC(10,4) DEFAULT 0,
                    created_at        TIMESTAMP DEFAULT NOW(),
                    UNIQUE(feed_id, guid)
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS media_digest_points (
                    id            VARCHAR(100) PRIMARY KEY,
                    episode_id    VARCHAR(100) REFERENCES media_episodes(id) ON DELETE CASCADE,
                    point_order   INT NOT NULL,
                    text          TEXT NOT NULL,
                    tickers       TEXT[] DEFAULT '{}',
                    sector_tags   TEXT[] DEFAULT '{}',
                    theme_tags    TEXT[] DEFAULT '{}',
                    timestamp_sec INT,
                    material      BOOLEAN DEFAULT FALSE,
                    cluster_id    VARCHAR(100),
                    created_at    TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS signals_watchlist (
                    id                VARCHAR(100) PRIMARY KEY,
                    kind              VARCHAR(20) NOT NULL,
                    value             TEXT NOT NULL,
                    associated_ticker VARCHAR(20),
                    muted             BOOLEAN DEFAULT FALSE,
                    note              TEXT,
                    created_at        TIMESTAMP DEFAULT NOW(),
                    UNIQUE(kind, value)
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS media_theme_clusters (
                    id              VARCHAR(100) PRIMARY KEY,
                    theme           TEXT NOT NULL,
                    summary         TEXT,
                    point_ids       TEXT[] DEFAULT '{}',
                    primary_tickers TEXT[] DEFAULT '{}',
                    week_start      DATE NOT NULL,
                    created_at      TIMESTAMP DEFAULT NOW()
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS notification_prefs (
                    key   VARCHAR(50) PRIMARY KEY,
                    value JSONB
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_episodes_status ON media_episodes(status)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_episodes_feed_published ON media_episodes(feed_id, published_at DESC)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_points_episode ON media_digest_points(episode_id, point_order)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_points_tickers_gin ON media_digest_points USING GIN(tickers)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_points_material ON media_digest_points(material, created_at DESC)')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_points_cluster ON media_digest_points(cluster_id)')

        print("Database tables initialized")
    except Exception as e:
        print(f"Database init error (may be normal on first run): {e}")

# Initialize database on startup
try:
    init_db()
except:
    pass  # Will init when DATABASE_URL is available

# Start APScheduler for media trackers (unless APSCHEDULER_DISABLED set)
if not os.environ.get('APSCHEDULER_DISABLED'):
    try:
        import scheduler as _media_scheduler
        _media_scheduler.start()
        print("Media tracker scheduler started")
    except Exception as e:
        print(f"Scheduler start failed (non-fatal): {e}")


# ============================================
# APP SETTINGS ENDPOINTS (API key persistence)
# ============================================

@app.route('/api/settings', methods=['GET'])
def get_settings():
    """Get all app settings"""
    try:
        with get_db() as (conn, cur):
            cur.execute('SELECT key, value FROM app_settings')
            rows = cur.fetchall()
        return jsonify({row['key']: row['value'] for row in rows})
    except Exception as e:
        print(f"Error getting settings: {e}")
        return jsonify({})

@app.route('/api/settings', methods=['POST'])
def save_settings():
    """Save app settings (upsert)"""
    try:
        data = request.json
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        with get_db(commit=True) as (conn, cur):
            for key, value in data.items():
                cur.execute('''
                    INSERT INTO app_settings (key, value, updated_at)
                    VALUES (%s, %s, CURRENT_TIMESTAMP)
                    ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value, updated_at = CURRENT_TIMESTAMP
                ''', (key, value))
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error saving settings: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# BACKGROUND ANALYSIS JOB SYSTEM
# ============================================

import uuid
from PyPDF2 import PdfReader, PdfWriter
import io

MAX_PDF_PAGES_PER_CHUNK = 95  # Claude API rejects >100 pages per document
MAX_PAGES_PER_BATCH = 100     # Initial batch hint; auto-split handles token limit errors dynamically


def _split_large_pdf(base64_data, max_pages=MAX_PDF_PAGES_PER_CHUNK):
    """Split a base64 PDF into chunks of ≤max_pages. Returns list of (base64_data, page_count)."""
    try:
        pdf_bytes = base64.b64decode(base64_data)
        reader = PdfReader(io.BytesIO(pdf_bytes))
        total = len(reader.pages)
        if total <= max_pages:
            return [(base64_data, total)]
        chunks = []
        for start in range(0, total, max_pages):
            writer = PdfWriter()
            end = min(start + max_pages, total)
            for p in range(start, end):
                writer.add_page(reader.pages[p])
            buf = io.BytesIO()
            writer.write(buf)
            chunks.append((base64.b64encode(buf.getvalue()).decode('utf-8'), end - start))
        return chunks
    except Exception as e:
        print(f"[_split_large_pdf] Error splitting PDF: {e}")
        return [(base64_data, 30)]


def _trim_analysis_for_context(analysis):
    """Remove non-analytical metadata from existing analysis before passing as LLM context.
    Preserves ALL analytical content (pillars, signposts, threats, sources with excerpts, conclusion).
    Only strips UI-only fields: history, documentHistory, updatedAt."""
    if not analysis or not isinstance(analysis, dict):
        return analysis
    import copy
    trimmed = copy.deepcopy(analysis)
    # Remove UI-only fields that have zero analytical value
    trimmed.pop('history', None)
    trimmed.pop('documentHistory', None)
    trimmed.pop('updatedAt', None)
    return trimmed


def _repair_truncated_json(text):
    """Attempt to repair truncated JSON from a max_tokens cutoff."""
    import re
    repaired = text
    # Strip trailing incomplete key:value pairs
    repaired = re.sub(r',\s*"[^"]*$', '', repaired)                    # trailing key without value
    repaired = re.sub(r',\s*"[^"]*":\s*"[^"]*$', '', repaired)        # trailing key:partial-string
    repaired = re.sub(r',\s*"[^"]*":\s*$', '', repaired)              # trailing key: (no value)
    repaired = re.sub(r',\s*$', '', repaired)                          # trailing comma
    # Close open structures
    open_brackets = repaired.count('[') - repaired.count(']')
    open_braces = repaired.count('{') - repaired.count('}')
    repaired += ']' * max(0, open_brackets)
    repaired += '}' * max(0, open_braces)
    try:
        result = json.loads(repaired)
        print(f"[_repair_truncated_json] Repair succeeded")
        return result
    except json.JSONDecodeError as e:
        raise Exception(f"JSON truncated and repair failed: {e}")


def _update_job(job_id, **kwargs):
    """Update analysis_jobs row with provided fields."""
    if not kwargs:
        return
    set_clauses = []
    values = []
    for key, value in kwargs.items():
        set_clauses.append(f"{key} = %s")
        values.append(value if not isinstance(value, dict) else json.dumps(value))
    set_clauses.append("updated_at = CURRENT_TIMESTAMP")
    values.append(job_id)
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute(f"UPDATE analysis_jobs SET {', '.join(set_clauses)} WHERE id = %s", values)
    except Exception as e:
        print(f"[analysis-job {job_id}] Failed to update: {e}")


def _update_pipeline_job(job_id, **kwargs):
    """Update a pipeline job's fields."""
    fields = []
    values = []
    for k, v in kwargs.items():
        fields.append(f'{k} = %s')
        values.append(v)
    fields.append('updated_at = NOW()')
    if kwargs.get('status') in ('complete', 'failed'):
        fields.append('completed_at = NOW()')
    values.append(job_id)
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute(f"UPDATE research_pipeline_jobs SET {', '.join(fields)} WHERE id = %s", values)
    except Exception as e:
        print(f"[pipeline-job {job_id}] Failed to update: {e}")


def _run_trading_agent(run_id, ticker, date_str, provider, model):
    """Run TradingAgents analysis in background thread."""
    try:
        def _append_log(agent, message):
            entry = json.dumps([{"ts": datetime.utcnow().isoformat(), "agent": agent, "message": str(message)[:500]}])
            with get_db(commit=True) as (conn, cur):
                cur.execute("UPDATE agent_runs SET logs = logs || %s::jsonb WHERE id = %s", (entry, run_id))

        _append_log("system", f"Starting TradingAgents: {ticker} @ {date_str} using {provider}/{model}")

        from tradingagents.graph.trading_graph import TradingAgentsGraph
        from tradingagents.default_config import DEFAULT_CONFIG

        config = DEFAULT_CONFIG.copy()
        config["llm_provider"] = provider
        config["deep_think_llm"] = model
        config["quick_think_llm"] = model

        _append_log("system", "Initializing agent graph...")
        ta = TradingAgentsGraph(debug=True, config=config)

        # Capture print output as logs
        import builtins
        original_print = builtins.print
        def capturing_print(*args, **kwargs):
            msg = " ".join(str(a) for a in args)
            if msg.strip():
                try:
                    _append_log("agent", msg[:500])
                except Exception:
                    pass
            original_print(*args, **kwargs)
        builtins.print = capturing_print

        try:
            _append_log("system", "Running analysis...")
            final_state, decision_signal = ta.propagate(ticker, date_str)
        finally:
            builtins.print = original_print

        # Extract full reports from final_state
        full_report = f"""# TradingAgents Analysis: {ticker} ({date_str})

## Final Decision: {decision_signal}

## Full Trade Decision
{final_state.get('final_trade_decision', 'N/A')}

## Investment Plan
{final_state.get('investment_plan', 'N/A')}

## Trader's Investment Decision
{final_state.get('trader_investment_plan', 'N/A')}

---

## Market Report
{final_state.get('market_report', 'N/A')}

## Sentiment Report
{final_state.get('sentiment_report', 'N/A')}

## News Report
{final_state.get('news_report', 'N/A')}

## Fundamentals Report
{final_state.get('fundamentals_report', 'N/A')}

---

## Investment Debate

### Bull Case
{final_state.get('investment_debate_state', {}).get('bull_history', 'N/A')}

### Bear Case
{final_state.get('investment_debate_state', {}).get('bear_history', 'N/A')}

### Judge Decision
{final_state.get('investment_debate_state', {}).get('judge_decision', 'N/A')}

---

## Risk Assessment

### Aggressive View
{final_state.get('risk_debate_state', {}).get('aggressive_history', 'N/A')}

### Conservative View
{final_state.get('risk_debate_state', {}).get('conservative_history', 'N/A')}

### Neutral View
{final_state.get('risk_debate_state', {}).get('neutral_history', 'N/A')}

### Risk Judge Decision
{final_state.get('risk_debate_state', {}).get('judge_decision', 'N/A')}
"""

        decision_text = f"{decision_signal}\n\n{final_state.get('final_trade_decision', '')}"
        _append_log("system", f"Analysis complete. Decision: {decision_signal}")

        # Estimate cost based on model
        cost_per_1k = {
            'claude-haiku-4-5-20251001': 0.001, 'claude-sonnet-4-6': 0.003, 'claude-opus-4-6': 0.015,
            'gpt-4.1-mini': 0.0004, 'gpt-4.1': 0.002, 'o4-mini': 0.001,
            'gemini-2.0-flash': 0.0001, 'gemini-2.5-pro': 0.005,
        }
        # Rough estimate from log count (each log ~500 tokens avg)
        est_tokens = len(logs) * 500 if 'logs' in dir() else 5000
        rate = cost_per_1k.get(model, 0.001)
        est_cost = (est_tokens / 1000) * rate

        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                UPDATE agent_runs SET status = 'complete', decision = %s, report = %s,
                estimated_cost = %s, completed_at = NOW() WHERE id = %s
            ''', (decision_text, full_report, est_cost, run_id))

        # Check for thesis conflict and send Telegram alert
        try:
            with get_db() as (_, cur2):
                cur2.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
                pa_row = cur2.fetchone()
            if pa_row and pa_row['analysis']:
                pa_analysis = pa_row['analysis'] if isinstance(pa_row['analysis'], dict) else json.loads(pa_row['analysis'])
                conclusion = (pa_analysis.get('conclusion', '') or '').lower()
                thesis_dir = 'bullish' if any(w in conclusion for w in ['own', 'buy', 'bullish', 'add']) else 'bearish' if any(w in conclusion for w in ['sell', 'avoid', 'underweight']) else 'neutral'
                agent_dir = decision_signal.upper()
                if (agent_dir == 'SELL' and thesis_dir == 'bullish') or (agent_dir == 'BUY' and thesis_dir == 'bearish'):
                    _notify_telegram(f"*THESIS CONFLICT:* TradingAgents says {agent_dir} on {ticker}, but your thesis is {thesis_dir}. Review?")
        except Exception:
            pass

        print(f"[trading-agent {run_id}] Complete: {ticker}")

    except Exception as e:
        print(f"[trading-agent {run_id}] Failed: {e}")
        import traceback
        traceback.print_exc()
        try:
            _append_log("error", str(e)[:500])
        except Exception:
            pass
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                UPDATE agent_runs SET status = 'error', decision = %s, completed_at = NOW() WHERE id = %s
            ''', (f"ERROR: {str(e)}", run_id))


def _run_pipeline_batch(batch_id, tickers, job_type, api_key):
    """Background thread: process a batch of tickers through the research pipeline sequentially."""
    print(f"[pipeline-batch {batch_id}] Starting batch for {len(tickers)} tickers: {tickers}")
    for ticker in tickers:
        try:
            with get_db() as (_, cur):
                cur.execute('SELECT id, status FROM research_pipeline_jobs WHERE batch_id = %s AND ticker = %s', (batch_id, ticker))
                job = cur.fetchone()
            if not job:
                print(f"[pipeline-batch {batch_id}] No job found for {ticker}, skipping")
                continue
            if job['status'] == 'cancelled':
                print(f"[pipeline-batch {batch_id}] Job {job['id']} for {ticker} was cancelled, skipping")
                continue
            _run_single_pipeline_job(job['id'], ticker, job_type, api_key)
        except Exception as e:
            print(f"[pipeline-batch {batch_id}] Error processing {ticker}: {e}")
    print(f"[pipeline-batch {batch_id}] Batch complete")


def _run_single_pipeline_job(job_id, ticker, job_type, api_key):
    """Process a single ticker through the research pipeline."""
    steps = [
        'Checking existing analysis',
        'Scanning for documents',
        'Processing documents',
        'Generating analysis',
        'Building deliverables',
        'Saving results',
        'Complete'
    ]
    try:
        _update_pipeline_job(job_id, status='running', current_step=steps[0], progress=0)

        # Step 1: Check existing analysis
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            existing = cur.fetchone()
        _update_pipeline_job(job_id, current_step=steps[1], progress=14)

        # Step 2: Scan for ALL documents (both Thesis and Research storage)
        thesis_docs = []
        research_docs = []
        with get_db() as (_, cur):
            cur.execute('SELECT id, filename, file_data, file_type, mime_type, file_size, metadata FROM document_files WHERE ticker = %s ORDER BY created_at DESC', (ticker,))
            thesis_docs = [dict(r) for r in cur.fetchall()]

        with get_db() as (_, cur):
            cur.execute('''
                SELECT rdf.id, rdf.filename, rdf.file_data, rdf.file_type, rdf.file_size,
                       rd.name as doc_name
                FROM research_document_files rdf
                JOIN research_documents rd ON rdf.document_id = rd.id
                JOIN research_categories rc ON rd.category_id = rc.id
                WHERE UPPER(rc.name) = %s
            ''', (ticker,))
            research_docs = [dict(r) for r in cur.fetchall()]

        all_available = []
        for d in thesis_docs:
            all_available.append({
                'id': f'thesis_{d["id"]}',
                'filename': d['filename'],
                'file_data': d.get('file_data', ''),
                'file_type': d.get('file_type', 'pdf'),
                'mime_type': d.get('mime_type', 'application/pdf'),
                'source': 'thesis',
                'weight': 1.0,
            })
        for d in research_docs:
            all_available.append({
                'id': f'research_{d["id"]}',
                'filename': d['filename'],
                'file_data': d.get('file_data', ''),
                'file_type': d.get('file_type', 'pdf'),
                'mime_type': d.get('mime_type', 'application/pdf'),
                'source': 'research',
                'weight': 1.0,
            })

        # If no docs in DB, check local manifest and request upload from local agent
        if not all_available:
            local_files = _local_file_manifest.get('manifest', {}).get(ticker, [])
            source_exts = {'.pdf', '.xlsx', '.xls', '.csv'}
            local_source_files = [f for f in local_files
                                  if f.get('extension', '').lower() in source_exts
                                  and f.get('folder', '') == 'main']
            if local_source_files:
                # Request urgent upload from local agent
                print(f'[pipeline {job_id}] No docs in DB for {ticker}, but {len(local_source_files)} found in iCloud manifest. Requesting upload...')
                _pending_doc_upload_requests[ticker] = {
                    'requested_at': datetime.utcnow(),
                    'job_id': job_id,
                }
                _update_pipeline_job(job_id, current_step=f'Waiting for {len(local_source_files)} iCloud docs to upload...', progress=18)

                # Poll for docs to appear in DB (local agent polls every 5s, upload takes a few seconds)
                max_wait = 60  # seconds
                poll_interval = 3
                waited = 0
                while waited < max_wait:
                    time.sleep(poll_interval)
                    waited += poll_interval
                    with get_db() as (_, cur):
                        cur.execute('SELECT COUNT(*) as cnt FROM document_files WHERE ticker = %s', (ticker,))
                        cnt = cur.fetchone()['cnt']
                    if cnt > 0:
                        print(f'[pipeline {job_id}] {cnt} docs now available in DB for {ticker} after {waited}s wait')
                        break

                # Re-query document_files after waiting
                with get_db() as (_, cur):
                    cur.execute('SELECT id, filename, file_data, file_type, mime_type, file_size, metadata FROM document_files WHERE ticker = %s ORDER BY created_at DESC', (ticker,))
                    thesis_docs = [dict(r) for r in cur.fetchall()]
                all_available = []
                for d in thesis_docs:
                    all_available.append({
                        'id': f'thesis_{d["id"]}',
                        'filename': d['filename'],
                        'file_data': d.get('file_data', ''),
                        'file_type': d.get('file_type', 'pdf'),
                        'mime_type': d.get('mime_type', 'application/pdf'),
                        'source': 'thesis',
                        'weight': 1.0,
                    })
                # Clean up request
                _pending_doc_upload_requests.pop(ticker, None)

            if not all_available and not existing:
                local_hint = f' ({len(local_source_files)} files seen in iCloud but upload timed out — is the local agent running?)' if local_source_files else ''
                raise Exception(f'No documents or existing analysis found for {ticker}{local_hint}')

        # Apply document config (inclusion/exclusion/weights)
        doc_config = {}
        with get_db() as (_, cur):
            cur.execute('SELECT steps_detail FROM research_pipeline_jobs WHERE id = %s', (job_id,))
            job_row = cur.fetchone()
            if job_row and job_row['steps_detail']:
                sd = job_row['steps_detail']
                if isinstance(sd, str):
                    try: sd = json.loads(sd)
                    except: sd = {}
                doc_config = sd.get('documentConfig', {})

        excluded_historical = doc_config.get('excludedHistoricalDocs', [])
        rebuild_from_scratch = doc_config.get('rebuildFromScratch', False)

        config_docs = doc_config.get('documents', [])
        config_map = {cd['id']: cd for cd in config_docs}

        selected_docs = []
        for doc in all_available:
            cfg = config_map.get(doc['id'], {})
            if cfg.get('included') is False:
                continue  # Explicitly excluded
            doc['weight'] = cfg.get('weight', doc['weight'])
            selected_docs.append(doc)

        _update_pipeline_job(job_id, current_step=f'Processing {len(selected_docs)} documents', progress=28)

        # Step 3: Ensure all documents exist in document_files for the analysis job
        # Research-sourced docs need to be copied to document_files so _run_analysis_job can find them
        for doc in selected_docs:
            if doc['source'] == 'research' and doc['file_data']:
                with get_db(commit=True) as (conn, cur):
                    cur.execute('''
                        INSERT INTO document_files (ticker, filename, file_data, file_type, mime_type, file_size, metadata, created_at)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, NOW())
                        ON CONFLICT (ticker, filename) DO NOTHING
                    ''', (ticker, doc['filename'], doc['file_data'], doc['file_type'], doc['mime_type'],
                          len(doc['file_data']) if doc['file_data'] else 0,
                          json.dumps({'source': 'research', 'copiedByPipeline': True})))

        _update_pipeline_job(job_id, current_step=steps[3], progress=42)

        # Step 4: Generate/update analysis via LLM
        # Use existing analysis job infrastructure
        doc_filenames = [d['filename'] for d in selected_docs]
        doc_details_list = [{'filename': d['filename'], 'weight': 1, 'isNew': True} for d in selected_docs]

        # Build weighting config from document config
        existing_weight = 0 if rebuild_from_scratch else doc_config.get('existingWeight', 70)
        new_weight = 100 - existing_weight
        weighting_config = {
            'mode': 'simple',
            'existingAnalysisWeight': existing_weight,
            'newDocsWeight': new_weight,
            'excludedHistoricalDocs': excluded_historical,
            'rebuildFromScratch': rebuild_from_scratch,
        }

        # For rebuild from scratch, don't pass existing analysis to the LLM
        effective_existing = None if rebuild_from_scratch else existing

        if not doc_filenames and not existing:
            raise Exception(f'No documents found for {ticker}. Ensure source files exist in iCloud and local agent is running.')

        if doc_filenames or (job_type == 'process' and existing):
            # Create an analysis sub-job
            sub_job_id = str(uuid.uuid4())
            with get_db(commit=True) as (conn, cur):
                cur.execute('''INSERT INTO analysis_jobs (id, ticker, status, progress, created_at, updated_at, api_key, request_payload)
                    VALUES (%s, %s, 'pending', '', NOW(), NOW(), %s, %s)''',
                    (sub_job_id, ticker, api_key, json.dumps({
                        'ticker': ticker,
                        'documentFilenames': doc_filenames,
                        'documentDetails': doc_details_list,
                        'existingAnalysis': {k: (v.isoformat() if hasattr(v, 'isoformat') else v) for k, v in dict(effective_existing).items()} if effective_existing else None,
                        'weightingConfig': weighting_config,
                    })))

            # Run the analysis synchronously within this thread
            _run_analysis_job(sub_job_id)

            # Wait for completion and check result
            with get_db() as (_, cur):
                cur.execute('SELECT status, result, error FROM analysis_jobs WHERE id = %s', (sub_job_id,))
                sub_result = cur.fetchone()

            if sub_result and sub_result['status'] == 'failed':
                raise Exception(f'Analysis failed: {sub_result.get("error", "Unknown error")}')

            # Save analysis result to portfolio_analyses
            if sub_result and sub_result['status'] == 'complete' and sub_result.get('result'):
                _update_pipeline_job(job_id, current_step='Saving analysis to portfolio', progress=57)
                result_data = sub_result['result'] if isinstance(sub_result['result'], dict) else json.loads(sub_result['result'])
                analysis_data = result_data.get('analysis', {})
                # Unwrap any depth of nested 'analysis' keys (LLM sometimes wraps its output)
                while 'analysis' in analysis_data and isinstance(analysis_data['analysis'], dict) and 'thesis' not in analysis_data:
                    analysis_data = analysis_data['analysis']
                analysis_changes = result_data.get('changes', [])
                doc_metadata = result_data.get('documentMetadata', [])
                # Company: prefer existing DB column, then analysis data, then fallback
                company = (existing['company'] if existing and existing.get('company') else '') or analysis_data.pop('company', '') or ticker
                # Preserve changes in analysis for the Thesis tab to display
                if analysis_changes:
                    analysis_data['_pipelineChanges'] = analysis_changes

                # Preserve history from existing analysis (skip on rebuild from scratch)
                if existing and not rebuild_from_scratch:
                    old_analysis = existing.get('analysis', {})
                    if isinstance(old_analysis, str):
                        try: old_analysis = json.loads(old_analysis)
                        except: old_analysis = {}
                    if not analysis_data.get('history'):
                        old_history = old_analysis.get('history') or []
                        if old_analysis.get('thesis') or old_analysis.get('signposts') or old_analysis.get('threats'):
                            old_history = list(old_history)
                            old_history.append({
                                'timestamp': old_analysis.get('updatedAt') or datetime.utcnow().isoformat(),
                                'thesis': old_analysis.get('thesis'),
                                'signposts': old_analysis.get('signposts'),
                                'threats': old_analysis.get('threats'),
                            })
                            old_history = old_history[-20:]
                        analysis_data['history'] = old_history
                    # Merge old documentHistory into new (don't lose docs from prior runs)
                    old_dh = old_analysis.get('documentHistory') or []
                    new_dh = analysis_data.get('documentHistory') or []
                    new_fnames = {d.get('filename') for d in new_dh}
                    for old_doc in old_dh:
                        if old_doc.get('filename') and old_doc['filename'] not in new_fnames:
                            new_dh.append(old_doc)
                            new_fnames.add(old_doc['filename'])
                    analysis_data['documentHistory'] = new_dh

                # Remove excluded historical docs from documentHistory
                if excluded_historical:
                    excluded_set = set(excluded_historical)
                    analysis_data['documentHistory'] = [
                        dh for dh in analysis_data.get('documentHistory', [])
                        if dh.get('filename') not in excluded_set
                    ]
                analysis_data['updatedAt'] = datetime.utcnow().isoformat()
                # Inject company + ticker into analysis JSON (frontend reads from here)
                analysis_data['company'] = company
                analysis_data['ticker'] = ticker

                with get_db(commit=True) as (conn, cur):
                    cur.execute('''
                        INSERT INTO portfolio_analyses (ticker, company, analysis, updated_at)
                        VALUES (%s, %s, %s, %s)
                        ON CONFLICT (ticker)
                        DO UPDATE SET company = EXCLUDED.company, analysis = EXCLUDED.analysis, updated_at = EXCLUDED.updated_at
                    ''', (ticker, company, json.dumps(analysis_data), datetime.utcnow()))
                print(f'[pipeline {job_id}] Saved analysis to portfolio_analyses for {ticker}')

                # Auto-snapshot for thesis evolution tracking
                try:
                    sc_data = None
                    with get_db() as (_, cur2):
                        cur2.execute('SELECT scorecard_data FROM thesis_scorecard_data WHERE ticker = %s', (ticker,))
                        sc_row = cur2.fetchone()
                    if sc_row:
                        sc_data = sc_row['scorecard_data']
                        if isinstance(sc_data, str):
                            try: sc_data = json.loads(sc_data)
                            except: sc_data = None
                    _create_thesis_snapshot(ticker, analysis_data, sc_data, 'pipeline')
                except Exception as snap_err:
                    print(f'[pipeline {job_id}] Snapshot error: {snap_err}')

        _update_pipeline_job(job_id, current_step=steps[4], progress=71)

        # Step 5: Build deliverables summary
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            final_analysis = cur.fetchone()

        result = {
            'ticker': ticker,
            'hasAnalysis': final_analysis is not None,
            'documentsProcessed': len(selected_docs),
            'completedAt': datetime.utcnow().isoformat(),
        }

        _update_pipeline_job(job_id, current_step=steps[5], progress=85)

        # Step 7-8-9: Generate thesis tiers if requested
        with get_db() as (_, cur):
            cur.execute('SELECT steps_detail FROM research_pipeline_jobs WHERE id = %s', (job_id,))
            sd_row = cur.fetchone()
        sd = sd_row['steps_detail'] if sd_row and sd_row['steps_detail'] else {}
        if isinstance(sd, str):
            try: sd = json.loads(sd)
            except: sd = {}
        generate_tiers = sd.get('generateTiers', 'detailed')

        if generate_tiers in ('full', 'all'):
            # Load the just-saved analysis
            with get_db() as (_, cur):
                cur.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
                ana_row = cur.fetchone()
            if ana_row:
                analysis_data = ana_row['analysis'] if isinstance(ana_row['analysis'], dict) else json.loads(ana_row['analysis'])

                # Generate Full tier
                _update_pipeline_job(job_id, current_step='Generating Full tier', progress=86)
                full_data = None
                try:
                    full_data = _build_full_tier(analysis_data, anthropic_key=api_key, gemini_key="")
                    with get_db(commit=True) as (conn, cur):
                        _snapshot_tier_history(cur, 'thesis_full', 'full_analysis', ticker)
                    with get_db(commit=True) as (conn, cur):
                        cur.execute('''
                            INSERT INTO thesis_full (ticker, full_analysis, updated_at)
                            VALUES (%s, %s, CURRENT_TIMESTAMP)
                            ON CONFLICT (ticker)
                            DO UPDATE SET full_analysis = EXCLUDED.full_analysis, updated_at = CURRENT_TIMESTAMP
                        ''', (ticker, json.dumps(full_data)))
                    print(f'[pipeline {job_id}] Full tier generated for {ticker}')
                except Exception as tier_err:
                    print(f'[pipeline {job_id}] Full tier generation failed for {ticker}: {tier_err}')

                if generate_tiers == 'all':
                    # Generate Condensed tier from Full (fall back to Detailed if Full failed)
                    _update_pipeline_job(job_id, current_step='Generating Condensed tier', progress=93)
                    try:
                        condensed_source = full_data if full_data else analysis_data
                        condensed_data = _build_condensed_tier(condensed_source, anthropic_key=api_key, gemini_key="")
                        with get_db(commit=True) as (conn, cur):
                            _snapshot_tier_history(cur, 'thesis_condensed', 'condensed_analysis', ticker)
                        with get_db(commit=True) as (conn, cur):
                            cur.execute('''
                                INSERT INTO thesis_condensed (ticker, condensed_analysis, updated_at)
                                VALUES (%s, %s, CURRENT_TIMESTAMP)
                                ON CONFLICT (ticker)
                                DO UPDATE SET condensed_analysis = EXCLUDED.condensed_analysis, updated_at = CURRENT_TIMESTAMP
                            ''', (ticker, json.dumps(condensed_data)))
                        print(f'[pipeline {job_id}] Condensed tier generated for {ticker}')
                    except Exception as tier_err:
                        print(f'[pipeline {job_id}] Condensed tier generation failed for {ticker}: {tier_err}')

        # Final step: Save results and mark complete
        _update_pipeline_job(job_id, status='complete', current_step='Complete', progress=100, result=json.dumps(result))
        print(f"[pipeline-job {job_id}] Completed successfully for {ticker}")
        doc_count = len(result.get('documentHistory', []))
        _notify_telegram(f"*Pipeline:* {ticker} thesis update complete ({doc_count} docs)")

    except Exception as e:
        print(f'[pipeline-job {job_id}] Failed for {ticker}: {e}')
        _update_pipeline_job(job_id, status='failed', error=str(e), current_step='Failed')
        _notify_telegram(f"*Pipeline:* {ticker} thesis update FAILED\n{str(e)[:200]}")


def _count_pdf_pages(base64_data):
    """Count pages in a base64-encoded PDF."""
    try:
        pdf_bytes = base64.b64decode(base64_data)
        reader = PdfReader(io.BytesIO(pdf_bytes))
        return len(reader.pages)
    except:
        return 30


def _build_analysis_content(batch_docs, all_docs, existing_analysis, historical_weights, weighting_config, batch_num, total_batches):
    """Build the Claude content array for analysis (documents + prompt)."""
    simple_mode = weighting_config.get('mode') == 'simple'
    existing_weight = weighting_config.get('existingAnalysisWeight', 70)
    new_docs_weight = weighting_config.get('newDocsWeight', 30)
    total_weight = sum(d.get('weight', 1) for d in all_docs) + sum(d.get('weight', 1) for d in (historical_weights or []))

    content = []

    # Document blocks
    for doc in batch_docs:
        doc_content = doc.get('file_data') or doc.get('fileData', '')
        doc_name = doc.get('filename', 'document.pdf')
        doc_type = doc.get('file_type') or doc.get('fileType', 'pdf')
        mime_type = doc.get('mime_type') or doc.get('mimeType', 'application/pdf')
        is_new = doc.get('isNew', True)

        if not doc_content:
            continue

        # Header with weight
        if simple_mode and existing_analysis and is_new:
            new_count = len([d for d in all_docs if d.get('isNew', True)])
            per_new = new_docs_weight / max(new_count, 1)
            header = f"\n=== NEW DOCUMENT ({round(per_new)}% update weight): {doc_name} ==="
        elif simple_mode and existing_analysis:
            header = f"\n=== EXISTING DOCUMENT (Reference): {doc_name} ==="
        else:
            doc_w = doc.get('weight', 1)
            pct = round((doc_w / total_weight) * 100) if total_weight > 0 else 0
            header = f"\n=== DOCUMENT: {doc_name} (Weight: {pct}%) ==="

        content.append({"type": "text", "text": header})

        if doc_type == 'pdf':
            content.append({"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": doc_content}})
        elif doc_type == 'image':
            content.append({"type": "image", "source": {"type": "base64", "media_type": mime_type or "image/png", "data": doc_content}})
        else:
            try:
                content.append({"type": "text", "text": base64.b64decode(doc_content).decode('utf-8')})
            except:
                continue

    if not content:
        return content

    # Build prompt
    STYLE_RULES = """IMPORTANT STYLE RULES:
- Do NOT reference any sellside broker names
- Do NOT reference specific analyst names
- Do NOT include specific broker price targets
- Write as independent analysis without attribution to sources"""

    JSON_SCHEMA = '{"ticker":"TICKER","company":"Name","thesis":{"summary":"...","pillars":[{"title":"...","description":"...","confidence":"High/Medium/Low","sources":[{"filename":"...","excerpt":"..."}]}]},"signposts":[{"metric":"...","target":"...","timeframe":"...","category":"Financial/Operational/Strategic/Market","confidence":"High/Medium/Low","sources":[]}],"threats":[{"threat":"...","likelihood":"...","impact":"...","triggerPoints":"...","sources":[]}],"documentMetadata":[{"filename":"...","docType":"broker_report","source":"...","publishDate":"YYYY-MM-DD","authors":[],"title":"..."}]}'

    LENGTH_RULES = """CONCISENESS RULES (the final output must fit in 2-3 printed pages):
- thesis.summary: 2-3 sentences MAX. One crisp paragraph capturing WHY we own it.
- pillars: 3-5 pillars. Each description is 1-2 sentences — punchy, not exhaustive.
- signposts: 4-6 signposts MAX. Combine related metrics rather than listing every variant.
- threats: 3-5 threats MAX. Each threat description is one sentence. triggerPoints is one sentence.
- conclusion: 1-2 sentences if included.
- Prioritize the MOST important and differentiated insights. Omit generic/obvious points.
- Do NOT pad with boilerplate or repeat the same point across pillars."""

    if existing_analysis:
        batch_note = f"\n\nNote: This is batch {batch_num} of {total_batches}. Incorporate these documents into the existing analysis.\n" if total_batches > 1 else ''
        if simple_mode:
            weight_instr = f"""CRITICAL WEIGHTING INSTRUCTION:
You MUST preserve {existing_weight}% of the existing analysis. The new documents can only contribute {new_docs_weight}% worth of changes.

What this means:
- KEEP {existing_weight}% of the existing thesis, pillars, signposts, and threats UNCHANGED
- Only make refinements or additions proportional to the {new_docs_weight}% new-document weight
- Do NOT fundamentally rewrite or replace the existing analysis
- Do NOT treat the new documents as "primary sources" — they are SUPPLEMENTARY to the existing analysis
- The new documents should ADD to or REFINE the existing analysis, not replace it

Example of correct behavior with {existing_weight}% existing / {new_docs_weight}% new:
- If existing thesis has 3 pillars, keep all 3, maybe slightly update wording or add a 4th minor pillar
- If existing has 5 signposts, keep them mostly intact, maybe add 1-2 new ones or update targets slightly
- Do NOT remove or majorly rewrite existing content unless it's factually contradicted by the new documents

In the "changes" array, describe what updates were made and why."""
        else:
            weight_instr = """DOCUMENT WEIGHTING:
- Each document has an assigned weight percentage shown at the start
- Give MORE emphasis to higher-weighted documents when forming conclusions
- Higher-weighted documents should have more influence on the thesis, signposts, and threats
- If documents conflict, prefer the view from the higher-weighted document"""
        # Build stale document exclusion instruction if any historical docs are excluded
        excluded_docs = weighting_config.get('excludedHistoricalDocs', [])
        stale_doc_instruction = ''
        if excluded_docs:
            doc_list_str = '\n'.join(f'  - {fname}' for fname in excluded_docs)
            stale_doc_instruction = f"""

CRITICAL — STALE DOCUMENT EXCLUSION:
The existing analysis was partly based on the following document(s) which are now STALE and must be DISREGARDED:
{doc_list_str}

You MUST actively revise the analysis to:
1. REMOVE or UPDATE any data points, statistics, or financial figures that originated from these stale documents
2. REVISE any arguments, conclusions, or thesis pillars that relied primarily on stale data
3. REMOVE signposts or threats whose basis was the stale document data
4. If a pillar used stale data as supporting evidence but has other valid support, keep the argument but remove the stale references
5. In the "changes" array, explicitly note what was removed or revised due to stale document exclusion
Do NOT simply ignore these documents — actively scrub their influence from the existing analysis."""

        # Trim verbose fields to prevent context bloat on multi-batch updates
        trimmed = _trim_analysis_for_context(existing_analysis)
        prompt = f"""Update this existing analysis with new information from the documents.{batch_note}

Existing Analysis:
{json.dumps(trimmed, indent=2)}

{weight_instr}
{stale_doc_instruction}

Review the new documents and:
1. Update or confirm the investment thesis
2. Add any new signposts or update existing ones — consolidate to 4-6 total
3. Add any new threats or update existing ones — consolidate to 3-5 total
4. Note what changed in the "changes" array
5. TRIM any bloat: if pillars > 5, merge or drop the weakest; if signposts > 6, combine related ones

{STYLE_RULES}

{LENGTH_RULES}

Return the updated analysis as JSON with the same structure plus a "changes" array.
Return ONLY valid JSON, no markdown, no explanation."""
    else:
        batch_note = f"\n\nNote: This is batch {batch_num} of {total_batches}. More documents will follow.\n" if total_batches > 1 else ''
        prompt = f"""Analyze these documents and create a comprehensive investment analysis.{batch_note}

Return a JSON object with this structure:
{JSON_SCHEMA}

{STYLE_RULES}

{LENGTH_RULES}

Focus on:
1. Investment Thesis with confidence and source citations
2. Signposts - specific KPIs, events, milestones
3. Threats - bear case scenarios with likelihood, impact, trigger points

Return ONLY valid JSON, no markdown, no explanation."""

    content.append({"type": "text", "text": prompt})
    return content


def _run_analysis_job(job_id):
    """Background thread: run full analysis pipeline."""
    try:
        # Load job
        with get_db() as (conn, cur):
            cur.execute('SELECT * FROM analysis_jobs WHERE id = %s', (job_id,))
            job = cur.fetchone()
        if not job:
            return

        api_key = job['api_key']
        payload = job['request_payload'] if isinstance(job['request_payload'], dict) else json.loads(job['request_payload'])
        ticker = job['ticker']
        existing_analysis = payload.get('existingAnalysis')
        historical_weights = payload.get('historicalWeights', [])
        weighting_config = payload.get('weightingConfig', {})
        doc_details = {d['filename']: d for d in payload.get('documentDetails', [])}

        _update_job(job_id, status='processing', progress='Loading documents...')

        # Load documents from DB
        doc_filenames = payload.get('documentFilenames', [])
        if doc_filenames:
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(doc_filenames))
                cur.execute(f'SELECT filename, file_data, file_type, mime_type, metadata FROM document_files WHERE ticker = %s AND filename IN ({placeholders})', [ticker] + doc_filenames)
                docs = cur.fetchall()
        else:
            docs = []

        # Also include any inline documents from the request (new uploads not yet in DB)
        inline_docs = payload.get('inlineDocuments', [])

        all_docs = []
        for doc in docs:
            d = dict(doc)
            detail = doc_details.get(d['filename'], {})
            d['weight'] = detail.get('weight', 1)
            d['isNew'] = detail.get('isNew', True)
            d['_pages'] = _count_pdf_pages(d['file_data']) if d.get('file_type') == 'pdf' else 0
            all_docs.append(d)

        for doc in inline_docs:
            doc['_pages'] = _count_pdf_pages(doc.get('fileData', '')) if doc.get('fileType') == 'pdf' else 0
            doc['file_data'] = doc.get('fileData', '')
            doc['file_type'] = doc.get('fileType', 'pdf')
            doc['mime_type'] = doc.get('mimeType', 'application/pdf')
            all_docs.append(doc)

        # Auto-split PDFs that exceed Claude's 100-page-per-document limit
        expanded_docs = []
        for doc in all_docs:
            if doc.get('file_type', doc.get('fileType', '')) == 'pdf' and doc['_pages'] > MAX_PDF_PAGES_PER_CHUNK:
                file_data = doc.get('file_data') or doc.get('fileData', '')
                chunks = _split_large_pdf(file_data)
                base_name = doc.get('filename', 'document.pdf')
                print(f"[analysis-job {job_id}] Split {base_name} ({doc['_pages']}p) into {len(chunks)} chunks")
                for ci, (chunk_data, chunk_pages) in enumerate(chunks, 1):
                    chunk_doc = dict(doc)
                    chunk_doc['file_data'] = chunk_data
                    chunk_doc['fileData'] = chunk_data
                    chunk_doc['_pages'] = chunk_pages
                    if len(chunks) > 1:
                        chunk_doc['filename'] = f"{base_name} (part {ci}/{len(chunks)})"
                    expanded_docs.append(chunk_doc)
            else:
                expanded_docs.append(doc)
        all_docs = expanded_docs

        total_pages = sum(d['_pages'] for d in all_docs)

        # Split into batches of <=MAX_PAGES_PER_BATCH pages
        if total_pages <= MAX_PAGES_PER_BATCH:
            batches = [all_docs]
        else:
            batches = []
            batch = []
            batch_pages = 0
            for doc in all_docs:
                if batch_pages + doc['_pages'] > MAX_PAGES_PER_BATCH and batch:
                    batches.append(batch)
                    batch = []
                    batch_pages = 0
                batch.append(doc)
                batch_pages += doc['_pages']
            if batch:
                batches.append(batch)

        _update_job(job_id, total_batches=len(batches), progress=f'{total_pages} pages, {len(batches)} batch(es)')

        # Process batches — 5 min timeout per batch, 1 retry on failure
        import httpx
        client = anthropic.Anthropic(
            api_key=api_key,
            timeout=httpx.Timeout(300.0, connect=30.0),  # 5 min total, 30s connect
        )
        current_analysis = existing_analysis
        all_changes = []
        total_usage = {'input_tokens': 0, 'output_tokens': 0}

        batch_idx = 0
        while batch_idx < len(batches):
            batch = batches[batch_idx]
            batch_pages = sum(d['_pages'] for d in batch)
            _update_job(job_id, batch_num=batch_idx + 1, chars_received=0,
                        progress=f'Processing batch {batch_idx + 1} of {len(batches)} ({batch_pages} pages)')

            content = _build_analysis_content(batch, all_docs, current_analysis, historical_weights, weighting_config, batch_idx + 1, len(batches))
            if not content:
                batch_idx += 1
                continue

            # Retry loop: try up to 2 attempts per batch
            max_attempts = 2
            batch_split = False
            for attempt in range(1, max_attempts + 1):
                try:
                    result_text = ""
                    usage_data = {}
                    if attempt > 1:
                        _update_job(job_id, chars_received=0,
                                    progress=f'Retrying batch {batch_idx + 1} of {len(batches)} (attempt {attempt})')
                        print(f"[analysis-job {job_id}] Retry batch {batch_idx + 1}, attempt {attempt}")

                    with client.messages.stream(
                        model="claude-sonnet-4-5-20250929",
                        max_tokens=64000,
                        messages=[{'role': 'user', 'content': content}],
                        system="You are an expert equity research analyst. Analyze documents thoroughly and provide institutional-quality investment analysis. Be concise: the final thesis should fit 2-3 printed pages (3-5 pillars, 4-6 signposts, 3-5 threats, each described in 1-2 sentences). Prioritize the most important insights. Always respond with valid JSON only."
                    ) as stream:
                        for text in stream.text_stream:
                            result_text += text
                            if len(result_text) % 500 < len(text):
                                _update_job(job_id, chars_received=len(result_text))
                        final_msg = stream.get_final_message()
                        usage_data = {
                            'input_tokens': final_msg.usage.input_tokens,
                            'output_tokens': final_msg.usage.output_tokens
                        }

                    # Parse JSON — with truncation repair
                    cleaned = result_text.strip()
                    if cleaned.startswith('```'):
                        cleaned = cleaned.split('\n', 1)[1] if '\n' in cleaned else cleaned[3:]
                    if cleaned.endswith('```'):
                        cleaned = cleaned.rsplit('\n', 1)[0]
                    if cleaned.startswith('json'):
                        cleaned = cleaned[4:].strip()

                    try:
                        analysis = json.loads(cleaned)
                    except json.JSONDecodeError as je:
                        print(f"[analysis-job {job_id}] JSON parse failed ({je}), attempting repair...")
                        analysis = _repair_truncated_json(cleaned)
                    break  # success — exit retry loop

                except (httpx.TimeoutException, httpx.ReadTimeout, httpx.ConnectTimeout) as e:
                    print(f"[analysis-job {job_id}] Batch {batch_idx + 1} timeout (attempt {attempt}): {e}")
                    if attempt == max_attempts:
                        raise Exception(f"Batch {batch_idx + 1} timed out after {max_attempts} attempts ({batch_pages} pages)")
                    import time
                    time.sleep(5)
                except json.JSONDecodeError as e:
                    print(f"[analysis-job {job_id}] Batch {batch_idx + 1} JSON parse error (attempt {attempt}): {e}")
                    if attempt == max_attempts:
                        raise Exception(f"Batch {batch_idx + 1} returned invalid JSON after {max_attempts} attempts")
                    import time
                    time.sleep(3)
                except Exception as e:
                    error_str = str(e)
                    # Auto-split: if prompt too long and batch has >1 doc, split in half and retry
                    if 'prompt is too long' in error_str and len(batch) > 1:
                        mid = len(batch) // 2
                        left, right = batch[:mid], batch[mid:]
                        left_pages = sum(d['_pages'] for d in left)
                        right_pages = sum(d['_pages'] for d in right)
                        print(f"[analysis-job {job_id}] Prompt too long ({batch_pages}p), auto-splitting into {left_pages}p + {right_pages}p")
                        batches[batch_idx:batch_idx + 1] = [left, right]
                        _update_job(job_id, total_batches=len(batches),
                                    progress=f'Auto-split: now {len(batches)} batches')
                        batch_split = True
                        break  # exit retry loop — while loop will retry with smaller first half
                    print(f"[analysis-job {job_id}] Batch {batch_idx + 1} error (attempt {attempt}): {e}")
                    if attempt == max_attempts:
                        raise
                    import time
                    time.sleep(5)

            if batch_split:
                continue  # retry same index with the smaller first-half batch

            changes = analysis.pop('changes', [])
            analysis.pop('documentMetadata', None)  # discard LLM metadata — we build it deterministically

            # Unwrap if LLM wrapped its output in an 'analysis' key
            while 'analysis' in analysis and isinstance(analysis['analysis'], dict) and 'thesis' not in analysis:
                analysis = analysis['analysis']
            current_analysis = analysis
            all_changes.extend(changes)
            total_usage['input_tokens'] += usage_data.get('input_tokens', 0)
            total_usage['output_tokens'] += usage_data.get('output_tokens', 0)
            batch_idx += 1

        # Build documentHistory: seed from prior runs, then add current run's docs
        if isinstance(current_analysis, dict):
            # 1. Start with documentHistory from existing analysis (preserves docs from prior runs)
            prior_dh = []
            if existing_analysis and isinstance(existing_analysis, dict):
                prior_dh = list(existing_analysis.get('documentHistory', []))
            # 2. Also keep any the LLM may have returned in current_analysis
            existing_dh = list(prior_dh)
            existing_fnames = {d.get('filename') for d in existing_dh}
            for dh in current_analysis.get('documentHistory', []):
                if dh.get('filename') and dh['filename'] not in existing_fnames:
                    existing_dh.append(dh)
                    existing_fnames.add(dh['filename'])
            # 3. Add current run's docs deterministically
            for doc in all_docs:
                fname = doc.get('filename', '')
                if fname and fname not in existing_fnames:
                    existing_dh.append({
                        'filename': fname,
                        'docType': doc.get('file_type', 'pdf'),
                        'pages': doc.get('_pages', 0),
                        'processedAt': datetime.utcnow().isoformat(),
                    })
                    existing_fnames.add(fname)
            current_analysis['documentHistory'] = existing_dh

        # Store result
        result = {
            'analysis': current_analysis,
            'changes': all_changes,
            'usage': total_usage
        }
        _update_job(job_id, status='complete', result=json.dumps(result), progress='Analysis complete', api_key='')
        print(f"[analysis-job {job_id}] Complete: {ticker}")

    except Exception as e:
        print(f"[analysis-job {job_id}] Failed: {e}")
        import traceback
        traceback.print_exc()
        _update_job(job_id, status='failed', error=str(e), api_key='')


@app.route('/api/analysis-job', methods=['POST'])
def create_analysis_job():
    """Create a background analysis job."""
    try:
        data = request.json
        api_key = data.get('apiKey', '')
        ticker = data.get('ticker', '').upper()
        if not api_key:
            return jsonify({'error': 'API key required'}), 400
        if not ticker:
            return jsonify({'error': 'Ticker required'}), 400

        job_id = str(uuid.uuid4())[:12]

        # Save any new inline documents to document_files first
        new_docs = data.get('newDocuments', [])
        if new_docs:
            with get_db(commit=True) as (conn, cur):
                for doc in new_docs:
                    cur.execute('''
                        INSERT INTO document_files (ticker, filename, file_data, file_type, mime_type, metadata, file_size)
                        VALUES (%s, %s, %s, %s, %s, %s, %s)
                        ON CONFLICT (ticker, filename) DO UPDATE SET
                            file_data = EXCLUDED.file_data,
                            file_type = EXCLUDED.file_type,
                            mime_type = EXCLUDED.mime_type,
                            metadata = EXCLUDED.metadata,
                            file_size = EXCLUDED.file_size
                    ''', (ticker, doc['filename'], doc.get('fileData', ''), doc.get('fileType', 'pdf'),
                          doc.get('mimeType', 'application/pdf'), json.dumps(doc.get('metadata', {})),
                          len(doc.get('fileData', ''))))

        # Build request payload (without large file data - docs are in DB)
        payload = {
            'documentFilenames': data.get('documentFilenames', []),
            'documentDetails': data.get('documentDetails', []),
            'existingAnalysis': data.get('existingAnalysis'),
            'historicalWeights': data.get('historicalWeights', []),
            'weightingConfig': data.get('weightingConfig', {})
        }

        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO analysis_jobs (id, ticker, status, api_key, request_payload)
                VALUES (%s, %s, 'pending', %s, %s)
            ''', (job_id, ticker, api_key, json.dumps(payload)))

        # Spawn background thread
        thread = threading.Thread(target=_run_analysis_job, args=(job_id,), daemon=True)
        thread.start()

        return jsonify({'jobId': job_id})
    except Exception as e:
        print(f"Error creating analysis job: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/analysis-job/<job_id>/status', methods=['GET'])
def get_analysis_job_status(job_id):
    """Poll for analysis job status."""
    try:
        with get_db() as (conn, cur):
            cur.execute('SELECT id, ticker, status, progress, batch_num, total_batches, chars_received, result, error, created_at FROM analysis_jobs WHERE id = %s', (job_id,))
            job = cur.fetchone()

        if not job:
            return jsonify({'error': 'Job not found'}), 404

        resp = {
            'status': job['status'],
            'progress': job['progress'],
            'batchNum': job['batch_num'],
            'totalBatches': job['total_batches'],
            'charsReceived': job['chars_received'],
            'elapsedSeconds': int((datetime.utcnow() - job['created_at']).total_seconds()) if job['created_at'] else 0
        }

        if job['status'] == 'complete' and job['result']:
            resp['result'] = job['result'] if isinstance(job['result'], dict) else json.loads(job['result'])
            # Clear API key now that result is delivered
            with get_db(commit=True) as (conn, cur):
                cur.execute("UPDATE analysis_jobs SET api_key = '' WHERE id = %s AND api_key != ''", (job_id,))

        if job['status'] == 'failed':
            resp['error'] = job['error']

        return jsonify(resp)
    except Exception as e:
        print(f"Error getting job status: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/analysis-jobs-batch', methods=['POST'])
def create_analysis_jobs_batch():
    """Create multiple independent background analysis jobs in one request."""
    try:
        data = request.json
        api_key = data.get('apiKey', '')
        if not api_key:
            return jsonify({'error': 'API key required'}), 400
        jobs_input = data.get('jobs', [])
        if not jobs_input:
            return jsonify({'error': 'No jobs provided'}), 400

        job_ids = {}
        for job_data in jobs_input:
            ticker = job_data.get('ticker', '').upper()
            if not ticker:
                continue

            job_id = str(uuid.uuid4())[:12]

            # Save any new documents to document_files
            new_docs = job_data.get('newDocuments', [])
            if new_docs:
                with get_db(commit=True) as (conn, cur):
                    for doc in new_docs:
                        cur.execute('''
                            INSERT INTO document_files (ticker, filename, file_data, file_type, mime_type, metadata, file_size)
                            VALUES (%s, %s, %s, %s, %s, %s, %s)
                            ON CONFLICT (ticker, filename) DO UPDATE SET
                                file_data = EXCLUDED.file_data,
                                file_type = EXCLUDED.file_type,
                                mime_type = EXCLUDED.mime_type,
                                metadata = EXCLUDED.metadata,
                                file_size = EXCLUDED.file_size
                        ''', (ticker, doc['filename'], doc.get('fileData', ''), doc.get('fileType', 'pdf'),
                              doc.get('mimeType', 'application/pdf'), json.dumps(doc.get('metadata', {})),
                              len(doc.get('fileData', ''))))

            # Build request payload (without large file data)
            payload = {
                'documentFilenames': job_data.get('documentFilenames', []),
                'documentDetails': job_data.get('documentDetails', []),
                'existingAnalysis': job_data.get('existingAnalysis'),
                'historicalWeights': job_data.get('historicalWeights', []),
                'weightingConfig': job_data.get('weightingConfig', {})
            }

            with get_db(commit=True) as (conn, cur):
                cur.execute('''
                    INSERT INTO analysis_jobs (id, ticker, status, api_key, request_payload)
                    VALUES (%s, %s, 'pending', %s, %s)
                ''', (job_id, ticker, api_key, json.dumps(payload)))

            # Spawn background thread
            thread = threading.Thread(target=_run_analysis_job, args=(job_id,), daemon=True)
            thread.start()
            job_ids[ticker] = job_id
            print(f"[batch-job] Started {ticker} → {job_id}")

        return jsonify({'jobIds': job_ids})
    except Exception as e:
        print(f"Error creating batch analysis jobs: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/analysis-jobs-batch-status', methods=['POST'])
def get_analysis_jobs_batch_status():
    """Poll multiple analysis job statuses in one request."""
    try:
        data = request.json
        job_ids = data.get('jobIds', [])
        if not job_ids:
            return jsonify({'jobs': {}})

        placeholders = ','.join(['%s'] * len(job_ids))
        with get_db() as (conn, cur):
            cur.execute(f'SELECT id, ticker, status, progress, batch_num, total_batches, chars_received, result, error, created_at FROM analysis_jobs WHERE id IN ({placeholders})', job_ids)
            rows = cur.fetchall()

        jobs = {}
        clear_ids = []
        stale_ids = []
        for job in rows:
            elapsed = int((datetime.utcnow() - job['created_at']).total_seconds()) if job['created_at'] else 0
            resp = {
                'ticker': job['ticker'],
                'status': job['status'],
                'progress': job['progress'],
                'batchNum': job['batch_num'],
                'totalBatches': job['total_batches'],
                'charsReceived': job['chars_received'],
                'elapsedSeconds': elapsed
            }
            # Stale job detection: if processing for >8 min, mark as failed
            if job['status'] == 'processing' and elapsed > 480:
                resp['status'] = 'failed'
                resp['error'] = f'Job timed out after {elapsed}s — the API call may have hung. Please retry.'
                stale_ids.append(job['id'])
                print(f"[analysis-job {job['id']}] Stale job detected for {job['ticker']} ({elapsed}s), marking failed")
            if job['status'] == 'complete' and job['result']:
                resp['result'] = job['result'] if isinstance(job['result'], dict) else json.loads(job['result'])
                clear_ids.append(job['id'])
            if job['status'] == 'failed':
                resp['error'] = job['error']
            jobs[job['id']] = resp

        # Mark stale jobs as failed in DB
        if stale_ids:
            stale_ph = ','.join(['%s'] * len(stale_ids))
            with get_db(commit=True) as (conn, cur):
                cur.execute(f"UPDATE analysis_jobs SET status = 'failed', error = 'Timed out', api_key = '' WHERE id IN ({stale_ph})", stale_ids)

        # Clear API keys for completed jobs
        if clear_ids:
            clear_ph = ','.join(['%s'] * len(clear_ids))
            with get_db(commit=True) as (conn, cur):
                cur.execute(f"UPDATE analysis_jobs SET api_key = '' WHERE id IN ({clear_ph}) AND api_key != ''", clear_ids)

        return jsonify({'jobs': jobs})
    except Exception as e:
        print(f"Error getting batch job status: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# PORTFOLIO ANALYSES ENDPOINTS
# ============================================

@app.route('/api/analyses', methods=['GET'])
def get_analyses():
    """Get all saved portfolio analyses"""
    try:
        cached = cache.get('analyses')
        if cached is not None:
            return jsonify(cached)

        with get_db() as (conn, cur):
            cur.execute('''
                SELECT ticker, company, analysis, updated_at
                FROM portfolio_analyses
                ORDER BY ticker ASC
            ''')
            rows = cur.fetchall()

        result = []
        for row in rows:
            result.append({
                'ticker': row['ticker'],
                'company': row['company'],
                'analysis': row['analysis'],
                'updated': row['updated_at'].isoformat() if row['updated_at'] else None
            })

        cache.set('analyses', result, ttl=300)
        return jsonify(result)
    except Exception as e:
        print(f"Error getting analyses: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/analysis/<ticker>', methods=['GET'])
def get_analysis(ticker):
    """Get a specific portfolio analysis by ticker"""
    try:
        with get_db() as (conn, cur):
            cur.execute('''
                SELECT ticker, company, analysis, updated_at
                FROM portfolio_analyses
                WHERE ticker = %s
            ''', (ticker.upper(),))
            row = cur.fetchone()

        if not row:
            return jsonify({'error': 'Analysis not found'}), 404

        # Normalize: unwrap any accidentally nested 'analysis' keys
        analysis = row['analysis']
        if isinstance(analysis, dict):
            while 'analysis' in analysis and isinstance(analysis['analysis'], dict) and 'thesis' not in analysis:
                analysis = analysis['analysis']

        return jsonify({
            'ticker': row['ticker'],
            'company': row['company'],
            'analysis': analysis,
            'updated': row['updated_at'].isoformat() if row['updated_at'] else None
        })
    except Exception as e:
        print(f"Error getting analysis: {e}")
        return jsonify({'error': str(e)}), 500

def _create_thesis_snapshot(ticker, analysis, scorecard_data, snapshot_type='analysis'):
    """Create a thesis snapshot capturing current state for evolution tracking."""
    try:
        if isinstance(analysis, str):
            try: analysis = json.loads(analysis)
            except: analysis = {}
        thesis = analysis.get('thesis', {})
        summary = thesis.get('summary', '')
        pillar_count = len(thesis.get('pillars', []))
        conclusion = analysis.get('conclusion', '')
        sp_statuses = []
        if scorecard_data and scorecard_data.get('signposts'):
            for sp in scorecard_data['signposts']:
                sp_statuses.append({'metric': sp.get('metric', ''), 'status': (sp.get('status', '') or '').lower()})
        rk_statuses = []
        if scorecard_data and scorecard_data.get('risks'):
            for rk in scorecard_data['risks']:
                rk_statuses.append({'risk': rk.get('riskFactor', ''), 'status': (rk.get('status', '') or '').lower()})
        g, y, r = _tally_statuses(scorecard_data) if scorecard_data else (0, 0, 0)
        total = g + y + r
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_snapshots
                (ticker, snapshot_type, thesis_summary, pillar_count,
                 signpost_statuses, risk_statuses, green_count, yellow_count,
                 red_count, total_items, conviction, raw_snapshot)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            ''', (ticker.upper(), snapshot_type, summary, pillar_count,
                  json.dumps(sp_statuses), json.dumps(rk_statuses),
                  g, y, r, total, conclusion,
                  json.dumps({'analysis': analysis, 'scorecard_data': scorecard_data})))
    except Exception as e:
        print(f"Error creating thesis snapshot for {ticker}: {e}")


@app.route('/api/save-analysis', methods=['POST'])
def save_analysis():
    """Save or update a portfolio analysis"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()
        company = data.get('companyName', data.get('company', ''))
        analysis = data.get('analysis', {})

        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400

        # Backend-side history & documentHistory preservation
        # If the incoming analysis lacks history/documentHistory, carry over from existing
        with get_db() as (conn, cur):
            cur.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            existing_row = cur.fetchone()
        if existing_row and existing_row['analysis']:
            old = existing_row['analysis']
            if isinstance(old, str):
                try: old = json.loads(old)
                except: old = {}

            # Preserve history: if incoming has none, snapshot old thesis + carry over
            if not analysis.get('history'):
                old_history = old.get('history') or []
                if old.get('thesis') or old.get('signposts') or old.get('threats'):
                    old_history = list(old_history)  # copy
                    old_history.append({
                        'timestamp': old.get('updatedAt') or datetime.utcnow().isoformat(),
                        'thesis': old.get('thesis'),
                        'signposts': old.get('signposts'),
                        'threats': old.get('threats'),
                    })
                    old_history = old_history[-20:]
                analysis['history'] = old_history

            # Preserve documentHistory: if incoming has none, carry over existing
            if not analysis.get('documentHistory'):
                analysis['documentHistory'] = old.get('documentHistory') or []

            # Set updatedAt if missing
            if not analysis.get('updatedAt'):
                analysis['updatedAt'] = datetime.utcnow().isoformat()

        with get_db(commit=True) as (conn, cur):
            # Upsert - insert or update
            cur.execute('''
                INSERT INTO portfolio_analyses (ticker, company, analysis, updated_at)
                VALUES (%s, %s, %s, %s)
                ON CONFLICT (ticker)
                DO UPDATE SET
                    company = EXCLUDED.company,
                    analysis = EXCLUDED.analysis,
                    updated_at = EXCLUDED.updated_at
                RETURNING ticker
            ''', (ticker, company, json.dumps(analysis), datetime.utcnow()))

            result = cur.fetchone()

        cache.invalidate('analyses')
        cache.invalidate('portfolio_dashboard')
        # Auto-snapshot for thesis evolution tracking
        try:
            sc_data = None
            with get_db() as (_, cur2):
                cur2.execute('SELECT scorecard_data FROM thesis_scorecard_data WHERE ticker = %s', (ticker,))
                sc_row = cur2.fetchone()
            if sc_row:
                sc_data = sc_row['scorecard_data']
                if isinstance(sc_data, str):
                    try: sc_data = json.loads(sc_data)
                    except: sc_data = None
            _create_thesis_snapshot(ticker, analysis, sc_data, 'analysis')
        except Exception as snap_err:
            print(f"Snapshot error: {snap_err}")
        return jsonify({'success': True, 'ticker': result['ticker']})
    except Exception as e:
        print(f"Error saving analysis: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/delete-analysis', methods=['POST'])
def delete_analysis():
    """Delete a portfolio analysis"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()

        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400

        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM portfolio_analyses WHERE ticker = %s', (ticker,))

        cache.invalidate('analyses')
        cache.invalidate('portfolio_dashboard')
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting analysis: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# PORTFOLIO DASHBOARD & THESIS SNAPSHOTS
# ============================================

@app.route('/api/portfolio/dashboard', methods=['GET'])
def get_portfolio_dashboard():
    """Return all tickers with scorecard statuses in a single call."""
    try:
        cached = cache.get('portfolio_dashboard')
        if cached is not None:
            return jsonify(cached)
        with get_db() as (_, cur):
            cur.execute('''
                SELECT pa.ticker, pa.company, pa.analysis, pa.updated_at, tsd.scorecard_data
                FROM portfolio_analyses pa
                LEFT JOIN thesis_scorecard_data tsd ON tsd.ticker = pa.ticker
                ORDER BY pa.ticker ASC
            ''')
            rows = cur.fetchall()
        result = []
        total_green = total_yellow = total_red = 0
        for row in rows:
            analysis = row['analysis'] or {}
            if isinstance(analysis, str):
                try: analysis = json.loads(analysis)
                except: analysis = {}
            scorecard_data = row['scorecard_data']
            if isinstance(scorecard_data, str):
                try: scorecard_data = json.loads(scorecard_data)
                except: scorecard_data = None
            g, y, r = _tally_statuses(scorecard_data)
            total_green += g; total_yellow += y; total_red += r
            thesis = analysis.get('thesis', {})
            summary = thesis.get('summary', '')
            first_sentence = (summary.split('. ')[0] + '.') if summary else ''
            sp_statuses = []
            if scorecard_data and scorecard_data.get('signposts'):
                for sp in scorecard_data['signposts']:
                    sp_statuses.append({'metric': sp.get('metric', ''), 'status': (sp.get('status', '') or '').lower()})
            rk_statuses = []
            if scorecard_data and scorecard_data.get('risks'):
                for rk in scorecard_data['risks']:
                    rk_statuses.append({'risk': rk.get('riskFactor', ''), 'status': (rk.get('status', '') or '').lower()})
            if r > 0: overall = 'red'
            elif y > 0: overall = 'yellow'
            elif g > 0: overall = 'green'
            else: overall = 'none'
            result.append({
                'ticker': row['ticker'], 'company': row['company'],
                'thesis_summary': first_sentence,
                'signpost_count': len(analysis.get('signposts', [])),
                'risk_count': len(analysis.get('threats', [])),
                'green': g, 'yellow': y, 'red': r,
                'signpost_statuses': sp_statuses, 'risk_statuses': rk_statuses,
                'overall': overall,
                'updated': row['updated_at'].isoformat() if row['updated_at'] else None,
                'has_scorecard': scorecard_data is not None
            })
        total_items = total_green + total_yellow + total_red
        health_score = round((total_green / total_items * 100) if total_items > 0 else 0, 1)
        dashboard = {
            'stocks': result,
            'summary': {
                'total_stocks': len(result), 'total_green': total_green,
                'total_yellow': total_yellow, 'total_red': total_red,
                'health_score': health_score
            }
        }
        cache.set('portfolio_dashboard', dashboard, ttl=120)
        return jsonify(dashboard)
    except Exception as e:
        print(f"Error getting portfolio dashboard: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-snapshots/<ticker>', methods=['GET'])
def get_thesis_snapshots(ticker):
    """Return all snapshots for a ticker, ordered by date."""
    try:
        limit = request.args.get('limit', 50, type=int)
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, ticker, snapshot_type, thesis_summary, pillar_count,
                       signpost_statuses, risk_statuses, green_count, yellow_count,
                       red_count, total_items, conviction, created_at
                FROM thesis_snapshots WHERE ticker = %s
                ORDER BY created_at DESC LIMIT %s
            ''', (ticker.upper(), limit))
            rows = cur.fetchall()
        result = []
        for row in rows:
            result.append({
                'id': row['id'], 'ticker': row['ticker'],
                'snapshot_type': row['snapshot_type'],
                'thesis_summary': row['thesis_summary'],
                'pillar_count': row['pillar_count'],
                'signpost_statuses': row['signpost_statuses'],
                'risk_statuses': row['risk_statuses'],
                'green': row['green_count'], 'yellow': row['yellow_count'], 'red': row['red_count'],
                'total_items': row['total_items'], 'conviction': row['conviction'],
                'date': row['created_at'].isoformat() if row['created_at'] else None
            })
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/portfolio/dashboard/timeline', methods=['GET'])
def get_portfolio_timeline():
    """Return aggregate G/Y/R snapshots over time for the whole portfolio."""
    try:
        days = request.args.get('days', 90, type=int)
        with get_db() as (_, cur):
            cur.execute('''
                SELECT DATE(created_at) as snap_date,
                       SUM(green_count) as total_green,
                       SUM(yellow_count) as total_yellow,
                       SUM(red_count) as total_red,
                       COUNT(DISTINCT ticker) as stock_count
                FROM (
                    SELECT DISTINCT ON (ticker, DATE(created_at))
                           ticker, green_count, yellow_count, red_count, created_at
                    FROM thesis_snapshots
                    WHERE created_at > NOW() - INTERVAL '%s days'
                    ORDER BY ticker, DATE(created_at), created_at DESC
                ) daily_snapshots
                GROUP BY DATE(created_at)
                ORDER BY snap_date ASC
            ''', (days,))
            rows = cur.fetchall()
        result = []
        for row in rows:
            total = row['total_green'] + row['total_yellow'] + row['total_red']
            result.append({
                'date': row['snap_date'].isoformat(),
                'green': row['total_green'], 'yellow': row['total_yellow'], 'red': row['total_red'],
                'total': total,
                'green_pct': round(row['total_green'] / total * 100, 1) if total > 0 else 0,
                'stock_count': row['stock_count']
            })
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============================================
# STOCK OVERVIEWS ENDPOINTS
# ============================================

@app.route('/api/overviews', methods=['GET'])
def get_overviews():
    """Get all saved stock overviews"""
    try:
        cached = cache.get('overviews')
        if cached is not None:
            return jsonify(cached)

        with get_db() as (conn, cur):
            cur.execute('''
                SELECT ticker, company_name, company_overview, business_model, business_mix,
                       opportunities, risks, conclusion, raw_content, history, updated_at
                FROM stock_overviews
                ORDER BY ticker ASC
            ''')
            rows = cur.fetchall()

        result = []
        for row in rows:
            result.append({
                'ticker': row['ticker'],
                'companyName': row['company_name'],
                'companyOverview': row['company_overview'],
                'businessModel': row['business_model'],
                'businessMix': row.get('business_mix', ''),
                'opportunities': row['opportunities'],
                'risks': row['risks'],
                'conclusion': row['conclusion'],
                'rawContent': row['raw_content'],
                'history': row['history'] or [],
                'updatedAt': row['updated_at'].isoformat() if row['updated_at'] else None
            })

        cache.set('overviews', result, ttl=300)
        return jsonify(result)
    except Exception as e:
        print(f"Error getting overviews: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/save-overview', methods=['POST'])
def save_overview():
    """Save or update a stock overview"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()
        
        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400

        with get_db(commit=True) as (conn, cur):
            # Upsert
            cur.execute('''
                INSERT INTO stock_overviews (
                    ticker, company_name, company_overview, business_model, business_mix,
                    opportunities, risks, conclusion, raw_content, history, updated_at
                )
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (ticker)
                DO UPDATE SET
                    company_name = EXCLUDED.company_name,
                    company_overview = EXCLUDED.company_overview,
                    business_model = EXCLUDED.business_model,
                    business_mix = EXCLUDED.business_mix,
                    opportunities = EXCLUDED.opportunities,
                    risks = EXCLUDED.risks,
                    conclusion = EXCLUDED.conclusion,
                    raw_content = EXCLUDED.raw_content,
                    history = EXCLUDED.history,
                    updated_at = EXCLUDED.updated_at
                RETURNING ticker
            ''', (
                ticker,
                data.get('companyName', ''),
                data.get('companyOverview', ''),
                data.get('businessModel', ''),
                data.get('businessMix', ''),
                data.get('opportunities', ''),
                data.get('risks', ''),
                data.get('conclusion', ''),
                data.get('rawContent', ''),
                json.dumps(data.get('history', [])),
                datetime.utcnow()
            ))

            result = cur.fetchone()

        cache.invalidate('overviews')
        return jsonify({'success': True, 'ticker': result['ticker']})
    except Exception as e:
        print(f"Error saving overview: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/delete-overview', methods=['POST'])
def delete_overview():
    """Delete a stock overview"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()

        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400

        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM stock_overviews WHERE ticker = %s', (ticker,))

        cache.invalidate('overviews')
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting overview: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# CHAT HISTORY ENDPOINTS
# ============================================

@app.route('/api/chats', methods=['GET'])
def get_chats():
    """Get all chat histories"""
    try:
        with get_db() as (conn, cur):
            cur.execute('''
                SELECT id, title, messages, updated_at
                FROM chat_histories
                ORDER BY updated_at DESC
            ''')
            rows = cur.fetchall()

        result = []
        for row in rows:
            result.append({
                'id': row['id'],
                'title': row['title'],
                'messages': row['messages'] or [],
                'updatedAt': row['updated_at'].isoformat() if row['updated_at'] else None
            })

        return jsonify(result)
    except Exception as e:
        print(f"Error getting chats: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/save-chat', methods=['POST'])
def save_chat():
    """Save or update a chat history"""
    try:
        data = request.json
        chat_id = data.get('id', '')
        
        if not chat_id:
            return jsonify({'error': 'Chat ID is required'}), 400

        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO chat_histories (id, title, messages, updated_at)
                VALUES (%s, %s, %s, %s)
                ON CONFLICT (id)
                DO UPDATE SET
                    title = EXCLUDED.title,
                    messages = EXCLUDED.messages,
                    updated_at = EXCLUDED.updated_at
                RETURNING id
            ''', (
                chat_id,
                data.get('title', 'New Chat'),
                json.dumps(data.get('messages', [])),
                datetime.utcnow()
            ))

            result = cur.fetchone()

        return jsonify({'success': True, 'id': result['id']})
    except Exception as e:
        print(f"Error saving chat: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/delete-chat', methods=['POST'])
def delete_chat():
    """Delete a chat history"""
    try:
        data = request.json
        chat_id = data.get('id', '')
        
        if not chat_id:
            return jsonify({'error': 'Chat ID is required'}), 400

        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM chat_histories WHERE id = %s', (chat_id,))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting chat: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# MEETING SUMMARY ENDPOINTS
# ============================================

@app.route('/api/summaries', methods=['GET'])
def get_summaries():
    """Get all meeting summaries"""
    try:
        cached = cache.get('summaries')
        if cached is not None:
            return jsonify(cached)

        with get_db() as (conn, cur):
            cur.execute('''
                SELECT id, title, raw_notes, summary, questions, assessment, topic, topic_type, source_type, source_files, doc_type, has_stored_files, categories, created_at
                FROM meeting_summaries
                ORDER BY created_at DESC
            ''')
            rows = cur.fetchall()

        result = []
        for row in rows:
            result.append({
                'id': row['id'],
                'title': row['title'],
                'rawNotes': row['raw_notes'],
                'summary': row['summary'],
                'questions': row['questions'],
                'assessment': row.get('assessment') or '',
                'topic': row.get('topic') or 'General',
                'topicType': row.get('topic_type') or 'other',
                'sourceType': row.get('source_type') or 'paste',
                'sourceFiles': row.get('source_files') or [],
                'docType': row.get('doc_type') or 'other',
                'hasStoredFiles': row.get('has_stored_files') or False,
                'categories': row.get('categories') or [],
                'createdAt': row['created_at'].isoformat() if row['created_at'] else None
            })

        cache.set('summaries', result, ttl=300)
        return jsonify(result)
    except Exception as e:
        print(f"Error getting summaries: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/save-summary', methods=['POST'])
def save_summary():
    """Save or update a meeting summary"""
    try:
        data = request.json
        summary_id = data.get('id', '')
        
        if not summary_id:
            return jsonify({'error': 'Summary ID is required'}), 400
        
        # Convert sourceFiles list to JSON
        source_files = data.get('sourceFiles', [])
        if isinstance(source_files, list):
            source_files = json.dumps(source_files)

        # Convert categories list to JSON
        categories = data.get('categories', [])
        if isinstance(categories, list):
            categories = json.dumps(categories)

        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO meeting_summaries (id, title, raw_notes, summary, questions, assessment, topic, topic_type, source_type, source_files, doc_type, categories, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (id)
                DO UPDATE SET
                    title = EXCLUDED.title,
                    raw_notes = EXCLUDED.raw_notes,
                    summary = EXCLUDED.summary,
                    questions = EXCLUDED.questions,
                    assessment = EXCLUDED.assessment,
                    topic = EXCLUDED.topic,
                    topic_type = EXCLUDED.topic_type,
                    source_type = EXCLUDED.source_type,
                    source_files = EXCLUDED.source_files,
                    doc_type = EXCLUDED.doc_type,
                    categories = EXCLUDED.categories
                RETURNING id
            ''', (
                summary_id,
                data.get('title', 'Meeting Summary'),
                data.get('rawNotes', ''),
                data.get('summary', ''),
                data.get('questions', ''),
                data.get('assessment', ''),
                data.get('topic', 'General'),
                data.get('topicType', 'other'),
                data.get('sourceType', 'paste'),
                source_files,
                data.get('docType', 'other'),
                categories,
                data.get('createdAt', datetime.utcnow().isoformat())
            ))

            result = cur.fetchone()

        cache.invalidate('summaries')
        return jsonify({'success': True, 'id': result['id']})
    except Exception as e:
        print(f"Error saving summary: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/delete-summary', methods=['POST'])
def delete_summary():
    """Delete a meeting summary"""
    try:
        data = request.json
        summary_id = data.get('id', '')

        if not summary_id:
            return jsonify({'error': 'Summary ID is required'}), 400

        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM summary_files WHERE summary_id = %s', (summary_id,))
            cur.execute('DELETE FROM meeting_summaries WHERE id = %s', (summary_id,))

        cache.invalidate('summaries')
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting summary: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# SUMMARY CATEGORIES ENDPOINTS
# ============================================

@app.route('/api/summary-categories', methods=['POST'])
def update_summary_categories():
    """Quick-update categories for a single summary"""
    try:
        data = request.json
        summary_id = data.get('summaryId', '')
        categories = data.get('categories', [])
        if not summary_id:
            return jsonify({'error': 'summaryId is required'}), 400
        with get_db(commit=True) as (conn, cur):
            cur.execute('UPDATE meeting_summaries SET categories = %s WHERE id = %s', (json.dumps(categories), summary_id))
        cache.invalidate('summaries')
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error updating categories: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/summary-categories/all', methods=['GET'])
def get_all_summary_categories():
    """Return deduplicated list of all category names across all summaries"""
    try:
        with get_db() as (_, cur):
            cur.execute("SELECT DISTINCT jsonb_array_elements_text(COALESCE(categories, '[]'::jsonb)) AS cat FROM meeting_summaries ORDER BY cat")
            rows = cur.fetchall()
        return jsonify([r['cat'] for r in rows])
    except Exception as e:
        print(f"Error getting categories: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# SUMMARY FILES ENDPOINTS
# ============================================

@app.route('/api/summary-files/<summary_id>', methods=['GET'])
def get_summary_files(summary_id):
    """Get stored files for a summary"""
    try:
        with get_db() as (conn, cur):
            cur.execute('''
                SELECT id, filename, file_type, file_data, file_size, created_at
                FROM summary_files
                WHERE summary_id = %s
                ORDER BY created_at ASC
            ''', (summary_id,))

            files = []
            for row in cur.fetchall():
                files.append({
                    'id': row['id'],
                    'filename': row['filename'],
                    'fileType': row['file_type'],
                    'fileData': row['file_data'],
                    'fileSize': row['file_size'],
                    'createdAt': row['created_at'].isoformat() if row['created_at'] else None
                })

        return jsonify(files)
    except Exception as e:
        print(f"Error getting summary files: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/summary-files/<summary_id>', methods=['POST'])
def save_summary_files(summary_id):
    """Save files for a summary"""
    try:
        data = request.json
        files = data.get('files', [])
        
        if not files:
            return jsonify({'error': 'No files provided'}), 400

        with get_db(commit=True) as (conn, cur):
            saved_count = 0
            for file_data in files:
                cur.execute('''
                    INSERT INTO summary_files (summary_id, filename, file_type, file_data, file_size)
                    VALUES (%s, %s, %s, %s, %s)
                ''', (
                    summary_id,
                    file_data.get('filename', 'document.pdf'),
                    file_data.get('fileType', 'application/pdf'),
                    file_data.get('fileData', ''),
                    file_data.get('fileSize', 0)
                ))
                saved_count += 1

            # Update has_stored_files flag on the summary
            cur.execute('''
                UPDATE meeting_summaries
                SET has_stored_files = TRUE
                WHERE id = %s
            ''', (summary_id,))

        return jsonify({'success': True, 'savedCount': saved_count})
    except Exception as e:
        print(f"Error saving summary files: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/summary-files/<summary_id>', methods=['DELETE'])
def delete_summary_files(summary_id):
    """Delete all stored files for a summary"""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM summary_files WHERE summary_id = %s', (summary_id,))

            # Update has_stored_files flag
            cur.execute('''
                UPDATE meeting_summaries
                SET has_stored_files = FALSE
                WHERE id = %s
            ''', (summary_id,))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting summary files: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-summary-section', methods=['POST'])
def email_summary_section():
    """Email a specific section of a summary (takeaways or questions)"""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    import re
    
    try:
        data = request.json
        email = data.get('email', '')
        subject = data.get('subject', 'Summary Notes')
        section = data.get('section', '')  # 'takeaways' or 'questions'
        content = data.get('content', '')
        title = data.get('title', 'Meeting Summary')
        topic = data.get('topic', 'General')
        smtp_config = data.get('smtpConfig', {})
        
        if not email or not content:
            return jsonify({'error': 'Email and content are required'}), 400
        
        # Get SMTP configuration
        use_gmail = smtp_config.get('use_gmail', True)
        gmail_user = smtp_config.get('gmail_user', '')
        gmail_password = smtp_config.get('gmail_app_password', '')
        from_email = smtp_config.get('from_email', gmail_user)
        
        if use_gmail and (not gmail_user or not gmail_password):
            return jsonify({'error': 'Gmail credentials required. Please set them in Settings.'}), 400
        
        # Convert HTML to plain text
        plain_text = re.sub(r'<[^>]+>', '', content)
        plain_text = plain_text.replace('&nbsp;', ' ').replace('&amp;', '&')
        
        # Format the section label
        section_label = "Key Takeaways" if section == 'takeaways' else "Follow-up Questions"
        header_color = "#0d9488" if section == 'takeaways' else "#d97706"
        
        # Build HTML email
        html_content = f"""
        <html>
        <head>
            <style>
                body {{
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                    line-height: 1.6;
                    color: #333;
                    max-width: 800px;
                    margin: 0 auto;
                    padding: 20px;
                }}
                h1 {{
                    color: {header_color};
                    border-bottom: 2px solid {header_color};
                    padding-bottom: 10px;
                }}
                h2 {{
                    color: #374151;
                    margin-top: 24px;
                }}
                h3 {{
                    color: #4b5563;
                }}
                ul, ol {{
                    padding-left: 24px;
                }}
                li {{
                    margin-bottom: 8px;
                }}
                strong {{
                    color: #111;
                }}
                .header {{
                    background: linear-gradient(135deg, {header_color} 0%, {'#0891b2' if section == 'takeaways' else '#ea580c'} 100%);
                    color: white;
                    padding: 20px;
                    border-radius: 8px;
                    margin-bottom: 24px;
                }}
                .topic-badge {{
                    display: inline-block;
                    background: rgba(255,255,255,0.2);
                    padding: 4px 12px;
                    border-radius: 12px;
                    font-size: 12px;
                    margin-top: 8px;
                }}
                .content {{
                    background: #f9fafb;
                    padding: 24px;
                    border-radius: 8px;
                    border: 1px solid #e5e7eb;
                }}
                .footer {{
                    margin-top: 24px;
                    padding-top: 16px;
                    border-top: 1px solid #e5e7eb;
                    font-size: 12px;
                    color: #6b7280;
                }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1 style="color: white; border: none; margin: 0;">{section_label}</h1>
                <p style="margin: 8px 0 0 0; opacity: 0.9;">{title}</p>
                <span class="topic-badge">{topic}</span>
            </div>
            <div class="content">
                {content}
            </div>
            <div class="footer">
                Generated by TDL Equity Analyzer
            </div>
        </body>
        </html>
        """
        
        # Create email message
        msg = MIMEMultipart('alternative')
        msg['From'] = from_email
        msg['To'] = email
        msg['Subject'] = subject
        
        # Attach both plain text and HTML versions
        msg.attach(MIMEText(plain_text, 'plain'))
        msg.attach(MIMEText(html_content, 'html'))
        
        # Send via Gmail SMTP
        if use_gmail:
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(gmail_user, gmail_password)
                server.send_message(msg)
        
        return jsonify({'success': True, 'message': 'Email sent successfully'})
        
    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'Gmail authentication failed. Check your email and app password.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        print(f"Error sending summary email: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-summary', methods=['POST'])
def generate_summary():
    """Generate a summary with explicit detail level control.
    Used for concise/detailed modes that need strong system prompts.
    Returns { summary, questions } matching the n8n webhook format."""
    try:
        data = request.json
        notes = data.get('notes', '')
        detail_level = data.get('detailLevel', 'standard')

        if not notes.strip():
            return jsonify({'error': 'No notes provided'}), 400

        keys = _get_api_keys(
            anthropic_api_key=data.get('apiKey', ''),
            gemini_api_key=data.get('geminiApiKey', '')
        )
        if not keys.get('gemini') and not keys.get('anthropic'):
            return jsonify({'error': 'No API key provided. Check your API keys in Settings.'}), 400

        # HTML format instructions matching the n8n webhook output that Charlie's CSS expects.
        # Charlie's .summary-content CSS auto-adds bullet styling to <p><strong>...</strong>...</p>
        html_format = """OUTPUT FORMAT — You MUST return raw HTML. No markdown. No code fences. No ```html.
Use EXACTLY this HTML structure:
- Section headers: <h2>Section Title</h2>
- Each key point as its own paragraph: <p><strong>Topic:</strong> Description text here.</p>
- For sub-points under a topic use: <ul><li>Sub-point text</li></ul>
- Bold key terms with <strong> tags
- Separate paragraphs with <p> tags. Do NOT use <br> tags.
- Do NOT wrap in code blocks. Start directly with HTML tags."""

        if detail_level == 'concise':
            summary_instruction = f"""Generate a SHORT, CONCISE summary of the following notes.

RULES:
- Maximum 300-500 words total
- Only the most critical 3-5 takeaways, decisions, and action items
- Omit minor details, examples, and supporting context
- Be ruthlessly brief

{html_format}"""
            questions_instruction = """Based on the notes, generate 2-3 key follow-up questions. Be brief.
Return raw HTML only. No markdown. No code fences. Use: <ol><li>Question?</li></ol>"""
        elif detail_level == 'detailed':
            summary_instruction = f"""Generate an EXTREMELY DETAILED and COMPREHENSIVE summary of the following notes.

RULES:
- Preserve virtually ALL content from the original
- Include every specific number, statistic, percentage, data point
- Include all names, quotes, attributions, and examples
- Include all anecdotes, supporting details, and context
- Your summary should be 3-5x longer than a normal summary
- Someone reading only your summary should miss almost nothing

{html_format}"""
            questions_instruction = """Based on the notes, generate 5-8 detailed follow-up questions. Be thorough.
Return raw HTML only. No markdown. No code fences. Use: <ol><li>Question?</li></ol>"""
        else:
            return jsonify({'error': 'Use the standard webhook for standard mode'}), 400

        # Assessment instruction — Claude's candid, opinionated analysis
        assessment_instruction = f"""You are a sharp, experienced advisor giving your CANDID, UNFILTERED assessment of this call, meeting, or conversation.

This is NOT a summary — the summary is generated separately. Your job is to provide your HONEST OPINION on how things went. Be direct, opinionated, and don't sugarcoat.

Cover whichever of these are relevant:
- **Overall assessment:** How did the call/meeting/conversation go? Was it productive, a waste of time, or somewhere in between?
- **Quality of answers:** Were the responses substantive and credible, or vague and evasive? Call out specific weak or strong answers.
- **Red flags / BS detection:** Did anyone dodge questions, give rehearsed non-answers, contradict themselves, or seem disingenuous? Be specific about what raised your suspicion and why.
- **Meeting flow and dynamics:** Was it well-structured? Did it go off-track? Was there tension or alignment? Who drove the conversation?
- **What was most effective:** What landed well? What was the strongest point made?
- **What could have been better:** What questions should have been asked but weren't? What was left on the table?
- **Credibility assessment:** Do you believe what was said? Rate the overall credibility of the key claims.
- **Bottom line:** One sentence on your overall take.

Be conversational and direct — write as if you're giving your honest debrief to a colleague after walking out of the meeting. Don't hedge. If something was weak, say it was weak. If someone was impressive, say so.

{html_format}"""

        def call_llm(system, user_text):
            """Call Gemini or Claude and return raw text response."""
            if keys.get('gemini'):
                try:
                    client = genai.Client(api_key=keys['gemini'])
                    response = client.models.generate_content(
                        model='gemini-2.5-flash',
                        contents=user_text,
                        config=genai_types.GenerateContentConfig(
                            system_instruction=system,
                            max_output_tokens=16384,
                            temperature=0.3,
                        )
                    )
                    return response.text.strip()
                except Exception as e:
                    print(f"Gemini summary failed: {e}")

            if keys.get('anthropic'):
                try:
                    client = anthropic.Anthropic(api_key=keys['anthropic'], timeout=120)
                    response = client.messages.create(
                        model='claude-haiku-4-5-20251001',
                        max_tokens=8192,
                        system=system,
                        messages=[{'role': 'user', 'content': user_text}],
                        temperature=0.3,
                    )
                    return response.content[0].text.strip()
                except Exception as e:
                    print(f"Anthropic summary failed: {e}")

            return None

        # Generate summary, questions, and assessment
        summary_html = call_llm(summary_instruction, notes)
        if not summary_html:
            return jsonify({'error': 'All LLM providers failed. Check your API keys.'}), 500

        questions_html = call_llm(questions_instruction, notes) or ''
        assessment_html = call_llm(assessment_instruction, notes) or ''

        # Strip code fences if LLM wraps output anyway
        def strip_fences(text):
            t = text.strip()
            if t.startswith('```'):
                t = t.split('\n', 1)[1] if '\n' in t else t[3:]
            if t.endswith('```'):
                t = t[:-3]
            return t.strip()

        return jsonify({
            'summary': strip_fences(summary_html),
            'questions': strip_fences(questions_html),
            'assessment': strip_fences(assessment_html)
        })

    except Exception as e:
        print(f"Generate summary error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-assessment', methods=['POST'])
def generate_assessment():
    """Generate only the assessment for notes (used when summary comes from N8N webhook)."""
    try:
        data = request.json
        notes = data.get('notes', '')

        if not notes.strip():
            return jsonify({'error': 'No notes provided'}), 400

        keys = _get_api_keys(
            anthropic_api_key=data.get('apiKey', ''),
            gemini_api_key=data.get('geminiApiKey', '')
        )
        if not keys.get('gemini') and not keys.get('anthropic'):
            return jsonify({'error': 'No API key provided. Check your API keys in Settings.'}), 400

        html_format = """OUTPUT FORMAT — You MUST return raw HTML. No markdown. No code fences. No ```html.
Use EXACTLY this HTML structure:
- Section headers: <h2>Section Title</h2>
- Each key point as its own paragraph: <p><strong>Topic:</strong> Description text here.</p>
- For sub-points under a topic use: <ul><li>Sub-point text</li></ul>
- Bold key terms with <strong> tags
- Separate paragraphs with <p> tags. Do NOT use <br> tags.
- Do NOT wrap in code blocks. Start directly with HTML tags."""

        assessment_instruction = f"""You are a sharp, experienced advisor giving your CANDID, UNFILTERED assessment of this call, meeting, or conversation.

This is NOT a summary — the summary is generated separately. Your job is to provide your HONEST OPINION on how things went. Be direct, opinionated, and don't sugarcoat.

Cover whichever of these are relevant:
- **Overall assessment:** How did the call/meeting/conversation go? Was it productive, a waste of time, or somewhere in between?
- **Quality of answers:** Were the responses substantive and credible, or vague and evasive? Call out specific weak or strong answers.
- **Red flags / BS detection:** Did anyone dodge questions, give rehearsed non-answers, contradict themselves, or seem disingenuous? Be specific about what raised your suspicion and why.
- **Meeting flow and dynamics:** Was it well-structured? Did it go off-track? Was there tension or alignment? Who drove the conversation?
- **What was most effective:** What landed well? What was the strongest point made?
- **What could have been better:** What questions should have been asked but weren't? What was left on the table?
- **Credibility assessment:** Do you believe what was said? Rate the overall credibility of the key claims.
- **Bottom line:** One sentence on your overall take.

Be conversational and direct — write as if you're giving your honest debrief to a colleague after walking out of the meeting. Don't hedge. If something was weak, say it was weak. If someone was impressive, say so.

{html_format}"""

        def call_llm(system, user_text):
            if keys.get('gemini'):
                try:
                    client = genai.Client(api_key=keys['gemini'])
                    response = client.models.generate_content(
                        model='gemini-2.5-flash',
                        contents=user_text,
                        config=genai_types.GenerateContentConfig(
                            system_instruction=system,
                            max_output_tokens=16384,
                            temperature=0.3,
                        )
                    )
                    return response.text.strip()
                except Exception as e:
                    print(f"Gemini assessment failed: {e}")
            if keys.get('anthropic'):
                try:
                    client = anthropic.Anthropic(api_key=keys['anthropic'], timeout=120)
                    response = client.messages.create(
                        model='claude-haiku-4-5-20251001',
                        max_tokens=8192,
                        system=system,
                        messages=[{'role': 'user', 'content': user_text}],
                        temperature=0.3,
                    )
                    return response.content[0].text.strip()
                except Exception as e:
                    print(f"Anthropic assessment failed: {e}")
            return None

        assessment_html = call_llm(assessment_instruction, notes) or ''

        def strip_fences(text):
            t = text.strip()
            if t.startswith('```'):
                t = t.split('\n', 1)[1] if '\n' in t else t[3:]
            if t.endswith('```'):
                t = t[:-3]
            return t.strip()

        return jsonify({'assessment': strip_fences(assessment_html)})

    except Exception as e:
        print(f"Generate assessment error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/extract-summary-text', methods=['POST'])
def extract_summary_text():
    """Extract text from uploaded files (PDF, DOCX, images, TXT) for summary generation.
    Accepts either multipart/form-data or JSON with base64-encoded files."""
    try:
        # Support both JSON (base64 files) and multipart/form-data uploads
        file_items = []  # list of (filename, file_content_bytes)

        if request.is_json or (request.content_type and 'application/json' in request.content_type):
            # JSON mode: files sent as base64
            data = request.get_json()
            if not data or 'files' not in data or not data['files']:
                return jsonify({'error': 'No files provided'}), 400
            api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '')
            for f in data['files']:
                file_content = base64.b64decode(f['fileData'])
                file_items.append((f['filename'], file_content))
            print(f"extract-summary-text (JSON): {len(file_items)} files, API key present: {bool(api_key)}")
        elif 'files' in request.files:
            # Multipart mode: standard file upload
            files = request.files.getlist('files')
            if not files or files[0].filename == '':
                return jsonify({'error': 'No files selected'}), 400
            api_key = os.environ.get('ANTHROPIC_API_KEY', '') or request.form.get('apiKey', '')
            for file in files:
                file_items.append((file.filename, file.read()))
            print(f"extract-summary-text (multipart): {len(file_items)} files, API key present: {bool(api_key)}")
        else:
            return jsonify({'error': 'No files provided'}), 400

        all_text = []
        first_filename = file_items[0][0]

        for orig_filename, file_content in file_items:
            filename = orig_filename.lower()
            extracted_text = ''
            
            try:
                # Handle PDF files
                if filename.endswith('.pdf'):
                    try:
                        import io
                        from PyPDF2 import PdfReader
                        pdf_reader = PdfReader(io.BytesIO(file_content))
                        for page in pdf_reader.pages:
                            text = page.extract_text()
                            if text:
                                extracted_text += text + '\n\n'
                    except ImportError:
                        # Fallback: try pdfplumber
                        try:
                            import pdfplumber
                            import io
                            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                                for page in pdf.pages:
                                    text = page.extract_text()
                                    if text:
                                        extracted_text += text + '\n\n'
                        except ImportError:
                            return jsonify({'error': 'PDF processing libraries not available'}), 500
                
                # Handle Word documents
                elif filename.endswith('.docx') or filename.endswith('.doc'):
                    try:
                        import io
                        from docx import Document
                        doc = Document(io.BytesIO(file_content))
                        for para in doc.paragraphs:
                            if para.text.strip():
                                extracted_text += para.text + '\n'
                        # Also extract from tables
                        for table in doc.tables:
                            for row in table.rows:
                                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                                if row_text:
                                    extracted_text += row_text + '\n'
                    except ImportError:
                        return jsonify({'error': 'Word document processing library not available'}), 500
                
                # Handle plain text files
                elif filename.endswith('.txt'):
                    try:
                        extracted_text = file_content.decode('utf-8')
                    except UnicodeDecodeError:
                        extracted_text = file_content.decode('latin-1')
                
                # Handle images (use Claude Vision API for OCR)
                elif filename.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp')):
                    # Use API key from form data or environment
                    if api_key:
                        try:
                            ext = filename.split('.')[-1].lower()
                            media_types = {'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'gif': 'image/gif', 'bmp': 'image/bmp', 'webp': 'image/webp'}
                            media_type = media_types.get(ext, 'image/jpeg')
                            image_base64 = base64.b64encode(file_content).decode('utf-8')

                            ocr_result = call_llm(
                                messages=[{
                                    'role': 'user',
                                    'content': [
                                        {
                                            'type': 'image',
                                            'source': {
                                                'type': 'base64',
                                                'media_type': media_type,
                                                'data': image_base64
                                            }
                                        },
                                        {
                                            'type': 'text',
                                            'text': 'Please extract ALL text from this image exactly as it appears. Preserve the original formatting, paragraphs, and structure. Output ONLY the extracted text, nothing else. If this is a screenshot of an article, extract the full article text.'
                                        }
                                    ]
                                }],
                                tier="fast",
                                max_tokens=8000,
                                timeout=60,
                                anthropic_api_key=api_key,
                            )
                            extracted_text = ocr_result["text"]
                            print(f"OCR extracted {len(extracted_text)} chars from {orig_filename} via {ocr_result['provider']}/{ocr_result['model']}")

                        except LLMError as llm_err:
                            print(f"OCR failed across all providers: {llm_err}")
                            extracted_text = f"[Image file: {orig_filename} - OCR failed: {str(llm_err)[:150]}]"
                        except Exception as ocr_error:
                            print(f"OCR error: {ocr_error}")
                            extracted_text = f"[Image file: {orig_filename} - OCR error: {str(ocr_error)[:150]}]"
                    else:
                        # Fallback to pytesseract if no API key
                        try:
                            import pytesseract
                            from PIL import Image
                            import io
                            image = Image.open(io.BytesIO(file_content))
                            extracted_text = pytesseract.image_to_string(image)
                        except ImportError:
                            extracted_text = f"[Image file: {orig_filename} - OCR not available. Please set your API key in Settings.]"
                        except Exception as ocr_error:
                            extracted_text = f"[Image file: {orig_filename} - Could not extract text: {str(ocr_error)}]"

                else:
                    # Try to read as text
                    try:
                        extracted_text = file_content.decode('utf-8')
                    except:
                        extracted_text = f"[Unsupported file type: {orig_filename}]"

                if extracted_text.strip():
                    # Add filename header if multiple files
                    if len(file_items) > 1:
                        all_text.append(f"=== {orig_filename} ===\n{extracted_text}")
                    else:
                        all_text.append(extracted_text)

            except Exception as file_error:
                print(f"Error processing file {orig_filename}: {file_error}")
                all_text.append(f"[Error processing {orig_filename}: {str(file_error)}]")
        
        combined_text = '\n\n'.join(all_text)
        
        if not combined_text.strip():
            return jsonify({'error': 'Could not extract any text from the uploaded files'}), 400
        
        # Check if OCR failed for images (all we got was placeholder text)
        if '[Image file:' in combined_text:
            # Count how many images failed
            failed_count = combined_text.count('[Image file:')
            if failed_count == len(file_items):
                return jsonify({
                    'error': f'Could not extract text from {failed_count} image(s). Please ensure your API key is set in Settings, or use PDFs/text files instead.'
                }), 400

        return jsonify({
            'success': True,
            'text': combined_text,
            'filename': first_filename,
            'fileCount': len(file_items)
        })
        
    except Exception as e:
        print(f"Error extracting text: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# AUDIO TRANSCRIPTION ENDPOINT
# ============================================

# In-memory store for async transcription jobs
_transcription_jobs = {}

# In-memory store for async infographic generation jobs
_infographic_jobs = {}


TRANSCRIPTION_PROMPT = """You are a professional transcriptionist. Produce an ABSOLUTE VERBATIM, word-for-word transcription of this ENTIRE audio recording from beginning to end. Do NOT omit, skip, summarize, or condense ANY portion.

CRITICAL RULES — YOU MUST FOLLOW ALL OF THESE:
1. Transcribe EVERY SINGLE WORD spoken in the recording, from the very first word to the very last word
2. Do NOT skip sections, do NOT summarize, do NOT paraphrase — write exactly what was said
3. Include ALL filler words: "um", "uh", "like", "you know", "I mean", "so", "right", etc.
4. Include false starts, repeated words, self-corrections, and stammering exactly as spoken
5. Identify different speakers (e.g., "Speaker 1:", "Speaker 2:", or use names if mentioned/identifiable)
6. Start each speaker's turn on a new line with their label
7. Use proper punctuation and paragraph breaks for readability
8. Do NOT add any commentary, headers, timestamps, introductions, or notes — output ONLY the transcription
9. If the audio is long, you MUST transcribe the entire thing — do NOT stop early or truncate
10. Your transcription should capture 100% of the spoken content"""

CHUNK_MINUTES = 15  # Transcribe in 15-minute segments for accuracy
OVERLAP_SECONDS = 30  # 30-second overlap between chunks to avoid missed content


def _transcribe_audio_content(client, audio_content, job_id, label=""):
    """Transcribe a single audio content (chunk or full file). Returns text or raises."""
    import time
    models_to_try = ['gemini-2.5-flash', 'gemini-2.0-flash']
    last_error = None

    for model_name in models_to_try:
        for attempt in range(2):
            try:
                print(f"[Job {job_id}] {label}Trying {model_name} (attempt {attempt + 1}/2)...")
                response = client.models.generate_content(
                    model=model_name,
                    contents=[audio_content, TRANSCRIPTION_PROMPT],
                    config=genai_types.GenerateContentConfig(max_output_tokens=65536)
                )
                try:
                    text = response.text
                except (ValueError, AttributeError) as text_err:
                    extracted = None
                    if hasattr(response, 'candidates') and response.candidates:
                        for c in response.candidates:
                            if hasattr(c, 'content') and c.content and hasattr(c.content, 'parts'):
                                for p in c.content.parts:
                                    if hasattr(p, 'text') and p.text:
                                        extracted = p.text
                                        break
                            if extracted:
                                break
                    if extracted:
                        text = extracted
                    else:
                        last_error = text_err
                        break
                if text and text.strip():
                    print(f"[Job {job_id}] {label}Success with {model_name}: {len(text)} chars")
                    return text
            except Exception as retry_err:
                last_error = retry_err
                err_str = str(retry_err)
                print(f"[Job {job_id}] {label}{model_name} attempt {attempt + 1} error: {err_str}")
                if '429' in err_str or 'RESOURCE_EXHAUSTED' in err_str:
                    if attempt == 0:
                        time.sleep(10)
                    else:
                        break
                else:
                    break
    raise Exception(f"All models failed. Last error: {last_error}")


def _get_audio_duration_ffprobe(filepath):
    """Get audio duration in seconds using ffprobe (no memory loading)."""
    import subprocess
    try:
        result = subprocess.run(
            ['ffprobe', '-v', 'quiet', '-show_entries', 'format=duration', '-of', 'csv=p=0', filepath],
            capture_output=True, text=True, timeout=30
        )
        return float(result.stdout.strip())
    except Exception:
        return None


def _split_audio_ffmpeg(filepath, chunk_sec, overlap_sec, output_dir):
    """Split audio into chunks using ffmpeg subprocess (disk-based, no memory). Returns list of (chunk_path, start_sec, end_sec)."""
    import subprocess

    duration = _get_audio_duration_ffprobe(filepath)
    if not duration:
        return []  # Can't determine duration, skip chunking

    chunks = []
    start = 0
    idx = 0
    while start < duration:
        end = min(start + chunk_sec, duration)
        chunk_path = os.path.join(output_dir, f"chunk_{idx:03d}.mp3")
        subprocess.run(
            ['ffmpeg', '-y', '-ss', str(start), '-t', str(chunk_sec), '-i', filepath, '-acodec', 'libmp3lame', '-q:a', '4', chunk_path],
            capture_output=True, timeout=120
        )
        if os.path.exists(chunk_path) and os.path.getsize(chunk_path) > 1000:
            chunks.append((chunk_path, start, end))
        idx += 1
        start += chunk_sec - overlap_sec

    return chunks


def _run_transcription(job_id, file_content, filename, mime_type, gemini_api_key, topic='', anthropic_api_key=''):
    """Background worker for audio transcription. Uses ffmpeg for chunking (disk-based, low memory)."""
    import time
    import io
    import tempfile
    import shutil

    tmp_dir = None
    try:
        file_size_mb = len(file_content) / (1024 * 1024)
        _transcription_jobs[job_id]['status'] = 'transcribing'
        print(f"[Job {job_id}] Starting transcription: {filename} ({file_size_mb:.1f}MB)")

        client = genai.Client(api_key=gemini_api_key)

        # Write file to temp disk to free memory
        tmp_dir = tempfile.mkdtemp(prefix='transcribe_')
        ext = os.path.splitext(filename)[1] or '.m4a'
        input_path = os.path.join(tmp_dir, f"input{ext}")
        with open(input_path, 'wb') as f:
            f.write(file_content)
        del file_content  # Free memory immediately

        # Check duration
        duration_sec = _get_audio_duration_ffprobe(input_path)
        chunk_threshold_sec = CHUNK_MINUTES * 60 * 1.5  # 22.5 min
        needs_chunking = duration_sec and duration_sec > chunk_threshold_sec

        if not duration_sec:
            # Heuristic fallback
            needs_chunking = file_size_mb > 25
            duration_sec = file_size_mb * 60  # rough estimate

        print(f"[Job {job_id}] Duration: {duration_sec:.0f}s ({duration_sec/60:.1f}min), chunking: {needs_chunking}")

        if needs_chunking:
            # === CHUNKED TRANSCRIPTION (disk-based, low memory) ===
            _transcription_jobs[job_id]['progress'] = 'Splitting audio...'
            chunks = _split_audio_ffmpeg(input_path, CHUNK_MINUTES * 60, OVERLAP_SECONDS, tmp_dir)

            if not chunks:
                print(f"[Job {job_id}] ffmpeg chunking failed, falling back to single-pass")
                needs_chunking = False
            else:
                print(f"[Job {job_id}] Split into {len(chunks)} chunks")

                # Delete input file to save disk space
                try:
                    os.remove(input_path)
                except Exception:
                    pass

                all_texts = []
                for idx, (chunk_path, start_sec, end_sec) in enumerate(chunks):
                    chunk_label = f"Chunk {idx + 1}/{len(chunks)}: "
                    _transcription_jobs[job_id]['progress'] = f"Chunk {idx + 1}/{len(chunks)} ({int(start_sec/60)}-{int(end_sec/60)} min)"

                    # Read chunk from disk (small — ~15min of mp3 ≈ 15MB)
                    with open(chunk_path, 'rb') as cf:
                        chunk_bytes = cf.read()
                    os.remove(chunk_path)  # Free disk immediately

                    chunk_size_mb = len(chunk_bytes) / (1024 * 1024)
                    uploaded = None
                    if chunk_size_mb > 20:
                        uploaded = client.files.upload(
                            file=io.BytesIO(chunk_bytes),
                            config=genai_types.UploadFileConfig(mime_type='audio/mpeg', display_name=f"{filename}_chunk{idx+1}")
                        )
                        del chunk_bytes
                        # Wait for file to be ready
                        wait_start = time.time()
                        while hasattr(uploaded, 'state') and str(uploaded.state) not in ('ACTIVE', 'State.ACTIVE', '2'):
                            if time.time() - wait_start > 120:
                                raise Exception('Gemini file processing timed out')
                            time.sleep(3)
                            uploaded = client.files.get(name=uploaded.name)
                        audio_content = uploaded
                    else:
                        audio_content = genai_types.Part.from_bytes(data=chunk_bytes, mime_type='audio/mpeg')
                        del chunk_bytes

                    try:
                        text = _transcribe_audio_content(client, audio_content, job_id, label=chunk_label)
                        all_texts.append(text)
                    finally:
                        if uploaded:
                            try:
                                client.files.delete(name=uploaded.name)
                            except Exception:
                                pass

                    if idx < len(chunks) - 1:
                        time.sleep(2)

                transcript_text = '\n\n'.join(all_texts)
                print(f"[Job {job_id}] All {len(chunks)} chunks done, total: {len(transcript_text)} chars")

        if not needs_chunking:
            # === SINGLE-PASS TRANSCRIPTION ===
            # Read from disk
            with open(input_path, 'rb') as f:
                audio_bytes = f.read()
            os.remove(input_path)

            uploaded_file = None
            audio_size_mb = len(audio_bytes) / (1024 * 1024)
            if audio_size_mb > 20:
                try:
                    _transcription_jobs[job_id]['progress'] = 'Uploading to Gemini...'
                    uploaded_file = client.files.upload(
                        file=io.BytesIO(audio_bytes),
                        config=genai_types.UploadFileConfig(mime_type=mime_type, display_name=filename)
                    )
                    del audio_bytes
                    print(f"[Job {job_id}] Uploaded to Gemini: {uploaded_file.name}, state: {uploaded_file.state}")
                    # Wait for file to be processed and ready
                    _transcription_jobs[job_id]['progress'] = 'Waiting for Gemini to process file...'
                    wait_start = time.time()
                    while hasattr(uploaded_file, 'state') and str(uploaded_file.state) not in ('ACTIVE', 'State.ACTIVE', '2'):
                        if time.time() - wait_start > 120:
                            raise Exception('Gemini file processing timed out after 2 minutes')
                        time.sleep(3)
                        uploaded_file = client.files.get(name=uploaded_file.name)
                        print(f"[Job {job_id}] File state: {uploaded_file.state}")
                    audio_content = uploaded_file
                    _transcription_jobs[job_id]['progress'] = 'Transcribing...'
                    print(f"[Job {job_id}] File ready, starting transcription")
                except Exception as upload_err:
                    print(f"[Job {job_id}] Upload failed: {upload_err}, using inline bytes")
                    audio_content = genai_types.Part.from_bytes(data=audio_bytes, mime_type=mime_type)
                    del audio_bytes
            else:
                audio_content = genai_types.Part.from_bytes(data=audio_bytes, mime_type=mime_type)
                del audio_bytes

            try:
                transcript_text = _transcribe_audio_content(client, audio_content, job_id)
            except Exception as e:
                _transcription_jobs[job_id] = {'status': 'error', 'error': str(e)}
                return
            finally:
                if uploaded_file:
                    try:
                        client.files.delete(name=uploaded_file.name)
                    except Exception:
                        pass

        if not transcript_text or not transcript_text.strip():
            _transcription_jobs[job_id] = {
                'status': 'error',
                'error': 'Gemini could not transcribe this audio. Try converting to MP3 or a shorter clip.'
            }
            return

        print(f"[Job {job_id}] Transcription complete: {len(transcript_text)} chars")

        # --- Transcript cleanup pass using Claude ---
        cleanup_key = anthropic_api_key or os.environ.get('ANTHROPIC_API_KEY', '')
        if cleanup_key:
            try:
                import anthropic as _anthropic_mod
                _transcription_jobs[job_id]['progress'] = 'Cleaning up transcript...'

                topic_context = ''
                if topic:
                    topic_context = (
                        f'\n\nMEETING CONTEXT: This transcript is from a meeting about "{topic}". '
                        f'Use your knowledge of this company/topic to identify and correct all misspelled '
                        f'product names, pipeline drugs, competitors, partners, subsidiaries, executives, '
                        f'and related industry terms. For example, if the topic is a pharma company, '
                        f'look up its actual drug names, pipeline assets, therapeutic areas, and competitors.'
                    )

                cleanup_prompt = (
                    "You are an expert transcript corrector for investment/finance meetings. "
                    "Speech-to-text models produce phonetic approximations of specialized terms. "
                    "Your job is to fix EVERY mistranscribed term in this transcript.\n\n"
                    "WHAT TO FIX:\n"
                    "- Drug/product names: The STT model gets these wrong almost every time. "
                    "Look up the REAL names of products, pipeline drugs, and competitors for this company. "
                    "(e.g., 'Kitruda'→'Keytruda', 'Sedera'→'Cidara', 'Elusitide'→'Elicitide')\n"
                    "- Company names and subsidiaries: Verify correct spellings\n"
                    "- Ticker symbols and acronyms: Must be UPPERCASE "
                    "(e.g., 'trp 2 ADC'→'TROP-2 ADC', 'pd l1'→'PD-L1', 'her 2'→'HER2')\n"
                    "- Scientific/medical terms: Correct drug mechanisms, biomarkers, disease names, "
                    "clinical trial phases\n"
                    "- People's names: CEO, CFO, executives mentioned in context\n"
                    "- Financial terms and metrics\n\n"
                    "RULES:\n"
                    "1. Fix ALL instances of each misspelled term consistently throughout the entire transcript\n"
                    "2. Do NOT change sentence structure, paraphrase, summarize, or remove any content\n"
                    "3. Do NOT add commentary, notes, or explanations — return ONLY the corrected transcript\n"
                    "4. Preserve all line breaks and formatting exactly as-is\n"
                    "5. When genuinely unsure about a term, keep the original"
                    f"{topic_context}"
                )
                _client = _anthropic_mod.Anthropic(api_key=cleanup_key, timeout=600)
                cleanup_resp = _client.messages.create(
                    model='claude-sonnet-4-20250514',
                    max_tokens=64000,
                    system=cleanup_prompt,
                    messages=[{'role': 'user', 'content': f"Please correct the specialized terms in this transcript:\n\n{transcript_text}"}],
                )
                cleaned = cleanup_resp.content[0].text.strip()
                if cleaned and len(cleaned) > len(transcript_text) * 0.5:
                    print(f"[Job {job_id}] Transcript cleaned: {len(transcript_text)} -> {len(cleaned)} chars")
                    transcript_text = cleaned
                else:
                    print(f"[Job {job_id}] Cleanup returned suspiciously short result, keeping original")
            except Exception as cleanup_err:
                print(f"[Job {job_id}] Transcript cleanup failed (keeping original): {cleanup_err}")

        _transcription_jobs[job_id] = {
            'status': 'done',
            'text': transcript_text,
            'filename': filename,
            'fileSizeMb': round(file_size_mb, 1),
            'charCount': len(transcript_text),
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"[Job {job_id}] Transcription failed: {e}")
        _transcription_jobs[job_id] = {'status': 'error', 'error': str(e)}
    finally:
        # Clean up temp directory
        if tmp_dir and os.path.exists(tmp_dir):
            try:
                shutil.rmtree(tmp_dir)
            except Exception:
                pass


@app.route('/api/transcribe-audio', methods=['POST'])
def transcribe_audio():
    """Start async audio transcription. Returns job_id immediately, poll /api/transcribe-audio/<job_id> for result."""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No audio file provided'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400

        gemini_api_key = os.environ.get('GEMINI_API_KEY', '') or request.form.get('geminiApiKey', '')
        if not gemini_api_key:
            return jsonify({'error': 'Gemini API key is required for audio transcription. Please add it in Settings.'}), 400

        allowed_extensions = ('.mp3', '.mp4', '.mpeg', '.mpga', '.m4a', '.wav', '.webm', '.ogg', '.flac')
        filename_lower = file.filename.lower()
        if not filename_lower.endswith(allowed_extensions):
            return jsonify({'error': f'Unsupported audio format. Supported: {", ".join(allowed_extensions)}'}), 400

        mime_map = {
            '.mp3': 'audio/mpeg', '.mp4': 'audio/mp4', '.mpeg': 'audio/mpeg',
            '.mpga': 'audio/mpeg', '.m4a': 'audio/mp4', '.wav': 'audio/wav',
            '.webm': 'audio/webm', '.ogg': 'audio/ogg', '.flac': 'audio/flac'
        }
        file_ext = '.' + filename_lower.rsplit('.', 1)[-1]
        mime_type = mime_map.get(file_ext, 'audio/mpeg')

        file_content = file.read()
        file_size_mb = len(file_content) / (1024 * 1024)

        # Create job and start background thread
        import uuid, threading
        job_id = str(uuid.uuid4())[:8]
        _transcription_jobs[job_id] = {'status': 'starting', 'filename': file.filename}

        topic = request.form.get('topic', '')
        anthropic_api_key = request.form.get('apiKey', '')
        thread = threading.Thread(target=_run_transcription, args=(job_id, file_content, file.filename, mime_type, gemini_api_key, topic, anthropic_api_key))
        thread.daemon = True
        thread.start()

        print(f"[Job {job_id}] Started transcription for {file.filename} ({file_size_mb:.1f}MB)")
        return jsonify({'success': True, 'job_id': job_id, 'filename': file.filename})

    except Exception as e:
        print(f"Error starting transcription: {e}")
        return jsonify({'error': f'Failed to start transcription: {str(e)}'}), 500


@app.route('/api/auto-process-audio', methods=['POST'])
def auto_process_audio():
    """Auto-process audio: transcribe + summarize + save. Used by local agent for iCloud auto-detection."""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No audio file provided'}), 400
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400

        gemini_api_key = os.environ.get('GEMINI_API_KEY', '') or request.form.get('geminiApiKey', '')
        if not gemini_api_key:
            return jsonify({'error': 'Gemini API key required'}), 400

        allowed_extensions = ('.mp3', '.mp4', '.mpeg', '.mpga', '.m4a', '.wav', '.webm', '.ogg', '.flac')
        if not file.filename.lower().endswith(allowed_extensions):
            return jsonify({'error': f'Unsupported audio format'}), 400

        file_content = file.read()
        filename = file.filename
        detail_level = request.form.get('detailLevel', 'standard')
        anthropic_api_key = request.form.get('apiKey', '') or os.environ.get('ANTHROPIC_API_KEY', '')

        mime_map = {'.mp3': 'audio/mpeg', '.mp4': 'audio/mp4', '.m4a': 'audio/mp4', '.wav': 'audio/wav', '.webm': 'audio/webm', '.ogg': 'audio/ogg', '.flac': 'audio/flac'}
        file_ext = '.' + filename.lower().rsplit('.', 1)[-1]
        mime_type = mime_map.get(file_ext, 'audio/mpeg')

        job_id = str(uuid.uuid4())[:8]
        _transcription_jobs[job_id] = {'status': 'starting', 'filename': filename, 'autoProcess': True}

        thread = threading.Thread(
            target=_run_auto_process_audio,
            args=(job_id, file_content, filename, mime_type, gemini_api_key, anthropic_api_key, detail_level),
            daemon=True
        )
        thread.start()

        print(f"[auto-audio {job_id}] Started for {filename}")
        return jsonify({'success': True, 'jobId': job_id, 'filename': filename})
    except Exception as e:
        print(f"Error auto-processing audio: {e}")
        return jsonify({'error': str(e)}), 500


def _run_auto_process_audio(job_id, file_content, filename, mime_type, gemini_api_key, anthropic_api_key, detail_level):
    """Background: transcribe audio, generate summary, save to DB, create alert."""
    try:
        # Step 1: Transcribe (reuse existing function)
        _run_transcription(job_id, file_content, filename, mime_type, gemini_api_key, '', anthropic_api_key)

        job = _transcription_jobs.get(job_id, {})
        transcript = job.get('transcript', '')
        if not transcript:
            print(f"[auto-audio {job_id}] Transcription failed or empty")
            return

        # Step 2: Generate summary
        _transcription_jobs[job_id]['status'] = 'summarizing'
        print(f"[auto-audio {job_id}] Transcription done ({len(transcript)} chars), generating summary...")

        html_format = """OUTPUT FORMAT — Return raw HTML. No markdown. No code fences.
Use: <h2>Section Title</h2>, <p><strong>Topic:</strong> Description.</p>, <ul><li>Sub-point</li></ul>"""

        summary_instruction = f"""Generate a clear, well-structured summary of the following transcript.
Include key points, decisions, action items, and important details.
{html_format}"""

        if detail_level == 'concise':
            summary_instruction = f"""Generate a SHORT, CONCISE summary. Maximum 300-500 words. Only critical takeaways.
{html_format}"""
        elif detail_level == 'detailed':
            summary_instruction = f"""Generate an EXTREMELY DETAILED summary. Preserve all content, numbers, quotes.
{html_format}"""

        keys = _get_api_keys(anthropic_api_key=anthropic_api_key, gemini_api_key=gemini_api_key)
        summary_result = call_llm(
            messages=[{"role": "user", "content": f"{summary_instruction}\n\nTRANSCRIPT:\n{transcript[:50000]}"}],
            system="You are a meeting notes analyst. Generate structured HTML summaries.",
            tier="standard",
            max_tokens=8192,
        )
        summary_html = summary_result.get('text', '')

        questions_result = call_llm(
            messages=[{"role": "user", "content": f"Based on this transcript, generate 3-5 key follow-up questions.\nReturn raw HTML: <ol><li>Question?</li></ol>\n\nTRANSCRIPT:\n{transcript[:20000]}"}],
            system="Generate insightful follow-up questions.",
            tier="fast",
            max_tokens=2048,
        )
        questions_html = questions_result.get('text', '')

        # Step 3: Save to DB
        title = os.path.splitext(filename)[0].replace('_', ' ').replace('-', ' ')
        summary_id = str(uuid.uuid4())
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO meeting_summaries (id, title, raw_notes, summary, questions, source_type, doc_type, created_at)
                VALUES (%s, %s, %s, %s, %s, 'audio_recording', 'audio_recording', NOW())
            ''', (summary_id, title, transcript, summary_html, questions_html))

        # Step 4: Create alert
        with get_db(commit=True) as (conn, cur):
            alert_id = str(uuid.uuid4())
            cur.execute('''
                INSERT INTO agent_alerts (id, alert_type, ticker, title, detail, status, created_at)
                VALUES (%s, 'audio_summary', '', %s, %s, 'new', NOW())
            ''', (alert_id, f'Summary generated: {title}',
                  json.dumps({'filename': filename, 'summaryId': summary_id, 'detailLevel': detail_level, 'transcriptLength': len(transcript)})))

        _transcription_jobs[job_id]['status'] = 'complete'
        _transcription_jobs[job_id]['summaryId'] = summary_id
        print(f"[auto-audio {job_id}] Complete: {title} saved as {summary_id}")

    except Exception as e:
        print(f"[auto-audio {job_id}] Failed: {e}")
        import traceback; traceback.print_exc()
        _transcription_jobs[job_id]['status'] = 'error'
        _transcription_jobs[job_id]['error'] = str(e)


@app.route('/api/transcribe-audio/<job_id>', methods=['GET'])
def transcribe_audio_status(job_id):
    """Poll for transcription job status."""
    job = _transcription_jobs.get(job_id)
    if not job:
        return jsonify({'error': 'Job not found'}), 404
    if job['status'] == 'done':
        result = dict(job)
        # Keep job for 10 min so user can return from background, then clean up
        if 'completed_at' not in job:
            job['completed_at'] = time.time()
        elif time.time() - job['completed_at'] > 600:
            del _transcription_jobs[job_id]
        return jsonify(result)
    if job['status'] == 'error':
        error = job.get('error', 'Unknown error')
        # Keep errors for 5 min
        if 'completed_at' not in job:
            job['completed_at'] = time.time()
        elif time.time() - job['completed_at'] > 300:
            del _transcription_jobs[job_id]
        return jsonify({'status': 'error', 'error': error})
    result = {'status': job['status']}
    if 'progress' in job:
        result['progress'] = job['progress']
    return jsonify(result)


@app.route('/api/transcription-jobs', methods=['GET'])
def list_transcription_jobs():
    """Debug: list all active transcription jobs."""
    jobs = {}
    for jid, job in _transcription_jobs.items():
        info = {'status': job.get('status'), 'progress': job.get('progress'), 'filename': job.get('filename')}
        if job.get('status') == 'done':
            info['charCount'] = job.get('charCount')
        if job.get('status') == 'error':
            info['error'] = job.get('error')
        jobs[jid] = info
    return jsonify({'jobs': jobs, 'count': len(jobs)})


@app.route('/api/text-to-docx', methods=['POST'])
def text_to_docx():
    """Convert transcript text to a .docx file and return as base64"""
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import base64
        import io

        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({'error': 'No text provided'}), 400

        text = data['text']
        title = data.get('title', 'Transcript')

        doc = Document()

        # Set default font
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)

        # Add title
        heading = doc.add_heading(title, level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

        # Add transcript content, preserving line breaks
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line)
            else:
                doc.add_paragraph('')

        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        docx_base64 = base64.b64encode(buffer.read()).decode('utf-8')

        return jsonify({
            'success': True,
            'fileData': docx_base64,
            'fileSize': len(buffer.getvalue())
        })

    except Exception as e:
        print(f"Error creating docx: {e}")
        return jsonify({'error': f'Failed to create document: {str(e)}'}), 500


# ============================================
# SUMMARY EXPORT ENDPOINTS (Word / PDF / Bulk)
# ============================================

def _html_to_docx_elements(doc, html_content):
    """Parse HTML content and add formatted elements to a python-docx Document."""
    from html.parser import HTMLParser
    from docx.shared import Pt, RGBColor

    class DocxHTMLParser(HTMLParser):
        def __init__(self, document):
            super().__init__()
            self.doc = document
            self.current_paragraph = None
            self.bold = False
            self.italic = False
            self.in_list = False
            self.list_type = 'ul'
            self.heading_level = 0

        def handle_starttag(self, tag, attrs):
            tag = tag.lower()
            if tag in ('h1', 'h2', 'h3', 'h4'):
                self.heading_level = int(tag[1])
                self.current_paragraph = self.doc.add_heading('', level=self.heading_level)
            elif tag == 'p':
                if self.in_list:
                    self.current_paragraph = self.doc.add_paragraph('', style='List Bullet' if self.list_type == 'ul' else 'List Number')
                else:
                    self.current_paragraph = self.doc.add_paragraph('')
            elif tag == 'ul':
                self.in_list = True
                self.list_type = 'ul'
            elif tag == 'ol':
                self.in_list = True
                self.list_type = 'ol'
            elif tag == 'li':
                style = 'List Bullet' if self.list_type == 'ul' else 'List Number'
                self.current_paragraph = self.doc.add_paragraph('', style=style)
            elif tag in ('strong', 'b'):
                self.bold = True
            elif tag in ('em', 'i'):
                self.italic = True
            elif tag == 'br':
                if self.current_paragraph:
                    self.current_paragraph.add_run('\n')

        def handle_endtag(self, tag):
            tag = tag.lower()
            if tag in ('h1', 'h2', 'h3', 'h4'):
                self.heading_level = 0
            elif tag in ('ul', 'ol'):
                self.in_list = False
            elif tag in ('strong', 'b'):
                self.bold = False
            elif tag in ('em', 'i'):
                self.italic = False

        def handle_data(self, data):
            text = data
            if not text.strip():
                return
            if self.current_paragraph is None:
                self.current_paragraph = self.doc.add_paragraph('')
            run = self.current_paragraph.add_run(text)
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.bold = self.bold
            run.italic = self.italic

        def handle_entityref(self, name):
            import html as html_mod
            self.handle_data(html_mod.unescape(f'&{name};'))

        def handle_charref(self, name):
            import html as html_mod
            self.handle_data(html_mod.unescape(f'&#{name};'))

    parser = DocxHTMLParser(doc)
    parser.feed(html_content)


def _generate_summary_docx_bytes(row, sections=None):
    """Generate docx bytes from a summary DB row dict.
    sections: list of section keys to include, e.g. ['takeaways','questions','assessment'].
              None means include all sections.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    include_all = sections is None
    sect_set = set(sections) if sections else set()

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0, 0, 0)

    title = row.get('title') or 'Summary'
    heading = doc.add_heading(title, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in heading.runs:
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 0, 0)

    created = row.get('created_at')
    if created:
        p = doc.add_paragraph(str(created)[:10])
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.color.rgb = RGBColor(128, 128, 128)

    summary_html = row.get('summary') or ''
    questions_html = row.get('questions') or ''
    assessment_html = row.get('assessment') or ''
    raw_notes = row.get('raw_notes') or ''
    source_type = row.get('source_type') or ''

    if summary_html and (include_all or 'takeaways' in sect_set):
        h = doc.add_heading('Key Takeaways', level=2)
        for r in h.runs:
            r.font.name = 'Calibri'
            r.font.color.rgb = RGBColor(0, 0, 0)
        _html_to_docx_elements(doc, summary_html)

    if questions_html and (include_all or 'questions' in sect_set):
        h = doc.add_heading('Follow-up Questions', level=2)
        for r in h.runs:
            r.font.name = 'Calibri'
            r.font.color.rgb = RGBColor(0, 0, 0)
        _html_to_docx_elements(doc, questions_html)

    if assessment_html and (include_all or 'assessment' in sect_set):
        h = doc.add_heading('Assessment', level=2)
        for r in h.runs:
            r.font.name = 'Calibri'
            r.font.color.rgb = RGBColor(0, 0, 0)
        _html_to_docx_elements(doc, assessment_html)

    if source_type == 'audio' and raw_notes and include_all:
        h = doc.add_heading('Full Transcript', level=2)
        for r in h.runs:
            r.font.name = 'Calibri'
            r.font.color.rgb = RGBColor(0, 0, 0)
        for line in raw_notes.split('\n'):
            p = doc.add_paragraph(line)
            for r in p.runs:
                r.font.name = 'Calibri'
                r.font.size = Pt(10)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _generate_summary_pdf_bytes(row, sections_filter=None):
    """Generate PDF bytes from a summary DB row dict.
    sections_filter: list of section keys to include, e.g. ['takeaways','questions','assessment'].
                     None means include all sections.
    """
    from xhtml2pdf import pisa

    include_all = sections_filter is None
    sect_set = set(sections_filter) if sections_filter else set()

    title = row.get('title') or 'Summary'
    created = row.get('created_at')
    date_str = str(created)[:10] if created else ''
    summary_html = row.get('summary') or ''
    questions_html = row.get('questions') or ''
    assessment_html = row.get('assessment') or ''
    raw_notes = row.get('raw_notes') or ''
    source_type = row.get('source_type') or ''

    sections = ''
    if summary_html and (include_all or 'takeaways' in sect_set):
        sections += f'<h2>Key Takeaways</h2>{summary_html}'
    if questions_html and (include_all or 'questions' in sect_set):
        sections += f'<h2>Follow-up Questions</h2>{questions_html}'
    if assessment_html and (include_all or 'assessment' in sect_set):
        sections += f'<h2>Assessment</h2>{assessment_html}'
    if source_type == 'audio' and raw_notes and include_all:
        escaped = raw_notes.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
        sections += f'<h2>Full Transcript</h2><p style="font-size:10pt;">{escaped}</p>'

    full_html = f"""<html><head><style>
        @page {{ margin: 1in; }}
        body {{ font-family: Calibri, Arial, Helvetica, sans-serif; font-size: 11pt; color: #000; line-height: 1.5; }}
        h1 {{ font-size: 18pt; margin-bottom: 4pt; }}
        h2 {{ font-size: 14pt; border-bottom: 1px solid #ccc; padding-bottom: 4pt; margin-top: 20pt; }}
        h3 {{ font-size: 12pt; }}
        ul, ol {{ padding-left: 20pt; }}
        li {{ margin-bottom: 6pt; }}
        p {{ margin-bottom: 8pt; }}
        .date {{ font-size: 9pt; color: #888; margin-bottom: 16pt; }}
    </style></head><body>
        <h1>{title}</h1>
        <p class="date">{date_str}</p>
        {sections}
    </body></html>"""

    buf = io.BytesIO()
    pisa.CreatePDF(full_html, dest=buf)
    return buf.getvalue()


@app.route('/api/summary-to-docx', methods=['POST'])
def summary_to_docx():
    """Export a summary as a formatted Word document."""
    try:
        data = request.get_json()
        summary_id = data.get('summaryId')
        if not summary_id:
            return jsonify({'error': 'No summary ID provided'}), 400

        with get_db() as (_, cur):
            cur.execute('SELECT * FROM meeting_summaries WHERE id = %s', (summary_id,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Summary not found'}), 404

        docx_bytes = _generate_summary_docx_bytes(dict(row))
        docx_b64 = base64.b64encode(docx_bytes).decode('utf-8')
        safe_title = re.sub(r'[^\w\s-]', '', row['title'] or 'Summary')[:50].strip().replace(' ', '_')

        return jsonify({
            'success': True,
            'fileData': docx_b64,
            'filename': f"{safe_title}.docx",
            'fileSize': len(docx_bytes)
        })
    except Exception as e:
        print(f"Error creating summary docx: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/summary-to-pdf', methods=['POST'])
def summary_to_pdf():
    """Export a summary as a styled PDF document."""
    try:
        data = request.get_json()
        summary_id = data.get('summaryId')
        if not summary_id:
            return jsonify({'error': 'No summary ID provided'}), 400

        with get_db() as (_, cur):
            cur.execute('SELECT * FROM meeting_summaries WHERE id = %s', (summary_id,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Summary not found'}), 404

        pdf_bytes = _generate_summary_pdf_bytes(dict(row))
        pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
        safe_title = re.sub(r'[^\w\s-]', '', row['title'] or 'Summary')[:50].strip().replace(' ', '_')

        return jsonify({
            'success': True,
            'fileData': pdf_b64,
            'filename': f"{safe_title}.pdf",
            'fileSize': len(pdf_bytes)
        })
    except Exception as e:
        print(f"Error creating summary PDF: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/summaries-bulk-export', methods=['POST'])
def summaries_bulk_export():
    """Export multiple summaries as a zip of Word or PDF files."""
    try:
        import zipfile

        data = request.get_json()
        summary_ids = data.get('summaryIds', [])
        export_format = data.get('format', 'docx')
        export_sections = data.get('sections')  # e.g. ['takeaways','questions','assessment'] or None for all

        if not summary_ids:
            return jsonify({'error': 'No summary IDs provided'}), 400

        with get_db() as (_, cur):
            placeholders = ','.join(['%s'] * len(summary_ids))
            cur.execute(f'SELECT * FROM meeting_summaries WHERE id IN ({placeholders})', summary_ids)
            rows = cur.fetchall()

        if not rows:
            return jsonify({'error': 'No summaries found'}), 404

        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            used_names = set()
            for row in rows:
                row_dict = dict(row)
                safe_title = re.sub(r'[^\w\s-]', '', row_dict.get('title') or 'Summary')[:50].strip().replace(' ', '_')
                filename = f"{safe_title}.{export_format}"
                # Deduplicate
                while filename in used_names:
                    filename = f"{safe_title}_{len(used_names)}.{export_format}"
                used_names.add(filename)

                if export_format == 'pdf':
                    file_bytes = _generate_summary_pdf_bytes(row_dict, sections_filter=export_sections)
                else:
                    file_bytes = _generate_summary_docx_bytes(row_dict, sections=export_sections)
                zf.writestr(filename, file_bytes)

        zip_bytes = zip_buf.getvalue()
        zip_b64 = base64.b64encode(zip_bytes).decode('utf-8')

        return jsonify({
            'success': True,
            'fileData': zip_b64,
            'filename': f"summaries-{len(rows)}-files.zip",
            'fileSize': len(zip_bytes),
            'fileCount': len(rows)
        })
    except Exception as e:
        print(f"Error bulk exporting summaries: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# ANALYSIS (THESIS) EXPORT ENDPOINTS
# ============================================

def _generate_analysis_docx_bytes(row):
    """Generate docx bytes from a portfolio_analyses row."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0, 0, 0)

    ticker = row.get('ticker', 'Stock')
    company = row.get('company', '')
    analysis = row.get('analysis') or {}
    if isinstance(analysis, str):
        import json as _json
        try: analysis = _json.loads(analysis)
        except: analysis = {}

    title_text = f"{ticker} — {company}" if company else ticker
    heading = doc.add_heading(title_text, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in heading.runs:
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 0, 0)

    updated = row.get('updated_at')
    if updated:
        p = doc.add_paragraph(str(updated)[:10])
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.color.rgb = RGBColor(128, 128, 128)

    is_condensed = analysis.get('_condensed', False)
    thesis = analysis.get('thesis', {})
    signposts = analysis.get('signposts', [])
    threats = analysis.get('threats', [])
    conclusion = analysis.get('conclusion', '')

    # Thesis section
    if thesis:
        h = doc.add_heading('Investment Thesis', level=2)
        for r in h.runs: r.font.name = 'Calibri'; r.font.color.rgb = RGBColor(0, 0, 0)
        if thesis.get('summary'):
            p = doc.add_paragraph(thesis['summary'])
            for r in p.runs: r.font.name = 'Calibri'
        for pillar in thesis.get('pillars', []):
            title = pillar.get('pillar', pillar.get('title', ''))
            desc = pillar.get('detail', pillar.get('description', ''))
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(f"{title}: ")
            run.bold = True
            run.font.name = 'Calibri'
            run2 = p.add_run(desc)
            run2.font.name = 'Calibri'

    # Signposts section
    if signposts:
        h = doc.add_heading('Signposts (What We\'re Watching)', level=2)
        for r in h.runs: r.font.name = 'Calibri'; r.font.color.rgb = RGBColor(0, 0, 0)
        for sp in signposts:
            metric = sp.get('metric', sp.get('signpost', ''))
            target = sp.get('target', '')
            timeframe = sp.get('timeframe', '')
            text = f"{metric}: {target}"
            if timeframe:
                text += f" ({timeframe})"
            p = doc.add_paragraph(text, style='List Bullet')
            for r in p.runs: r.font.name = 'Calibri'

    # Threats section
    if threats:
        h = doc.add_heading('Thesis Threats', level=2)
        for r in h.runs: r.font.name = 'Calibri'; r.font.color.rgb = RGBColor(0, 0, 0)
        for threat in threats:
            threat_desc = threat.get('threat', '')
            triggers = threat.get('triggerPoints', '')
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(threat_desc)
            run.bold = True
            run.font.name = 'Calibri'
            details = []
            if triggers: details.append(f"Watch for: {triggers}")
            if details:
                run2 = p.add_run(f"\n{' | '.join(details)}")
                run2.font.name = 'Calibri'
                run2.font.size = Pt(10)
                run2.font.color.rgb = RGBColor(100, 100, 100)

    # Conclusion
    if conclusion:
        h = doc.add_heading('Conclusion', level=2)
        for r in h.runs: r.font.name = 'Calibri'; r.font.color.rgb = RGBColor(0, 0, 0)
        p = doc.add_paragraph(conclusion)
        for r in p.runs: r.font.name = 'Calibri'

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _generate_analysis_pdf_bytes(row):
    """Generate PDF bytes from a portfolio_analyses row."""
    from xhtml2pdf import pisa

    ticker = row.get('ticker', 'Stock')
    company = row.get('company', '')
    analysis = row.get('analysis') or {}
    if isinstance(analysis, str):
        import json as _json
        try: analysis = _json.loads(analysis)
        except: analysis = {}
    updated = row.get('updated_at')
    date_str = str(updated)[:10] if updated else ''

    is_condensed = analysis.get('_condensed', False)
    thesis = analysis.get('thesis', {})
    signposts = analysis.get('signposts', [])
    threats = analysis.get('threats', [])
    conclusion = analysis.get('conclusion', '')

    sections = ''
    if thesis:
        sections += '<h2>Investment Thesis</h2>'
        if thesis.get('summary'):
            sections += f'<p>{thesis["summary"]}</p>'
        for pillar in thesis.get('pillars', []):
            t = pillar.get('pillar', pillar.get('title', ''))
            d = pillar.get('detail', pillar.get('description', ''))
            sections += f'<li><strong>{t}:</strong> {d}</li>'
        if thesis.get('pillars'):
            sections = sections  # already has <li> items

    if signposts:
        sections += '<h2>Signposts (What We\'re Watching)</h2><ul>'
        for sp in signposts:
            m = sp.get('metric', sp.get('signpost', ''))
            t = sp.get('target', '')
            tf = sp.get('timeframe', '')
            sections += f'<li><strong>{m}:</strong> {t}'
            if tf: sections += f' <em>({tf})</em>'
            sections += '</li>'
        sections += '</ul>'

    if threats:
        sections += '<h2>Thesis Threats</h2><ul>'
        for threat in threats:
            td = threat.get('threat', '')
            lk = threat.get('likelihood', '')
            imp = threat.get('impact', '')
            tr = threat.get('triggerPoints', '')
            sections += f'<li><strong>{td}</strong>'
            if not is_condensed and (lk or imp): sections += f'<br/><span style="color:#666;font-size:0.9em;">Likelihood: {lk} | Impact: {imp}</span>'
            if tr: sections += f'<br/><span style="color:#666;font-size:0.9em;">Watch for: {tr}</span>'
            sections += '</li>'
        sections += '</ul>'

    if conclusion:
        esc = conclusion.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
        sections += f'<h2>Conclusion</h2><p>{esc}</p>'

    title_text = f"{ticker} — {company}" if company else ticker
    full_html = f"""<html><head><style>
        @page {{ margin: 1in; }}
        body {{ font-family: Calibri, Arial, Helvetica, sans-serif; font-size: 11pt; color: #000; line-height: 1.5; }}
        h1 {{ font-size: 18pt; margin-bottom: 4pt; color: #1a365d; border-bottom: 2px solid #2c5282; padding-bottom: 10px; }}
        h2 {{ font-size: 14pt; color: #2c5282; margin-top: 20pt; }}
        ul {{ padding-left: 20pt; }}
        li {{ margin-bottom: 6pt; }}
        p {{ margin-bottom: 8pt; }}
        .date {{ font-size: 9pt; color: #888; margin-bottom: 16pt; }}
    </style></head><body>
        <h1>{title_text}</h1>
        <p class="date">{date_str}</p>
        {sections}
    </body></html>"""

    buf = io.BytesIO()
    pisa.CreatePDF(full_html, dest=buf)
    return buf.getvalue()


def _resolve_thesis_tier(data):
    """Resolve thesisTier from request data with backward compat."""
    thesis_tier = data.get('thesisTier', 'detailed')
    if not thesis_tier or thesis_tier == 'detailed':
        if data.get('useCondensed'):
            thesis_tier = 'condensed'
    return thesis_tier


def _overlay_thesis_tier(row_dict, thesis_tier):
    """Overlay full or condensed thesis data onto analysis row dict."""
    if thesis_tier == 'detailed' or not thesis_tier:
        return
    table = 'thesis_full' if thesis_tier == 'full' else 'thesis_condensed'
    column = 'full_analysis' if thesis_tier == 'full' else 'condensed_analysis'
    ticker = row_dict.get('ticker', '')
    with get_db() as (_, cur):
        cur.execute(f'SELECT {column} FROM {table} WHERE ticker = %s', (ticker,))
        tier_row = cur.fetchone()
    if not tier_row:
        return
    tier_data = tier_row[column]
    if isinstance(tier_data, str):
        tier_data = json.loads(tier_data)
    orig = row_dict.get('analysis') or {}
    if isinstance(orig, str):
        orig = json.loads(orig)
    orig['thesis'] = tier_data.get('thesis', orig.get('thesis', {}))
    orig['signposts'] = tier_data.get('signposts', orig.get('signposts', []))
    orig['threats'] = tier_data.get('threats', orig.get('threats', []))
    orig['conclusion'] = tier_data.get('conclusion', orig.get('conclusion', ''))
    if thesis_tier == 'condensed':
        orig['_condensed'] = True
    row_dict['analysis'] = orig


def _thesis_tier_suffix(thesis_tier):
    """Return filename suffix for thesis tier."""
    return {'full': '_Full', 'condensed': '_Condensed'}.get(thesis_tier, '')


@app.route('/api/analysis-to-docx', methods=['POST'])
def analysis_to_docx():
    """Export a portfolio analysis as a Word document."""
    try:
        data = request.get_json()
        ticker = data.get('ticker')
        thesis_tier = _resolve_thesis_tier(data)
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Analysis not found'}), 404
        row_dict = dict(row)
        _overlay_thesis_tier(row_dict, thesis_tier)
        docx_bytes = _generate_analysis_docx_bytes(row_dict)
        suffix = _thesis_tier_suffix(thesis_tier)
        docx_b64 = base64.b64encode(docx_bytes).decode('utf-8')
        return jsonify({'success': True, 'fileData': docx_b64, 'filename': f"{ticker}_Thesis{suffix}.docx", 'fileSize': len(docx_bytes)})
    except Exception as e:
        print(f"Error creating analysis docx: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/analysis-to-pdf', methods=['POST'])
def analysis_to_pdf():
    """Export a portfolio analysis as a PDF."""
    try:
        data = request.get_json()
        ticker = data.get('ticker')
        thesis_tier = _resolve_thesis_tier(data)
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Analysis not found'}), 404
        row_dict = dict(row)
        _overlay_thesis_tier(row_dict, thesis_tier)
        pdf_bytes = _generate_analysis_pdf_bytes(row_dict)
        suffix = _thesis_tier_suffix(thesis_tier)
        pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
        return jsonify({'success': True, 'fileData': pdf_b64, 'filename': f"{ticker}_Thesis{suffix}.pdf", 'fileSize': len(pdf_bytes)})
    except Exception as e:
        print(f"Error creating analysis pdf: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyses-bulk-export', methods=['POST'])
def analyses_bulk_export():
    """Export multiple analyses as a zip of Word or PDF files."""
    try:
        import zipfile
        data = request.get_json()
        tickers = data.get('tickers', [])
        export_format = data.get('format', 'docx')
        thesis_tier = _resolve_thesis_tier(data)
        if not tickers:
            return jsonify({'error': 'No tickers provided'}), 400
        with get_db() as (_, cur):
            placeholders = ','.join(['%s'] * len(tickers))
            cur.execute(f'SELECT * FROM portfolio_analyses WHERE ticker IN ({placeholders})', tickers)
            rows = cur.fetchall()
        if not rows:
            return jsonify({'error': 'No analyses found'}), 404
        # Load tier overlay data if requested
        tier_map = {}
        if thesis_tier == 'full':
            with get_db() as (_, cur2):
                cur2.execute(f'SELECT ticker, full_analysis FROM thesis_full WHERE ticker IN ({placeholders})', tickers)
                for cr in cur2.fetchall():
                    ca = cr['full_analysis']
                    if isinstance(ca, str):
                        ca = json.loads(ca)
                    tier_map[cr['ticker']] = ca
        elif thesis_tier == 'condensed':
            with get_db() as (_, cur2):
                cur2.execute(f'SELECT ticker, condensed_analysis FROM thesis_condensed WHERE ticker IN ({placeholders})', tickers)
                for cr in cur2.fetchall():
                    ca = cr['condensed_analysis']
                    if isinstance(ca, str):
                        ca = json.loads(ca)
                    tier_map[cr['ticker']] = ca
        suffix = _thesis_tier_suffix(thesis_tier)
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            for row in rows:
                row_dict = dict(row)
                t = row_dict.get('ticker', 'Stock')
                if thesis_tier in ('full', 'condensed') and t in tier_map:
                    orig = row_dict.get('analysis') or {}
                    if isinstance(orig, str):
                        orig = json.loads(orig)
                    tier_data = tier_map[t]
                    orig['thesis'] = tier_data.get('thesis', orig.get('thesis', {}))
                    orig['signposts'] = tier_data.get('signposts', orig.get('signposts', []))
                    orig['threats'] = tier_data.get('threats', orig.get('threats', []))
                    orig['conclusion'] = tier_data.get('conclusion', orig.get('conclusion', ''))
                    if thesis_tier == 'condensed':
                        orig['_condensed'] = True
                    row_dict['analysis'] = orig
                fn = f"{t}_Thesis{suffix}.{export_format}"
                if export_format == 'pdf':
                    zf.writestr(fn, _generate_analysis_pdf_bytes(row_dict))
                else:
                    zf.writestr(fn, _generate_analysis_docx_bytes(row_dict))
        zip_bytes = zip_buf.getvalue()
        label = {'full': 'full-theses', 'condensed': 'condensed-theses'}.get(thesis_tier, 'theses')
        return jsonify({'success': True, 'fileData': base64.b64encode(zip_bytes).decode('utf-8'), 'filename': f"{label}-{len(rows)}-files.zip", 'fileSize': len(zip_bytes), 'fileCount': len(rows)})
    except Exception as e:
        print(f"Error bulk exporting analyses: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-analyses-bulk', methods=['POST'])
def email_analyses_bulk():
    """Email multiple analyses in one combined email."""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    try:
        data = request.get_json()
        tickers = data.get('tickers', [])
        recipient = data.get('email')
        smtp_config = data.get('smtpConfig', {})
        thesis_tier = _resolve_thesis_tier(data)
        if not tickers or not recipient:
            return jsonify({'error': 'Tickers and email required'}), 400

        with get_db() as (_, cur):
            placeholders = ','.join(['%s'] * len(tickers))
            cur.execute(f'SELECT * FROM portfolio_analyses WHERE ticker IN ({placeholders})', tickers)
            rows = cur.fetchall()
        if not rows:
            return jsonify({'error': 'No analyses found'}), 404

        # Load tier overlay data if requested
        tier_map = {}
        if thesis_tier == 'full':
            with get_db() as (_, cur2):
                cur2.execute(f'SELECT ticker, full_analysis FROM thesis_full WHERE ticker IN ({placeholders})', tickers)
                for cr in cur2.fetchall():
                    ca = cr['full_analysis']
                    if isinstance(ca, str):
                        ca = json.loads(ca)
                    tier_map[cr['ticker']] = ca
        elif thesis_tier == 'condensed':
            with get_db() as (_, cur2):
                cur2.execute(f'SELECT ticker, condensed_analysis FROM thesis_condensed WHERE ticker IN ({placeholders})', tickers)
                for cr in cur2.fetchall():
                    ca = cr['condensed_analysis']
                    if isinstance(ca, str):
                        ca = json.loads(ca)
                    tier_map[cr['ticker']] = ca

        # Sort by ticker for consistency
        rows = sorted(rows, key=lambda r: r['ticker'])
        ticker_list = ', '.join(r['ticker'] for r in rows)
        tier_prefix = {'full': 'Full ', 'condensed': 'Condensed '}.get(thesis_tier, '')
        subject = f"{tier_prefix}Investment Theses: {ticker_list}"

        html_body = '<html><body style="font-family: Arial, sans-serif; line-height: 1.6; color: #333; max-width: 700px;">'
        plain_text = ''
        for idx, row in enumerate(rows):
            row_dict = dict(row)
            ticker = row_dict.get('ticker', 'Stock')
            company = row_dict.get('company', '')
            analysis = row_dict.get('analysis') or {}
            if isinstance(analysis, str):
                import json as _json
                try: analysis = _json.loads(analysis)
                except: analysis = {}

            # Overlay tier data if available
            if thesis_tier in ('full', 'condensed') and ticker in tier_map:
                tier_data = tier_map[ticker]
                if isinstance(tier_data, str):
                    try: tier_data = json.loads(tier_data)
                    except: tier_data = {}
                if isinstance(tier_data, dict):
                    analysis['thesis'] = tier_data.get('thesis', analysis.get('thesis', {}))
                    analysis['signposts'] = tier_data.get('signposts', analysis.get('signposts', []))
                    analysis['threats'] = tier_data.get('threats', analysis.get('threats', []))

            def _ensure_dict(val, fallback_key=None):
                if isinstance(val, str):
                    try: return json.loads(val)
                    except: return {fallback_key: val} if fallback_key else {}
                return val if isinstance(val, (dict, list)) else {}

            thesis = _ensure_dict(analysis.get('thesis', {}), 'summary')
            signposts = analysis.get('signposts', [])
            if isinstance(signposts, str):
                try: signposts = json.loads(signposts)
                except: signposts = []
            threats = analysis.get('threats', [])
            if isinstance(threats, str):
                try: threats = json.loads(threats)
                except: threats = []

            if idx > 0:
                html_body += '<hr style="border:none;border-top:2px solid #ccc;margin:30px 0;"/>'
                plain_text += '\n' + '='*60 + '\n\n'

            html_body += f'<h1 style="color:#1a365d;border-bottom:2px solid #2c5282;padding-bottom:10px;">{ticker} — {company}</h1>'
            plain_text += f"{ticker} — {company}\n{'='*40}\n\n"

            if thesis:
                html_body += '<h2 style="color: #2c5282; margin-top: 25px;">1. Investment Thesis</h2>'
                html_body += f'<p style="margin-left: 20px;">{thesis.get("summary","")}</p>'
                plain_text += f"1. INVESTMENT THESIS\n{thesis.get('summary','')}\n\n"
                if thesis.get('pillars'):
                    html_body += '<ul style="margin-left: 20px;">'
                    for pillar in thesis['pillars']:
                        if isinstance(pillar, str): pillar = {'pillar': pillar}
                        t = pillar.get('pillar', pillar.get('title', ''))
                        d = pillar.get('detail', pillar.get('description', ''))
                        html_body += f'<li style="margin-bottom: 8px;"><strong>{t}:</strong> {d}</li>'
                        plain_text += f"  - {t}: {d}\n"
                    html_body += '</ul>'
                    plain_text += '\n'

            if signposts:
                html_body += '<h2 style="color: #2c5282; margin-top: 25px;">2. Signposts (What We\'re Watching)</h2><ul style="margin-left: 20px;">'
                plain_text += "2. SIGNPOSTS\n"
                for sp in signposts:
                    if isinstance(sp, str): sp = {'metric': sp}
                    m = sp.get('metric', sp.get('signpost', ''))
                    tgt = sp.get('target', '')
                    tf = sp.get('timeframe', '')
                    line = f"{m}: {tgt}"
                    if tf: line += f" ({tf})"
                    html_body += f'<li style="margin-bottom: 8px;"><strong>{m}:</strong> {tgt}'
                    if tf: html_body += f' <em>({tf})</em>'
                    html_body += '</li>'
                    plain_text += f"  - {line}\n"
                html_body += '</ul>'
                plain_text += '\n'

            if threats:
                html_body += '<h2 style="color: #2c5282; margin-top: 25px;">3. Thesis Threats (Where We Can Be Wrong)</h2><ul style="margin-left: 20px;">'
                plain_text += "3. THESIS THREATS\n"
                for threat in threats:
                    if isinstance(threat, str): threat = {'threat': threat}
                    td = threat.get('threat', '')
                    triggers = threat.get('triggerPoints', '')
                    html_body += f'<li style="margin-bottom: 10px;"><strong>{td}</strong>'
                    if triggers:
                        html_body += f'<br><span style="color: #666; font-size: 0.9em;">Watch for: {triggers}</span>'
                    html_body += '</li>'
                    plain_text += f"  - {td}\n"
                    if triggers:
                        plain_text += f"    Watch for: {triggers}\n"
                html_body += '</ul>'
                plain_text += '\n'

        html_body += '</body></html>'

        use_gmail = smtp_config.get('use_gmail', True)
        gmail_user = smtp_config.get('gmail_user', '')
        gmail_password = smtp_config.get('gmail_app_password', '')
        from_email = smtp_config.get('from_email', gmail_user)

        msg = MIMEMultipart('alternative')
        msg['Subject'] = subject
        msg['From'] = from_email
        msg['To'] = recipient
        msg.attach(MIMEText(plain_text, 'plain'))
        msg.attach(MIMEText(html_body, 'html'))

        if use_gmail:
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(gmail_user, gmail_password)
                server.send_message(msg)
        else:
            smtp_server = smtp_config.get('smtp_server', 'smtp.gmail.com')
            smtp_port = smtp_config.get('smtp_port', 587)
            with smtplib.SMTP(smtp_server, smtp_port) as server:
                server.starttls()
                server.login(from_email, smtp_config.get('password', ''))
                server.send_message(msg)

        return jsonify({'success': True, 'message': f'Emailed {len(rows)} theses to {recipient}'})
    except Exception as e:
        print(f"Error bulk emailing analyses: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# OVERVIEW EXPORT ENDPOINTS
# ============================================

def _generate_overview_docx_bytes(row):
    """Generate docx bytes from a stock_overviews row."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0, 0, 0)

    ticker = row.get('ticker', 'Stock')
    company_name = row.get('company_name', '')
    title_text = f"{ticker} — {company_name}" if company_name else ticker
    heading = doc.add_heading(title_text, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in heading.runs:
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 0, 0)

    updated = row.get('updated_at')
    if updated:
        p = doc.add_paragraph(str(updated)[:10])
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.color.rgb = RGBColor(128, 128, 128)

    overview_sections = [
        ('Company Overview', row.get('company_overview', '')),
        ('Business Model', row.get('business_model', '')),
        ('Business Mix', row.get('business_mix', '')),
        ('Opportunities', row.get('opportunities', '')),
        ('Risks', row.get('risks', '')),
        ('Conclusion', row.get('conclusion', '')),
    ]

    for section_title, content in overview_sections:
        if not content:
            continue
        h = doc.add_heading(section_title, level=2)
        for r in h.runs: r.font.name = 'Calibri'; r.font.color.rgb = RGBColor(0, 0, 0)
        for line in content.split('\n'):
            p = doc.add_paragraph(line)
            for r in p.runs: r.font.name = 'Calibri'

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _generate_overview_pdf_bytes(row):
    """Generate PDF bytes from a stock_overviews row."""
    from xhtml2pdf import pisa

    ticker = row.get('ticker', 'Stock')
    company_name = row.get('company_name', '')
    title_text = f"{ticker} — {company_name}" if company_name else ticker
    updated = row.get('updated_at')
    date_str = str(updated)[:10] if updated else ''

    overview_sections = [
        ('Company Overview', row.get('company_overview', '')),
        ('Business Model', row.get('business_model', '')),
        ('Business Mix', row.get('business_mix', '')),
        ('Opportunities', row.get('opportunities', '')),
        ('Risks', row.get('risks', '')),
        ('Conclusion', row.get('conclusion', '')),
    ]

    sections = ''
    for section_title, content in overview_sections:
        if not content:
            continue
        esc = content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
        sections += f'<h2>{section_title}</h2><p>{esc}</p>'

    full_html = f"""<html><head><style>
        @page {{ margin: 1in; }}
        body {{ font-family: Calibri, Arial, Helvetica, sans-serif; font-size: 11pt; color: #000; line-height: 1.5; }}
        h1 {{ font-size: 18pt; margin-bottom: 4pt; color: #1a365d; border-bottom: 2px solid #2c5282; padding-bottom: 10px; }}
        h2 {{ font-size: 14pt; color: #2c5282; margin-top: 20pt; }}
        p {{ margin-bottom: 8pt; }}
        .date {{ font-size: 9pt; color: #888; margin-bottom: 16pt; }}
    </style></head><body>
        <h1>{title_text}</h1>
        <p class="date">{date_str}</p>
        {sections}
    </body></html>"""

    buf = io.BytesIO()
    pisa.CreatePDF(full_html, dest=buf)
    return buf.getvalue()


@app.route('/api/overview-to-docx', methods=['POST'])
def overview_to_docx():
    """Export a stock overview as a Word document."""
    try:
        data = request.get_json()
        ticker = data.get('ticker')
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM stock_overviews WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Overview not found'}), 404
        docx_bytes = _generate_overview_docx_bytes(dict(row))
        docx_b64 = base64.b64encode(docx_bytes).decode('utf-8')
        return jsonify({'success': True, 'fileData': docx_b64, 'filename': f"{ticker}_Overview.docx", 'fileSize': len(docx_bytes)})
    except Exception as e:
        print(f"Error creating overview docx: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/overview-to-pdf', methods=['POST'])
def overview_to_pdf():
    """Export a stock overview as a PDF."""
    try:
        data = request.get_json()
        ticker = data.get('ticker')
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM stock_overviews WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Overview not found'}), 404
        pdf_bytes = _generate_overview_pdf_bytes(dict(row))
        pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
        return jsonify({'success': True, 'fileData': pdf_b64, 'filename': f"{ticker}_Overview.pdf", 'fileSize': len(pdf_bytes)})
    except Exception as e:
        print(f"Error creating overview pdf: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/overviews-bulk-export', methods=['POST'])
def overviews_bulk_export():
    """Export multiple overviews as a zip of Word or PDF files."""
    try:
        import zipfile
        data = request.get_json()
        tickers = data.get('tickers', [])
        export_format = data.get('format', 'docx')
        if not tickers:
            return jsonify({'error': 'No tickers provided'}), 400
        with get_db() as (_, cur):
            placeholders = ','.join(['%s'] * len(tickers))
            cur.execute(f'SELECT * FROM stock_overviews WHERE ticker IN ({placeholders})', tickers)
            rows = cur.fetchall()
        if not rows:
            return jsonify({'error': 'No overviews found'}), 404
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            for row in rows:
                row_dict = dict(row)
                t = row_dict.get('ticker', 'Stock')
                fn = f"{t}_Overview.{export_format}"
                if export_format == 'pdf':
                    zf.writestr(fn, _generate_overview_pdf_bytes(row_dict))
                else:
                    zf.writestr(fn, _generate_overview_docx_bytes(row_dict))
        zip_bytes = zip_buf.getvalue()
        return jsonify({'success': True, 'fileData': base64.b64encode(zip_bytes).decode('utf-8'), 'filename': f"overviews-{len(rows)}-files.zip", 'fileSize': len(zip_bytes), 'fileCount': len(rows)})
    except Exception as e:
        print(f"Error bulk exporting overviews: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# THESIS FORMAT STUDIO ENDPOINTS
# ============================================

@app.route('/api/thesis-format/scorecard-data/<ticker>', methods=['GET'])
def get_scorecard_data(ticker):
    """Get saved scorecard data for a ticker."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT scorecard_data FROM thesis_scorecard_data WHERE ticker = %s', (ticker.upper(),))
            row = cur.fetchone()
        if not row:
            return jsonify({'scorecard_data': None})
        return jsonify({'scorecard_data': row['scorecard_data']})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/scorecard-data', methods=['POST'])
def save_scorecard_data():
    """Save scorecard data (traffic-light thresholds) for a ticker."""
    try:
        data = request.get_json()
        ticker = data.get('ticker', '').upper()
        scorecard_data = data.get('scorecard_data', {})
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400
        with get_db() as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_scorecard_data (ticker, scorecard_data, updated_at)
                VALUES (%s, %s, CURRENT_TIMESTAMP)
                ON CONFLICT (ticker) DO UPDATE SET scorecard_data = %s, updated_at = CURRENT_TIMESTAMP
            ''', (ticker, json.dumps(scorecard_data), json.dumps(scorecard_data)))
            conn.commit()
        cache.invalidate('portfolio_dashboard')
        # Auto-snapshot for thesis evolution tracking
        try:
            with get_db() as (_, cur2):
                cur2.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
                a_row = cur2.fetchone()
            a_data = a_row['analysis'] if a_row else {}
            _create_thesis_snapshot(ticker, a_data, scorecard_data, 'scorecard')
        except Exception as snap_err:
            print(f"Snapshot error: {snap_err}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


def _fmt_escape(text):
    """HTML-escape text for format templates."""
    if not text:
        return ''
    return str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def _normalize_status(status):
    """Convert numeric score (1-10) or legacy color to canonical green/yellow/red."""
    s = (status or '').strip().lower()
    if s in ('green', 'yellow', 'red'):
        return s
    try:
        n = int(s)
    except (ValueError, TypeError):
        return ''
    if n >= 7:
        return 'green'
    elif n >= 4:
        return 'yellow'
    elif n >= 1:
        return 'red'
    return ''


def _status_color(status):
    """Return CSS color for a traffic-light status."""
    s = _normalize_status(status)
    if s == 'green':
        return '#16a34a'
    elif s == 'yellow':
        return '#ca8a04'
    elif s == 'red':
        return '#dc2626'
    return '#94a3b8'


def _status_bg(status):
    """Return CSS background for a traffic-light status."""
    s = _normalize_status(status)
    if s == 'green':
        return '#dcfce7'
    elif s == 'yellow':
        return '#fef9c3'
    elif s == 'red':
        return '#fee2e2'
    return '#f1f5f9'


def _status_dot(status):
    """Return a colored circle HTML for status."""
    color = _status_color(status)
    return f'<span style="display:inline-block;width:12px;height:12px;border-radius:50%;background:{color};margin-right:4px;vertical-align:middle;"></span>'


def _parse_analysis_data(row):
    """Extract and parse analysis data from a portfolio_analyses row."""
    analysis = row.get('analysis') or {}
    if isinstance(analysis, str):
        try:
            analysis = json.loads(analysis)
        except:
            analysis = {}
    ticker = row.get('ticker', 'Stock')
    company = row.get('company', '')
    updated = row.get('updated_at')
    date_str = str(updated)[:10] if updated else ''
    thesis = analysis.get('thesis', {})
    signposts = analysis.get('signposts', [])
    threats = analysis.get('threats', [])
    conclusion = analysis.get('conclusion', '')
    is_condensed = analysis.get('_condensed', False)
    # Strip confidence from pillars when condensed
    if is_condensed:
        for p in thesis.get('pillars', []):
            p.pop('confidence', None)
    # Always strip likelihood/impact from threats
    for t in threats:
        t.pop('likelihood', None)
        t.pop('impact', None)
    return {
        'ticker': ticker, 'company': company, 'date_str': date_str,
        'thesis': thesis, 'signposts': signposts, 'threats': threats,
        'conclusion': conclusion, 'analysis': analysis,
        '_condensed': is_condensed,
    }


def _snapshot_tier_history(cur, table, analysis_col, ticker, max_versions=20):
    """Push current tier data into history array before overwriting."""
    cur.execute(f'SELECT {analysis_col}, history FROM {table} WHERE ticker = %s', (ticker,))
    row = cur.fetchone()
    if not row or not row[analysis_col]:
        return
    current = row[analysis_col]
    if isinstance(current, str):
        current = json.loads(current)
    history = row.get('history') or []
    if isinstance(history, str):
        history = json.loads(history)
    snapshot = {
        'timestamp': datetime.utcnow().isoformat(),
        'thesis': current.get('thesis'),
        'signposts': current.get('signposts'),
        'threats': current.get('threats'),
        'conclusion': current.get('conclusion', ''),
    }
    history.append(snapshot)
    history = history[-max_versions:]
    cur.execute(f'UPDATE {table} SET history = %s WHERE ticker = %s', (json.dumps(history), ticker))


def _build_pillar_data(thesis, scorecard_data):
    """Build pillar list, filtering excluded items and appending manual additions."""
    raw = thesis.get('pillars', [])
    if not scorecard_data or not scorecard_data.get('pillars'):
        return raw
    sc_pillars = scorecard_data['pillars']
    # Build lookup of scorecard pillars by title
    sc_map = {}
    for sp in sc_pillars:
        sc_map[sp.get('title', '')] = sp
    result = []
    # Walk raw pillars, skip excluded
    seen = set()
    for p in raw:
        title = p.get('pillar', p.get('title', ''))
        sc = sc_map.get(title, {})
        if sc.get('included') is False:
            continue
        result.append(p)
        seen.add(title)
    # Append manual pillars not already in raw
    for sp in sc_pillars:
        if sp.get('isManual') and sp.get('included') is not False and sp.get('title', '') not in seen:
            result.append({'pillar': sp.get('title', ''), 'detail': sp.get('description', '')})
    return result


def _build_signpost_data(signposts, scorecard_data):
    """Build enriched signpost data merging analysis signposts with scorecard data."""
    sc_signposts = {}
    if scorecard_data and scorecard_data.get('signposts'):
        for sp in scorecard_data['signposts']:
            sc_signposts[sp.get('metric', '')] = sp
    rows = []
    for sp in signposts:
        key = sp.get('metric', sp.get('signpost', ''))
        sc = sc_signposts.get(key, {})
        if sc.get('included') is False:
            continue
        rows.append({
            'metric': _fmt_escape(key),
            'target': _fmt_escape(sp.get('target', '')),
            'timeframe': _fmt_escape(sp.get('timeframe', '')),
            'ltGoal': _fmt_escape(sc.get('ltGoal', sp.get('target', ''))),
            'latest': _fmt_escape(sc.get('latest', '')),
            'status': sc.get('status', ''),
            'green': _fmt_escape(sc.get('greenThreshold', '')),
            'yellow': _fmt_escape(sc.get('yellowThreshold', '')),
            'red': _fmt_escape(sc.get('redThreshold', '')),
        })
    # Append manual signposts
    if scorecard_data and scorecard_data.get('signposts'):
        for sp in scorecard_data['signposts']:
            if sp.get('isManual') and sp.get('included') is not False:
                rows.append({
                    'metric': _fmt_escape(sp.get('metric', '')),
                    'target': '',
                    'timeframe': '',
                    'ltGoal': _fmt_escape(sp.get('ltGoal', '')),
                    'latest': _fmt_escape(sp.get('latest', '')),
                    'status': sp.get('status', ''),
                    'green': _fmt_escape(sp.get('greenThreshold', '')),
                    'yellow': _fmt_escape(sp.get('yellowThreshold', '')),
                    'red': _fmt_escape(sp.get('redThreshold', '')),
                })
    return rows


def _build_risk_data(threats, scorecard_data):
    """Build enriched risk data merging analysis threats with scorecard data."""
    sc_risks = {}
    if scorecard_data and scorecard_data.get('risks'):
        for r in scorecard_data['risks']:
            sc_risks[r.get('riskFactor', '')] = r
    rows = []
    for threat in threats:
        key = threat.get('threat', '')
        sc = sc_risks.get(key, {})
        if sc.get('included') is False:
            continue
        rows.append({
            'threat': _fmt_escape(key),
            'triggers': _fmt_escape(threat.get('triggerPoints', '')),
            'status': sc.get('status', ''),
            'statusNote': _fmt_escape(sc.get('statusNote', '')),
            'green': _fmt_escape(sc.get('greenDescription', '')),
            'yellow': _fmt_escape(sc.get('yellowDescription', '')),
            'red': _fmt_escape(sc.get('redDescription', '')),
        })
    # Append manual risks
    if scorecard_data and scorecard_data.get('risks'):
        for r in scorecard_data['risks']:
            if r.get('isManual') and r.get('included') is not False:
                rows.append({
                    'threat': _fmt_escape(r.get('riskFactor', '')),
                    'triggers': '',
                    'status': r.get('status', ''),
                    'statusNote': _fmt_escape(r.get('statusNote', '')),
                    'green': _fmt_escape(r.get('greenDescription', '')),
                    'yellow': _fmt_escape(r.get('yellowDescription', '')),
                    'red': _fmt_escape(r.get('redDescription', '')),
                })
    return rows


def _tally_statuses(scorecard_data):
    """Count green/yellow/red across signposts and risks (excluding items with included=False)."""
    g = y = r = 0
    if not scorecard_data:
        return g, y, r
    for sp in scorecard_data.get('signposts', []):
        if sp.get('included') is False:
            continue
        s = _normalize_status(sp.get('status', ''))
        if s == 'green': g += 1
        elif s == 'yellow': y += 1
        elif s == 'red': r += 1
    for rk in scorecard_data.get('risks', []):
        if rk.get('included') is False:
            continue
        s = _normalize_status(rk.get('status', ''))
        if s == 'green': g += 1
        elif s == 'yellow': y += 1
        elif s == 'red': r += 1
    return g, y, r


def _generate_executive_brief_pdf(row, scorecard_data=None):
    """Format 1: Executive Brief — clean navy professional layout. Numbered pillars, clean signpost/risk tables."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    title = f"{ticker} — {company}" if company else ticker
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    has_sc = scorecard_data and (scorecard_data.get('signposts') or scorecard_data.get('risks'))

    pillars_html = ''
    for i, p in enumerate(_build_pillar_data(thesis, scorecard_data), 1):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        pillars_html += f'<tr><td style="padding:10px 14px;border-bottom:1px solid #e2e8f0;vertical-align:top;width:32px;color:#1e3a5f;font-weight:bold;font-size:12pt;">{i}.</td><td style="padding:10px 14px;border-bottom:1px solid #e2e8f0;vertical-align:top;"><b style="color:#1e3a5f;font-size:10pt;">{t}</b><br/><span style="color:#475569;font-size:9.5pt;">{desc}</span></td></tr>'

    # Signpost table — include Latest/Status columns only when scorecard data exists
    sp_rows = ''
    for s in sp_data:
        st = s['status']
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        extra = ''
        if has_sc:
            extra = f'<td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;color:#1e293b;text-align:center;font-weight:bold;width:12%;">{s["latest"] or "—"}</td><td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;width:10%;">{st.upper() if st else "—"}</td>'
        sp_rows += f'<tr><td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;font-weight:600;color:#1e293b;width:30%;">{s["metric"]}</td><td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;color:#475569;width:30%;">{s["target"]}</td><td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;color:#475569;width:18%;">{s["timeframe"]}</td>{extra}</tr>'

    sp_header_extra = ''
    if has_sc:
        sp_header_extra = '<th style="padding:8px 12px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:12%;">Latest</th><th style="padding:8px 12px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:10%;">Status</th>'

    # Risk table — trigger points folded under risk name as subtitle, not a separate column
    risk_rows = ''
    for r in rk_data:
        st = r['status']
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        trigger_sub = f'<br/><span style="font-size:8.5pt;color:#64748b;font-weight:normal;">{r["triggers"]}</span>' if r['triggers'] else ''
        status_col = ''
        if has_sc:
            status_col = f'<td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;width:10%;">{st.upper() if st else "—"}</td>'
        risk_rows += f'<tr><td style="padding:8px 12px;border-bottom:1px solid #e2e8f0;font-weight:600;color:#1e293b;width:70%;">{r["threat"]}{trigger_sub}</td>{status_col}</tr>'

    risk_header_status = ''
    if has_sc:
        risk_header_status = '<th style="padding:8px 12px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:10%;">Status</th>'

    risk_lk_imp_headers = ''
    risk_threat_width = '70%'

    html = f"""<html><head><style>
        @page {{ margin: 0.75in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.5; }}
        table {{ border-collapse: collapse; width: 100%; }}
        h2, h3 {{ page-break-after: avoid; }}
        table {{ page-break-inside: avoid; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <p style="font-size:22pt;font-weight:bold;color:#1e293b;margin:0 0 2px 0;">{_fmt_escape(title)}</p>
        <table style="margin-bottom:20px;"><tr><td style="border-top:3px solid #1e3a5f;padding-top:6px;"><span style="font-size:9pt;color:#64748b;">Investment Thesis  |  {date_str}</span></td></tr></table>

        <table style="margin-bottom:6px;"><tr><td style="background:#1e3a5f;padding:8px 16px;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHY DO WE OWN IT?</b></td></tr></table>
        <table style="margin-bottom:4px;"><tr><td style="padding:4px 0;"><b style="color:#1e293b;font-size:10pt;">INVESTMENT THESIS</b></td></tr></table>
        <p style="color:#334155;font-size:10pt;line-height:1.7;margin:0 0 16px 12px;">{_fmt_escape(thesis.get('summary', ''))}</p>
        <table style="margin-bottom:20px;">{pillars_html}</table>

        {'<table style="margin-bottom:6px;"><tr><td style="background:#1e3a5f;padding:8px 16px;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHAT ARE WE WATCHING?</b></td></tr></table><table style="margin-bottom:20px;"><thead><tr><th style="padding:8px 12px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:30%;">Measure</th><th style="padding:8px 12px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:30%;">Target</th><th style="padding:8px 12px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:18%;">Timeframe</th>' + sp_header_extra + '</tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}

        {'<table style="margin-bottom:6px;"><tr><td style="background:#1e3a5f;padding:8px 16px;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHAT ARE THE RISKS?</b></td></tr></table><table style="margin-bottom:20px;"><thead><tr><th style="padding:8px 12px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border-bottom:2px solid #1e3a5f;width:' + risk_threat_width + ';">Risk Factor</th>' + risk_lk_imp_headers + risk_header_status + '</tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}

        {('<p style="font-size:10pt;color:#1e3a5f;font-weight:bold;border-bottom:2px solid #1e3a5f;padding-bottom:4px;margin:4px 0 8px 0;">Conclusion</p><p style="margin:0;color:#334155;font-size:10pt;line-height:1.65;">' + _fmt_escape(conclusion) + '</p>') if conclusion else ''}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_scorecard_pdf(row, scorecard_data=None):
    """Format 2: Thesis Scorecard — green headers, G/Y/R threshold columns for signposts and risks."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)

    # Pillars as compact blocks
    pillars_html = ''
    for p in _build_pillar_data(thesis, scorecard_data):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        pillars_html += f'<tr><td style="padding:8px 14px;border-bottom:1px solid #d1d5db;vertical-align:top;width:100%;"><b style="color:#1e293b;">{t}.</b> <span style="color:#475569;font-size:9.5pt;">{desc}</span></td></tr>'

    def _tl_cell(color_name, desc, current):
        """Render a G/Y/R threshold cell, highlighted if active."""
        is_active = _normalize_status(current) == color_name
        if is_active:
            return f'<td style="padding:6px 8px;border:1px solid #d1d5db;background:{_status_bg(color_name)};text-align:center;font-size:8pt;font-weight:bold;color:{_status_color(color_name)};width:14%;">{desc if desc else color_name.upper()}</td>'
        return f'<td style="padding:6px 8px;border:1px solid #d1d5db;text-align:center;font-size:8pt;color:#94a3b8;width:14%;">{desc if desc else ""}</td>'

    # Signpost rows with G/Y/R threshold columns
    sp_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        sp_rows += f'<tr><td style="padding:7px 10px;border:1px solid #d1d5db;font-weight:600;color:#1e293b;font-size:9.5pt;width:22%;">{s["metric"]}</td><td style="padding:7px 10px;border:1px solid #d1d5db;color:#475569;font-size:9pt;text-align:center;width:22%;">{s["ltGoal"]}</td><td style="padding:7px 10px;border:1px solid #d1d5db;color:#1e293b;font-size:9pt;text-align:center;font-weight:bold;width:14%;">{s["latest"] or "—"}</td>{_tl_cell("green", s["green"], st)}{_tl_cell("yellow", s["yellow"], st)}{_tl_cell("red", s["red"], st)}</tr>'

    # Risk rows with G/Y/R threshold columns
    risk_rows = ''
    for r in rk_data:
        st = r['status'] or ''
        norm_st = _normalize_status(st)
        st_display = r['statusNote'] if r['statusNote'] else (norm_st.upper() if norm_st else '—')
        st_bg = _status_bg(st) if st else '#f8fafc'
        st_color = _status_color(st) if st else '#94a3b8'
        risk_rows += f'<tr><td style="padding:7px 10px;border:1px solid #d1d5db;font-weight:600;color:#1e293b;font-size:9.5pt;width:22%;">{r["threat"]}</td><td style="padding:7px 10px;border:1px solid #d1d5db;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:9pt;width:16%;">{st_display}</td>{_tl_cell("green", r["green"], st)}{_tl_cell("yellow", r["yellow"], st)}{_tl_cell("red", r["red"], st)}</tr>'

    html = f"""<html><head><style>
        @page {{ margin: 0.65in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.5; }}
        table {{ border-collapse: collapse; width: 100%; }}
    </style></head><body>
        <p style="font-size:24pt;font-weight:bold;color:#1e293b;margin:0 0 4px 0;">Equity Thesis: {_fmt_escape(ticker)}</p>
        <p style="font-size:11pt;color:#475569;margin:0 0 4px 0;">{_fmt_escape(company)}</p>
        <table style="margin-bottom:22px;"><tr><td style="border-top:3px solid #166534;padding-top:4px;width:100%;"><span style="font-size:8pt;color:#64748b;">{date_str}</span></td></tr></table>

        <table style="margin-bottom:6px;"><tr><td style="background:#166534;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHY DO WE OWN IT?</b></td></tr></table>
        <table style="margin-bottom:4px;"><tr><td style="padding:4px 0;width:100%;"><b style="color:#1e293b;font-size:10pt;">INVESTMENT THESIS</b></td></tr></table>
        <p style="color:#334155;font-size:10pt;line-height:1.65;margin:0 0 10px 12px;">{_fmt_escape(thesis.get('summary', ''))}</p>
        <table style="margin-bottom:22px;">{pillars_html}</table>

        {'<table style="margin-bottom:8px;"><tr><td style="background:#166534;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHAT ARE WE LOOKING FOR?</b></td></tr></table><table style="border:1px solid #d1d5db;margin-bottom:22px;"><thead><tr><th style="padding:8px 10px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:22%;">Measure</th><th style="padding:8px 10px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:22%;">LT Goal</th><th style="padding:8px 10px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:14%;">Latest</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#16a34a;border:1px solid #d1d5db;width:14%;">GREEN</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#1e293b;background:#facc15;border:1px solid #d1d5db;width:14%;">YELLOW</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#dc2626;border:1px solid #d1d5db;width:14%;">RED</th></tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#166534;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHAT ARE THE RISKS?</b></td></tr></table><table style="border:1px solid #d1d5db;margin-bottom:22px;"><thead><tr><th style="padding:8px 10px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:22%;">Risk Factor</th><th style="padding:8px 10px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:16%;">Status</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#16a34a;border:1px solid #d1d5db;width:20%;">GREEN</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#1e293b;background:#facc15;border:1px solid #d1d5db;width:20%;">YELLOW</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#dc2626;border:1px solid #d1d5db;width:22%;">RED</th></tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_onepager_pdf(row, scorecard_data=None):
    """Format 3: Two-Column One-Pager — compact layout fitting everything on one page."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)

    pillars_html = ''
    for i, p in enumerate(_build_pillar_data(thesis, scorecard_data), 1):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        pillars_html += f'<tr><td style="padding:3px 0 3px 8px;border-bottom:1px solid #e2e8f0;vertical-align:top;"><b style="color:#1e3a5f;font-size:8.5pt;">{i}. {t}:</b> <span style="color:#475569;font-size:8pt;">{desc}</span></td></tr>'

    sp_rows = ''
    for s in sp_data:
        st = s['status']
        st_color = _status_color(st) if st else '#94a3b8'
        sp_rows += f'<tr><td style="padding:3px 6px;border-bottom:1px solid #e2e8f0;font-size:8pt;font-weight:600;color:#1e293b;">{s["metric"]}</td><td style="padding:3px 6px;border-bottom:1px solid #e2e8f0;font-size:8pt;color:#475569;">{s["target"]}</td><td style="padding:3px 6px;border-bottom:1px solid #e2e8f0;font-size:8pt;text-align:center;font-weight:bold;color:{st_color};">{st.upper() if st else "—"}</td></tr>'

    risk_rows = ''
    for r in rk_data:
        st = r['status']
        st_color = _status_color(st) if st else '#94a3b8'
        risk_rows += f'<tr><td style="padding:3px 6px;border-bottom:1px solid #e2e8f0;font-size:8pt;font-weight:600;color:#1e293b;">{r["threat"]}</td><td style="padding:3px 6px;border-bottom:1px solid #e2e8f0;font-size:8pt;text-align:center;font-weight:bold;color:{st_color};">{st.upper() if st else "—"}</td></tr>'

    html = f"""<html><head><style>
        @page {{ margin: 0.45in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 9pt; color: #1e293b; line-height: 1.35; }}
        table {{ border-collapse: collapse; width: 100%; }}
    </style></head><body>
        <table style="margin-bottom:8px;"><tr>
            <td style="background:#0f172a;padding:10px 16px;width:55%;"><span style="font-size:18pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span> <span style="font-size:10pt;color:#94a3b8;"> {_fmt_escape(company)}</span></td>
            <td style="background:#0f172a;padding:10px 16px;text-align:right;"><span style="font-size:8pt;color:#64748b;">INVESTMENT THESIS | {date_str}</span></td>
        </tr></table>

        <table style="border:1px solid #cbd5e1;margin-bottom:10px;"><tr><td style="padding:10px 14px;background:#f8fafc;"><b style="color:#1e293b;font-size:9pt;">THESIS:</b> <span style="color:#334155;font-size:9pt;">{_fmt_escape(thesis.get('summary', ''))}</span></td></tr></table>

        <table><tr>
            <td style="vertical-align:top;width:50%;padding-right:8px;">
                <table style="margin-bottom:4px;"><tr><td style="background:#1e3a5f;padding:5px 10px;"><b style="color:#ffffff;font-size:8pt;letter-spacing:1px;">KEY PILLARS</b></td></tr></table>
                <table style="margin-bottom:8px;">{pillars_html}</table>
                {'<table style="margin-bottom:4px;"><tr><td style="background:#1e3a5f;padding:5px 10px;"><b style="color:#ffffff;font-size:8pt;letter-spacing:1px;">SIGNPOSTS</b></td></tr></table><table style="border:1px solid #e2e8f0;"><thead><tr style="background:#f1f5f9;"><th style="padding:3px 6px;text-align:left;font-size:7pt;color:#475569;border-bottom:1px solid #cbd5e1;">METRIC</th><th style="padding:3px 6px;text-align:left;font-size:7pt;color:#475569;border-bottom:1px solid #cbd5e1;">TARGET</th><th style="padding:3px 6px;text-align:center;font-size:7pt;color:#475569;border-bottom:1px solid #cbd5e1;">STATUS</th></tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}
            </td>
            <td style="vertical-align:top;width:50%;padding-left:8px;border-left:1px solid #cbd5e1;">
                <table style="margin-bottom:4px;"><tr><td style="background:#7f1d1d;padding:5px 10px;"><b style="color:#ffffff;font-size:8pt;letter-spacing:1px;">RISKS TO THESIS</b></td></tr></table>
                <table style="border:1px solid #e2e8f0;margin-bottom:8px;"><thead><tr style="background:#fef2f2;"><th style="padding:3px 6px;text-align:left;font-size:7pt;color:#7f1d1d;border-bottom:1px solid #fecaca;">RISK</th>{'<th style="padding:3px 6px;text-align:center;font-size:7pt;color:#7f1d1d;border-bottom:1px solid #fecaca;">L/I</th>' if not is_condensed else ''}<th style="padding:3px 6px;text-align:center;font-size:7pt;color:#7f1d1d;border-bottom:1px solid #fecaca;">STATUS</th></tr></thead><tbody>{risk_rows}</tbody></table>
                {('<table style="margin-bottom:4px;"><tr><td style="background:#0f172a;padding:5px 10px;"><b style="color:#ffffff;font-size:8pt;letter-spacing:1px;">CONCLUSION</b></td></tr></table><p style="color:#475569;font-size:8.5pt;line-height:1.4;margin:0 0 0 4px;">' + _fmt_escape(conclusion) + '</p>') if conclusion else ''}
            </td>
        </tr></table>
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_board_pdf(row, scorecard_data=None):
    """Format 4: Conviction Dashboard — conviction gauge with G/Y/R counts, color-coded pillar cards, status tables."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r = _tally_statuses(scorecard_data)
    total = g + y + r
    pct = round(g / total * 100) if total > 0 else 0
    conv_color = '#16a34a' if pct >= 70 else '#ca8a04' if pct >= 40 else '#dc2626'
    conv_label = 'HIGH CONVICTION' if pct >= 70 else 'MEDIUM CONVICTION' if pct >= 40 else 'LOW CONVICTION'

    gauge_html = ''
    if total > 0:
        gauge_html = f'<table style="border:2px solid #e2e8f0;margin-bottom:18px;"><tr><td style="padding:16px 24px;text-align:center;width:40%;"><span style="font-size:40pt;font-weight:bold;color:{conv_color};">{pct}%</span><br/><span style="font-size:10pt;color:#475569;font-weight:700;letter-spacing:1px;">{conv_label}</span></td><td style="padding:16px;text-align:center;background:#dcfce7;width:20%;border-left:2px solid #e2e8f0;"><span style="font-size:24pt;font-weight:bold;color:#16a34a;">{g}</span><br/><span style="font-size:8pt;color:#166534;font-weight:700;">GREEN</span></td><td style="padding:16px;text-align:center;background:#fef9c3;width:20%;border-left:2px solid #e2e8f0;"><span style="font-size:24pt;font-weight:bold;color:#ca8a04;">{y}</span><br/><span style="font-size:8pt;color:#854d0e;font-weight:700;">YELLOW</span></td><td style="padding:16px;text-align:center;background:#fee2e2;width:20%;border-left:2px solid #e2e8f0;"><span style="font-size:24pt;font-weight:bold;color:#dc2626;">{r}</span><br/><span style="font-size:8pt;color:#991b1b;font-weight:700;">RED</span></td></tr></table>'

    pillars_html = ''
    is_condensed = d.get('_condensed', False)
    for p in _build_pillar_data(thesis, scorecard_data):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        conf = (p.get('confidence', '') or '').lower()
        bc = '#16a34a' if conf == 'high' else '#ca8a04' if conf == 'medium' else '#dc2626' if conf == 'low' else '#cbd5e1'
        pillars_html += f'<table style="margin-bottom:6px;"><tr><td style="width:5px;background:{bc};"></td><td style="padding:8px 14px;border:1px solid #e2e8f0;"><b style="color:#0f172a;font-size:10pt;">{t}</b><br/><span style="color:#475569;font-size:9pt;">{desc}</span></td></tr></table>'

    sp_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        left_border = f'border-left:4px solid {st_color};' if st else ''
        sp_rows += f'<tr style="background:{st_bg};"><td style="padding:7px 12px;border:1px solid #e2e8f0;{left_border}font-weight:600;color:#0f172a;font-size:9pt;width:35%;">{s["metric"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;color:#475569;font-size:9pt;text-align:center;width:30%;">{s["target"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;color:#0f172a;font-size:9pt;text-align:center;font-weight:bold;width:20%;">{s["latest"] or "—"}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;font-weight:bold;color:{st_color};font-size:9pt;width:15%;">{st.upper() if st else "—"}</td></tr>'

    risk_rows = ''
    for rk in rk_data:
        st = rk['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        left_border = f'border-left:4px solid {st_color};' if st else ''
        if is_condensed:
            risk_rows += f'<tr style="background:{st_bg};"><td style="padding:7px 12px;border:1px solid #e2e8f0;{left_border}font-weight:600;color:#0f172a;font-size:9pt;width:75%;">{rk["threat"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;font-weight:bold;color:{st_color};font-size:9pt;width:25%;">{st.upper() if st else "—"}</td></tr>'
        else:
            risk_rows += f'<tr style="background:{st_bg};"><td style="padding:7px 12px;border:1px solid #e2e8f0;{left_border}font-weight:600;color:#0f172a;font-size:9pt;width:35%;">{rk["threat"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;color:#475569;font-size:9pt;text-align:center;width:20%;">{rk["likelihood"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;color:#475569;font-size:9pt;text-align:center;width:20%;">{rk["impact"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;font-weight:bold;color:{st_color};font-size:9pt;width:25%;">{st.upper() if st else "—"}</td></tr>'

    html = f"""<html><head><style>
        @page {{ margin: 0.6in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.45; }}
        table {{ border-collapse: collapse; width: 100%; }}
        h2, h3 {{ page-break-after: avoid; }}
        table {{ page-break-inside: avoid; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:14px;"><tr>
            <td style="background:#0f172a;padding:18px 24px;width:60%;"><span style="font-size:28pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span><br/><span style="font-size:11pt;color:#94a3b8;">{_fmt_escape(company)}</span></td>
            <td style="background:#0f172a;padding:18px 24px;text-align:right;vertical-align:bottom;"><span style="font-size:8pt;color:#64748b;">CONVICTION DASHBOARD | {date_str}</span></td>
        </tr></table>

        {gauge_html}

        <p style="font-size:11pt;color:#0f172a;font-weight:bold;border-bottom:2px solid #0f172a;padding-bottom:4px;margin:0 0 8px 0;">INVESTMENT THESIS</p>
        <p style="color:#334155;font-size:10pt;line-height:1.65;margin:0 0 12px 0;">{_fmt_escape(thesis.get('summary', ''))}</p>
        {pillars_html}

        {'<p style="font-size:11pt;color:#0f172a;font-weight:bold;border-bottom:2px solid #0f172a;padding-bottom:4px;margin:16px 0 8px 0;">SIGNPOSTS</p><table style="border:1px solid #e2e8f0;margin-bottom:14px;"><thead><tr style="background:#f1f5f9;"><th style="padding:7px 12px;text-align:left;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">METRIC</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">TARGET</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">LATEST</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">STATUS</th></tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}

        {'<p style="font-size:11pt;color:#0f172a;font-weight:bold;border-bottom:2px solid #0f172a;padding-bottom:4px;margin:0 0 8px 0;">RISKS TO THESIS</p><table style="border:1px solid #e2e8f0;margin-bottom:14px;"><thead><tr style="background:#f1f5f9;"><th style="padding:7px 12px;text-align:left;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">RISK FACTOR</th>' + ('' if is_condensed else '<th style="padding:7px 12px;text-align:center;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">LIKELIHOOD</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">IMPACT</th>') + '<th style="padding:7px 12px;text-align:center;font-size:8pt;color:#475569;font-weight:700;border:1px solid #e2e8f0;">STATUS</th></tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_ic_memo_pdf(row, scorecard_data=None):
    """Format 5: Investment Committee Memo — formal, structured presentation for IC review."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r = _tally_statuses(scorecard_data)
    total = g + y + r

    # Status summary line
    status_line = ''
    if total > 0:
        status_line = f'<table style="border:2px solid #e2e8f0;margin-bottom:20px;"><tr><td style="padding:10px 16px;width:50%;"><b style="font-size:9pt;color:#475569;">THESIS STATUS SUMMARY</b></td><td style="padding:10px 12px;text-align:center;background:#dcfce7;border-left:1px solid #e2e8f0;"><b style="color:#16a34a;font-size:14pt;">{g}</b> <span style="font-size:8pt;color:#166534;">GREEN</span></td><td style="padding:10px 12px;text-align:center;background:#fef9c3;border-left:1px solid #e2e8f0;"><b style="color:#ca8a04;font-size:14pt;">{y}</b> <span style="font-size:8pt;color:#854d0e;">YELLOW</span></td><td style="padding:10px 12px;text-align:center;background:#fee2e2;border-left:1px solid #e2e8f0;"><b style="color:#dc2626;font-size:14pt;">{r}</b> <span style="font-size:8pt;color:#991b1b;">RED</span></td></tr></table>'

    pillars_html = ''
    for i, p in enumerate(_build_pillar_data(thesis, scorecard_data), 1):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        conf = p.get('confidence', '')
        conf_html = f' <span style="color:#64748b;font-size:8pt;">({_fmt_escape(conf)} confidence)</span>' if conf else ''
        pillars_html += f'<tr><td style="padding:8px 14px;border-bottom:1px solid #e2e8f0;vertical-align:top;"><b style="color:#1e293b;">Pillar {i}: {t}</b>{conf_html}<br/><span style="color:#475569;font-size:9.5pt;">{desc}</span></td></tr>'

    has_sc = scorecard_data and (scorecard_data.get('signposts') or scorecard_data.get('risks'))

    # Column widths that always sum to 100%
    sp_measure_w = '43%' if has_sc else '60%'
    sp_goal_w = '27%' if has_sc else '40%'

    sp_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        st_color = _status_color(st) if st else '#94a3b8'
        extra = ''
        if has_sc:
            extra = f'<td style="padding:7px 12px;border-bottom:1px solid #e2e8f0;color:#1e293b;text-align:center;font-weight:bold;width:15%;">{s["latest"] or "—"}</td><td style="padding:7px 12px;border-bottom:1px solid #e2e8f0;text-align:center;font-weight:bold;color:{st_color};width:15%;">{st.upper() if st else "—"}</td>'
        sp_rows += f'<tr><td style="padding:7px 12px;border-bottom:1px solid #e2e8f0;font-weight:600;color:#1e293b;width:{sp_measure_w};">{s["metric"]}</td><td style="padding:7px 12px;border-bottom:1px solid #e2e8f0;color:#475569;width:{sp_goal_w};">{s["ltGoal"]}</td>{extra}</tr>'

    sp_hdr_extra = ''
    if has_sc:
        sp_hdr_extra = '<th style="padding:7px 12px;text-align:center;font-size:9pt;color:#475569;font-weight:700;border-bottom:2px solid #475569;width:15%;">Latest</th><th style="padding:7px 12px;text-align:center;font-size:9pt;color:#475569;font-weight:700;border-bottom:2px solid #475569;width:15%;">Status</th>'

    # Risk column widths
    if has_sc:
        rk_threat_w = '80%'
        rk_status_w = '20%'
    else:
        rk_threat_w = '100%'
        rk_status_w = '0%'

    risk_rows = ''
    for rk in rk_data:
        st = rk['status'] or ''
        st_color = _status_color(st) if st else '#94a3b8'
        trigger_sub = f'<br/><span style="font-size:8.5pt;color:#64748b;font-weight:normal;">{rk["triggers"]}</span>' if rk['triggers'] else ''
        status_col = ''
        if has_sc:
            status_col = f'<td style="padding:7px 12px;border-bottom:1px solid #e2e8f0;text-align:center;font-weight:bold;color:{st_color};width:{rk_status_w};">{st.upper() if st else "—"}</td>'
        risk_rows += f'<tr><td style="padding:7px 12px;border-bottom:1px solid #e2e8f0;font-weight:600;color:#1e293b;width:{rk_threat_w};">{rk["threat"]}{trigger_sub}</td>{status_col}</tr>'

    risk_hdr_status = ''
    if has_sc:
        risk_hdr_status = f'<th style="padding:7px 12px;text-align:center;font-size:9pt;color:#475569;font-weight:700;border-bottom:2px solid #475569;width:{rk_status_w};">Status</th>'
    risk_lk_imp_headers_ic = ''
    risk_threat_width_ic = rk_threat_w

    html = f"""<html><head><style>
        @page {{ margin: 0.75in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.5; }}
        table {{ border-collapse: collapse; width: 100%; }}
        h2, h3 {{ page-break-after: avoid; }}
        table {{ page-break-inside: avoid; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:20px;"><tr><td style="padding:0 0 8px 0;border-bottom:3px solid #1e293b;width:100%;"><span style="font-size:22pt;font-weight:bold;color:#0f172a;">{_fmt_escape(ticker)}</span><span style="font-size:14pt;color:#475569;"> — {_fmt_escape(company)}</span><br/><span style="font-size:9pt;color:#94a3b8;">{date_str}</span></td></tr></table>

        {status_line}

        <p style="font-size:11pt;font-weight:bold;color:#1e293b;margin:0 0 6px 0;border-bottom:1px solid #cbd5e1;padding-bottom:4px;">I. INVESTMENT THESIS</p>
        <p style="color:#334155;font-size:10pt;line-height:1.7;margin:0 0 14px 0;">{_fmt_escape(thesis.get('summary', ''))}</p>
        <table style="margin-bottom:18px;">{pillars_html}</table>

        {'<p style="font-size:11pt;font-weight:bold;color:#1e293b;margin:0 0 6px 0;border-bottom:1px solid #cbd5e1;padding-bottom:4px;">II. KEY SIGNPOSTS</p><table style="margin-bottom:18px;"><thead><tr><th style="padding:7px 12px;text-align:left;font-size:9pt;color:#475569;font-weight:700;border-bottom:2px solid #475569;width:' + sp_measure_w + ';">Measure</th><th style="padding:7px 12px;text-align:left;font-size:9pt;color:#475569;font-weight:700;border-bottom:2px solid #475569;width:' + sp_goal_w + ';">LT Goal</th>' + sp_hdr_extra + '</tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}

        {'<p style="font-size:11pt;font-weight:bold;color:#1e293b;margin:0 0 6px 0;border-bottom:1px solid #cbd5e1;padding-bottom:4px;">III. RISK ASSESSMENT</p><table style="margin-bottom:18px;"><thead><tr><th style="padding:7px 12px;text-align:left;font-size:9pt;color:#475569;font-weight:700;border-bottom:2px solid #475569;width:' + risk_threat_width_ic + ';">Risk Factor</th>' + risk_lk_imp_headers_ic + risk_hdr_status + '</tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}

        {('<p style="font-size:11pt;font-weight:bold;color:#1e293b;margin:0 0 6px 0;border-bottom:1px solid #cbd5e1;padding-bottom:4px;">IV. CONCLUSION</p><p style="color:#334155;font-size:10pt;line-height:1.65;margin:0;">' + _fmt_escape(conclusion) + '</p>') if conclusion else ''}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_risk_focus_pdf(row, scorecard_data=None):
    """Format 6: Risk-Focused Report — emphasizes risks and signpost monitoring with G/Y/R thresholds."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r_count = _tally_statuses(scorecard_data)
    total = g + y + r_count

    # Summary bar
    summary_html = ''
    if total > 0:
        pct_g = round(g / total * 100)
        pct_y = round(y / total * 100)
        pct_r = round(r_count / total * 100)
        summary_html = f'<table style="border:1px solid #d1d5db;margin-bottom:16px;"><tr><td style="padding:8px 14px;width:30%;"><b style="font-size:9pt;color:#475569;">STATUS OVERVIEW</b><br/><span style="font-size:8pt;color:#94a3b8;">{total} monitored items</span></td><td style="padding:8px 14px;text-align:center;width:23%;background:#dcfce7;border-left:1px solid #d1d5db;"><b style="color:#16a34a;font-size:16pt;">{g}</b><span style="font-size:8pt;color:#166534;"> ({pct_g}%)</span></td><td style="padding:8px 14px;text-align:center;width:23%;background:#fef9c3;border-left:1px solid #d1d5db;"><b style="color:#ca8a04;font-size:16pt;">{y}</b><span style="font-size:8pt;color:#854d0e;"> ({pct_y}%)</span></td><td style="padding:8px 14px;text-align:center;width:24%;background:#fee2e2;border-left:1px solid #d1d5db;"><b style="color:#dc2626;font-size:16pt;">{r_count}</b><span style="font-size:8pt;color:#991b1b;"> ({pct_r}%)</span></td></tr></table>'

    # Thesis summary (brief)
    thesis_brief = f'<table style="border:1px solid #cbd5e1;margin-bottom:16px;"><tr><td style="padding:10px 14px;background:#f8fafc;"><b style="color:#334155;font-size:9pt;">THESIS SUMMARY:</b> <span style="color:#475569;font-size:9pt;">{_fmt_escape(thesis.get("summary", ""))}</span></td></tr></table>'

    def _tl_cell_rf(color_name, desc, current):
        """Render a G/Y/R threshold cell for risk focus template."""
        is_active = _normalize_status(current) == color_name
        if is_active:
            return f'<td style="padding:6px 8px;border:1px solid #d1d5db;background:{_status_bg(color_name)};text-align:center;font-size:8pt;font-weight:bold;color:{_status_color(color_name)};width:14%;">{desc if desc else color_name.upper()}</td>'
        return f'<td style="padding:6px 8px;border:1px solid #d1d5db;text-align:center;font-size:8pt;color:#94a3b8;width:14%;">{desc if desc else ""}</td>'

    # RISKS first (this template emphasizes risks)
    risk_rows = ''
    for rk in rk_data:
        st = rk['status'] or ''
        norm_st = _normalize_status(st)
        st_display = rk['statusNote'] if rk['statusNote'] else (norm_st.upper() if norm_st else '—')
        st_bg = _status_bg(st) if st else '#f8fafc'
        st_color = _status_color(st) if st else '#94a3b8'
        risk_rows += f'<tr><td style="padding:7px 10px;border:1px solid #d1d5db;font-weight:600;color:#1e293b;font-size:9.5pt;width:22%;">{rk["threat"]}</td><td style="padding:7px 10px;border:1px solid #d1d5db;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:9pt;width:16%;">{st_display}</td>{_tl_cell_rf("green", rk["green"], st)}{_tl_cell_rf("yellow", rk["yellow"], st)}{_tl_cell_rf("red", rk["red"], st)}</tr>'

    # Signposts with thresholds
    sp_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        sp_rows += f'<tr><td style="padding:7px 10px;border:1px solid #d1d5db;font-weight:600;color:#1e293b;font-size:9.5pt;width:22%;">{s["metric"]}</td><td style="padding:7px 10px;border:1px solid #d1d5db;color:#475569;font-size:9pt;text-align:center;width:22%;">{s["ltGoal"]}</td><td style="padding:7px 10px;border:1px solid #d1d5db;color:#1e293b;font-size:9pt;text-align:center;font-weight:bold;width:14%;">{s["latest"] or "—"}</td>{_tl_cell_rf("green", s["green"], st)}{_tl_cell_rf("yellow", s["yellow"], st)}{_tl_cell_rf("red", s["red"], st)}</tr>'

    html = f"""<html><head><style>
        @page {{ margin: 0.65in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.5; }}
        table {{ border-collapse: collapse; width: 100%; }}
    </style></head><body>
        <p style="font-size:22pt;font-weight:bold;color:#1e293b;margin:0 0 2px 0;">Risk Monitor: {_fmt_escape(ticker)}</p>
        <p style="font-size:11pt;color:#475569;margin:0 0 4px 0;">{_fmt_escape(company)}</p>
        <table style="margin-bottom:20px;"><tr><td style="border-top:3px solid #7f1d1d;padding-top:4px;width:100%;"><span style="font-size:8pt;color:#64748b;">{date_str}</span></td></tr></table>

        {summary_html}
        {thesis_brief}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#7f1d1d;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">RISK ASSESSMENT</b></td></tr></table><table style="border:1px solid #d1d5db;margin-bottom:22px;"><thead><tr><th style="padding:8px 10px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:22%;">Risk Factor</th><th style="padding:8px 10px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:16%;">Status</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#16a34a;border:1px solid #d1d5db;width:20%;">GREEN</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#1e293b;background:#facc15;border:1px solid #d1d5db;width:20%;">YELLOW</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#dc2626;border:1px solid #d1d5db;width:22%;">RED</th></tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#1e3a5f;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">SIGNPOST MONITORING</b></td></tr></table><table style="border:1px solid #d1d5db;margin-bottom:20px;"><thead><tr><th style="padding:8px 10px;text-align:left;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:22%;">Measure</th><th style="padding:8px 10px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:22%;">LT Goal</th><th style="padding:8px 10px;text-align:center;font-size:9pt;color:#1e293b;font-weight:bold;border:1px solid #d1d5db;width:14%;">Latest</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#16a34a;border:1px solid #d1d5db;width:14%;">GREEN</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#1e293b;background:#facc15;border:1px solid #d1d5db;width:14%;">YELLOW</th><th style="padding:8px 10px;text-align:center;font-size:9pt;font-weight:bold;color:#ffffff;background:#dc2626;border:1px solid #d1d5db;width:14%;">RED</th></tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_catalyst_timeline_pdf(row, scorecard_data=None):
    """Catalyst Timeline — visual timeline of upcoming catalysts derived from signpost timeframes."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)

    # Thesis summary box
    thesis_html = f'<table style="border:1px solid #99f6e4;margin-bottom:18px;"><tr><td style="padding:10px 16px;background:#f0fdfa;"><b style="color:#0d9488;font-size:9pt;">INVESTMENT THESIS:</b> <span style="color:#334155;font-size:9.5pt;">{_fmt_escape(thesis.get("summary", ""))}</span></td></tr></table>'

    # Sort signposts by timeframe for timeline grouping
    tf_order = {'near-term': 0, 'q1': 1, 'q2': 2, 'q3': 3, 'q4': 4, 'h1': 5, 'h2': 6, '1 year': 7, '2 year': 8, '3 year': 9, 'long-term': 10}
    def _tf_sort_key(s):
        tf_lower = (s.get('timeframe') or '').strip().lower()
        for k, v in tf_order.items():
            if k in tf_lower:
                return v
        return 50
    sorted_sp = sorted(sp_data, key=_tf_sort_key)

    # Build timeline entries
    timeline_rows = ''
    for s in sorted_sp:
        st = s['status'] or ''
        st_bg = _status_bg(st) if st else '#f0fdfa'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = norm_st.upper() if norm_st else '—'
        tf_display = s['timeframe'] if s['timeframe'] else '—'
        latest_display = s['latest'] if s['latest'] else '—'
        timeline_rows += f'<tr style="page-break-inside:avoid;"><td style="padding:10px 14px;border:1px solid #e2e8f0;background:#f0fdfa;text-align:center;vertical-align:top;width:22%;"><b style="color:#0d9488;font-size:9.5pt;">{tf_display}</b></td><td style="padding:10px 14px;border:1px solid #e2e8f0;width:50%;"><b style="color:#0f172a;font-size:9.5pt;">{s["metric"]}</b><br/><span style="color:#64748b;font-size:8.5pt;">Target: {s["target"]}</span></td><td style="padding:10px 14px;border:1px solid #e2e8f0;text-align:center;width:13%;"><span style="color:#0f172a;font-weight:bold;font-size:9pt;">{latest_display}</span></td><td style="padding:10px 14px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};width:15%;"><b style="color:{st_color};font-size:8.5pt;">{st_label}</b></td></tr>'

    # Risk watchlist — compact
    risk_rows = ''
    for rk in rk_data:
        st = rk['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = rk['statusNote'] if rk['statusNote'] else (norm_st.upper() if norm_st else '—')
        risk_rows += f'<tr><td style="padding:6px 12px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:60%;">{rk["threat"]}</td><td style="padding:6px 12px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:20%;">{st_label}</td><td style="padding:6px 12px;border:1px solid #e2e8f0;color:#64748b;font-size:8pt;width:20%;">{rk["triggers"][:50] + "..." if len(rk["triggers"]) > 50 else rk["triggers"]}</td></tr>'

    html = f"""<html><head><style>
        @page {{ margin: 0.65in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.45; }}
        table {{ border-collapse: collapse; width: 100%; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:14px;"><tr>
            <td style="background:#0d9488;padding:18px 24px;width:60%;"><span style="font-size:26pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span><br/><span style="font-size:11pt;color:#ccfbf1;">{_fmt_escape(company)}</span></td>
            <td style="background:#0d9488;padding:18px 24px;text-align:right;vertical-align:bottom;width:40%;"><span style="font-size:8pt;color:#ccfbf1;">CATALYST TIMELINE | {date_str}</span></td>
        </tr></table>

        {thesis_html}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#134e4a;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">CATALYST TIMELINE</b></td></tr></table><table style="border:1px solid #e2e8f0;margin-bottom:20px;"><thead><tr style="background:#f0fdfa;"><th style="padding:8px 14px;text-align:center;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:22%;">TIMEFRAME</th><th style="padding:8px 14px;text-align:left;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:50%;">CATALYST / METRIC</th><th style="padding:8px 14px;text-align:center;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:13%;">LATEST</th><th style="padding:8px 14px;text-align:center;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:15%;">STATUS</th></tr></thead><tbody>' + timeline_rows + '</tbody></table>' if signposts else ''}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#134e4a;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">RISK WATCHLIST</b></td></tr></table><table style="border:1px solid #e2e8f0;margin-bottom:18px;"><thead><tr style="background:#f0fdfa;"><th style="padding:6px 12px;text-align:left;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:60%;">RISK FACTOR</th><th style="padding:6px 12px;text-align:center;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:20%;">STATUS</th><th style="padding:6px 12px;text-align:left;font-size:8pt;color:#0d9488;font-weight:700;border:1px solid #e2e8f0;width:20%;">TRIGGER</th></tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_valuation_summary_pdf(row, scorecard_data=None):
    """Valuation Summary — valuation-focused layout presenting thesis from a value perspective."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)

    # Investment Case box
    case_html = f'<table style="border:1px solid #a7f3d0;margin-bottom:18px;"><tr><td style="padding:12px 16px;background:#ecfdf5;"><b style="color:#064e3b;font-size:10pt;">INVESTMENT CASE</b><br/><span style="color:#334155;font-size:9.5pt;line-height:1.6;">{_fmt_escape(thesis.get("summary", ""))}</span></td></tr></table>'

    # Value Drivers — numbered pillars
    pillar_rows = ''
    for i, p in enumerate(_build_pillar_data(thesis, scorecard_data), 1):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        conf = (p.get('confidence', '') or '').lower()
        conf_color = '#16a34a' if conf == 'high' else '#ca8a04' if conf == 'medium' else '#dc2626' if conf == 'low' else '#64748b'
        conf_label = conf.upper() if conf else ''
        conf_html = f'<span style="font-size:7.5pt;color:{conf_color};font-weight:bold;"> [{conf_label}]</span>' if conf_label else ''
        pillar_rows += f'<tr><td style="padding:8px 12px;border-bottom:1px solid #d1fae5;vertical-align:top;width:30px;color:#064e3b;font-weight:bold;font-size:12pt;">{i}.</td><td style="padding:8px 12px;border-bottom:1px solid #d1fae5;vertical-align:top;width:100%;"><b style="color:#064e3b;font-size:10pt;">{t}</b>{conf_html}<br/><span style="color:#475569;font-size:9pt;">{desc}</span></td></tr>'

    # Valuation Checkpoints — signpost table
    sp_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = norm_st.upper() if norm_st else '—'
        sp_rows += f'<tr><td style="padding:7px 12px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:30%;">{s["metric"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;color:#475569;font-size:9pt;text-align:center;width:25%;">{s["target"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;color:#0f172a;font-size:9pt;text-align:center;font-weight:bold;width:20%;">{s["latest"] or "—"}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:15%;">{st_label}</td></tr>'

    # Valuation Risks — compact
    risk_rows = ''
    for rk in rk_data:
        st = rk['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = rk['statusNote'] if rk['statusNote'] else (norm_st.upper() if norm_st else '—')
        trigger_sub = f'<br/><span style="font-size:8pt;color:#64748b;font-weight:normal;">{rk["triggers"]}</span>' if rk['triggers'] else ''
        risk_rows += f'<tr><td style="padding:7px 12px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:70%;">{rk["threat"]}{trigger_sub}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:30%;">{st_label}</td></tr>'

    # Valuation Assessment — conclusion
    conclusion_html = ''
    if conclusion:
        conclusion_html = f'<table style="border:1px solid #a7f3d0;margin-top:4px;"><tr><td style="padding:12px 16px;background:#ecfdf5;"><b style="color:#064e3b;font-size:10pt;">VALUATION ASSESSMENT</b><br/><span style="color:#334155;font-size:9.5pt;line-height:1.6;">{_fmt_escape(conclusion)}</span></td></tr></table>'

    html = f"""<html><head><style>
        @page {{ margin: 0.65in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.45; }}
        table {{ border-collapse: collapse; width: 100%; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:14px;"><tr>
            <td style="background:#064e3b;padding:18px 24px;width:60%;"><span style="font-size:26pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span><br/><span style="font-size:11pt;color:#a7f3d0;">{_fmt_escape(company)}</span></td>
            <td style="background:#064e3b;padding:18px 24px;text-align:right;vertical-align:bottom;width:40%;"><span style="font-size:8pt;color:#a7f3d0;">VALUATION FRAMEWORK | {date_str}</span></td>
        </tr></table>

        {case_html}

        <table style="margin-bottom:8px;"><tr><td style="background:#065f46;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">VALUE DRIVERS</b></td></tr></table>
        <table style="margin-bottom:18px;">{pillar_rows}</table>

        {'<table style="margin-bottom:8px;"><tr><td style="background:#065f46;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">VALUATION CHECKPOINTS</b></td></tr></table><table style="border:1px solid #e2e8f0;margin-bottom:18px;"><thead><tr style="background:#ecfdf5;"><th style="padding:7px 12px;text-align:left;font-size:8pt;color:#064e3b;font-weight:700;border:1px solid #e2e8f0;width:30%;">METRIC</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#064e3b;font-weight:700;border:1px solid #e2e8f0;width:25%;">TARGET</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#064e3b;font-weight:700;border:1px solid #e2e8f0;width:20%;">LATEST</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#064e3b;font-weight:700;border:1px solid #e2e8f0;width:15%;">STATUS</th></tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#065f46;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">VALUATION RISKS</b></td></tr></table><table style="border:1px solid #e2e8f0;margin-bottom:18px;"><thead><tr style="background:#ecfdf5;"><th style="padding:7px 12px;text-align:left;font-size:8pt;color:#064e3b;font-weight:700;border:1px solid #e2e8f0;width:70%;">RISK FACTOR</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#064e3b;font-weight:700;border:1px solid #e2e8f0;width:30%;">STATUS</th></tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}

        {conclusion_html}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_bull_bear_pdf(row, scorecard_data=None):
    """Bull/Bear Framework — side-by-side bull case vs bear case presentation."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r_count = _tally_statuses(scorecard_data)
    total = g + y + r_count

    # Build bull case items from pillars + green signposts
    bull_items = ''
    for p in _build_pillar_data(thesis, scorecard_data):
        t = _fmt_escape(p.get('pillar', p.get('title', '')))
        desc = _fmt_escape(p.get('detail', p.get('description', '')))
        bull_items += f'<tr><td style="padding:6px 10px;border-bottom:1px solid #bbf7d0;vertical-align:top;width:100%;"><b style="color:#166534;font-size:9pt;">{t}</b><br/><span style="color:#334155;font-size:8.5pt;">{desc}</span></td></tr>'

    # Add green-status signposts to bull side
    for s in sp_data:
        st = s['status'] or ''
        norm_st = _normalize_status(st)
        if norm_st == 'green':
            bull_items += f'<tr><td style="padding:5px 10px;border-bottom:1px solid #bbf7d0;vertical-align:top;width:100%;background:#f0fdf4;"><span style="color:#16a34a;font-size:8pt;font-weight:bold;">SIGNPOST:</span> <span style="color:#166534;font-size:8.5pt;font-weight:600;">{s["metric"]}</span><br/><span style="color:#475569;font-size:8pt;">Latest: {s["latest"] or "—"} | Target: {s["target"]}</span></td></tr>'

    # Build bear case items from risks + red/yellow signposts
    bear_items = ''
    for rk in rk_data:
        st = rk['status'] or ''
        st_color = _status_color(st) if st else '#991b1b'
        bear_items += f'<tr><td style="padding:6px 10px;border-bottom:1px solid #fecaca;vertical-align:top;width:100%;"><b style="color:#991b1b;font-size:9pt;">{rk["threat"]}</b><br/><span style="color:#334155;font-size:8.5pt;">{rk["triggers"]}</span></td></tr>'

    # Add red/yellow signposts to bear side
    for s in sp_data:
        st = s['status'] or ''
        norm_st = _normalize_status(st)
        if norm_st in ('red', 'yellow'):
            label_color = '#dc2626' if norm_st == 'red' else '#ca8a04'
            bg = '#fef2f2' if norm_st == 'red' else '#fefce8'
            bear_items += f'<tr><td style="padding:5px 10px;border-bottom:1px solid #fecaca;vertical-align:top;width:100%;background:{bg};"><span style="color:{label_color};font-size:8pt;font-weight:bold;">SIGNPOST:</span> <span style="color:#991b1b;font-size:8.5pt;font-weight:600;">{s["metric"]}</span><br/><span style="color:#475569;font-size:8pt;">Latest: {s["latest"] or "—"} | Target: {s["target"]}</span></td></tr>'

    # Conviction verdict
    pct = round(g / total * 100) if total > 0 else 0
    conv_color = '#16a34a' if pct >= 70 else '#ca8a04' if pct >= 40 else '#dc2626'
    conv_label = 'HIGH CONVICTION' if pct >= 70 else 'MEDIUM CONVICTION' if pct >= 40 else 'LOW CONVICTION'

    verdict_html = ''
    if conclusion:
        gauge_part = ''
        if total > 0:
            gauge_part = f'<table style="border:1px solid #e2e8f0;margin-bottom:10px;"><tr><td style="padding:12px 20px;text-align:center;width:30%;"><span style="font-size:32pt;font-weight:bold;color:{conv_color};">{pct}%</span><br/><span style="font-size:9pt;color:#475569;font-weight:700;letter-spacing:1px;">{conv_label}</span></td><td style="padding:12px 20px;text-align:center;background:#dcfce7;width:23%;border-left:1px solid #e2e8f0;"><b style="color:#16a34a;font-size:18pt;">{g}</b><br/><span style="font-size:8pt;color:#166534;font-weight:700;">GREEN</span></td><td style="padding:12px 20px;text-align:center;background:#fef9c3;width:23%;border-left:1px solid #e2e8f0;"><b style="color:#ca8a04;font-size:18pt;">{y}</b><br/><span style="font-size:8pt;color:#854d0e;font-weight:700;">YELLOW</span></td><td style="padding:12px 20px;text-align:center;background:#fee2e2;width:24%;border-left:1px solid #e2e8f0;"><b style="color:#dc2626;font-size:18pt;">{r_count}</b><br/><span style="font-size:8pt;color:#991b1b;font-weight:700;">RED</span></td></tr></table>'
        verdict_html = f'{gauge_part}<table style="border:1px solid #e2e8f0;"><tr><td style="padding:10px 16px;background:#f8fafc;"><b style="color:#0f172a;font-size:10pt;">VERDICT</b><br/><span style="color:#334155;font-size:9.5pt;line-height:1.6;">{_fmt_escape(conclusion)}</span></td></tr></table>'

    html = f"""<html><head><style>
        @page {{ margin: 0.6in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.45; }}
        table {{ border-collapse: collapse; width: 100%; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:14px;"><tr>
            <td style="background:#0f172a;padding:18px 24px;width:100%;text-align:center;"><span style="font-size:26pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span><span style="font-size:11pt;color:#94a3b8;"> &mdash; Bull/Bear Framework</span><br/><span style="font-size:8pt;color:#64748b;">{_fmt_escape(company)} | {date_str}</span></td>
        </tr></table>

        <table style="margin-bottom:20px;"><tr>
            <td style="vertical-align:top;width:49%;padding-right:4px;">
                <table><tr><td style="background:#166534;padding:8px 14px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">BULL CASE</b></td></tr></table>
                <table style="border:1px solid #bbf7d0;">{bull_items if bull_items else '<tr><td style="padding:10px;color:#64748b;font-size:9pt;width:100%;">No bull case items</td></tr>'}</table>
            </td>
            <td style="vertical-align:top;width:2%;"></td>
            <td style="vertical-align:top;width:49%;padding-left:4px;">
                <table><tr><td style="background:#991b1b;padding:8px 14px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">BEAR CASE</b></td></tr></table>
                <table style="border:1px solid #fecaca;">{bear_items if bear_items else '<tr><td style="padding:10px;color:#64748b;font-size:9pt;width:100%;">No bear case items</td></tr>'}</table>
            </td>
        </tr></table>

        {verdict_html}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_portfolio_context_pdf(row, scorecard_data=None):
    """Portfolio Context — position-level view for portfolio review."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r_count = _tally_statuses(scorecard_data)
    total = g + y + r_count

    # Thesis Health summary bar
    health_html = ''
    if total > 0:
        pct_g = round(g / total * 100)
        pct_y = round(y / total * 100)
        pct_r = round(r_count / total * 100)
        overall = 'HEALTHY' if pct_g >= 60 else 'MIXED' if pct_g >= 30 else 'AT RISK'
        overall_color = '#16a34a' if pct_g >= 60 else '#ca8a04' if pct_g >= 30 else '#dc2626'
        health_html = f'<table style="border:2px solid #e2e8f0;margin-bottom:16px;"><tr><td style="padding:10px 16px;width:34%;border-right:1px solid #e2e8f0;"><b style="font-size:9pt;color:#475569;">THESIS HEALTH</b><br/><span style="font-size:16pt;font-weight:bold;color:{overall_color};">{overall}</span></td><td style="padding:10px 14px;text-align:center;width:22%;background:#dcfce7;border-right:1px solid #e2e8f0;"><b style="color:#16a34a;font-size:18pt;">{g}</b><br/><span style="font-size:8pt;color:#166534;font-weight:700;">GREEN ({pct_g}%)</span></td><td style="padding:10px 14px;text-align:center;width:22%;background:#fef9c3;border-right:1px solid #e2e8f0;"><b style="color:#ca8a04;font-size:18pt;">{y}</b><br/><span style="font-size:8pt;color:#854d0e;font-weight:700;">YELLOW ({pct_y}%)</span></td><td style="padding:10px 14px;text-align:center;width:22%;background:#fee2e2;"><b style="color:#dc2626;font-size:18pt;">{r_count}</b><br/><span style="font-size:8pt;color:#991b1b;font-weight:700;">RED ({pct_r}%)</span></td></tr></table>'

    # Compact thesis
    thesis_html = f'<table style="border:1px solid #cbd5e1;margin-bottom:16px;"><tr><td style="padding:10px 14px;background:#f8fafc;width:100%;"><b style="color:#1e293b;font-size:9pt;">THESIS:</b> <span style="color:#475569;font-size:9pt;">{_fmt_escape(thesis.get("summary", ""))}</span></td></tr></table>'

    # Combined monitoring table — signposts + risks in one table
    monitor_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = norm_st.upper() if norm_st else '—'
        monitor_rows += f'<tr><td style="padding:7px 12px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:50%;">{s["metric"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;color:#0ea5e9;font-size:8pt;font-weight:700;width:20%;">SIGNPOST</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:30%;">{st_label}</td></tr>'

    for rk in rk_data:
        st = rk['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = rk['statusNote'] if rk['statusNote'] else (norm_st.upper() if norm_st else '—')
        monitor_rows += f'<tr><td style="padding:7px 12px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:50%;">{rk["threat"]}</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;color:#ef4444;font-size:8pt;font-weight:700;width:20%;">RISK</td><td style="padding:7px 12px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:30%;">{st_label}</td></tr>'

    # Position Outlook
    outlook_html = ''
    if conclusion:
        outlook_html = f'<table style="border:1px solid #bae6fd;margin-top:4px;"><tr><td style="padding:12px 16px;background:#f0f9ff;width:100%;"><b style="color:#1e293b;font-size:10pt;">POSITION OUTLOOK</b><br/><span style="color:#334155;font-size:9.5pt;line-height:1.6;">{_fmt_escape(conclusion)}</span></td></tr></table>'

    html = f"""<html><head><style>
        @page {{ margin: 0.65in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.45; }}
        table {{ border-collapse: collapse; width: 100%; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:14px;"><tr>
            <td style="background:#1e293b;padding:18px 24px;width:60%;"><span style="font-size:26pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span><br/><span style="font-size:11pt;color:#94a3b8;">{_fmt_escape(company)}</span></td>
            <td style="background:#1e293b;padding:18px 24px;text-align:right;vertical-align:bottom;width:40%;"><span style="font-size:8pt;color:#94a3b8;">POSITION REVIEW | {date_str}</span></td>
        </tr></table>

        {health_html}
        {thesis_html}

        <table style="margin-bottom:8px;"><tr><td style="background:#0c4a6e;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">KEY MONITORING POINTS</b></td></tr></table>
        <table style="border:1px solid #e2e8f0;margin-bottom:18px;"><thead><tr style="background:#f0f9ff;"><th style="padding:7px 12px;text-align:left;font-size:8pt;color:#1e293b;font-weight:700;border:1px solid #e2e8f0;width:50%;">ITEM</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#1e293b;font-weight:700;border:1px solid #e2e8f0;width:20%;">TYPE</th><th style="padding:7px 12px;text-align:center;font-size:8pt;color:#1e293b;font-weight:700;border:1px solid #e2e8f0;width:30%;">STATUS</th></tr></thead><tbody>{monitor_rows}</tbody></table>

        {outlook_html}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


def _generate_earnings_snapshot_pdf(row, scorecard_data=None):
    """Earnings Snapshot — quick-reference card for earnings season."""
    from xhtml2pdf import pisa
    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r_count = _tally_statuses(scorecard_data)
    total = g + y + r_count
    pct = round(g / total * 100) if total > 0 else 0
    pillar_count = len(_build_pillar_data(thesis, scorecard_data))

    # Compact thesis — truncated
    summary_text = _fmt_escape(thesis.get('summary', ''))
    if len(summary_text) > 300:
        summary_text = summary_text[:297] + '...'
    thesis_html = f'<table style="border:1px solid #c4b5fd;margin-bottom:16px;"><tr><td style="padding:8px 14px;background:#f5f3ff;width:100%;"><b style="color:#4c1d95;font-size:9pt;">THESIS:</b> <span style="color:#334155;font-size:9pt;">{summary_text}</span></td></tr></table>'

    # What To Watch — signpost table (compact)
    sp_rows = ''
    for s in sp_data:
        st = s['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = norm_st.upper() if norm_st else '—'
        sp_rows += f'<tr><td style="padding:6px 10px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:30%;">{s["metric"]}</td><td style="padding:6px 10px;border:1px solid #e2e8f0;color:#475569;font-size:8.5pt;text-align:center;width:25%;">{s["target"]}</td><td style="padding:6px 10px;border:1px solid #e2e8f0;color:#0f172a;font-size:9pt;text-align:center;font-weight:bold;width:20%;">{s["latest"] or "—"}</td><td style="padding:6px 10px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:15%;">{st_label}</td></tr>'

    # Key Risks Into Print
    risk_rows = ''
    for rk in rk_data:
        st = rk['status'] or ''
        st_bg = _status_bg(st) if st else '#ffffff'
        st_color = _status_color(st) if st else '#94a3b8'
        norm_st = _normalize_status(st)
        st_label = rk['statusNote'] if rk['statusNote'] else (norm_st.upper() if norm_st else '—')
        risk_rows += f'<tr><td style="padding:6px 10px;border:1px solid #e2e8f0;font-weight:600;color:#1e293b;font-size:9pt;width:70%;">{rk["threat"]}</td><td style="padding:6px 10px;border:1px solid #e2e8f0;text-align:center;background:{st_bg};color:{st_color};font-weight:bold;font-size:8.5pt;width:30%;">{st_label}</td></tr>'

    # Quick reference box
    conv_color = '#16a34a' if pct >= 70 else '#ca8a04' if pct >= 40 else '#dc2626'
    conv_label = 'HIGH' if pct >= 70 else 'MEDIUM' if pct >= 40 else 'LOW'
    ref_html = ''
    if total > 0:
        ref_html = f'<table style="border:2px solid #c4b5fd;margin-top:4px;"><tr><td style="padding:12px 18px;text-align:center;background:#f5f3ff;width:34%;"><span style="font-size:8pt;color:#4c1d95;font-weight:700;">CONVICTION</span><br/><span style="font-size:22pt;font-weight:bold;color:{conv_color};">{pct}%</span><br/><span style="font-size:8pt;color:{conv_color};font-weight:700;">{conv_label}</span></td><td style="padding:12px 18px;text-align:center;width:22%;border-left:1px solid #c4b5fd;"><span style="font-size:8pt;color:#4c1d95;font-weight:700;">PILLARS</span><br/><span style="font-size:20pt;font-weight:bold;color:#4c1d95;">{pillar_count}</span></td><td style="padding:12px 18px;text-align:center;width:22%;border-left:1px solid #c4b5fd;"><span style="font-size:8pt;color:#4c1d95;font-weight:700;">SIGNPOSTS</span><br/><span style="font-size:20pt;font-weight:bold;color:#4c1d95;">{len(sp_data)}</span></td><td style="padding:12px 18px;text-align:center;width:22%;border-left:1px solid #c4b5fd;"><span style="font-size:8pt;color:#4c1d95;font-weight:700;">RISKS</span><br/><span style="font-size:20pt;font-weight:bold;color:#4c1d95;">{len(rk_data)}</span></td></tr></table>'

    html = f"""<html><head><style>
        @page {{ margin: 0.6in 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.45; }}
        table {{ border-collapse: collapse; width: 100%; }}
        tr {{ page-break-inside: avoid; }}
    </style></head><body>
        <table style="margin-bottom:14px;"><tr>
            <td style="background:#4c1d95;padding:18px 24px;width:60%;"><span style="font-size:26pt;font-weight:bold;color:#ffffff;">{_fmt_escape(ticker)}</span><br/><span style="font-size:11pt;color:#ddd6fe;">{_fmt_escape(company)}</span></td>
            <td style="background:#4c1d95;padding:18px 24px;text-align:right;vertical-align:bottom;width:40%;"><span style="font-size:8pt;color:#ddd6fe;">EARNINGS SNAPSHOT | {date_str}</span></td>
        </tr></table>

        {thesis_html}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#5b21b6;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">WHAT TO WATCH THIS QUARTER</b></td></tr></table><table style="border:1px solid #e2e8f0;margin-bottom:18px;"><thead><tr style="background:#f5f3ff;"><th style="padding:6px 10px;text-align:left;font-size:8pt;color:#4c1d95;font-weight:700;border:1px solid #e2e8f0;width:30%;">METRIC</th><th style="padding:6px 10px;text-align:center;font-size:8pt;color:#4c1d95;font-weight:700;border:1px solid #e2e8f0;width:25%;">TARGET</th><th style="padding:6px 10px;text-align:center;font-size:8pt;color:#4c1d95;font-weight:700;border:1px solid #e2e8f0;width:20%;">LATEST</th><th style="padding:6px 10px;text-align:center;font-size:8pt;color:#4c1d95;font-weight:700;border:1px solid #e2e8f0;width:15%;">STATUS</th></tr></thead><tbody>' + sp_rows + '</tbody></table>' if signposts else ''}

        {'<table style="margin-bottom:8px;"><tr><td style="background:#5b21b6;padding:8px 16px;width:100%;"><b style="color:#ffffff;font-size:10pt;letter-spacing:1px;">KEY RISKS INTO PRINT</b></td></tr></table><table style="border:1px solid #e2e8f0;margin-bottom:18px;"><thead><tr style="background:#f5f3ff;"><th style="padding:6px 10px;text-align:left;font-size:8pt;color:#4c1d95;font-weight:700;border:1px solid #e2e8f0;width:70%;">RISK FACTOR</th><th style="padding:6px 10px;text-align:center;font-size:8pt;color:#4c1d95;font-weight:700;border:1px solid #e2e8f0;width:30%;">STATUS</th></tr></thead><tbody>' + risk_rows + '</tbody></table>' if threats else ''}

        {ref_html}
    </body></html>"""
    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    return buf.getvalue()


# ============================================
# PPTX SLIDE DECK GENERATOR
# ============================================

def _pptx_status_color(status):
    """Return RGBColor for a status string."""
    from pptx.dml.color import RGBColor
    s = _normalize_status(status)
    if s == 'green': return RGBColor(0x16, 0xA3, 0x4A)
    elif s == 'yellow': return RGBColor(0xCA, 0x8A, 0x04)
    elif s == 'red': return RGBColor(0xDC, 0x26, 0x26)
    return RGBColor(0x94, 0xA3, 0xB8)

def _pptx_status_bg(status):
    """Return RGBColor background for a status cell."""
    from pptx.dml.color import RGBColor
    s = _normalize_status(status)
    if s == 'green': return RGBColor(0xDC, 0xFC, 0xE7)
    elif s == 'yellow': return RGBColor(0xFE, 0xF9, 0xC3)
    elif s == 'red': return RGBColor(0xFE, 0xE2, 0xE2)
    return None

def _generate_thesis_pptx_bytes(row, scorecard_data=None, fmt='executive'):
    """Generate a 5-slide PowerPoint deck for a stock thesis."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import io

    d = _parse_analysis_data(row)
    ticker, company, date_str = d['ticker'], d['company'], d['date_str']
    thesis, signposts, threats, conclusion = d['thesis'], d['signposts'], d['threats'], d['conclusion']
    is_condensed = d.get('_condensed', False)
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)
    g, y, r = _tally_statuses(scorecard_data)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(0x1E, 0x29, 0x3B)
    DARK = RGBColor(0x0F, 0x17, 0x2A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEAL = RGBColor(0x14, 0xB8, 0xA6)
    SLATE = RGBColor(0x94, 0xA3, 0xB8)
    LIGHT = RGBColor(0xCB, 0xD5, 0xE1)

    def add_bg(slide, color=NAVY):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()

    def add_bar(slide, text, left, top, width):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(36))
        bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
        tf = bar.text_frame; tf.margin_left = Pt(10)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = text; run.font.size = Pt(14); run.font.color.rgb = WHITE; run.font.bold = True

    def style_cell(cell, font_size=Pt(11), font_color=WHITE, bold=False, fill_color=None, align=PP_ALIGN.LEFT):
        if fill_color:
            cell.fill.solid(); cell.fill.fore_color.rgb = fill_color
        cell.margin_left = Pt(6); cell.margin_right = Pt(6)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = align
            for run in p.runs:
                run.font.size = font_size; run.font.color.rgb = font_color; run.font.bold = bold

    # ---- SLIDE 1: Title ----
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1)
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f"{ticker} — {company}" if company else ticker
    run.font.size = Pt(44); run.font.color.rgb = WHITE; run.font.bold = True
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
    r2 = p2.add_run(); r2.text = "Investment Thesis"; r2.font.size = Pt(24); r2.font.color.rgb = TEAL
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.333), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()
    dtb = s1.shapes.add_textbox(Inches(9), Inches(6.5), Inches(3.5), Inches(0.5))
    dp = dtb.text_frame.paragraphs[0]; dp.alignment = PP_ALIGN.RIGHT
    dr = dp.add_run(); dr.text = date_str; dr.font.size = Pt(14); dr.font.color.rgb = SLATE

    # ---- SLIDE 2: Thesis + Pillars ----
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2)
    add_bar(s2, "WHY DO WE OWN IT?", Inches(0.5), Inches(0.4), Inches(12.333))
    stb = s2.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.9), Inches(1.8))
    stf = stb.text_frame; stf.word_wrap = True
    sp = stf.paragraphs[0]; sr = sp.add_run()
    sr.text = thesis.get('summary', ''); sr.font.size = Pt(15); sr.font.color.rgb = WHITE
    y_pos = Inches(3.2)
    pillars = _build_pillar_data(thesis, scorecard_data)
    for i, pil in enumerate(pillars, 1):
        ptitle = pil.get('pillar', pil.get('title', ''))
        pdesc = pil.get('detail', pil.get('description', ''))
        pb = s2.shapes.add_textbox(Inches(0.7), y_pos, Inches(11.9), Inches(0.7))
        ptf = pb.text_frame; ptf.word_wrap = True
        pp = ptf.paragraphs[0]
        r1 = pp.add_run(); r1.text = f"{i}. {ptitle}: "; r1.font.size = Pt(14); r1.font.color.rgb = TEAL; r1.font.bold = True
        r2 = pp.add_run(); r2.text = pdesc[:200] if pdesc else ''; r2.font.size = Pt(13); r2.font.color.rgb = LIGHT
        y_pos += Inches(0.75)

    # ---- SLIDE 3: Signposts Table ----
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3)
    add_bar(s3, "WHAT ARE WE WATCHING?", Inches(0.5), Inches(0.4), Inches(12.333))
    if sp_data:
        rows_n = len(sp_data) + 1
        tbl = s3.shapes.add_table(rows_n, 4, Inches(0.5), Inches(1.3), Inches(12.333), Inches(min(rows_n * 0.55, 5.5))).table
        tbl.columns[0].width = Inches(4.5); tbl.columns[1].width = Inches(3.0); tbl.columns[2].width = Inches(2.833); tbl.columns[3].width = Inches(2.0)
        for j, h in enumerate(['Metric', 'LT Goal', 'Latest', 'Status']):
            tbl.cell(0, j).text = h; style_cell(tbl.cell(0, j), Pt(12), WHITE, True, TEAL, PP_ALIGN.CENTER if j == 3 else PP_ALIGN.LEFT)
        for i, s in enumerate(sp_data):
            ri = i + 1
            tbl.cell(ri, 0).text = s.get('metric', ''); style_cell(tbl.cell(ri, 0), fill_color=DARK)
            tbl.cell(ri, 1).text = s.get('ltGoal', ''); style_cell(tbl.cell(ri, 1), fill_color=DARK)
            tbl.cell(ri, 2).text = s.get('latest', '') or '—'; style_cell(tbl.cell(ri, 2), fill_color=DARK)
            st = (s.get('status', '') or '').lower()
            tbl.cell(ri, 3).text = st.upper() if st else '—'
            sbg = _pptx_status_bg(st)
            style_cell(tbl.cell(ri, 3), font_color=_pptx_status_color(st) if not sbg else RGBColor(0x0F, 0x17, 0x2A), fill_color=sbg or DARK, align=PP_ALIGN.CENTER, bold=True)

    # ---- SLIDE 4: Risk Assessment ----
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4)
    add_bar(s4, "WHAT COULD GO WRONG?", Inches(0.5), Inches(0.4), Inches(12.333))
    if rk_data:
        rows_n = len(rk_data) + 1
        tbl = s4.shapes.add_table(rows_n, 2, Inches(0.5), Inches(1.3), Inches(12.333), Inches(min(rows_n * 0.55, 5.5))).table
        tbl.columns[0].width = Inches(9.5); tbl.columns[1].width = Inches(2.833)
        for j, h in enumerate(['Risk Factor', 'Status']):
            tbl.cell(0, j).text = h; style_cell(tbl.cell(0, j), Pt(12), WHITE, True, TEAL, PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT)
        for i, rk in enumerate(rk_data):
            ri = i + 1
            tbl.cell(ri, 0).text = rk.get('threat', ''); style_cell(tbl.cell(ri, 0), fill_color=DARK)
            st = (rk.get('status', '') or '').lower()
            tbl.cell(ri, 1).text = st.upper() if st else '—'
            sbg = _pptx_status_bg(st)
            style_cell(tbl.cell(ri, 1), font_color=_pptx_status_color(st) if not sbg else RGBColor(0x0F, 0x17, 0x2A), fill_color=sbg or DARK, align=PP_ALIGN.CENTER, bold=True)

    # ---- SLIDE 5: Conclusion + Conviction ----
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5)
    add_bar(s5, "CONCLUSION & CONVICTION", Inches(0.5), Inches(0.4), Inches(12.333))
    # Conviction boxes
    total = g + y + r
    bx = Inches(0.7); by = Inches(1.3)
    for label, count, color, bg_c in [('GREEN', g, RGBColor(0x16,0xA3,0x4A), RGBColor(0x05,0x2E,0x16)),
                                       ('YELLOW', y, RGBColor(0xCA,0x8A,0x04), RGBColor(0x42,0x2D,0x09)),
                                       ('RED', r, RGBColor(0xDC,0x26,0x26), RGBColor(0x45,0x0A,0x0A))]:
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(3.5), Inches(1.2))
        box.fill.solid(); box.fill.fore_color.rgb = bg_c; box.line.fill.background()
        btf = box.text_frame; btf.word_wrap = True
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = str(count); br.font.size = Pt(40); br.font.color.rgb = color; br.font.bold = True
        bp2 = btf.add_paragraph(); bp2.alignment = PP_ALIGN.CENTER
        br2 = bp2.add_run(); br2.text = label; br2.font.size = Pt(12); br2.font.color.rgb = SLATE
        bx += Inches(4.1)
    # Health score
    if total > 0:
        pct = round(g / total * 100)
        htb = s5.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(12), Inches(0.6))
        hp = htb.text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        hr1 = hp.add_run(); hr1.text = f"Health Score: {pct}%"; hr1.font.size = Pt(20); hr1.font.color.rgb = TEAL; hr1.font.bold = True
    # Conclusion
    conclusion_text = conclusion if isinstance(conclusion, str) else str(conclusion or '')
    if conclusion_text:
        ctb = s5.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(11.9), Inches(3))
        ctf = ctb.text_frame; ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cr = cp.add_run(); cr.text = conclusion_text[:800]; cr.font.size = Pt(15); cr.font.color.rgb = WHITE

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================================
# THESIS CONDENSED ENDPOINTS
# ============================================

def _build_condensed_tier(source_data, anthropic_key="", gemini_key=""):
    """Build a condensed thesis tier from source data (Full or Detailed). Returns condensed dict.

    ALL sections go to the LLM for compression — summary, pillars, signposts, threats, conclusion.
    """
    thesis = source_data.get('thesis', {})
    signposts = source_data.get('signposts', [])
    threats = source_data.get('threats', [])
    conclusion = source_data.get('conclusion', '')

    source_json = json.dumps({
        'summary': thesis.get('summary', ''),
        'pillars': thesis.get('pillars', []),
        'signposts': signposts,
        'threats': threats,
        'conclusion': conclusion,
    }, indent=2)

    prompt = f"""Compress this investment thesis into a short, punchy version. Rules:

1. Summary: 2-3 tight sentences. Preserve the core WHY.
2. Pillars: keep 4-5 strongest. 1 sentence description each. Keep title exactly as-is. Drop confidence and sources.
3. Signposts: keep 3-5 most actionable. Keep the metric name exactly as-is. Tighten the target to 1 sentence. Drop confidence, sources, and timeframe.
4. Threats: keep 3-4 highest-impact. Keep the threat title exactly as-is. Tighten triggerPoints to 1 sentence.
5. Conclusion: 1-2 sentences max.

Return ONLY valid JSON:
{{
  "thesis": {{
    "summary": "...",
    "pillars": [{{"title": "exact original title", "description": "1 sentence"}}]
  }},
  "signposts": [{{"metric": "exact original metric name", "target": "1 sentence target"}}],
  "threats": [{{"threat": "exact original threat title", "triggerPoints": "1 sentence"}}],
  "conclusion": "1-2 sentences"
}}

Source data:
{source_json}"""

    result = call_llm(
        messages=[{"role": "user", "content": prompt}],
        system="You are a concise financial analyst. Return only valid JSON, no markdown fences.",
        tier="fast",
        max_tokens=4096,
        anthropic_api_key=anthropic_key,
        gemini_api_key=gemini_key,
    )

    try:
        condensed = _extract_json(result['text'])
    except json.JSONDecodeError:
        print(f"[_build_condensed_tier] Raw LLM response:\n{result['text'][:2000]}")
        raise

    # Validate structure
    if 'thesis' not in condensed:
        condensed = {'thesis': condensed, 'signposts': [], 'threats': [], 'conclusion': conclusion}
    if 'signposts' not in condensed:
        condensed['signposts'] = []
    if 'threats' not in condensed:
        condensed['threats'] = []
    if 'conclusion' not in condensed:
        condensed['conclusion'] = conclusion

    return condensed


@app.route('/api/thesis-condensed/generate', methods=['POST'])
def generate_condensed_thesis():
    """Generate a condensed version of a thesis using LLM."""
    try:
        data = request.get_json()
        ticker = data.get('ticker', '').upper()
        anthropic_key = data.get('apiKey', '')
        gemini_key = data.get('geminiApiKey', '')
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400

        # Load analysis — prefer thesis_full (curated) if it exists, else raw detailed
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No analysis found for this ticker'}), 404

        analysis = row['analysis'] if isinstance(row['analysis'], dict) else json.loads(row['analysis'])

        # Check for curated full version to source from
        with get_db() as (_, cur):
            cur.execute('SELECT full_analysis FROM thesis_full WHERE ticker = %s', (ticker,))
            full_row = cur.fetchone()
        if full_row:
            full_data = full_row['full_analysis']
            if isinstance(full_data, str):
                full_data = json.loads(full_data)
            analysis['thesis'] = full_data.get('thesis', analysis.get('thesis', {}))
            analysis['signposts'] = full_data.get('signposts', analysis.get('signposts', []))
            analysis['threats'] = full_data.get('threats', analysis.get('threats', []))
            analysis['conclusion'] = full_data.get('conclusion', analysis.get('conclusion', ''))

        condensed = _build_condensed_tier(analysis, anthropic_key, gemini_key)

        # Snapshot existing tier before overwrite
        with get_db(commit=True) as (conn, cur):
            _snapshot_tier_history(cur, 'thesis_condensed', 'condensed_analysis', ticker)

        # Upsert into DB
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_condensed (ticker, condensed_analysis, updated_at)
                VALUES (%s, %s, CURRENT_TIMESTAMP)
                ON CONFLICT (ticker)
                DO UPDATE SET condensed_analysis = EXCLUDED.condensed_analysis, updated_at = CURRENT_TIMESTAMP
                RETURNING ticker
            ''', (ticker, json.dumps(condensed)))

        # Fetch updated history and timestamp
        with get_db() as (_, cur):
            cur.execute('SELECT history, updated_at FROM thesis_condensed WHERE ticker = %s', (ticker,))
            meta = cur.fetchone()

        return jsonify({
            'success': True, 'ticker': ticker, 'condensed': condensed,
            'history': (meta.get('history') or []) if meta else [],
            'updatedAt': meta['updated_at'].isoformat() if meta and meta.get('updated_at') else None,
        })
    except json.JSONDecodeError as je:
        print(f"Error parsing condensed thesis JSON: {je}")
        return jsonify({'error': 'Failed to parse condensed thesis from Claude'}), 500
    except Exception as e:
        print(f"Error generating condensed thesis: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/thesis-condensed/<ticker>', methods=['GET'])
def get_condensed_thesis(ticker):
    """Get stored condensed thesis for a ticker."""
    try:
        ticker = ticker.upper()
        with get_db() as (_, cur):
            cur.execute('SELECT condensed_analysis, history, created_at, updated_at FROM thesis_condensed WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No condensed thesis found'}), 404
        condensed = row['condensed_analysis']
        if isinstance(condensed, str):
            condensed = json.loads(condensed)
        return jsonify({
            'ticker': ticker,
            'condensed': condensed,
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
            'updatedAt': row['updated_at'].isoformat() if row['updated_at'] else None,
            'history': row.get('history') or [],
        })
    except Exception as e:
        print(f"Error getting condensed thesis: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/thesis-condensed/<ticker>', methods=['DELETE'])
def delete_condensed_thesis(ticker):
    """Delete stored condensed thesis."""
    try:
        ticker = ticker.upper()
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM thesis_condensed WHERE ticker = %s', (ticker,))
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting condensed thesis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-condensed/<ticker>', methods=['PUT'])
def update_condensed_thesis(ticker):
    """Save edited condensed thesis."""
    try:
        ticker = ticker.upper()
        data = request.get_json()
        condensed = data.get('condensed_analysis')
        if not condensed:
            return jsonify({'error': 'condensed_analysis required'}), 400
        # Snapshot existing tier before overwrite
        with get_db(commit=True) as (conn, cur):
            _snapshot_tier_history(cur, 'thesis_condensed', 'condensed_analysis', ticker)
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_condensed (ticker, condensed_analysis, updated_at)
                VALUES (%s, %s, %s)
                ON CONFLICT (ticker)
                DO UPDATE SET condensed_analysis = EXCLUDED.condensed_analysis, updated_at = EXCLUDED.updated_at
            ''', (ticker, json.dumps(condensed), datetime.utcnow()))
        with get_db() as (_, cur):
            cur.execute('SELECT history, updated_at FROM thesis_condensed WHERE ticker = %s', (ticker,))
            meta = cur.fetchone()
        return jsonify({
            'success': True,
            'history': (meta.get('history') or []) if meta else [],
            'updatedAt': meta['updated_at'].isoformat() if meta and meta.get('updated_at') else None,
        })
    except Exception as e:
        print(f"Error saving condensed thesis: {e}")
        return jsonify({'error': str(e)}), 500


# --- Thesis Full (curated selection) endpoints ---

def _build_full_tier(analysis, anthropic_key="", gemini_key=""):
    """Build a curated Full thesis tier from detailed analysis. Returns full_data dict.

    Selects top items and allows light prose tightening while preserving all fields.
    """
    thesis = analysis.get('thesis', {})
    signposts = analysis.get('signposts', [])
    threats = analysis.get('threats', [])
    conclusion = analysis.get('conclusion', '')

    thesis_text = json.dumps(thesis, indent=2)
    signposts_text = json.dumps(signposts, indent=2)
    threats_text = json.dumps(threats, indent=2)

    prompt = f"""You are curating an investment thesis by selecting the most important items. You may lightly tighten prose but do NOT compress or rewrite substantially.

Rules:
1. Select 4-5 most important pillars. Keep title exactly as-is. You may tighten the description to remove filler, but preserve meaning. Keep confidence and sources unchanged.
2. Select 5-7 most actionable signposts. You may tighten the text slightly. Keep all fields (confidence, sources, timeframe).
3. Select 4-5 highest-risk threats. Drop likelihood and impact fields entirely, but KEEP the triggerPoints field (this is the "Watch for" text — it MUST be preserved). You may tighten the threat sentence.
4. Summary: keep as-is or remove filler only. Do not change meaning.
5. Conclusion: keep exactly as-is.

Return ONLY valid JSON with this exact structure:
{{
  "thesis": {{
    "summary": "original or lightly tightened summary",
    "pillars": [selected pillars with ALL original fields]
  }},
  "signposts": [selected signposts with ALL original fields],
  "threats": [selected threats — each MUST include "threat" and "triggerPoints" fields],
  "conclusion": "exact original conclusion"
}}

Original thesis:
{thesis_text}

Original signposts:
{signposts_text}

Original threats:
{threats_text}

Original conclusion:
{json.dumps(conclusion)}"""

    result = call_llm(
        messages=[{"role": "user", "content": prompt}],
        system="You are a financial analyst curating the most important items from a thesis. Return only valid JSON, no markdown fences. Light tightening is OK; do not compress or rewrite.",
        tier="fast",
        max_tokens=16384,
        anthropic_api_key=anthropic_key,
        gemini_api_key=gemini_key,
    )

    print(f"[_build_full_tier] LLM response length: {len(result['text'])} chars, provider: {result.get('provider')}, model: {result.get('model')}")
    try:
        full_data = _extract_json(result['text'])
    except json.JSONDecodeError:
        print(f"[_build_full_tier] FAILED TO PARSE. Raw LLM response (first 3000 chars):\n{result['text'][:3000]}")
        print(f"[_build_full_tier] Raw LLM response (last 500 chars):\n{result['text'][-500:]}")
        raise

    # Validate structure
    if 'thesis' not in full_data:
        full_data = {'thesis': full_data, 'signposts': signposts[:7], 'threats': threats[:5], 'conclusion': conclusion}
    if 'signposts' not in full_data:
        full_data['signposts'] = signposts[:7]
    if 'threats' not in full_data:
        full_data['threats'] = threats[:5]
    if 'conclusion' not in full_data:
        full_data['conclusion'] = conclusion

    # Restore triggerPoints if LLM dropped them — match by threat name to originals
    orig_threats_map = {t.get('threat', ''): t for t in threats}
    for ft in full_data.get('threats', []):
        if not ft.get('triggerPoints'):
            orig = orig_threats_map.get(ft.get('threat', ''), {})
            if orig.get('triggerPoints'):
                ft['triggerPoints'] = orig['triggerPoints']

    return full_data


@app.route('/api/thesis-full/generate', methods=['POST'])
def generate_full_thesis():
    """Generate a curated Full version by selecting top items from detailed analysis."""
    try:
        data = request.get_json()
        ticker = data.get('ticker', '').upper()
        anthropic_key = data.get('apiKey', '')
        gemini_key = data.get('geminiApiKey', '')
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400

        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No analysis found for this ticker'}), 404

        analysis = row['analysis'] if isinstance(row['analysis'], dict) else json.loads(row['analysis'])

        full_data = _build_full_tier(analysis, anthropic_key, gemini_key)

        # Snapshot existing tier before overwrite
        with get_db(commit=True) as (conn, cur):
            _snapshot_tier_history(cur, 'thesis_full', 'full_analysis', ticker)

        # Upsert into DB
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_full (ticker, full_analysis, updated_at)
                VALUES (%s, %s, CURRENT_TIMESTAMP)
                ON CONFLICT (ticker)
                DO UPDATE SET full_analysis = EXCLUDED.full_analysis, updated_at = CURRENT_TIMESTAMP
                RETURNING ticker
            ''', (ticker, json.dumps(full_data)))

        # Invalidate stale condensed version (should be re-derived from Full)
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM thesis_condensed WHERE ticker = %s', (ticker,))

        # Fetch updated history and timestamp
        with get_db() as (_, cur):
            cur.execute('SELECT history, updated_at FROM thesis_full WHERE ticker = %s', (ticker,))
            meta = cur.fetchone()

        return jsonify({
            'success': True, 'ticker': ticker, 'full': full_data,
            'history': (meta.get('history') or []) if meta else [],
            'updatedAt': meta['updated_at'].isoformat() if meta and meta.get('updated_at') else None,
        })
    except json.JSONDecodeError as je:
        print(f"Error parsing full thesis JSON: {je}")
        return jsonify({'error': 'Failed to parse full thesis from LLM'}), 500
    except Exception as e:
        print(f"Error generating full thesis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-full/<ticker>', methods=['GET'])
def get_full_thesis(ticker):
    """Get stored full (curated) thesis for a ticker."""
    try:
        ticker = ticker.upper()
        with get_db() as (_, cur):
            cur.execute('SELECT full_analysis, history, created_at, updated_at FROM thesis_full WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No full thesis found'}), 404
        full_data = row['full_analysis']
        if isinstance(full_data, str):
            full_data = json.loads(full_data)
        return jsonify({
            'ticker': ticker,
            'full': full_data,
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
            'updatedAt': row['updated_at'].isoformat() if row['updated_at'] else None,
            'history': row.get('history') or [],
        })
    except Exception as e:
        print(f"Error getting full thesis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-full/<ticker>', methods=['DELETE'])
def delete_full_thesis(ticker):
    """Delete stored full thesis."""
    try:
        ticker = ticker.upper()
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM thesis_full WHERE ticker = %s', (ticker,))
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting full thesis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-full/<ticker>', methods=['PUT'])
def update_full_thesis(ticker):
    """Save edited full thesis."""
    try:
        ticker = ticker.upper()
        data = request.get_json()
        full_data = data.get('full_analysis')
        if not full_data:
            return jsonify({'error': 'full_analysis required'}), 400
        # Snapshot existing tier before overwrite
        with get_db(commit=True) as (conn, cur):
            _snapshot_tier_history(cur, 'thesis_full', 'full_analysis', ticker)
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_full (ticker, full_analysis, updated_at)
                VALUES (%s, %s, %s)
                ON CONFLICT (ticker)
                DO UPDATE SET full_analysis = EXCLUDED.full_analysis, updated_at = EXCLUDED.updated_at
            ''', (ticker, json.dumps(full_data), datetime.utcnow()))
        with get_db() as (_, cur):
            cur.execute('SELECT history, updated_at FROM thesis_full WHERE ticker = %s', (ticker,))
            meta = cur.fetchone()
        return jsonify({
            'success': True,
            'history': (meta.get('history') or []) if meta else [],
            'updatedAt': meta['updated_at'].isoformat() if meta and meta.get('updated_at') else None,
        })
    except Exception as e:
        print(f"Error saving full thesis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-tiers/generate-all', methods=['POST'])
def generate_all_thesis_tiers():
    """Generate Full then Condensed thesis tiers sequentially."""
    try:
        data = request.get_json()
        ticker = data.get('ticker', '').upper()
        anthropic_key = data.get('apiKey', '')
        gemini_key = data.get('geminiApiKey', '')
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400

        # Load detailed analysis
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No analysis found for this ticker'}), 404

        analysis = row['analysis'] if isinstance(row['analysis'], dict) else json.loads(row['analysis'])

        # Step 1: Build Full tier
        full_data = _build_full_tier(analysis, anthropic_key, gemini_key)

        # Snapshot existing full tier before overwrite
        with get_db(commit=True) as (conn, cur):
            _snapshot_tier_history(cur, 'thesis_full', 'full_analysis', ticker)

        # Upsert Full into DB
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO thesis_full (ticker, full_analysis, updated_at)
                VALUES (%s, %s, CURRENT_TIMESTAMP)
                ON CONFLICT (ticker)
                DO UPDATE SET full_analysis = EXCLUDED.full_analysis, updated_at = CURRENT_TIMESTAMP
            ''', (ticker, json.dumps(full_data)))

        # Invalidate stale condensed
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM thesis_condensed WHERE ticker = %s', (ticker,))

        # Step 2: Build Condensed tier from Full
        result = {'success': True, 'ticker': ticker, 'full': full_data}
        try:
            condensed_data = _build_condensed_tier(full_data, anthropic_key, gemini_key)

            # Snapshot existing condensed tier before overwrite
            with get_db(commit=True) as (conn, cur):
                _snapshot_tier_history(cur, 'thesis_condensed', 'condensed_analysis', ticker)

            with get_db(commit=True) as (conn, cur):
                cur.execute('''
                    INSERT INTO thesis_condensed (ticker, condensed_analysis, updated_at)
                    VALUES (%s, %s, CURRENT_TIMESTAMP)
                    ON CONFLICT (ticker)
                    DO UPDATE SET condensed_analysis = EXCLUDED.condensed_analysis, updated_at = CURRENT_TIMESTAMP
                ''', (ticker, json.dumps(condensed_data)))

            result['condensed'] = condensed_data
        except Exception as ce:
            print(f"Condensed generation failed (Full succeeded): {ce}")
            result['partialSuccess'] = True
            result['condensedError'] = str(ce)

        # Fetch history and timestamps for both tiers
        with get_db() as (_, cur):
            cur.execute('SELECT history, updated_at FROM thesis_full WHERE ticker = %s', (ticker,))
            frow = cur.fetchone()
            result['fullHistory'] = (frow.get('history') or []) if frow else []
            result['fullUpdatedAt'] = frow['updated_at'].isoformat() if frow and frow.get('updated_at') else None
            cur.execute('SELECT history, updated_at FROM thesis_condensed WHERE ticker = %s', (ticker,))
            crow = cur.fetchone()
            result['condensedHistory'] = (crow.get('history') or []) if crow else []
            result['condensedUpdatedAt'] = crow['updated_at'].isoformat() if crow and crow.get('updated_at') else None

        return jsonify(result)
    except json.JSONDecodeError as je:
        print(f"Error parsing thesis JSON in generate-all: {je}")
        return jsonify({'error': 'Failed to parse thesis from LLM'}), 500
    except Exception as e:
        print(f"Error in generate-all thesis tiers: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/generate', methods=['POST'])
def generate_thesis_format():
    """Generate a formatted thesis document in a specific template."""
    try:
        data = request.get_json()
        ticker = data.get('ticker', '').upper()
        fmt = data.get('format', 'executive')  # executive, scorecard, onepager, board
        output_type = data.get('outputType', 'pdf')  # pdf or docx
        thesis_tier = _resolve_thesis_tier(data)
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400

        # Fetch analysis
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No analysis found for this ticker'}), 404

        # Fetch scorecard data
        scorecard_data = None
        with get_db() as (_, cur):
            cur.execute('SELECT scorecard_data FROM thesis_scorecard_data WHERE ticker = %s', (ticker,))
            sc_row = cur.fetchone()
        if sc_row:
            scorecard_data = sc_row['scorecard_data']
            if isinstance(scorecard_data, str):
                try:
                    scorecard_data = json.loads(scorecard_data)
                except:
                    scorecard_data = None

        row_dict = dict(row)
        _overlay_thesis_tier(row_dict, thesis_tier)

        # Generate based on format
        format_names = {
            'executive': 'Executive_Brief',
            'scorecard': 'Scorecard',
            'onepager': 'One_Pager',
            'board': 'Conviction_Card',
            'ic_memo': 'IC_Memo',
            'risk_focus': 'Risk_Monitor',
            'catalyst': 'Catalyst_Timeline',
            'valuation': 'Valuation_Summary',
            'bull_bear': 'Bull_Bear_Framework',
            'portfolio': 'Portfolio_Context',
            'earnings': 'Earnings_Snapshot',
        }
        format_name = format_names.get(fmt, 'Formatted')

        if output_type == 'pdf':
            generators = {
                'executive': _generate_executive_brief_pdf,
                'scorecard': _generate_scorecard_pdf,
                'onepager': _generate_onepager_pdf,
                'board': _generate_board_pdf,
                'ic_memo': _generate_ic_memo_pdf,
                'risk_focus': _generate_risk_focus_pdf,
                'catalyst': _generate_catalyst_timeline_pdf,
                'valuation': _generate_valuation_summary_pdf,
                'bull_bear': _generate_bull_bear_pdf,
                'portfolio': _generate_portfolio_context_pdf,
                'earnings': _generate_earnings_snapshot_pdf,
            }
            gen_func = generators.get(fmt, _generate_executive_brief_pdf)
            file_bytes = gen_func(row_dict, scorecard_data)
            filename = f"{ticker}_{format_name}.pdf"
            mime = 'application/pdf'
        elif output_type == 'pptx':
            file_bytes = _generate_thesis_pptx_bytes(row_dict, scorecard_data, fmt)
            filename = f"{ticker}_{format_name}.pptx"
            mime = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        else:
            file_bytes = _generate_analysis_docx_bytes(row_dict)
            filename = f"{ticker}_{format_name}.docx"
            mime = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'

        b64 = base64.b64encode(file_bytes).decode('utf-8')

        # Auto-save to history
        saved_id = None
        try:
            with get_db(commit=True) as (conn, cur):
                cur.execute('''
                    INSERT INTO thesis_format_history (ticker, format_type, output_type, file_data, filename, file_size)
                    VALUES (%s, %s, %s, %s, %s, %s)
                    RETURNING id
                ''', (ticker, fmt, output_type, b64, filename, len(file_bytes)))
                saved_id = cur.fetchone()['id']
        except Exception as save_err:
            print(f"Failed to save format to history: {save_err}")

        return jsonify({'success': True, 'fileData': b64, 'filename': filename, 'fileSize': len(file_bytes), 'savedId': saved_id})
    except Exception as e:
        print(f"Error generating thesis format: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/history/<ticker>', methods=['GET'])
def get_format_history(ticker):
    """List all format versions for a ticker (metadata only)."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, ticker, format_type, output_type, filename, file_size, created_at
                FROM thesis_format_history
                WHERE ticker = %s
                ORDER BY created_at DESC
            ''', (ticker.upper(),))
            rows = cur.fetchall()
        return jsonify([{
            'id': r['id'],
            'ticker': r['ticker'],
            'formatType': r['format_type'],
            'outputType': r['output_type'],
            'filename': r['filename'],
            'fileSize': r['file_size'],
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
        } for r in rows])
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/history/<int:history_id>/download', methods=['GET'])
def download_format_history(history_id):
    """Get file data for a specific format version."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT file_data, filename, output_type FROM thesis_format_history WHERE id = %s', (history_id,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Not found'}), 404
        mime_map = {'pdf': 'application/pdf', 'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}
        mime = mime_map.get(row['output_type'], 'application/octet-stream')
        return jsonify({'fileData': row['file_data'], 'filename': row['filename'], 'mimeType': mime})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/history/<int:history_id>', methods=['DELETE'])
def delete_format_history(history_id):
    """Delete a specific format version."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM thesis_format_history WHERE id = %s', (history_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-format', methods=['POST'])
def email_format():
    """Email a thesis format document as attachment."""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    from email.mime.base import MIMEBase
    from email import encoders

    try:
        data = request.get_json()
        file_data = data.get('fileData', '')
        filename = data.get('filename', 'document.pdf')
        ticker = data.get('ticker', '')
        output_type = data.get('outputType', 'pdf')
        recipient = data.get('email', '')
        subject = data.get('customSubject', f'{ticker} — Thesis Format')
        smtp_config = data.get('smtpConfig', {})

        if not recipient:
            return jsonify({'error': 'Recipient email is required'}), 400
        if not file_data:
            return jsonify({'error': 'No file data'}), 400

        use_gmail = smtp_config.get('use_gmail', True)
        gmail_user = smtp_config.get('gmail_user', '')
        gmail_password = smtp_config.get('gmail_app_password', '')
        from_email = smtp_config.get('from_email', gmail_user)

        if use_gmail and (not gmail_user or not gmail_password):
            return jsonify({'error': 'Gmail credentials required'}), 400

        msg = MIMEMultipart()
        msg['From'] = from_email
        msg['To'] = recipient
        msg['Subject'] = subject

        body = f"Please find the attached {ticker} thesis format document."
        msg.attach(MIMEText(body, 'plain'))

        mime_type = 'application/pdf' if output_type == 'pdf' else 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        attachment = MIMEBase(*mime_type.split('/'))
        attachment.set_payload(base64.b64decode(file_data))
        encoders.encode_base64(attachment)
        attachment.add_header('Content-Disposition', f'attachment; filename="{filename}"')
        msg.attach(attachment)

        if use_gmail:
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(gmail_user, gmail_password)
                server.send_message(msg)

        return jsonify({'success': True, 'message': f'Format emailed to {recipient}'})

    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'Gmail authentication failed.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============================================
# THESIS INFOGRAPHIC GENERATION
# ============================================

def _build_thesis_infographic_prompt(d, scorecard_data, style_prompt, mode, slide_num=None, detail='full', show_risk_detail=False, color_scheme=None, include_company=False, style_key=None):
    """Build a Gemini image-generation prompt for a thesis infographic slide."""
    ticker = d['ticker']
    company = d['company']
    thesis = d['thesis']
    signposts = d['signposts']
    threats = d['threats']
    conclusion = d['conclusion']
    sp_data = _build_signpost_data(signposts, scorecard_data)
    rk_data = _build_risk_data(threats, scorecard_data)

    header = ticker + (f" — {company}" if company and include_company else "")

    # Build thesis section text
    summary = thesis.get('summary', '')
    pillars_list = _build_pillar_data(thesis, scorecard_data)
    pillars_text = ""

    if detail == 'simple':
        # Simple: first sentence only for summary
        if '.' in summary:
            summary = summary[:summary.index('.') + 1]
        # Simple: pillar titles only, max 3
        for i, p in enumerate(pillars_list[:3], 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pillars_text += f"\n  {i}. {ptitle}"
    elif detail == 'summary':
        for i, p in enumerate(pillars_list, 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pillars_text += f"\n  {i}. {ptitle}"
    else:
        for i, p in enumerate(pillars_list, 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pdesc = p.get('detail', p.get('description', ''))
            pillars_text += f"\n  {i}. {ptitle}: {pdesc}"

    # Build signpost section text
    signpost_text = ""
    for s in sp_data:
        if detail == 'simple':
            # Simple: metric name only, no status
            line = f"  - {s['metric']}"
        elif detail == 'summary':
            line = f"  - {s['metric']}: {s['latest'] or s['ltGoal']}"
            if s['status']:
                line += f" [{s['status'].upper()}]"
        else:
            line = f"  - {s['metric']}: Target = {s['ltGoal']}"
            if s['latest']:
                line += f", Latest = {s['latest']}"
            if s['status']:
                line += f" [{s['status'].upper()}]"
        signpost_text += line + "\n"

    # Build threats section text
    is_condensed = d.get('_condensed', False)
    threats_text = ""
    for r in rk_data:
        if detail == 'simple':
            # Simple: threat name only, no status
            line = f"  - {r['threat']}"
        elif is_condensed:
            # Condensed: threat name only, no likelihood/impact
            line = f"  - {r['threat']}"
        elif show_risk_detail:
            if detail == 'summary':
                line = f"  - {r['threat']} (Impact: {r['impact']})"
            else:
                line = f"  - {r['threat']}: Likelihood = {r['likelihood']}, Impact = {r['impact']}"
        else:
            line = f"  - {r['threat']}"
        if detail != 'simple' and r['status']:
            line += f" [{r['status'].upper()}]"
        threats_text += line + "\n"

    conclusion_text = conclusion if isinstance(conclusion, str) else str(conclusion)

    # For simple mode, omit conclusion entirely
    if detail == 'simple':
        conclusion_text = ''
    # For summary mode, truncate to first sentence
    elif detail == 'summary':
        if '.' in summary:
            summary = summary[:summary.index('.') + 1]
        if '.' in conclusion_text:
            conclusion_text = conclusion_text[:conclusion_text.index('.') + 1]

    parts = []
    parts.append(f"Generate a visually compelling 16:9 LANDSCAPE infographic image.\n")
    # Append color scheme suffix if provided
    parts.append(f"VISUAL STYLE (follow exactly):\n{style_prompt}")
    if color_scheme and color_scheme != 'default' and style_key:
        sdef = THESIS_INFOGRAPHIC_STYLES.get(style_key, {})
        for cs in sdef.get('colorSchemes', []):
            if cs['id'] == color_scheme and cs.get('promptSuffix'):
                parts.append(cs['promptSuffix'])
                break
    parts.append(
        "CRITICAL RULES:\n"
        "- ALL text MUST be in English\n"
        "- DO NOT use quotation marks around any text — no double quotes, no curly quotes, no single quotes for emphasis\n"
        "- DO NOT include any watermarks, AI generation notices, footer text, branding lines, or attribution text (e.g. no 'Financial Feature | Analysis | 2024' footers)\n"
        "- Use 16:9 widescreen aspect ratio\n"
        "- Include the ticker symbol prominently but do NOT add labels like 'Ticker Symbol' underneath it\n"
        "- All text must be perfectly clear, legible, and properly rendered\n"
        "- Letters must NOT be distorted or artistically modified — B must not look like 8, O must not look like 0\n"
        "- Do NOT use decorative drop caps or oversized first-letters at the start of paragraphs\n"
        "- Use standard clean fonts — no decorative fonts that sacrifice readability\n"
        "- Ensure sufficient contrast between text and background\n"
        "- DO NOT include any target prices, price targets, or stock price predictions in the infographic\n"
        "- DO NOT use quotation marks anywhere in the infographic — no opening/closing quotes around thesis summaries, conclusions, or any text\n"
        "- DO NOT place emoji icons, colored circles, or traffic-light dots next to individual signpost or risk items — list them as clean text only\n"
        "- If status information is provided in brackets like [GREEN], render it as a small subtle text label, NOT as a colored circle or emoji"
    )

    if detail == 'simple':
        summary_layout = (
            "\nThis is a SIMPLE AT-A-GLANCE version — absolute minimum text. "
            "Show only section headers and item names as clean text lists. No paragraphs, no data values, no descriptions, no status indicators. "
            "Think executive cheat sheet — names only."
        )
    elif detail == 'summary':
        summary_layout = (
            "\nThis is a SUMMARY/PRESENTATION version — use MINIMAL text. "
            "Use LARGE bold numbers and metrics as focal points. Use short bullet points, NOT paragraphs. "
            "Emphasize visual hierarchy with oversized key figures and clean section headers. "
            "Think executive presentation — the audience will be viewing from a distance."
        )
    else:
        summary_layout = ""

    if mode == '1':
        parts.append(f"\nTITLE: {header} — Investment Thesis\n")
        parts.append(f"INVESTMENT THESIS:\n{summary}")
        if pillars_text:
            parts.append(f"\nKEY PILLARS:{pillars_text}")
        if signpost_text:
            parts.append(f"\nSIGNPOSTS TO MONITOR:\n{signpost_text}")
        if threats_text:
            parts.append(f"\nKEY RISKS:\n{threats_text}")
        if conclusion_text:
            parts.append(f"\nCONCLUSION: {conclusion_text}")
        if detail in ('summary', 'simple'):
            parts.append(summary_layout)
        else:
            parts.append("\nLayout this as a dense but readable infographic with all sections visible. Use visual hierarchy: large thesis header, medium section headers, compact data tables/lists for signposts and risks.")
    elif mode == '3':
        if slide_num == 1:
            parts.append(f"\nTITLE: {header} — Investment Thesis\n")
            parts.append(f"THESIS SUMMARY:\n{summary}")
            if pillars_text:
                parts.append(f"\nKEY PILLARS:{pillars_text}")
            if conclusion_text:
                parts.append(f"\nCONCLUSION: {conclusion_text}")
            if detail in ('summary', 'simple'):
                parts.append(summary_layout + " Slide 1 of 3.")
            else:
                parts.append("\nFocus this slide on the investment thesis narrative. Make the summary prominent, pillars clearly numbered, conclusion at bottom. Slide 1 of 3.")
        elif slide_num == 2:
            parts.append(f"\nTITLE: {header} — Key Signposts\n")
            parts.append(f"SIGNPOSTS TO MONITOR:\n{signpost_text}")
            if detail in ('summary', 'simple'):
                parts.append(summary_layout + " Slide 2 of 3.")
            else:
                parts.append("\nPresent each signpost as a clean monitoring dashboard item showing the metric name, target value, and latest reading. Use clean text layout — no colored circles, emojis, or traffic-light dots next to items. Slide 2 of 3.")
        elif slide_num == 3:
            parts.append(f"\nTITLE: {header} — Risk Assessment\n")
            parts.append(f"KEY RISKS:\n{threats_text}")
            if detail in ('summary', 'simple'):
                parts.append(summary_layout + " Slide 3 of 3.")
            else:
                if show_risk_detail:
                    parts.append("\nPresent each risk as a clean text card showing the risk name, likelihood, and impact level. Use visual weight (font size, boldness) to show severity — no colored circles or emojis. Slide 3 of 3.")
                else:
                    parts.append("\nPresent each risk as a clean text list with the risk name. Use visual weight (font size, boldness) to show severity — no colored circles or emojis. Slide 3 of 3.")

    return "\n".join(parts)


def _strip_quotes_from_data(d):
    """Remove all quotation marks from text fields in the data dictionary."""
    import copy
    d = copy.deepcopy(d)
    QUOTES = ['\u201c', '\u201d', '\u2018', '\u2019', '"']
    def strip_q(text):
        if not isinstance(text, str):
            return text
        for q in QUOTES:
            text = text.replace(q, '')
        return text
    if d.get('thesis') and isinstance(d['thesis'], dict):
        d['thesis']['summary'] = strip_q(d['thesis'].get('summary', ''))
        for p in d['thesis'].get('pillars', []):
            for k in ('pillar', 'title', 'detail', 'description'):
                if k in p:
                    p[k] = strip_q(p[k])
    d['conclusion'] = strip_q(d.get('conclusion', ''))
    for sp in d.get('signposts', []):
        for k in ('metric', 'signpost', 'target', 'timeframe'):
            if k in sp:
                sp[k] = strip_q(sp[k])
    for t in d.get('threats', []):
        for k in ('threat', 'triggerPoints'):
            if k in t:
                t[k] = strip_q(t[k])
    return d


def _apply_pillow_edit(d, edit_prompt):
    """Apply text edits from a prompt to the data dictionary for Pillow re-rendering.

    Supports:
      - "remove X" / "delete X" → removes X from all text fields
      - "change X to Y" / "replace X with Y" → substitutes X→Y in all text fields
      - Direct find→replace pairs separated by newlines
    """
    import re as _re
    if not edit_prompt or not edit_prompt.strip():
        return d

    # Collect all substitution pairs: (old, new)
    subs = []
    for line in edit_prompt.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        # "remove X" / "delete X"
        m = _re.match(r'^(?:remove|delete)\s+["\']?(.+?)["\']?\s*$', line, _re.IGNORECASE)
        if m:
            subs.append((m.group(1), ''))
            continue
        # "change X to Y" / "replace X with Y"
        m = _re.match(r'^(?:change|replace)\s+["\']?(.+?)["\']?\s+(?:to|with)\s+["\']?(.+?)["\']?\s*$', line, _re.IGNORECASE)
        if m:
            subs.append((m.group(1), m.group(2)))
            continue
        # "X -> Y" or "X → Y"
        for arrow in [' -> ', ' → ', '→', '->']:
            if arrow in line:
                parts = line.split(arrow, 1)
                if len(parts) == 2 and parts[0].strip():
                    subs.append((parts[0].strip().strip('"\''), parts[1].strip().strip('"\'')))
                break

    if not subs:
        # Treat entire prompt as a single removal if short, otherwise skip
        if len(edit_prompt.strip()) < 60:
            subs.append((edit_prompt.strip(), ''))
        else:
            return d

    def apply_subs(text):
        if not isinstance(text, str):
            return text
        for old, new in subs:
            text = text.replace(old, new)
        return text

    # Apply to all text fields in the data dictionary
    import copy
    d = copy.deepcopy(d)
    if d.get('thesis'):
        if isinstance(d['thesis'], dict):
            d['thesis']['summary'] = apply_subs(d['thesis'].get('summary', ''))
            for p in d['thesis'].get('pillars', []):
                for k in ('pillar', 'title', 'detail', 'description'):
                    if k in p:
                        p[k] = apply_subs(p[k])
    d['conclusion'] = apply_subs(d.get('conclusion', ''))
    for sp in d.get('signposts', []):
        for k in ('metric', 'signpost', 'target', 'timeframe'):
            if k in sp:
                sp[k] = apply_subs(sp[k])
    for t in d.get('threats', []):
        for k in ('threat', 'triggerPoints'):
            if k in t:
                t[k] = apply_subs(t[k])
    return d


def _run_thesis_infographic(job_id, d, scorecard_data, style_key, mode, gemini_key, detail='full', edit_prompt=None, parent_id=None, show_risk_detail=False, color_scheme=None, include_company=False, template_id=None):
    """Background worker to generate thesis infographic images via Gemini or Pillow."""
    import time as _time
    job = _infographic_jobs[job_id]
    style_def = THESIS_INFOGRAPHIC_STYLES.get(style_key, THESIS_INFOGRAPHIC_STYLES['professional'])
    engine = style_def.get('engine', 'gemini')

    # Load template reference image if specified
    ref_image = None
    if template_id:
        try:
            with get_db() as (conn, cur):
                cur.execute('SELECT reference_image FROM thesis_infographic_templates WHERE id = %s', (template_id,))
                tpl = cur.fetchone()
                if tpl:
                    ref_image = tpl['reference_image']
        except Exception as e:
            print(f"Failed to load template {template_id}: {e}")

    try:
        # Strip quotation marks from all text fields (per user preference)
        d = _strip_quotes_from_data(d)

        if engine == 'pillow':
            # Deterministic Pillow rendering — no Gemini API needed
            job['current'] = 1
            job['progress'] = 50
            # Apply text edits if provided (quick fix for Pillow infographics)
            if edit_prompt:
                d = _apply_pillow_edit(d, edit_prompt)
            pillow_renderers = {
                'precision': _generate_precision_infographic,
                'analyst_brief': _generate_analyst_brief_infographic,
                'quad_grid': _generate_quad_grid_infographic,
            }
            renderer = pillow_renderers.get(style_key, _generate_precision_infographic)
            images = renderer(d, scorecard_data, mode, detail, show_risk_detail=show_risk_detail, color_scheme=color_scheme, include_company=include_company)
            job['images'] = images
        elif mode == '1':
            style_prompt = style_def['prompt']
            job['current'] = 1
            prompt = _build_thesis_infographic_prompt(d, scorecard_data, style_prompt, '1', detail=detail, show_risk_detail=show_risk_detail, color_scheme=color_scheme, include_company=include_company, style_key=style_key)
            if edit_prompt:
                prompt += f"\n\nADDITIONAL INSTRUCTIONS (edit request): {edit_prompt}"
            img = _generate_slide_image(prompt, gemini_key, reference_image=ref_image)
            if img:
                job['images'].append(img)
                job['progress'] = 100
            else:
                job['error'] = 'Image generation failed after retries'
                job['status'] = 'error'
                return
        else:
            style_prompt = style_def['prompt']
            # 3-slide mode
            for i in range(1, 4):
                job['current'] = i
                job['progress'] = int((i - 1) / 3 * 100)
                prompt = _build_thesis_infographic_prompt(d, scorecard_data, style_prompt, '3', slide_num=i, detail=detail, show_risk_detail=show_risk_detail, color_scheme=color_scheme, include_company=include_company, style_key=style_key)
                if edit_prompt:
                    prompt += f"\n\nADDITIONAL INSTRUCTIONS (edit request): {edit_prompt}"
                img = _generate_slide_image(prompt, gemini_key, reference_image=ref_image)
                if img:
                    job['images'].append(img)
                else:
                    job['images'].append(None)  # partial failure
                job['progress'] = int(i / 3 * 100)
                if i < 3:
                    _time.sleep(1)  # small delay between slides

            # Check if all failed
            if all(x is None for x in job['images']):
                job['error'] = 'All image generations failed'
                job['status'] = 'error'
                return

        job['status'] = 'done'
        job['progress'] = 100

        # Persist to history
        try:
            with get_db(commit=True) as (conn, cur):
                cur.execute('''
                    INSERT INTO thesis_infographic_history
                    (ticker, mode, detail, style, slide_images, edit_prompt, parent_id, color_scheme)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id
                ''', (
                    job['ticker'], job['mode'], job.get('detail', 'full'),
                    job['style'], json.dumps(job['images']),
                    edit_prompt, parent_id, job.get('color_scheme', 'default')
                ))
                job['saved_id'] = cur.fetchone()['id']
        except Exception as e:
            print(f"Failed to save infographic to history: {e}")

    except Exception as e:
        print(f"Infographic generation error for job {job_id}: {e}")
        job['error'] = str(e)
        job['status'] = 'error'


def _resolve_infographic_params(mode, detail, style, gemini_key, color_scheme,
                                show_risk_detail, include_company, template_id,
                                thesis_tier_override=None):
    """Resolve infographic params, applying template overrides if needed.
    Returns (mode, detail, style, gemini_key, color_scheme, show_risk_detail, include_company, thesis_tier, error_msg)."""
    # If template specified, load its params as defaults
    if template_id:
        try:
            with get_db() as (_, cur):
                cur.execute('SELECT mode, detail, style, color_scheme, show_risk_detail, include_company FROM thesis_infographic_templates WHERE id = %s', (template_id,))
                tpl = cur.fetchone()
                if tpl:
                    mode = tpl['mode']
                    detail = tpl['detail']
                    style = tpl['style']
                    color_scheme = tpl['color_scheme'] or 'default'
                    show_risk_detail = tpl['show_risk_detail'] or False
                    include_company = tpl['include_company'] or False
        except Exception as e:
            print(f"Template load warning: {e}")

    if style not in THESIS_INFOGRAPHIC_STYLES:
        return None, None, None, None, None, None, None, None, f'Unknown style: {style}'
    style_def = THESIS_INFOGRAPHIC_STYLES[style]
    engine = style_def.get('engine', 'gemini')
    if engine != 'pillow' and not gemini_key:
        return None, None, None, None, None, None, None, None, 'Gemini API key required for infographic generation'

    # Map detail level to thesis tier
    detail_to_tier = {'full': 'detailed', 'summary': 'full', 'simple': 'condensed'}
    thesis_tier = thesis_tier_override or detail_to_tier.get(detail, 'detailed')

    return mode, detail, style, gemini_key, color_scheme, show_risk_detail, include_company, thesis_tier, None


def _launch_infographic_job(ticker, mode, detail, style, gemini_key, color_scheme,
                            show_risk_detail, include_company, thesis_tier, template_id=None):
    """Fetch data and launch a single infographic job for a ticker.
    Returns (job_id, total, error_msg). On error, job_id is None."""
    import threading

    # Fetch analysis
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
        row = cur.fetchone()
    if not row:
        return None, 0, f'No analysis found for {ticker}'

    # Fetch scorecard data
    scorecard_data = None
    with get_db() as (_, cur):
        cur.execute('SELECT scorecard_data FROM thesis_scorecard_data WHERE ticker = %s', (ticker,))
        sc_row = cur.fetchone()
    if sc_row:
        scorecard_data = sc_row['scorecard_data']
        if isinstance(scorecard_data, str):
            try:
                scorecard_data = json.loads(scorecard_data)
            except:
                scorecard_data = None

    row_dict = dict(row)

    # Overlay thesis tier data if not detailed
    try:
        _overlay_thesis_tier(row_dict, thesis_tier)
    except Exception as ce:
        print(f"Thesis tier overlay warning: {ce}")

    d = _parse_analysis_data(row_dict)
    total = 1 if mode == '1' else 3
    job_id = f"infog_{ticker}_{int(time.time()*1000)}"

    _infographic_jobs[job_id] = {
        'status': 'running',
        'ticker': ticker,
        'mode': mode,
        'detail': detail,
        'style': style,
        'color_scheme': color_scheme,
        'progress': 0,
        'current': 0,
        'total': total,
        'images': [],
        'error': None,
        'created_at': time.time(),
    }

    t = threading.Thread(
        target=_run_thesis_infographic,
        args=(job_id, d, scorecard_data, style, mode, gemini_key, detail),
        kwargs={'show_risk_detail': show_risk_detail, 'color_scheme': color_scheme, 'include_company': include_company, 'template_id': template_id},
        daemon=True
    )
    t.start()
    return job_id, total, None


@app.route('/api/thesis-format/infographic', methods=['POST'])
def start_thesis_infographic():
    """Start async thesis infographic generation."""
    try:
        data = request.get_json()
        ticker = data.get('ticker', '').upper()
        if not ticker:
            return jsonify({'error': 'No ticker provided'}), 400

        mode, detail, style, gemini_key, color_scheme, show_risk_detail, include_company, thesis_tier, err = \
            _resolve_infographic_params(
                data.get('mode', '1'), data.get('detail', 'full'), data.get('style', 'professional'),
                data.get('geminiApiKey', ''), data.get('colorScheme', 'default'),
                data.get('showRiskDetail', False), data.get('includeCompanyName', False),
                data.get('templateId'), thesis_tier_override=data.get('thesisTier'))
        if err:
            return jsonify({'error': err}), 400

        job_id, total, launch_err = _launch_infographic_job(
            ticker, mode, detail, style, gemini_key, color_scheme,
            show_risk_detail, include_company, thesis_tier, template_id=data.get('templateId'))
        if launch_err:
            return jsonify({'error': launch_err}), 404

        return jsonify({'jobId': job_id, 'total': total})
    except Exception as e:
        print(f"Error starting infographic: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/batch', methods=['POST'])
def start_thesis_infographic_batch():
    """Start batch infographic generation for multiple tickers."""
    try:
        data = request.get_json()
        tickers = data.get('tickers', [])
        if not tickers:
            return jsonify({'error': 'No tickers provided'}), 400
        if len(tickers) > 30:
            return jsonify({'error': 'Maximum 30 tickers per batch'}), 400

        mode, detail, style, gemini_key, color_scheme, show_risk_detail, include_company, thesis_tier, err = \
            _resolve_infographic_params(
                data.get('mode', '1'), data.get('detail', 'full'), data.get('style', 'professional'),
                data.get('geminiApiKey', ''), data.get('colorScheme', 'default'),
                data.get('showRiskDetail', False), data.get('includeCompanyName', False),
                data.get('templateId'), thesis_tier_override=data.get('thesisTier'))
        if err:
            return jsonify({'error': err}), 400

        job_ids = {}
        errors = {}
        for t in tickers:
            tk = t.upper().strip()
            if not tk:
                continue
            job_id, total, launch_err = _launch_infographic_job(
                tk, mode, detail, style, gemini_key, color_scheme,
                show_risk_detail, include_company, thesis_tier, template_id=data.get('templateId'))
            if launch_err:
                errors[tk] = launch_err
            else:
                job_ids[tk] = job_id

        return jsonify({'jobIds': job_ids, 'errors': errors})
    except Exception as e:
        print(f"Error starting batch infographic: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/batch-status', methods=['POST'])
def poll_thesis_infographic_batch():
    """Poll status of multiple infographic generation jobs."""
    data = request.get_json()
    job_ids = data.get('jobIds', {})  # {TICKER: job_id, ...}

    results = {}
    for ticker, job_id in job_ids.items():
        job = _infographic_jobs.get(job_id)
        if not job:
            results[ticker] = {'status': 'not_found'}
            continue
        results[ticker] = {
            'status': job['status'],
            'progress': job['progress'],
            'error': job.get('error'),
        }
        if job['status'] == 'done':
            results[ticker]['images'] = job['images']
            if job.get('saved_id'):
                results[ticker]['savedId'] = job['saved_id']

    return jsonify(results)


@app.route('/api/thesis-format/infographic/<job_id>', methods=['GET'])
def poll_thesis_infographic(job_id):
    """Poll status of an infographic generation job."""
    # Auto-cleanup old jobs
    now = time.time()
    stale = [k for k, v in _infographic_jobs.items()
             if (v['status'] == 'done' and now - v['created_at'] > 600)
             or (v['status'] == 'error' and now - v['created_at'] > 300)]
    for k in stale:
        del _infographic_jobs[k]

    job = _infographic_jobs.get(job_id)
    if not job:
        return jsonify({'error': 'Job not found or expired'}), 404

    result = {
        'status': job['status'],
        'progress': job['progress'],
        'current': job['current'],
        'total': job['total'],
        'ticker': job['ticker'],
        'mode': job['mode'],
        'detail': job.get('detail', 'full'),
    }
    if job['status'] == 'done':
        result['images'] = job['images']
        if job.get('saved_id'):
            result['savedId'] = job['saved_id']
    if job['status'] == 'error':
        result['error'] = job['error']

    return jsonify(result)


# ============================================
# THESIS INFOGRAPHIC HISTORY ENDPOINTS
# ============================================

@app.route('/api/thesis-format/infographic/history/<ticker>', methods=['GET'])
def get_infographic_history(ticker):
    """List all infographic versions for a ticker (metadata only, no images)."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, ticker, mode, detail, style, color_scheme, edit_prompt, parent_id, created_at,
                       jsonb_array_length(slide_images) as image_count
                FROM thesis_infographic_history
                WHERE ticker = %s
                ORDER BY created_at DESC
            ''', (ticker.upper(),))
            rows = cur.fetchall()
        return jsonify([{
            'id': r['id'],
            'ticker': r['ticker'],
            'mode': r['mode'],
            'detail': r['detail'],
            'style': r['style'],
            'colorScheme': r.get('color_scheme', 'default'),
            'editPrompt': r['edit_prompt'],
            'parentId': r['parent_id'],
            'imageCount': r['image_count'],
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
        } for r in rows])
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/history/<int:history_id>/images', methods=['GET'])
def get_infographic_history_images(history_id):
    """Get full image data for a specific infographic version."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, ticker, mode, detail, style, color_scheme, slide_images, edit_prompt, parent_id, created_at
                FROM thesis_infographic_history WHERE id = %s
            ''', (history_id,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Not found'}), 404
        images = row['slide_images']
        if isinstance(images, str):
            images = json.loads(images)
        return jsonify({
            'id': row['id'],
            'ticker': row['ticker'],
            'mode': row['mode'],
            'detail': row['detail'],
            'style': row['style'],
            'images': images,
            'editPrompt': row['edit_prompt'],
            'parentId': row['parent_id'],
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/history/<int:history_id>', methods=['DELETE'])
def delete_infographic_history(history_id):
    """Delete a specific infographic version."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('UPDATE thesis_infographic_history SET parent_id = NULL WHERE parent_id = %s', (history_id,))
            cur.execute('DELETE FROM thesis_infographic_history WHERE id = %s', (history_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/templates', methods=['GET'])
def list_infographic_templates():
    """List all saved infographic templates."""
    try:
        with get_db() as (conn, cur):
            cur.execute('''
                SELECT id, name, mode, detail, style, color_scheme, show_risk_detail,
                       include_company, reference_image, source_ticker, created_at
                FROM thesis_infographic_templates ORDER BY created_at DESC
            ''')
            rows = cur.fetchall()
        return jsonify([{
            'id': r['id'], 'name': r['name'], 'mode': r['mode'], 'detail': r['detail'],
            'style': r['style'], 'colorScheme': r['color_scheme'],
            'showRiskDetail': r['show_risk_detail'], 'includeCompany': r['include_company'],
            'referenceImage': r['reference_image'], 'sourceTicker': r['source_ticker'],
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None
        } for r in rows])
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/templates', methods=['POST'])
def save_infographic_template():
    """Save an infographic history version as a reusable template."""
    data = request.get_json()
    history_id = data.get('historyId')
    name = (data.get('name') or '').strip()
    if not history_id or not name:
        return jsonify({'error': 'historyId and name are required'}), 400
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('SELECT * FROM thesis_infographic_history WHERE id = %s', (history_id,))
            hist = cur.fetchone()
            if not hist:
                return jsonify({'error': 'History version not found'}), 404
            images = json.loads(hist['slide_images']) if isinstance(hist['slide_images'], str) else hist['slide_images']
            ref_image = images[0] if images else None
            if not ref_image:
                return jsonify({'error': 'No image found in history version'}), 400
            cur.execute('''
                INSERT INTO thesis_infographic_templates
                (name, mode, detail, style, color_scheme, reference_image, source_ticker)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                RETURNING id
            ''', (name, hist['mode'], hist['detail'], hist['style'], hist.get('color_scheme', 'default') or 'default', ref_image, hist['ticker']))
            new_id = cur.fetchone()['id']
        return jsonify({'id': new_id, 'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/templates/<int:template_id>', methods=['DELETE'])
def delete_infographic_template(template_id):
    """Delete a saved infographic template."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM thesis_infographic_templates WHERE id = %s', (template_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/edit', methods=['POST'])
def edit_thesis_infographic():
    """Start async edit of an existing infographic version."""
    try:
        data = request.get_json()
        parent_id = data.get('parentId')
        edit_prompt = data.get('editPrompt', '')
        gemini_key = data.get('geminiApiKey', '')

        if not parent_id:
            return jsonify({'error': 'parentId required'}), 400
        if not edit_prompt:
            return jsonify({'error': 'editPrompt required'}), 400

        # Fetch parent version
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM thesis_infographic_history WHERE id = %s', (parent_id,))
            parent = cur.fetchone()
        if not parent:
            return jsonify({'error': 'Parent version not found'}), 404

        # Check if Gemini key is needed (not required for Pillow-based styles)
        style_key_parent = parent['style']
        parent_engine = THESIS_INFOGRAPHIC_STYLES.get(style_key_parent, {}).get('engine', 'gemini')
        if parent_engine != 'pillow' and not gemini_key:
            return jsonify({'error': 'Gemini API key required for this style'}), 400

        ticker = parent['ticker']
        mode = parent['mode']
        detail_level = parent['detail']
        style_key = parent['style']
        color_scheme = parent.get('color_scheme', 'default') or 'default'

        # Map detail level to thesis tier for data sourcing
        detail_to_tier = {'full': 'detailed', 'summary': 'full', 'simple': 'condensed'}
        thesis_tier = detail_to_tier.get(detail_level, 'detailed')

        # Fetch analysis data
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'No analysis found'}), 404

        scorecard_data = None
        with get_db() as (_, cur):
            cur.execute('SELECT scorecard_data FROM thesis_scorecard_data WHERE ticker = %s', (ticker,))
            sc_row = cur.fetchone()
        if sc_row:
            scorecard_data = sc_row['scorecard_data']
            if isinstance(scorecard_data, str):
                try:
                    scorecard_data = json.loads(scorecard_data)
                except:
                    scorecard_data = None

        row_dict = dict(row)
        try:
            _overlay_thesis_tier(row_dict, thesis_tier)
        except Exception as ce:
            print(f"Edit thesis tier overlay warning: {ce}")
        d = _parse_analysis_data(row_dict)
        total = 1 if mode == '1' else 3
        job_id = f"infog_edit_{ticker}_{int(time.time()*1000)}"

        _infographic_jobs[job_id] = {
            'status': 'running',
            'ticker': ticker,
            'mode': mode,
            'detail': detail_level,
            'style': style_key,
            'color_scheme': color_scheme,
            'progress': 0,
            'current': 0,
            'total': total,
            'images': [],
            'error': None,
            'created_at': time.time(),
        }

        import threading
        t = threading.Thread(
            target=_run_thesis_infographic,
            args=(job_id, d, scorecard_data, style_key, mode, gemini_key, detail_level, edit_prompt, parent_id),
            kwargs={'color_scheme': color_scheme},
            daemon=True
        )
        t.start()

        return jsonify({'jobId': job_id, 'total': total})
    except Exception as e:
        print(f"Error starting infographic edit: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis-format/infographic/export-pdf', methods=['POST'])
def export_infographic_pdf():
    """Combine infographic images into a single PDF."""
    try:
        import img2pdf
        data = request.get_json()
        images = data.get('images', [])
        ticker = data.get('ticker', 'Infographic')
        if not images or all(x is None for x in images):
            return jsonify({'error': 'No images provided'}), 400
        img_bytes_list = []
        for img in images:
            if img:
                img_bytes_list.append(base64.b64decode(img))
        if not img_bytes_list:
            return jsonify({'error': 'No valid images'}), 400
        pdf_bytes = img2pdf.convert(img_bytes_list)
        pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
        return jsonify({'success': True, 'fileData': pdf_b64, 'filename': f'{ticker}_Thesis_Infographic.pdf'})
    except Exception as e:
        print(f"Error exporting infographic PDF: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-infographic', methods=['POST'])
def email_infographic():
    """Email infographic images as inline attachments."""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    from email.mime.image import MIMEImage

    try:
        data = request.get_json()
        images = data.get('images', [])
        ticker = data.get('ticker', '')
        mode = data.get('mode', '1')
        detail = data.get('detail', 'full')
        style = data.get('style', '')
        recipient = data.get('email', '')
        subject = data.get('customSubject', f'{ticker} — Thesis Infographic')
        smtp_config = data.get('smtpConfig', {})

        if not recipient:
            return jsonify({'error': 'Recipient email is required'}), 400
        if not images or all(x is None for x in images):
            return jsonify({'error': 'No images to send'}), 400

        use_gmail = smtp_config.get('use_gmail', True)
        gmail_user = smtp_config.get('gmail_user', '')
        gmail_password = smtp_config.get('gmail_app_password', '')
        from_email = smtp_config.get('from_email', gmail_user)

        if use_gmail and (not gmail_user or not gmail_password):
            return jsonify({'error': 'Gmail credentials required'}), 400

        # Build HTML with inline CID images
        slide_labels = ['Investment Thesis', 'Key Signposts', 'Risk Assessment']
        detail_label = 'Summary' if detail == 'summary' else 'Full Detail'
        html_parts = [
            f'<div style="font-family:Arial,sans-serif;max-width:800px;margin:0 auto;">',
            f'<h2 style="color:#1e293b;border-bottom:2px solid #6366f1;padding-bottom:8px;">{ticker} — Thesis Infographic</h2>',
            f'<p style="color:#64748b;font-size:13px;">{mode}-slide &middot; {detail_label} &middot; {style}</p>',
        ]
        cid_images = []
        for idx, img in enumerate(images):
            if not img:
                continue
            cid = f'infographic_{idx}'
            cid_images.append((cid, img))
            label = slide_labels[idx] if mode == '3' and idx < len(slide_labels) else ''
            if label:
                html_parts.append(f'<h3 style="color:#475569;margin-top:20px;">{label}</h3>')
            html_parts.append(f'<img src="cid:{cid}" style="max-width:100%;border:1px solid #e2e8f0;border-radius:8px;margin-bottom:16px;" />')

        html_parts.append('</div>')
        html_body = '\n'.join(html_parts)

        msg = MIMEMultipart('related')
        msg['From'] = from_email
        msg['To'] = recipient
        msg['Subject'] = subject

        html_part = MIMEText(html_body, 'html')
        msg.attach(html_part)

        for cid, img_b64 in cid_images:
            img_data = base64.b64decode(img_b64)
            mime_img = MIMEImage(img_data, _subtype='png')
            mime_img.add_header('Content-ID', f'<{cid}>')
            mime_img.add_header('Content-Disposition', 'inline', filename=f'{ticker}_{cid}.png')
            msg.attach(mime_img)

        if use_gmail:
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(gmail_user, gmail_password)
                server.send_message(msg)

        return jsonify({'success': True, 'message': f'Infographic emailed to {recipient}'})

    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'Gmail authentication failed. Check your email and app password.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        print(f"Error emailing infographic: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# DOCUMENT STORAGE ENDPOINTS
# ============================================

@app.route('/api/documents/<ticker>', methods=['GET'])
def get_documents(ticker):
    """Get all stored documents for a ticker"""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT filename, file_type, mime_type, metadata, file_size, created_at
                FROM document_files
                WHERE ticker = %s
                ORDER BY created_at DESC
            ''', (ticker.upper(),))
            docs = cur.fetchall()

        return jsonify({
            'documents': [{
                'filename': d['filename'],
                'fileType': d['file_type'],
                'mimeType': d['mime_type'],
                'metadata': d['metadata'] or {},
                'fileSize': d['file_size'],
                'createdAt': d['created_at'].isoformat() if d['created_at'] else None,
                'stored': True
            } for d in docs]
        })
    except Exception as e:
        print(f"Error getting documents: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/documents/<ticker>/content', methods=['GET'])
def get_documents_with_content(ticker):
    """Get all stored documents with file content for re-analysis"""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT filename, file_data, file_type, mime_type, metadata
                FROM document_files
                WHERE ticker = %s
            ''', (ticker.upper(),))
            docs = cur.fetchall()

        return jsonify({
            'documents': [{
                'filename': d['filename'],
                'fileData': d['file_data'],
                'fileType': d['file_type'],
                'mimeType': d['mime_type'],
                'metadata': d['metadata'] or {},
                'stored': True
            } for d in docs]
        })
    except Exception as e:
        print(f"Error getting document content: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/documents/save', methods=['POST'])
def save_documents():
    """Save document files to database for a ticker"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()
        documents = data.get('documents', [])
        
        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400
        
        if not documents:
            return jsonify({'error': 'No documents provided'}), 400

        with get_db(commit=True) as (_, cur):
            saved_count = 0
            for doc in documents:
                filename = doc.get('filename')
                file_data = doc.get('fileData')
                file_type = doc.get('fileType', 'pdf')
                mime_type = doc.get('mimeType', 'application/pdf')
                metadata = doc.get('metadata', {})

                if not filename or not file_data:
                    continue

                # Calculate approximate file size (base64 is ~1.33x original)
                file_size = len(file_data) * 3 // 4

                cur.execute('''
                    INSERT INTO document_files (ticker, filename, file_data, file_type, mime_type, metadata, file_size)
                    VALUES (%s, %s, %s, %s, %s, %s, %s)
                    ON CONFLICT (ticker, filename)
                    DO UPDATE SET
                        file_data = EXCLUDED.file_data,
                        file_type = EXCLUDED.file_type,
                        mime_type = EXCLUDED.mime_type,
                        metadata = EXCLUDED.metadata,
                        file_size = EXCLUDED.file_size
                ''', (ticker, filename, file_data, file_type, mime_type, json.dumps(metadata), file_size))
                saved_count += 1

        return jsonify({'success': True, 'savedCount': saved_count})
    except Exception as e:
        print(f"Error saving documents: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/documents/delete', methods=['POST'])
def delete_document():
    """Delete a specific document for a ticker"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()
        filename = data.get('filename')
        
        if not ticker or not filename:
            return jsonify({'error': 'Ticker and filename are required'}), 400

        with get_db(commit=True) as (_, cur):
            cur.execute('''
                DELETE FROM document_files
                WHERE ticker = %s AND filename = %s
            ''', (ticker, filename))
            deleted = cur.rowcount > 0

        return jsonify({'success': True, 'deleted': deleted})
    except Exception as e:
        print(f"Error deleting document: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/documents/delete-all', methods=['POST'])
def delete_all_documents():
    """Delete all documents for a ticker"""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper()
        
        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400

        with get_db(commit=True) as (_, cur):
            cur.execute('DELETE FROM document_files WHERE ticker = %s', (ticker,))
            deleted_count = cur.rowcount

        return jsonify({'success': True, 'deletedCount': deleted_count})
    except Exception as e:
        print(f"Error deleting all documents: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/documents/storage-stats', methods=['GET'])
def get_storage_stats():
    """Get storage statistics"""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT
                    ticker,
                    COUNT(*) as doc_count,
                    SUM(file_size) as total_size
                FROM document_files
                GROUP BY ticker
                ORDER BY total_size DESC
            ''')
            stats = cur.fetchall()

            cur.execute('SELECT SUM(file_size) as total FROM document_files')
            total = cur.fetchone()

        return jsonify({
            'byTicker': [{
                'ticker': s['ticker'],
                'docCount': s['doc_count'],
                'totalSize': s['total_size'] or 0
            } for s in stats],
            'totalSize': total['total'] or 0 if total else 0
        })
    except Exception as e:
        print(f"Error getting storage stats: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# ANTHROPIC API PROXY ENDPOINTS
# ============================================

@app.route('/api/chat', methods=['POST'])
def chat():
    """Proxy chat requests to Anthropic API"""
    try:
        data = request.json
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('api_key', '')
        messages = data.get('messages', [])
        system = data.get('system', '')

        if not api_key:
            return jsonify({'error': 'No API key provided. Please add your API key in Settings.'}), 400

        result = call_llm(
            messages=messages,
            system=system,
            tier="standard",
            max_tokens=4096,
            timeout=120,
            anthropic_api_key=api_key,
        )

        return jsonify({
            'response': result["text"],
            'usage': result["usage"]
        })

    except LLMError as e:
        return jsonify({'error': str(e)}), 502
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyze-multi', methods=['POST'])
def analyze_multi():
    """Analyze multiple PDF documents and generate investment thesis"""
    try:
        data = request.json
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '')
        documents = data.get('documents', [])
        existing_analysis = data.get('existingAnalysis')
        historical_weights = data.get('historicalWeights', [])
        weighting_config = data.get('weightingConfig', {})

        if not api_key:
            return jsonify({'error': 'No API key provided. Please add your API key in Settings.'}), 400
        
        if not documents:
            return jsonify({'error': 'No documents provided'}), 400
        
        # Filter enabled documents
        enabled_docs = [d for d in documents if d.get('enabled', True)]
        
        if not enabled_docs:
            return jsonify({'error': 'No enabled documents'}), 400
        
        # Build the content array for Claude
        content = []
        
        # Check if using simple weighting mode
        simple_mode = weighting_config.get('mode') == 'simple'
        existing_weight = weighting_config.get('existingAnalysisWeight', 70) if simple_mode else None
        new_docs_weight = weighting_config.get('newDocsWeight', 30) if simple_mode else None
        
        # Build weighting information string
        weight_info = ""
        
        # Pre-categorize documents for simple mode
        truly_new_docs = [d for d in enabled_docs if d.get('isNew', True)]
        stored_existing_docs = [d for d in enabled_docs if not d.get('isNew', True)]
        
        if simple_mode and existing_analysis:
            # Simple mode: clear instruction about preservation vs updates
            weight_info = f"""=== ANALYSIS UPDATE MODE: SIMPLE WEIGHTING ===

PRESERVATION RATIO: {existing_weight}% existing / {new_docs_weight}% new

This means:
- PRESERVE {existing_weight}% of the existing thesis, pillars, signposts, and threats
- Allow only {new_docs_weight}% worth of modifications from the new document(s)
- The new document(s) are SUPPLEMENTARY, not replacements

"""
            if stored_existing_docs:
                weight_info += "EXISTING DOCUMENTS (re-uploaded for reference, part of the preserved analysis):\n"
                for doc in stored_existing_docs:
                    doc_name = doc.get('filename', 'document.pdf')
                    weight_info += f"- {doc_name}\n"
                weight_info += "\n"
            
            if truly_new_docs:
                per_new_doc_weight = new_docs_weight / len(truly_new_docs)
                weight_info += f"NEW DOCUMENTS (sharing the {new_docs_weight}% update allocation):\n"
                for doc in truly_new_docs:
                    doc_name = doc.get('filename', 'document.pdf')
                    weight_info += f"- {doc_name} ({round(per_new_doc_weight)}% weight)\n"
                weight_info += "\n"
            
            weight_info += f"""Remember: With {existing_weight}% preservation, you should keep most existing content intact.
Only add minor refinements, new data points, or small additions from the new document(s).
Do NOT rewrite or fundamentally change the existing analysis.

"""
        elif simple_mode and not existing_analysis:
            # Simple mode but new analysis - just list documents
            weight_info = "DOCUMENT WEIGHTING:\n\n"
            weight_info += "NEW DOCUMENTS (being analyzed now):\n"
            for doc in enabled_docs:
                doc_name = doc.get('filename', 'document.pdf')
                weight_info += f"- {doc_name}\n"
            weight_info += "\n"
        elif not simple_mode:
            # Advanced mode: per-document weights
            # Calculate total weight including both new and historical docs
            new_doc_weight = sum(doc.get('weight', 1) for doc in enabled_docs)
            hist_doc_weight = sum(hw.get('weight', 1) for hw in historical_weights)
            total_weight = new_doc_weight + hist_doc_weight
            
            # Historical documents (from existing analysis)
            if historical_weights:
                weight_info += "PREVIOUSLY ANALYZED DOCUMENTS (their insights are in the existing analysis):\n"
                for hw in historical_weights:
                    hw_name = hw.get('filename', 'document')
                    hw_weight = hw.get('weight', 1)
                    hw_pct = round((hw_weight / total_weight) * 100) if total_weight > 0 else 0
                    weight_info += f"- {hw_name}: {hw_pct}% weight\n"
                weight_info += "\n"
            
            # New documents being analyzed now
            weight_info += "NEW DOCUMENTS (being analyzed now):\n"
            for doc in enabled_docs:
                doc_name = doc.get('filename', 'document.pdf')
                doc_weight = doc.get('weight', 1)
                doc_pct = round((doc_weight / total_weight) * 100) if total_weight > 0 else 0
                weight_info += f"- {doc_name}: {doc_pct}% weight\n"
            
            weight_info += "\nWhen synthesizing the analysis:\n"
            weight_info += "- Give MORE emphasis to higher-weighted documents\n"
            weight_info += "- If updating existing analysis, respect the weights of previously analyzed documents\n"
            weight_info += "- Higher-weighted historical docs = keep more of their conclusions in the existing analysis\n"
        
        content.append({
            "type": "text",
            "text": weight_info
        })
        
        # Calculate total weight for document headers (use simple mode weight or calculated weight)
        if simple_mode and existing_analysis:
            # In simple mode, only truly NEW documents share the new_docs_weight
            # Stored existing documents are re-uploaded for context but shouldn't count as "new"
            truly_new_docs = [d for d in enabled_docs if d.get('isNew', True)]
            stored_docs = [d for d in enabled_docs if not d.get('isNew', True)]
            
            per_new_doc_weight = new_docs_weight / len(truly_new_docs) if truly_new_docs else 0
        else:
            new_doc_weight = sum(doc.get('weight', 1) for doc in enabled_docs)
            hist_doc_weight = sum(hw.get('weight', 1) for hw in historical_weights)
            total_weight = new_doc_weight + hist_doc_weight
        
        # Add each document
        for doc in enabled_docs:
            doc_content = doc.get('fileData', '')
            doc_name = doc.get('filename', 'document.pdf')
            doc_type = doc.get('fileType', 'pdf')
            mime_type = doc.get('mimeType', 'application/pdf')
            is_new = doc.get('isNew', True)
            
            if simple_mode and existing_analysis:
                if is_new:
                    doc_pct = round(per_new_doc_weight)
                    doc_header = f"\n=== NEW DOCUMENT (Supplementary - {doc_pct}% update weight): {doc_name} ==="
                else:
                    # Stored existing document - re-uploaded for reference, part of the existing analysis
                    doc_header = f"\n=== EXISTING DOCUMENT (Reference - part of {existing_weight}% preserved analysis): {doc_name} ==="
            else:
                doc_weight = doc.get('weight', 1)
                doc_pct = round((doc_weight / total_weight) * 100) if total_weight > 0 else 0
                doc_header = f"\n=== DOCUMENT: {doc_name} (Weight: {doc_pct}%) ==="
            
            if not doc_content:
                continue
            
            # Add document header with weight
            content.append({
                "type": "text",
                "text": doc_header
            })
                
            if doc_type == 'pdf':
                content.append({
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": doc_content
                    }
                })
            elif doc_type == 'image':
                content.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": mime_type or "image/png",
                        "data": doc_content
                    }
                })
            else:
                try:
                    decoded_text = base64.b64decode(doc_content).decode('utf-8')
                    content.append({
                        "type": "text",
                        "text": decoded_text
                    })
                except:
                    continue
        
        if not content:
            return jsonify({'error': 'No valid documents to analyze'}), 400
        
        # Add the analysis prompt
        analysis_prompt = """Analyze these broker research documents and create a concise, high-conviction investment analysis.

CONCISENESS RULES (the final output MUST fit in 2-3 printed pages):
- thesis.summary: 2-3 sentences MAX. One crisp paragraph capturing WHY we own it.
- pillars: 3-5 pillars ONLY. Each description is 1-2 sentences — punchy, not exhaustive.
- signposts: 4-6 signposts MAX. Combine related metrics rather than listing every variant.
- threats: 3-5 threats MAX. Each threat is one sentence. triggerPoints is one sentence.
- Prioritize the MOST important and differentiated insights. Omit generic/obvious points.
- Do NOT pad with boilerplate or repeat the same point across pillars.

Return a JSON object with this exact structure:
{
    "ticker": "STOCK_TICKER",
    "company": "Company Name",
    "thesis": {
        "summary": "2-3 sentence investment thesis summary",
        "pillars": [
            {"title": "Pillar Title", "description": "1-2 sentence explanation", "confidence": "High/Medium/Low", "sources": [{"filename": "Document name", "excerpt": "Brief supporting quote"}]}
        ]
    },
    "signposts": [
        {"metric": "Key metric name", "target": "Target value", "timeframe": "When to expect", "category": "Financial/Operational/Strategic/Market", "confidence": "High/Medium/Low", "sources": [{"filename": "Document name", "excerpt": "Brief supporting quote"}]}
    ],
    "threats": [
        {"threat": "Risk factor in one sentence", "triggerPoints": "One sentence on what to watch for", "sources": [{"filename": "Document name", "excerpt": "Brief supporting quote"}]}
    ],
    "documentMetadata": [
        {"filename": "exact_filename.pdf", "docType": "broker_report", "source": "Citi", "publishDate": "YYYY-MM-DD", "authors": ["Analyst Name"], "title": "Report Title"},
        {"filename": "transcript.pdf", "docType": "earnings_call", "source": "Company Name", "publishDate": "YYYY-MM-DD", "quarter": "Q3 2025", "title": "Q3 2025 Earnings Call"},
        {"filename": "email_screenshot.png", "docType": "email", "source": "Sender Name/Org", "publishDate": "YYYY-MM-DD", "title": "Email Subject"}
    ]
}

DOCUMENT METADATA EXTRACTION (CRITICAL):
For EACH document provided, identify the document type and extract appropriate metadata:

**For Broker Reports:**
- "docType": "broker_report"
- "source": Investment bank/broker name (e.g., "Citi", "Morgan Stanley", "Goldman Sachs", "Wolfe Research")
- "publishDate": Report date in YYYY-MM-DD format
- "authors": Array of analyst names
- "title": Report title/headline

**For Earnings Call Transcripts:**
- "docType": "earnings_call"  
- "source": Company name (e.g., "Union Pacific", "Apple Inc.")
- "publishDate": Call date in YYYY-MM-DD format
- "quarter": Fiscal quarter (e.g., "Q3 2025", "FY 2025")
- "title": e.g., "Q3 2025 Earnings Call Transcript"

**For SEC Filings (10-K, 10-Q, 8-K):**
- "docType": "sec_filing"
- "source": Company name
- "publishDate": Filing date in YYYY-MM-DD format
- "filingType": "10-K", "10-Q", "8-K", etc.
- "title": Filing description

**For Emails/Email Screenshots:**
- "docType": "email"
- "source": Sender name or organization
- "publishDate": Email date in YYYY-MM-DD format
- "title": Email subject line if visible

**For Company Presentations:**
- "docType": "presentation"
- "source": Company name or presenting organization
- "publishDate": Presentation date in YYYY-MM-DD format
- "title": Presentation title

**For News Articles:**
- "docType": "news"
- "source": Publication name (e.g., "Wall Street Journal", "Reuters")
- "publishDate": Article date in YYYY-MM-DD format
- "authors": Array of journalist names if visible
- "title": Article headline

**For Screenshots/Images (if content type unclear):**
- "docType": "screenshot"
- "source": Infer source if visible in image
- "publishDate": Infer date if visible, otherwise null
- "title": Brief description of content

Always include "filename" with the exact filename provided.

IMPORTANT STYLE RULES:
- Do NOT reference any sellside broker names (e.g., "Goldman Sachs believes...", "According to Morgan Stanley...")
- Do NOT reference specific analyst names
- Do NOT include specific broker price targets
- Write as independent analysis that synthesizes the information without attribution to sources in the prose
- The output should read like original independent research, not a summary of broker views

DOCUMENT WEIGHTING:
- Each document has an assigned weight percentage shown at the start
- Give MORE emphasis to higher-weighted documents when forming conclusions
- Higher-weighted documents should have more influence on the thesis, signposts, and threats
- If documents conflict, prefer the view from the higher-weighted document

Focus on:
1. Why own this stock? (Investment Thesis) - include confidence level and source citations
2. What are we looking for? (Signposts - specific KPIs, events, milestones with metric names)
3. Where can we be wrong? (Threats - bear case scenarios with trigger points to watch)

For each pillar, signpost, and threat, include:
- "sources": Array of source documents that support this point, with filename and a brief excerpt
- Use the actual document filenames provided in the analysis

Return ONLY valid JSON, no markdown, no explanation."""

        if existing_analysis:
            # Build weighting instruction specific to the mode
            if simple_mode:
                weighting_instruction = f"""
CRITICAL WEIGHTING INSTRUCTION (SIMPLE MODE):
You MUST preserve {existing_weight}% of the existing analysis. The new documents can only contribute {new_docs_weight}% worth of changes.

What this means:
- KEEP {existing_weight}% of the existing thesis, pillars, signposts, and threats UNCHANGED
- Only make MINOR refinements or additions based on the new document(s)
- Do NOT fundamentally rewrite or replace the existing analysis
- Do NOT treat the new document as a "primary source" - it is a SUPPLEMENTARY source
- The new document should ADD to or SLIGHTLY REFINE the existing analysis, not replace it

Example of correct behavior with {existing_weight}% existing / {new_docs_weight}% new:
- If existing thesis has 3 pillars, keep all 3, maybe slightly update wording or add a 4th minor pillar
- If existing has 5 signposts, keep them mostly intact, maybe add 1-2 new ones or update targets slightly
- Do NOT remove or majorly rewrite existing content unless it's factually contradicted

In the "changes" array, describe what minor updates were made, NOT that you've rewritten the analysis.
"""
            else:
                weighting_instruction = """
DOCUMENT WEIGHTING:
- Each document has an assigned weight percentage shown at the start
- Give MORE emphasis to higher-weighted documents when forming conclusions
- Higher-weighted documents should have more influence on the thesis, signposts, and threats
- If documents conflict, prefer the view from the higher-weighted document
"""
            
            analysis_prompt = f"""Update this existing analysis with new information from the documents.

Existing Analysis:
{json.dumps(existing_analysis, indent=2)}

{weighting_instruction}

CONCISENESS RULES (the final output MUST fit in 2-3 printed pages):
- thesis.summary: 2-3 sentences MAX.
- pillars: 3-5 pillars ONLY. Each description is 1-2 sentences.
- signposts: 4-6 signposts MAX. Combine related metrics rather than listing every variant.
- threats: 3-5 threats MAX. Each threat + triggerPoints is one sentence each.
- If the existing analysis is bloated (too many pillars/signposts/threats), CONSOLIDATE it now.
  Merge overlapping pillars. Combine related signposts. Drop the least important threats.
- Prioritize the MOST important insights. Omit generic/obvious points.

Review the new documents and:
1. Update or confirm the investment thesis (respecting the weighting above)
2. Add any new signposts or update existing ones — consolidate to 4-6 total
3. Add any new threats or update existing ones — consolidate to 3-5 total
4. Note what has changed in the "changes" array
5. Update sources for each point based on all documents analyzed
6. Extract metadata for ALL documents (both new and from existing analysis)

DOCUMENT METADATA EXTRACTION:
For EACH document (new AND previously analyzed), extract metadata:
**Broker Reports:** docType="broker_report", source=Broker name, publishDate, authors, title
**Earnings Calls:** docType="earnings_call", source=Company, publishDate, quarter, title
**SEC Filings:** docType="sec_filing", source=Company, publishDate, filingType, title
**Emails:** docType="email", source=Sender, publishDate, title=Subject
**Presentations:** docType="presentation", source=Company/Org, publishDate, title
**News:** docType="news", source=Publication, publishDate, authors, title
**Screenshots:** docType="screenshot", source=Inferred, publishDate=if visible, title=description

IMPORTANT STYLE RULES:
- Do NOT reference any sellside broker names or specific analyst names
- Do NOT include specific broker price targets
- Write as independent analysis without attribution to sources in the prose

For each pillar, signpost, and threat, include "sources" and "confidence".

Return the updated analysis as JSON with the same structure (including "documentMetadata" array), plus a "changes" array.

Return ONLY valid JSON, no markdown, no explanation."""

        content.append({
            "type": "text",
            "text": analysis_prompt
        })
        
        # Use Anthropic SDK streaming directly (no fallback chain) to:
        # 1. Prevent Render 502 timeouts via heartbeat spaces during streaming
        # 2. Avoid fallback to OpenAI/Gemini which can't properly handle PDF documents
        # JSON.parse() ignores leading whitespace, so the frontend parses normally.
        def generate():
            try:
                import httpx as _httpx
                client = anthropic.Anthropic(api_key=api_key, timeout=_httpx.Timeout(300.0, connect=30.0))
                result_text = ""
                kwargs = {
                    "model": "claude-sonnet-4-5-20250929",
                    "max_tokens": 64000,
                    "messages": [{'role': 'user', 'content': content}],
                    "system": "You are an expert equity research analyst. Provide institutional-quality investment analysis that is CONCISE: 2-3 printed pages max. Limit to 3-5 pillars, 4-6 signposts, 3-5 threats, each described in 1-2 sentences. Prioritize the most important insights and consolidate related points. Always respond with valid JSON only.",
                }
                usage_data = {}
                with client.messages.stream(**kwargs) as stream:
                    for text in stream.text_stream:
                        result_text += text
                        yield b' '
                    final_msg = stream.get_final_message()
                    usage_data = {
                        "input_tokens": final_msg.usage.input_tokens,
                        "output_tokens": final_msg.usage.output_tokens,
                    }
            except Exception as e:
                print(f"[analyze-multi] Anthropic streaming failed: {type(e).__name__}: {e}")
                yield json.dumps({'error': f'Analysis failed: {str(e)}'}).encode()
                return

            try:
                cleaned = result_text.strip()
                if cleaned.startswith('```'):
                    cleaned = cleaned.split('\n', 1)[1]
                if cleaned.endswith('```'):
                    cleaned = cleaned.rsplit('\n', 1)[0]
                if cleaned.startswith('json'):
                    cleaned = cleaned[4:].strip()

                try:
                    analysis = json.loads(cleaned)
                except json.JSONDecodeError:
                    print(f"[analyze-multi] JSON parse failed, attempting repair...")
                    analysis = _repair_truncated_json(cleaned)
                changes = analysis.pop('changes', [])
                document_metadata = analysis.pop('documentMetadata', [])

                yield json.dumps({
                    'analysis': analysis,
                    'changes': changes,
                    'documentMetadata': document_metadata,
                    'usage': usage_data
                }).encode()
            except json.JSONDecodeError as e:
                yield json.dumps({
                    'error': f'Failed to parse analysis: {str(e)}',
                    'raw_response': result_text
                }).encode()
            except Exception as e:
                yield json.dumps({'error': f'Server error: {str(e)}'}).encode()

        return Response(generate(), mimetype='application/json')
    except Exception as e:
        import traceback
        print(f"Error in analyze-multi: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Server error: {str(e)}'}), 500


@app.route('/api/parse', methods=['POST'])
def parse():
    """Use Claude to intelligently parse stock analysis into sections"""
    try:
        data = request.json
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('api_key', '')
        content = data.get('content', '')

        if not api_key:
            return jsonify({'error': 'No API key provided. Please add your API key in Settings.'}), 400

        result = call_llm(
            messages=[{'role': 'user', 'content': content}],
            system='You are a precise JSON extractor. Extract content into the exact JSON format requested. Return ONLY valid JSON with no markdown formatting, no code blocks, no explanation - just the raw JSON object.',
            tier="fast",
            max_tokens=4096,
            timeout=120,
            anthropic_api_key=api_key,
        )

        return jsonify({
            'response': result["text"],
            'usage': result["usage"]
        })

    except LLMError as e:
        return jsonify({'error': str(e)}), 502
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============================================
# EMAIL ENDPOINTS
# ============================================

@app.route('/api/email', methods=['POST'])
def send_email():
    """Send email via SMTP"""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    
    try:
        data = request.json
        smtp_server = data.get('smtp_server')
        smtp_port = data.get('smtp_port', 587)
        email = data.get('email')
        password = data.get('password')
        recipient = data.get('recipient')
        subject = data.get('subject')
        body = data.get('body')
        
        if not all([smtp_server, email, password, recipient, subject, body]):
            return jsonify({'error': 'Missing required email fields'}), 400
        
        msg = MIMEMultipart()
        msg['From'] = email
        msg['To'] = recipient
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'plain'))
        
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(email, password)
            server.send_message(msg)
        
        return jsonify({'success': True, 'message': 'Email sent successfully'})
        
    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'SMTP authentication failed. Check your email and password/app password.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-overview', methods=['POST'])
def send_overview_email():
    """Send Overview email via SMTP with HTML support"""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    
    try:
        data = request.json
        ticker = data.get('ticker', '')
        company_name = data.get('companyName', '')
        html_body = data.get('htmlBody', '')
        recipient = data.get('email')
        subject = data.get('customSubject', f'{ticker} - Stock Overview')
        smtp_config = data.get('smtpConfig', {})
        
        use_gmail = smtp_config.get('use_gmail', True)
        gmail_user = smtp_config.get('gmail_user', '')
        gmail_password = smtp_config.get('gmail_app_password', '')
        from_email = smtp_config.get('from_email', gmail_user)
        
        if not recipient:
            return jsonify({'error': 'Recipient email is required'}), 400
        
        if use_gmail and (not gmail_user or not gmail_password):
            return jsonify({'error': 'Gmail credentials required'}), 400
        
        msg = MIMEMultipart('alternative')
        msg['From'] = from_email
        msg['To'] = recipient
        msg['Subject'] = subject
        
        plain_text = html_body.replace('<h1>', '\n').replace('</h1>', '\n' + '='*50 + '\n')
        plain_text = plain_text.replace('<h2>', '\n\n').replace('</h2>', '\n' + '-'*30 + '\n')
        plain_text = plain_text.replace('<p>', '').replace('</p>', '\n')
        plain_text = plain_text.replace('<br>', '\n').replace('<em>', '').replace('</em>', '')
        
        msg.attach(MIMEText(plain_text, 'plain'))
        msg.attach(MIMEText(html_body, 'html'))
        
        if use_gmail:
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(gmail_user, gmail_password)
                server.send_message(msg)
        
        return jsonify({'success': True, 'message': 'Overview email sent successfully'})
        
    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'Gmail authentication failed. Check your email and app password.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-analysis', methods=['POST'])
def send_analysis_email():
    """Send Analysis email via SMTP with HTML formatting"""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    
    try:
        data = request.json
        analysis = data.get('analysis', {})
        recipient = data.get('email')
        smtp_config = data.get('smtpConfig', {})
        
        ticker = analysis.get('ticker', 'Stock')
        company = analysis.get('company', '')
        
        # Default subject if not provided
        default_subject = f"{ticker} - Investment Analysis"
        subject = data.get('customSubject') or default_subject
        
        use_gmail = smtp_config.get('use_gmail', True)
        gmail_user = smtp_config.get('gmail_user', '')
        gmail_password = smtp_config.get('gmail_app_password', '')
        from_email = smtp_config.get('from_email', gmail_user)
        
        if not recipient:
            return jsonify({'error': 'Recipient email is required'}), 400
        
        if use_gmail and (not gmail_user or not gmail_password):
            return jsonify({'error': 'Gmail credentials required'}), 400
        
        thesis = analysis.get('thesis', {})
        signposts = analysis.get('signposts', [])
        threats = analysis.get('threats', [])
        
        # Build HTML email
        html_body = f"""
<html>
<body style="font-family: Arial, sans-serif; line-height: 1.6; color: #333; max-width: 700px;">
    <h1 style="color: #1a365d; border-bottom: 2px solid #2c5282; padding-bottom: 10px;">{ticker} - {company}</h1>
    
    <h2 style="color: #2c5282; margin-top: 25px;">1. Investment Thesis</h2>
    <p style="margin-left: 20px;">{thesis.get('summary', 'N/A')}</p>
"""
        
        if thesis.get('pillars'):
            html_body += '<ul style="margin-left: 20px;">'
            for pillar in thesis['pillars']:
                title = pillar.get('pillar', pillar.get('title', ''))
                desc = pillar.get('detail', pillar.get('description', ''))
                html_body += f'<li style="margin-bottom: 8px;"><strong>{title}:</strong> {desc}</li>'
            html_body += '</ul>'
        
        html_body += '<h2 style="color: #2c5282; margin-top: 25px;">2. Signposts (What We\'re Watching)</h2>'
        html_body += '<ul style="margin-left: 20px;">'
        for sp in signposts:
            metric = sp.get('metric', sp.get('signpost', ''))
            target = sp.get('target', '')
            timeframe = sp.get('timeframe', '')
            html_body += f'<li style="margin-bottom: 8px;"><strong>{metric}:</strong> {target}'
            if timeframe:
                html_body += f' <em>({timeframe})</em>'
            html_body += '</li>'
        html_body += '</ul>'
        
        html_body += '<h2 style="color: #2c5282; margin-top: 25px;">3. Thesis Threats (Where We Can Be Wrong)</h2>'
        html_body += '<ul style="margin-left: 20px;">'
        for threat in threats:
            threat_desc = threat.get('threat', '')
            triggers = threat.get('triggerPoints', '')
            html_body += f'<li style="margin-bottom: 10px;"><strong>{threat_desc}</strong>'
            if triggers:
                html_body += f'<br><span style="color: #666; font-size: 0.9em;">Watch for: {triggers}</span>'
            html_body += '</li>'
        html_body += '</ul>'
        
        html_body += """
</body>
</html>
"""
        
        # Plain text version
        plain_text = f"{ticker} - {company}\n\n"
        plain_text += "1. INVESTMENT THESIS\n"
        plain_text += f"{thesis.get('summary', 'N/A')}\n\n"
        
        if thesis.get('pillars'):
            for pillar in thesis['pillars']:
                title = pillar.get('pillar', pillar.get('title', ''))
                desc = pillar.get('detail', pillar.get('description', ''))
                plain_text += f"  - {title}: {desc}\n"
        
        plain_text += "\n2. SIGNPOSTS\n"
        for sp in signposts:
            metric = sp.get('metric', sp.get('signpost', ''))
            target = sp.get('target', '')
            plain_text += f"  - {metric}: {target}\n"
        
        plain_text += "\n3. THESIS THREATS\n"
        for threat in threats:
            plain_text += f"  - {threat.get('threat', '')}\n"
        
        msg = MIMEMultipart('alternative')
        msg['From'] = from_email
        msg['To'] = recipient
        msg['Subject'] = subject
        
        msg.attach(MIMEText(plain_text, 'plain'))
        msg.attach(MIMEText(html_body, 'html'))
        
        if use_gmail:
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(gmail_user, gmail_password)
                server.send_message(msg)
        
        return jsonify({'success': True, 'message': 'Analysis email sent successfully'})
        
    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'Gmail authentication failed. Check your email and app password.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============================================
# PDF EXTRACTION ENDPOINT
# ============================================

@app.route('/api/extract-pdf', methods=['POST'])
def extract_pdf():
    """
    Extract text from uploaded PDF file.
    Used by Research tab for document analysis.
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Read PDF and extract text
        from PyPDF2 import PdfReader
        import io
        
        pdf_bytes = file.read()
        pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
        
        text_content = []
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_content.append(page_text)
        
        full_text = '\n\n'.join(text_content)
        
        if not full_text.strip():
            return jsonify({'error': 'Could not extract text from PDF. It may be scanned/image-based.'}), 400
        
        return jsonify({
            'text': full_text,
            'pages': len(pdf_reader.pages),
            'filename': file.filename
        })
        
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return jsonify({'error': f'Failed to extract PDF: {str(e)}'}), 500


# ============================================
# RESEARCH PIPELINE
# ============================================

PIPELINE_SECTOR_MAP = {
    'Pharma': ['MRK', 'LLY', 'JNJ', 'PFE', 'BMY', 'ABBV'],
    'Biotech': ['GILD', 'REGN', 'VRTX'],
    'MedTech & Life Sci': ['MDT', 'BDX', 'BSX', 'ABT', 'TMO', 'DHR', 'DGX'],
    'Managed Care & Distro': ['UNH', 'CI', 'CVS', 'COR', 'CAH'],
    'Industrials': ['RTX', 'GD', 'CARR', 'ETN', 'DE', 'PH', 'DOV', 'MMM', 'RSG', 'UNP', 'NSC'],
    'REITs': ['VTR', 'PLD', 'AVB', 'AMT'],
}
# Reverse lookup: ticker -> sector
PIPELINE_TICKER_SECTOR = {}
for _sec, _tks in PIPELINE_SECTOR_MAP.items():
    for _tk in _tks:
        PIPELINE_TICKER_SECTOR[_tk] = _sec

RESEARCH_NOTE_PLAYBOOK = """
## NOTE FORMAT RULES

### Structure (adapt emphasis per stock, but include all sections):
1. Executive Summary / Investment Thesis (bull + bear in 3-4 bullets each)
2. Business Overview & Segment Breakdown
3. Key Revenue & Earnings Drivers
4. What the Street Is Debating (the 2-3 key open questions)
5. Catalyst Calendar (earnings, FDA dates, contract renewals, etc.)
6. Valuation Context (vs. history, vs. peers)
7. Risks
8. Bottom Line: Own / Avoid / Revisit at $X

### Sector-specific additions:
- Pharma/Biotech: Patent cliffs, pipeline table, LOE timeline
- MedTech: Procedure volume trends, ASP dynamics, new product cycles
- Managed Care: Membership trends, MLR, PBM reform risk, star ratings
- Distribution: Drug pricing dynamics, biosimilar opportunity, generic deflation
- Industrials: Cycle positioning, book-to-bill, aftermarket mix, margin expansion
- REITs: Same-store NOI, occupancy, cap rates, lease spreads, FFO/AFFO

### STRICT RULES:
- NEVER reference specific analyst names, firms, or broker ratings
- Synthesize data points from reports but attribute nothing to specific brokers
- Hard facts (reported financials, FDA approvals, deals, guidance) = state directly
- Sellside opinions: include if valuable but do NOT attribute
- NEVER attribute headline YoY EPS growth to a single narrative driver without decomposing the bridge
- Cross-check every narrative claim against data tables
- When FY EPS includes >$0.50/share in non-recurring items, flag explicitly
- Distinguish reported vs underlying/organic growth rates

### TONE:
- Write as if the analyst is authoring the note to their PM
- Confident, concise, first-person where appropriate
- No hedging language
- Lead with conclusion, support with evidence
- Use precise numbers, no rounding
- No emojis, no filler
"""


@app.route('/api/pipeline/start', methods=['POST'])
def pipeline_start():
    """Start a research pipeline batch for one or more tickers."""
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'Request body required'}), 400

        tickers = data.get('tickers', [])
        job_type = data.get('jobType', 'process')
        api_key = data.get('apiKey', '')
        document_config = data.get('documentConfig', {})
        generate_tiers = data.get('generateTiers', 'detailed')

        if not tickers:
            return jsonify({'error': 'At least one ticker is required'}), 400
        if not api_key:
            return jsonify({'error': 'API key is required'}), 400

        # Determine total steps based on tier generation
        total_steps = 7
        if generate_tiers == 'full':
            total_steps = 9
        elif generate_tiers == 'all':
            total_steps = 10

        batch_id = str(uuid.uuid4())
        jobs = []

        with get_db(commit=True) as (conn, cur):
            for ticker in tickers:
                job_id = str(uuid.uuid4())
                doc_config = document_config.get(ticker.upper(), {}) if document_config else {}
                cur.execute('''
                    INSERT INTO research_pipeline_jobs (id, batch_id, ticker, job_type, status, progress, current_step, total_steps, steps_detail)
                    VALUES (%s, %s, %s, %s, 'queued', 0, 'Queued', %s, %s)
                ''', (job_id, batch_id, ticker.upper(), job_type, total_steps, json.dumps({'documentConfig': doc_config, 'generateTiers': generate_tiers})))
                jobs.append({'id': job_id, 'ticker': ticker.upper(), 'status': 'queued'})

        # Spawn background thread for the batch
        threading.Thread(
            target=_run_pipeline_batch,
            args=(batch_id, [t.upper() for t in tickers], job_type, api_key),
            daemon=True
        ).start()

        return jsonify({'batchId': batch_id, 'jobs': jobs})

    except Exception as e:
        print(f"Pipeline start error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/jobs', methods=['GET'])
def pipeline_jobs():
    """List pipeline jobs with optional filters."""
    try:
        status_filter = request.args.get('status')
        ticker_filter = request.args.get('ticker')
        limit = int(request.args.get('limit', 50))

        conditions = []
        params = []

        if status_filter:
            conditions.append('status = %s')
            params.append(status_filter)
        if ticker_filter:
            conditions.append('ticker = %s')
            params.append(ticker_filter.upper())

        where_clause = f"WHERE {' AND '.join(conditions)}" if conditions else ''
        params.append(limit)

        with get_db() as (_, cur):
            cur.execute(f'''
                SELECT id, batch_id, ticker, job_type, status, progress, current_step, total_steps,
                       steps_detail, error, created_at, updated_at, completed_at
                FROM research_pipeline_jobs
                {where_clause}
                ORDER BY created_at DESC
                LIMIT %s
            ''', params)
            rows = cur.fetchall()

        jobs = []
        for row in rows:
            job = dict(row)
            # Serialize timestamps
            for ts_field in ('created_at', 'updated_at', 'completed_at'):
                if job.get(ts_field):
                    job[ts_field] = job[ts_field].isoformat()
            # Parse JSONB fields
            if isinstance(job.get('steps_detail'), str):
                job['steps_detail'] = json.loads(job['steps_detail'])
            jobs.append(job)

        return jsonify({'jobs': jobs})

    except Exception as e:
        print(f"Pipeline jobs list error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/job/<job_id>/status', methods=['GET'])
def pipeline_job_status(job_id):
    """Get status of a single pipeline job."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, batch_id, ticker, job_type, status, progress, current_step, total_steps,
                       steps_detail, result, error, created_at, updated_at, completed_at
                FROM research_pipeline_jobs WHERE id = %s
            ''', (job_id,))
            row = cur.fetchone()

        if not row:
            return jsonify({'error': 'Job not found'}), 404

        job = dict(row)
        for ts_field in ('created_at', 'updated_at', 'completed_at'):
            if job.get(ts_field):
                job[ts_field] = job[ts_field].isoformat()
        if isinstance(job.get('steps_detail'), str):
            job['steps_detail'] = json.loads(job['steps_detail'])
        if isinstance(job.get('result'), str):
            job['result'] = json.loads(job['result'])

        return jsonify(job)

    except Exception as e:
        print(f"Pipeline job status error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/batch-status', methods=['POST'])
def pipeline_batch_status():
    """Get status for multiple pipeline jobs at once (efficient bulk poll)."""
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'Request body required'}), 400

        job_ids = data.get('jobIds', [])
        if not job_ids:
            return jsonify({'jobs': {}})

        placeholders = ','.join(['%s'] * len(job_ids))
        with get_db() as (_, cur):
            cur.execute(f'''
                SELECT id, batch_id, ticker, job_type, status, progress, current_step, total_steps,
                       steps_detail, result, error, created_at, updated_at, completed_at
                FROM research_pipeline_jobs WHERE id IN ({placeholders})
            ''', job_ids)
            rows = cur.fetchall()

        jobs = {}
        for row in rows:
            job = dict(row)
            for ts_field in ('created_at', 'updated_at', 'completed_at'):
                if job.get(ts_field):
                    job[ts_field] = job[ts_field].isoformat()
            if isinstance(job.get('steps_detail'), str):
                job['steps_detail'] = json.loads(job['steps_detail'])
            if isinstance(job.get('result'), str):
                job['result'] = json.loads(job['result'])
            jobs[job['id']] = job

        return jsonify({'jobs': jobs})

    except Exception as e:
        print(f"Pipeline batch status error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/job/<job_id>', methods=['DELETE'])
def pipeline_job_delete(job_id):
    """Delete or cancel a pipeline job."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT id, status FROM research_pipeline_jobs WHERE id = %s', (job_id,))
            job = cur.fetchone()

        if not job:
            return jsonify({'error': 'Job not found'}), 404

        if job['status'] == 'running':
            # Can't delete a running job, mark as cancelled instead
            with get_db(commit=True) as (conn, cur):
                cur.execute("UPDATE research_pipeline_jobs SET status = 'cancelled', updated_at = NOW() WHERE id = %s", (job_id,))
            return jsonify({'message': 'Job cancelled', 'id': job_id, 'status': 'cancelled'})
        else:
            # Queued, complete, or failed — delete the row
            with get_db(commit=True) as (conn, cur):
                cur.execute('DELETE FROM research_pipeline_jobs WHERE id = %s', (job_id,))
            return jsonify({'message': 'Job deleted', 'id': job_id})

    except Exception as e:
        print(f"Pipeline job delete error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/universe', methods=['GET'])
def pipeline_universe():
    """Return the stock universe with processing status."""
    sort_by = request.args.get('sort', 'sector')  # 'alpha', 'updated', 'sector'
    try:
        with get_db() as (_, cur):
            # Get all tickers from portfolio_analyses
            cur.execute('''
                SELECT pa.ticker, pa.company, pa.updated_at as last_analysis_date,
                       CASE WHEN tsd.ticker IS NOT NULL THEN true ELSE false END as has_scorecard
                FROM portfolio_analyses pa
                LEFT JOIN thesis_scorecard_data tsd ON pa.ticker = tsd.ticker
                ORDER BY pa.ticker
            ''')
            analyses = cur.fetchall()

            # Get latest pipeline job per ticker
            cur.execute('''
                SELECT DISTINCT ON (ticker) ticker, status as last_status, completed_at as last_processed, updated_at
                FROM research_pipeline_jobs
                ORDER BY ticker, created_at DESC
            ''')
            pipeline_map = {row['ticker']: row for row in cur.fetchall()}

            # Get document counts per ticker (from both thesis and research storage)
            cur.execute('SELECT ticker, COUNT(*) as cnt FROM document_files GROUP BY ticker')
            doc_counts = {row['ticker']: row['cnt'] for row in cur.fetchall()}

            # Get which tickers have research notes
            cur.execute('SELECT DISTINCT ticker FROM research_notes')
            tickers_with_notes = {row['ticker'] for row in cur.fetchall()}

            # Get custom sector assignments from ticker_settings
            cur.execute('SELECT ticker, sector FROM ticker_settings WHERE sector IS NOT NULL')
            custom_sectors = {row['ticker']: row['sector'] for row in cur.fetchall()}

        universe = []
        for a in analyses:
            entry = {
                'ticker': a['ticker'],
                'company': a['company'],
                'lastAnalysisDate': a['last_analysis_date'].isoformat() if a['last_analysis_date'] else None,
                'hasScorecard': a['has_scorecard'],
                'lastProcessed': None,
                'lastStatus': None,
            }
            # Sector assignment: custom overrides first, then hardcoded map, then 'Other'
            entry['sector'] = custom_sectors.get(a['ticker']) or PIPELINE_TICKER_SECTOR.get(a['ticker'], 'Other')
            # Calculate staleness: days since last analysis
            if a['last_analysis_date']:
                days_old = (datetime.utcnow() - a['last_analysis_date']).days
                entry['daysOld'] = days_old
                entry['freshness'] = 'fresh' if days_old <= 7 else 'stale' if days_old <= 30 else 'outdated'
            else:
                entry['daysOld'] = None
                entry['freshness'] = 'never'
            entry['docCount'] = doc_counts.get(a['ticker'], 0)
            entry['hasNote'] = a['ticker'] in tickers_with_notes
            pj = pipeline_map.get(a['ticker'])
            if pj:
                entry['lastProcessed'] = pj['last_processed'].isoformat() if pj['last_processed'] else None
                entry['lastStatus'] = pj['last_status']
            universe.append(entry)

        # Apply sorting
        if sort_by == 'alpha':
            universe.sort(key=lambda x: x['ticker'])
        elif sort_by == 'updated':
            universe.sort(key=lambda x: x.get('lastAnalysisDate') or '', reverse=True)
        # 'sector' is the default — already grouped by sector in the frontend

        # Build combined sector list (hardcoded + custom)
        all_sectors = set(PIPELINE_SECTOR_MAP.keys())
        for s in custom_sectors.values():
            if s:
                all_sectors.add(s)

        return jsonify({
            'universe': universe,
            'total': len(universe),
            'sectors': sorted(all_sectors),
            'sectorMap': PIPELINE_SECTOR_MAP
        })

    except Exception as e:
        print(f"Pipeline universe error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/ticker/add', methods=['POST'])
def add_ticker():
    """Manually add a ticker to the universe."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper().strip()
    company = data.get('company', '').strip()
    sector = data.get('sector', '').strip()

    if not ticker:
        return jsonify({'error': 'Ticker required'}), 400

    # Create entry in portfolio_analyses if it doesn't exist
    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO portfolio_analyses (ticker, company, analysis, updated_at)
            VALUES (%s, %s, %s, NOW())
            ON CONFLICT (ticker) DO UPDATE SET
                company = COALESCE(NULLIF(%s, ''), portfolio_analyses.company),
                updated_at = NOW()
        ''', (ticker, company or ticker, json.dumps({}), company))

    # Save sector assignment if provided
    if sector:
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO ticker_settings (ticker, sector, custom_company, updated_at)
                VALUES (%s, %s, %s, NOW())
                ON CONFLICT (ticker) DO UPDATE SET
                    sector = COALESCE(NULLIF(%s, ''), ticker_settings.sector),
                    custom_company = COALESCE(NULLIF(%s, ''), ticker_settings.custom_company),
                    updated_at = NOW()
            ''', (ticker, sector, company, sector, company))

    return jsonify({'success': True, 'ticker': ticker})


@app.route('/api/ticker/sector', methods=['POST'])
def update_ticker_sector():
    """Update a ticker's sector assignment."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper().strip()
    sector = data.get('sector', '').strip()

    if not ticker or not sector:
        return jsonify({'error': 'Ticker and sector required'}), 400

    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO ticker_settings (ticker, sector, updated_at)
            VALUES (%s, %s, NOW())
            ON CONFLICT (ticker) DO UPDATE SET sector = %s, updated_at = NOW()
        ''', (ticker, sector, sector))

    return jsonify({'success': True, 'ticker': ticker, 'sector': sector})


@app.route('/api/ticker/sectors', methods=['GET'])
def get_all_sectors():
    """Return all available sectors (hardcoded + custom)."""
    sectors = set(PIPELINE_SECTOR_MAP.keys())
    with get_db() as (_, cur):
        cur.execute('SELECT DISTINCT sector FROM ticker_settings WHERE sector IS NOT NULL')
        for row in cur.fetchall():
            if row['sector']:
                sectors.add(row['sector'])
    return jsonify({'sectors': sorted(sectors)})


@app.route('/api/ticker/delete', methods=['POST'])
def delete_ticker():
    """Delete a ticker from the universe (archives, doesn't hard-delete)."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper().strip()
    if not ticker:
        return jsonify({'error': 'Ticker required'}), 400

    try:
        with get_db(commit=True) as (conn, cur):
            # Remove from portfolio_analyses
            cur.execute('DELETE FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            # Remove from ticker_settings
            cur.execute('DELETE FROM ticker_settings WHERE ticker = %s', (ticker,))
            # Remove scorecard data
            cur.execute('DELETE FROM thesis_scorecard_data WHERE ticker = %s', (ticker,))

        cache.invalidate('analyses')
        cache.invalidate('portfolio_dashboard')
        return jsonify({'success': True, 'ticker': ticker})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/thesis/diff/<ticker>', methods=['GET'])
def thesis_diff(ticker):
    """Get diff between current thesis and previous version."""
    ticker = ticker.upper()
    with get_db() as (_, cur):
        cur.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
        row = cur.fetchone()
    if not row or not row['analysis']:
        return jsonify({'error': 'No analysis found'}), 404

    analysis = row['analysis'] if isinstance(row['analysis'], dict) else json.loads(row['analysis'])
    history = analysis.get('history', [])

    if not history:
        return jsonify({'hasDiff': False, 'message': 'No prior versions'})

    prev = history[-1]  # Most recent prior version
    current = {
        'summary': analysis.get('thesis', {}).get('summary', ''),
        'pillars': analysis.get('thesis', {}).get('pillars', []),
        'signposts': analysis.get('signposts', []),
        'threats': analysis.get('threats', []),
        'conclusion': analysis.get('conclusion', ''),
    }
    previous = {
        'summary': prev.get('thesis', {}).get('summary', '') if isinstance(prev.get('thesis'), dict) else '',
        'pillars': prev.get('thesis', {}).get('pillars', []) if isinstance(prev.get('thesis'), dict) else [],
        'signposts': prev.get('signposts', []),
        'threats': prev.get('threats', []),
        'conclusion': '',
    }

    changes = []
    # Compare summary
    if current['summary'] != previous['summary']:
        changes.append({'section': 'Summary', 'before': previous['summary'][:500], 'after': current['summary'][:500]})

    # Compare pillars
    curr_pillar_titles = [p.get('title', p.get('pillar', '')) for p in current['pillars']]
    prev_pillar_titles = [p.get('title', p.get('pillar', '')) for p in previous['pillars']]
    added_pillars = [t for t in curr_pillar_titles if t not in prev_pillar_titles]
    removed_pillars = [t for t in prev_pillar_titles if t not in curr_pillar_titles]
    if added_pillars:
        changes.append({'section': 'Pillars Added', 'before': '', 'after': ', '.join(added_pillars)})
    if removed_pillars:
        changes.append({'section': 'Pillars Removed', 'before': ', '.join(removed_pillars), 'after': ''})

    # Compare signposts
    curr_sp = [s.get('metric', s.get('signpost', '')) for s in current['signposts']]
    prev_sp = [s.get('metric', s.get('signpost', '')) for s in previous['signposts']]
    added_sp = [s for s in curr_sp if s not in prev_sp]
    removed_sp = [s for s in prev_sp if s not in curr_sp]
    if added_sp:
        changes.append({'section': 'Signposts Added', 'before': '', 'after': ', '.join(added_sp)})
    if removed_sp:
        changes.append({'section': 'Signposts Removed', 'before': ', '.join(removed_sp), 'after': ''})

    # Compare threats
    curr_th = [t.get('threat', '') for t in current['threats']]
    prev_th = [t.get('threat', '') for t in previous['threats']]
    added_th = [t for t in curr_th if t not in prev_th]
    removed_th = [t for t in prev_th if t not in curr_th]
    if added_th:
        changes.append({'section': 'Threats Added', 'before': '', 'after': ', '.join(added_th)})
    if removed_th:
        changes.append({'section': 'Threats Removed', 'before': ', '.join(removed_th), 'after': ''})

    # Build full side-by-side sections for the diff modal
    def format_pillars(pillars):
        return '\n\n'.join([f"**{p.get('title', p.get('pillar', ''))}**\n{p.get('detail', p.get('description', ''))}" for p in pillars]) if pillars else ''

    def format_signposts(items):
        return '\n'.join([f"- {s.get('metric', s.get('signpost', ''))}: {s.get('target', s.get('threshold', ''))}" for s in items]) if items else ''

    def format_threats(items):
        return '\n'.join([f"- {t.get('threat', '')}: {t.get('detail', t.get('mitigation', ''))}" for t in items]) if items else ''

    sections = [
        {'name': 'Summary', 'previous': previous['summary'], 'current': current['summary'], 'changed': current['summary'] != previous['summary']},
        {'name': 'Pillars', 'previous': format_pillars(previous['pillars']), 'current': format_pillars(current['pillars']), 'changed': curr_pillar_titles != prev_pillar_titles},
        {'name': 'Signposts', 'previous': format_signposts(previous['signposts']), 'current': format_signposts(current['signposts']), 'changed': curr_sp != prev_sp},
        {'name': 'Threats', 'previous': format_threats(previous['threats']), 'current': format_threats(current['threats']), 'changed': curr_th != prev_th},
    ]

    return jsonify({
        'hasDiff': len(changes) > 0,
        'changes': changes,
        'sections': sections,
        'previousTimestamp': prev.get('timestamp', ''),
        'currentTimestamp': analysis.get('updated', ''),
        'totalChanges': len(changes),
    })


@app.route('/api/pipeline/refresh-sector', methods=['POST'])
def pipeline_refresh_sector():
    """Trigger pipeline jobs for all stocks in a sector."""
    data = request.get_json()
    sector = data.get('sector', '')
    job_type = data.get('jobType', 'update')
    api_key = data.get('apiKey', '')
    generate_tiers = data.get('generateTiers', 'detailed')
    document_config = data.get('documentConfig', {})

    tickers = PIPELINE_SECTOR_MAP.get(sector, [])
    if not tickers:
        return jsonify({'error': f'Unknown sector: {sector}'}), 400

    # Determine total steps based on tier generation
    total_steps = 7
    if generate_tiers == 'full':
        total_steps = 9
    elif generate_tiers == 'all':
        total_steps = 10

    # Filter to only tickers that exist in portfolio_analyses
    with get_db() as (_, cur):
        cur.execute('SELECT ticker FROM portfolio_analyses WHERE ticker = ANY(%s)', (tickers,))
        existing = [r['ticker'] for r in cur.fetchall()]

    if not existing:
        return jsonify({'error': f'No analyzed stocks found in sector {sector}'}), 404

    batch_id = str(uuid.uuid4())
    jobs = []
    with get_db(commit=True) as (conn, cur):
        for ticker in existing:
            job_id = str(uuid.uuid4())
            # Include per-ticker document config if available
            ticker_doc_config = document_config.get(ticker, {})
            steps_detail = {'generateTiers': generate_tiers}
            if ticker_doc_config:
                steps_detail['documentConfig'] = ticker_doc_config
            cur.execute('''
                INSERT INTO research_pipeline_jobs (id, batch_id, ticker, job_type, status, progress, current_step, total_steps, steps_detail)
                VALUES (%s, %s, %s, %s, 'queued', 0, 'Queued', %s, %s)
            ''', (job_id, batch_id, ticker, job_type, total_steps, json.dumps(steps_detail)))
            jobs.append({'id': job_id, 'ticker': ticker, 'status': 'queued'})

    threading.Thread(target=_run_pipeline_batch, args=(batch_id, existing, job_type, api_key), daemon=True).start()

    return jsonify({'batchId': batch_id, 'sector': sector, 'jobs': jobs, 'tickerCount': len(existing)})


@app.route('/api/pipeline/documents/<ticker>', methods=['GET'])
def pipeline_documents(ticker):
    """Return all documents for a ticker from both Thesis and Research storage."""
    ticker = ticker.upper()
    try:
        docs = []
        # Source 1: document_files (Thesis tab)
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, filename, file_type, mime_type, file_size, created_at, metadata
                FROM document_files WHERE ticker = %s ORDER BY created_at DESC
            ''', (ticker,))
            for row in cur.fetchall():
                docs.append({
                    'id': f'thesis_{row["id"]}',
                    'source': 'thesis',
                    'filename': row['filename'],
                    'fileType': row.get('file_type', ''),
                    'fileSize': row.get('file_size', 0),
                    'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
                    'weight': 1.0,
                })

        # Source 2: research_document_files (Research tab) — join through research_documents → research_categories
        with get_db() as (_, cur):
            cur.execute('''
                SELECT rdf.id, rdf.filename, rdf.file_type, rdf.file_size, rdf.created_at,
                       rd.name as doc_name, rd.doc_type, rd.published_date
                FROM research_document_files rdf
                JOIN research_documents rd ON rdf.document_id = rd.id
                JOIN research_categories rc ON rd.category_id = rc.id
                WHERE UPPER(rc.name) = %s
                ORDER BY rdf.created_at DESC
            ''', (ticker,))
            for row in cur.fetchall():
                docs.append({
                    'id': f'research_{row["id"]}',
                    'source': 'research',
                    'filename': row['filename'],
                    'fileType': row.get('file_type', ''),
                    'fileSize': row.get('file_size', 0),
                    'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
                    'docName': row.get('doc_name', ''),
                    'docType': row.get('doc_type', ''),
                    'weight': 1.0,
                })

        # Cross-reference with documentHistory to show which docs were used in thesis
        used_in_thesis = set()
        with get_db() as (_, cur):
            cur.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            pa_row = cur.fetchone()
        if pa_row and pa_row['analysis']:
            analysis = pa_row['analysis'] if isinstance(pa_row['analysis'], dict) else json.loads(pa_row['analysis'])
            for dh in analysis.get('documentHistory', []):
                used_in_thesis.add(dh.get('filename', ''))

        # Cross-reference with research_notes metadata to show which docs were used in notes
        used_in_note = set()
        with get_db() as (_, cur):
            cur.execute('SELECT metadata FROM research_notes WHERE ticker = %s ORDER BY created_at DESC LIMIT 1', (ticker,))
            note_row = cur.fetchone()
        if note_row and note_row['metadata']:
            meta = note_row['metadata'] if isinstance(note_row['metadata'], dict) else json.loads(note_row['metadata'] or '{}')
            for fn in meta.get('documentFilenames', []):
                used_in_note.add(fn)

        for doc in docs:
            doc['usedInAnalysis'] = doc['filename'] in used_in_thesis
            doc['usedInThesis'] = doc['filename'] in used_in_thesis
            doc['usedInNote'] = doc['filename'] in used_in_note

        return jsonify({'ticker': ticker, 'documents': docs, 'total': len(docs), 'usedInThesisCount': len(used_in_thesis), 'usedInNoteCount': len(used_in_note)})
    except Exception as e:
        print(f'Error fetching pipeline documents for {ticker}: {e}')
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/history', methods=['GET'])
def pipeline_history():
    """Return pipeline job history with analytics."""
    limit = int(request.args.get('limit', 100))
    ticker = request.args.get('ticker', '')

    try:
        with get_db() as (_, cur):
            if ticker:
                cur.execute('''
                    SELECT id, batch_id, ticker, job_type, status, progress, current_step, error,
                           created_at, updated_at, completed_at
                    FROM research_pipeline_jobs
                    WHERE ticker = %s
                    ORDER BY created_at DESC LIMIT %s
                ''', (ticker.upper(), limit))
            else:
                cur.execute('''
                    SELECT id, batch_id, ticker, job_type, status, progress, current_step, error,
                           created_at, updated_at, completed_at
                    FROM research_pipeline_jobs
                    ORDER BY created_at DESC LIMIT %s
                ''', (limit,))
            jobs = cur.fetchall()

            # Analytics
            cur.execute('''
                SELECT
                    COUNT(*) as total_jobs,
                    COUNT(*) FILTER (WHERE status = 'complete') as completed,
                    COUNT(*) FILTER (WHERE status = 'failed') as failed,
                    COUNT(*) FILTER (WHERE status = 'running') as running,
                    COUNT(*) FILTER (WHERE status = 'queued') as queued,
                    COUNT(DISTINCT ticker) as unique_tickers,
                    AVG(EXTRACT(EPOCH FROM (completed_at - created_at))) FILTER (WHERE status = 'complete') as avg_duration_seconds
                FROM research_pipeline_jobs
            ''')
            analytics = cur.fetchone()

        return jsonify({
            'jobs': [{
                'id': j['id'],
                'batchId': j['batch_id'],
                'ticker': j['ticker'],
                'jobType': j['job_type'],
                'status': j['status'],
                'progress': j['progress'],
                'currentStep': j['current_step'],
                'error': j['error'],
                'createdAt': j['created_at'].isoformat() if j['created_at'] else None,
                'completedAt': j['completed_at'].isoformat() if j['completed_at'] else None,
                'duration': round((j['completed_at'] - j['created_at']).total_seconds()) if j['completed_at'] and j['created_at'] else None,
            } for j in jobs],
            'analytics': {
                'totalJobs': analytics['total_jobs'],
                'completed': analytics['completed'],
                'failed': analytics['failed'],
                'running': analytics['running'],
                'queued': analytics['queued'],
                'uniqueTickers': analytics['unique_tickers'],
                'avgDurationSeconds': round(analytics['avg_duration_seconds']) if analytics['avg_duration_seconds'] else None,
            }
        })
    except Exception as e:
        print(f'Error fetching pipeline history: {e}')
        return jsonify({'error': str(e)}), 500


@app.route('/api/pipeline/history/clear', methods=['POST'])
def pipeline_clear_history():
    """Clear completed/failed pipeline jobs."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute("DELETE FROM research_pipeline_jobs WHERE status IN ('complete', 'failed')")
            deleted = cur.rowcount
        return jsonify({'deleted': deleted})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============================================
# RESEARCH NOTE GENERATION
# ============================================

@app.route('/api/notes/generate', methods=['POST'])
def generate_research_note():
    """Generate a full equity research note for a ticker."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper()
    api_key = data.get('apiKey', '')
    mode = data.get('mode', 'new')  # 'new' or 'update'
    file_selection = data.get('fileSelection', [])
    reprocess = data.get('reprocess', False)

    if not ticker:
        return jsonify({'error': 'No ticker provided'}), 400
    if not api_key:
        return jsonify({'error': 'API key required'}), 400

    # Create a pipeline job for the local agent to pick up
    # NOTE: Do NOT spawn a backend thread — note generation runs on the local agent
    # which has access to iCloud files and can call Claude directly
    job_id = str(uuid.uuid4())
    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO research_pipeline_jobs (id, batch_id, ticker, job_type, status, progress, current_step, total_steps, steps_detail)
            VALUES (%s, %s, %s, 'note', 'queued', 0, 'Waiting for local agent', 6, %s)
        ''', (job_id, str(uuid.uuid4()), ticker, json.dumps({'mode': mode, 'fileSelection': file_selection, 'reprocess': reprocess})))

    return jsonify({'jobId': job_id, 'ticker': ticker})


@app.route('/api/notes/<ticker>', methods=['GET'])
def get_research_note(ticker):
    """Get the latest research note for a ticker."""
    ticker = ticker.upper()
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_notes WHERE ticker = %s ORDER BY created_at DESC LIMIT 1', (ticker,))
        note = cur.fetchone()
    if not note:
        return jsonify({'error': 'No note found'}), 404

    return jsonify({
        'id': note['id'],
        'ticker': note['ticker'],
        'version': note['version'],
        'noteMarkdown': note['note_markdown'],
        'sourcesMarkdown': note['sources_markdown'],
        'changelogMarkdown': note['changelog_markdown'],
        'noteDocx': note['note_docx'],
        'charts': note['charts'] if isinstance(note['charts'], list) else json.loads(note['charts'] or '[]'),
        'metadata': note['metadata'] if isinstance(note['metadata'], dict) else json.loads(note['metadata'] or '{}'),
        'createdAt': note['created_at'].isoformat() if note['created_at'] else None,
    })


@app.route('/api/notes/<ticker>/pdf', methods=['GET'])
def get_research_note_pdf(ticker):
    """Generate a PDF from the latest research note for a ticker."""
    from xhtml2pdf import pisa
    ticker = ticker.upper()
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_notes WHERE ticker = %s ORDER BY created_at DESC LIMIT 1', (ticker,))
        note = cur.fetchone()
    if not note:
        return jsonify({'error': 'No note found'}), 404

    note_md = note['note_markdown'] or ''
    version = note.get('version', '1.0')

    # Convert markdown to simple HTML for PDF
    import re
    html_body = note_md
    # Headers — -pdf-keep-with-next prevents orphaned headings at page bottom
    html_body = re.sub(r'^### (.+)$', r'<h3 style="color:#1e293b;font-size:11pt;margin:12px 0 6px 0;-pdf-keep-with-next:true;">\1</h3>', html_body, flags=re.MULTILINE)
    html_body = re.sub(r'^## (\d+\.\s+)?(.+)$', r'<h2 style="color:#1e293b;font-size:13pt;font-weight:bold;border-bottom:2px solid #1e3a5f;padding-bottom:4px;margin:18px 0 8px 0;-pdf-keep-with-next:true;">\1\2</h2>', html_body, flags=re.MULTILINE)
    html_body = re.sub(r'^# (.+)$', r'<h1 style="color:#1e293b;font-size:18pt;margin:0 0 4px 0;-pdf-keep-with-next:true;">\1</h1>', html_body, flags=re.MULTILINE)
    # Bold and italic
    html_body = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', html_body)
    html_body = re.sub(r'\*(.+?)\*', r'<i>\1</i>', html_body)
    # Bullets
    html_body = re.sub(r'^- (.+)$', r'<p style="margin:2px 0 2px 20px;font-size:10pt;color:#334155;">&bull; \1</p>', html_body, flags=re.MULTILINE)
    # Markdown tables -> HTML tables (BEFORE horizontal rule replacement to avoid breaking separator rows)
    def _md_table_to_html(match):
        lines = match.group(0).strip().split('\n')
        # Filter out separator rows (|---|---|---| or |:--|:--:|--:|)
        rows = [l for l in lines if l.strip() and not re.match(r'^\|[\s\-:|]+\|$', l.strip())]
        if not rows:
            return ''
        html = '<table style="border-collapse:collapse;width:100%;margin:8px 0;font-size:9pt;page-break-inside:avoid;">'
        for ri, row in enumerate(rows):
            cells = [c.strip().replace('**', '') for c in row.strip('|').split('|')]
            tag = 'th' if ri == 0 else 'td'
            style = 'padding:5px 8px;border:1px solid #d1d5db;color:#1e293b;word-wrap:break-word;' + ('background:#f1f5f9;font-weight:bold;' if ri == 0 else '')
            html += '<tr>' + ''.join(f'<{tag} style="{style}">{c}</{tag}>' for c in cells) + '</tr>'
        html += '</table>'
        return html
    html_body = re.sub(r'(\|.+\|[\n\r]*)+', _md_table_to_html, html_body)
    # Horizontal rules (only standalone --- lines, not inside other content)
    html_body = re.sub(r'^---+$', '<hr style="border:none;border-top:1px solid #cbd5e1;margin:12px 0;">', html_body, flags=re.MULTILINE)
    # Paragraphs (lines that aren't already HTML)
    lines = html_body.split('\n')
    processed = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            processed.append('')
        elif stripped.startswith('<'):
            processed.append(stripped)
        else:
            processed.append(f'<p style="font-size:10pt;color:#334155;line-height:1.6;margin:4px 0;">{stripped}</p>')
    html_body = '\n'.join(processed)

    # Embed chart images if available
    charts_html = ''
    charts_data = note.get('charts', [])
    if isinstance(charts_data, str):
        try: charts_data = json.loads(charts_data)
        except: charts_data = []
    for chart in (charts_data or []):
        if chart.get('data'):
            chart_type = chart.get('type', 'chart').title()
            charts_html += f'<p style="margin:16px 0 4px 0;font-weight:bold;font-size:11pt;color:#1e293b;">{ticker} {chart_type} Breakdown</p>'
            charts_html += f'<img src="data:image/png;base64,{chart["data"]}" style="width:100%;max-width:500px;margin:0 0 16px 0;" />'

    html = f"""<html><head><style>
        @page {{ margin: 0.7in; size: letter; }}
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 10pt; color: #1e293b; line-height: 1.5; }}
        h1, h2, h3 {{ -pdf-keep-with-next: true; }}
        table {{ border-collapse: collapse; width: 100%; page-break-inside: avoid; }}
        img {{ max-width: 100%; }}
    </style></head><body>
    {html_body}
    {charts_html}
    </body></html>"""

    buf = io.BytesIO()
    pisa.CreatePDF(html, dest=buf)
    pdf_b64 = base64.b64encode(buf.getvalue()).decode('utf-8')

    return jsonify({
        'success': True,
        'fileData': pdf_b64,
        'filename': f'{ticker}_Note.pdf',
        'fileSize': len(buf.getvalue()),
    })


@app.route('/api/notes/<ticker>/history', methods=['GET'])
def get_research_note_history(ticker):
    """Get all note versions for a ticker."""
    ticker = ticker.upper()
    with get_db() as (_, cur):
        cur.execute('''
            SELECT id, ticker, version, metadata, created_at
            FROM research_notes WHERE ticker = %s ORDER BY created_at DESC
        ''', (ticker,))
        notes = cur.fetchall()
    return jsonify({
        'notes': [{
            'id': n['id'], 'version': n['version'],
            'metadata': n['metadata'] if isinstance(n['metadata'], dict) else json.loads(n['metadata'] or '{}'),
            'createdAt': n['created_at'].isoformat() if n['created_at'] else None,
        } for n in notes]
    })


def _generate_research_note(job_id, ticker, api_key, mode='new'):
    """Generate a full equity research note in a background thread."""
    try:
        _update_pipeline_job(job_id, status='running', current_step='Loading data', progress=0)

        # Step 1: Load existing analysis (thesis data)
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
            analysis_row = cur.fetchone()

        analysis = {}
        company = ticker
        if analysis_row:
            analysis = analysis_row['analysis'] if isinstance(analysis_row['analysis'], dict) else json.loads(analysis_row['analysis'] or '{}')
            company = analysis_row.get('company', ticker) or ticker

        # Step 2: Load source documents
        _update_pipeline_job(job_id, current_step='Reading source documents', progress=10)
        with get_db() as (_, cur):
            cur.execute('SELECT filename, file_data, file_type FROM document_files WHERE ticker = %s', (ticker,))
            docs = cur.fetchall()

        if not docs and not analysis:
            raise Exception(f'No documents or analysis found for {ticker}')

        # Build document content for LLM
        doc_contents = []
        for doc in docs:
            if doc.get('file_type') == 'pdf' and doc.get('file_data'):
                doc_contents.append({
                    'type': 'document',
                    'source': {'type': 'base64', 'media_type': 'application/pdf', 'data': doc['file_data']},
                })
                doc_contents.append({'type': 'text', 'text': f'[Document: {doc["filename"]}]'})

        # Step 3: Load existing note if updating
        existing_note = None
        if mode == 'update':
            with get_db() as (_, cur):
                cur.execute('SELECT * FROM research_notes WHERE ticker = %s ORDER BY created_at DESC LIMIT 1', (ticker,))
                existing_note = cur.fetchone()

        _update_pipeline_job(job_id, current_step='Generating research note', progress=20)

        # Step 4: Generate the note via LLM
        sector = PIPELINE_TICKER_SECTOR.get(ticker, 'Other')

        existing_thesis_context = ""
        if analysis:
            thesis = analysis.get('thesis', {})
            signposts = analysis.get('signposts', [])
            threats = analysis.get('threats', [])
            conclusion = analysis.get('conclusion', '')
            existing_thesis_context = f"""
EXISTING STRUCTURED THESIS (use as foundation, expand with document details):
Summary: {thesis.get('summary', '')}
Pillars: {json.dumps(thesis.get('pillars', []), indent=2)}
Signposts: {json.dumps(signposts, indent=2)}
Threats: {json.dumps(threats, indent=2)}
Conclusion: {conclusion}
"""

        update_context = ""
        if existing_note and mode == 'update':
            update_context = f"""
EXISTING NOTE (update this, don't rewrite from scratch):
{existing_note['note_markdown'][:8000]}

Update the note with new information from the source documents. Add a "What's Changed" section at the top. Update any data points, estimates, or catalysts that have changed.
"""

        prompt = f"""You are a senior equity research analyst writing a comprehensive investment research note.

TICKER: {ticker}
COMPANY: {company}
SECTOR: {sector}

{RESEARCH_NOTE_PLAYBOOK}

{existing_thesis_context}

{update_context}

Using the source documents provided, write a complete equity research note in markdown format.

The note should be 8-12 pages when printed, with these sections:
1. Executive Summary / Investment Thesis
2. Business Overview & Segment Breakdown
3. Key Revenue & Earnings Drivers
4. What the Street Is Debating
5. Catalyst Calendar
6. Valuation Context
7. Risks
8. Bottom Line

Also provide:
- Revenue segment data for pie chart (JSON array: [{{"segment": "name", "revenue": number_in_millions}}])
- Profit segment data for pie chart (JSON array: [{{"segment": "name", "profit": number_in_millions}}])

Return your response in this exact format:

===NOTE_START===
[full markdown note here]
===NOTE_END===

===SOURCES_START===
[sources document: section-by-section, which reports informed each claim, broker name + date + page ref]
===SOURCES_END===

===REVENUE_CHART_DATA===
[JSON array of revenue segments]
===REVENUE_CHART_END===

===PROFIT_CHART_DATA===
[JSON array of profit segments]
===PROFIT_CHART_END===
"""

        # Build messages with documents
        messages = [{"role": "user", "content": doc_contents + [{"type": "text", "text": prompt}]}] if doc_contents else [{"role": "user", "content": prompt}]

        result = call_llm(
            messages=messages,
            system="You are a senior equity research analyst. Write thorough, data-driven research notes. Be precise with numbers. No sellside attribution.",
            tier="advanced",
            max_tokens=16384,
            anthropic_api_key=api_key,
        )

        _update_pipeline_job(job_id, current_step='Parsing note', progress=60)

        response_text = result['text']

        # Parse sections
        note_md = ''
        sources_md = ''
        revenue_data = []
        profit_data = []

        note_match = re.search(r'===NOTE_START===\s*(.*?)\s*===NOTE_END===', response_text, re.DOTALL)
        if note_match:
            note_md = note_match.group(1).strip()
        else:
            # Fallback: treat entire response as note
            note_md = response_text

        sources_match = re.search(r'===SOURCES_START===\s*(.*?)\s*===SOURCES_END===', response_text, re.DOTALL)
        if sources_match:
            sources_md = sources_match.group(1).strip()

        rev_match = re.search(r'===REVENUE_CHART_DATA===\s*(.*?)\s*===REVENUE_CHART_END===', response_text, re.DOTALL)
        if rev_match:
            try:
                revenue_data = json.loads(rev_match.group(1).strip())
            except Exception:
                pass

        profit_match = re.search(r'===PROFIT_CHART_DATA===\s*(.*?)\s*===PROFIT_CHART_END===', response_text, re.DOTALL)
        if profit_match:
            try:
                profit_data = json.loads(profit_match.group(1).strip())
            except Exception:
                pass

        # Step 5: Generate charts
        _update_pipeline_job(job_id, current_step='Generating charts', progress=70)
        charts = []

        try:
            if revenue_data:
                rev_chart = _generate_donut_chart(ticker, 'Revenue', revenue_data, 'revenue')
                if rev_chart:
                    charts.append({'type': 'revenue', 'data': rev_chart, 'filename': f'{ticker}_Revenue_Breakdown.png'})

            if profit_data:
                prof_chart = _generate_donut_chart(ticker, 'Profit', profit_data, 'profit')
                if prof_chart:
                    charts.append({'type': 'profit', 'data': prof_chart, 'filename': f'{ticker}_Profit_Breakdown.png'})
        except Exception as chart_err:
            print(f'[note-gen {job_id}] Chart generation error: {chart_err}')

        # Step 6: Generate DOCX
        _update_pipeline_job(job_id, current_step='Building Word document', progress=80)
        docx_b64 = ''
        try:
            docx_b64 = _generate_note_docx(ticker, company, note_md, charts)
        except Exception as docx_err:
            print(f'[note-gen {job_id}] DOCX generation error: {docx_err}')

        # Determine version
        version = '1.0'
        if mode == 'update' and existing_note:
            old_ver = existing_note.get('version', '1.0')
            parts = old_ver.split('.')
            try:
                parts[-1] = str(int(parts[-1]) + 1)
                version = '.'.join(parts)
            except Exception:
                version = old_ver + '.1'

        # Generate changelog
        now_str = datetime.utcnow().strftime('%Y-%m-%d')
        changelog_md = f"# {ticker} Changelog\n\n## Version {version} — {now_str}\n- {'Updated' if mode == 'update' else 'Initial'} research note\n- {len(docs)} source documents processed\n- Generated via Charlie Pipeline\n"
        if existing_note and existing_note.get('changelog_markdown'):
            old_changelog = existing_note['changelog_markdown']
            if '\n' in old_changelog:
                changelog_md += '\n' + old_changelog.split('\n', 1)[1]

        # Save to DB
        _update_pipeline_job(job_id, current_step='Saving note', progress=90)
        note_id = str(uuid.uuid4())
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO research_notes (id, ticker, version, note_markdown, sources_markdown, changelog_markdown, note_docx, charts, metadata)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
            ''', (note_id, ticker, version, note_md, sources_md, changelog_md, docx_b64,
                  json.dumps([{'type': c['type'], 'filename': c['filename']} for c in charts]),
                  json.dumps({
                      'mode': mode,
                      'documentsProcessed': len(docs),
                      'documentFilenames': [d['filename'] for d in docs],
                      'provider': result.get('provider', ''),
                      'model': result.get('model', ''),
                      'charCount': len(note_md),
                  })))

        _update_pipeline_job(job_id, status='complete', current_step='Complete', progress=100,
                            result=json.dumps({'noteId': note_id, 'version': version, 'ticker': ticker}))
        print(f'[note-gen {job_id}] Complete: {ticker} v{version}')
        _notify_telegram(f"*Pipeline:* {ticker} note generated (v{version}, {len(note_md):,} chars)")

    except Exception as e:
        print(f'[note-gen {job_id}] Failed: {e}')
        import traceback
        traceback.print_exc()
        _update_pipeline_job(job_id, status='failed', error=str(e), current_step='Failed')
        _notify_telegram(f"*Pipeline:* {ticker} note generation FAILED\n{str(e)[:200]}")


def _generate_donut_chart(ticker, chart_type, data, value_key):
    """Generate a donut chart as base64 PNG."""
    try:
        import numpy as np
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.patheffects as pe

        labels = [d.get('segment', '') for d in data]
        values = [d.get(value_key, d.get('revenue', d.get('profit', 0))) for d in data]

        if not values or sum(values) == 0:
            return None

        total = sum(values)
        colors = ['#5DADE2', '#F7DC6F', '#F1948A', '#7DCEA0', '#BB8FCE', '#85C1E9', '#F8C471', '#82E0AA']

        fig, ax = plt.subplots(figsize=(10, 8), facecolor='white')
        wedges, _ = ax.pie(values, labels=None, colors=colors[:len(values)],
                           startangle=90, wedgeprops={'width': 0.58, 'edgecolor': 'none', 'linewidth': 0})

        # Center text
        centre_circle = plt.Circle((0, 0), 0.30, fc='white')
        ax.add_artist(centre_circle)
        total_str = f'\${total/1000:.1f}B' if total >= 1000 else f'\${total:.0f}M'
        ax.text(0, 0.05, 'Total', ha='center', va='center', fontsize=12, color='#333333', fontweight='normal')
        ax.text(0, -0.08, total_str, ha='center', va='center', fontsize=16, color='#333333', fontweight='bold')

        # Inside labels
        for i, (wedge, value) in enumerate(zip(wedges, values)):
            pct = value / total * 100
            ang = (wedge.theta2 - wedge.theta1) / 2. + wedge.theta1
            x = 0.70 * np.cos(np.deg2rad(ang))
            y = 0.70 * np.sin(np.deg2rad(ang))
            val_str = f'\${value/1000:.1f}B' if value >= 1000 else f'\${value:.0f}M'
            fontsize = 11 if pct > 15 else 10 if pct > 8 else 9
            if pct >= 3:
                ax.text(x, y, f'{val_str}\n({pct:.1f}%)', ha='center', va='center',
                        fontsize=fontsize, fontweight='bold', color='white',
                        path_effects=[pe.withStroke(linewidth=2, foreground='black')])

        # Outside labels
        for i, (wedge, label) in enumerate(zip(wedges, labels)):
            ang = (wedge.theta2 - wedge.theta1) / 2. + wedge.theta1
            x = 1.15 * np.cos(np.deg2rad(ang))
            y = 1.15 * np.sin(np.deg2rad(ang))
            ax.text(x, y, label, ha='center', va='center', fontsize=10, fontweight='bold', color='#333333')

        ax.set_title(f'{ticker} — {chart_type} Breakdown', fontsize=14, fontweight='bold', color='#333333', pad=20)
        ax.set_aspect('equal')
        plt.tight_layout()

        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        buf.seek(0)
        return base64.b64encode(buf.read()).decode('ascii')
    except Exception as e:
        print(f'Chart generation error for {ticker} {chart_type}: {e}')
        return None


def _generate_note_docx(ticker, company, markdown_text, charts):
    """Generate a Word document from markdown note text with embedded charts."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)
        font.color.rgb = RGBColor(0, 0, 0)

        # Title
        title = doc.add_heading(f'{ticker} — {company}', level=1)
        title.runs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Subtitle
        sub = doc.add_paragraph(f'Equity Research Note — {datetime.utcnow().strftime("%B %Y")}')
        sub.style.font.size = Pt(10)
        sub.style.font.color.rgb = RGBColor(100, 100, 100)

        # Parse markdown and add to doc
        lines = markdown_text.split('\n')
        i = 0

        while i < len(lines):
            line = lines[i]

            # Headings
            if line.startswith('### '):
                h = doc.add_heading(line[4:].strip(), level=3)
                h.paragraph_format.keep_with_next = True
                h.runs[0].font.color.rgb = RGBColor(0, 0, 0)
            elif line.startswith('## '):
                section_title = line[3:].strip()
                h = doc.add_heading(section_title, level=2)
                h.paragraph_format.keep_with_next = True
                h.runs[0].font.color.rgb = RGBColor(0, 0, 0)
            elif line.startswith('# '):
                h = doc.add_heading(line[2:].strip(), level=1)
                h.paragraph_format.keep_with_next = True
                h.runs[0].font.color.rgb = RGBColor(0, 0, 0)
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:].strip(), style='List Bullet')
            elif re.match(r'^\d+\.\s', line):
                doc.add_paragraph(line.split('. ', 1)[-1].strip(), style='List Number')
            elif line.strip() == '':
                pass  # Skip blank lines
            else:
                # Handle bold/italic text
                p = doc.add_paragraph()
                _add_formatted_text(p, line)

            i += 1

        # Add charts at end
        if charts:
            ch = doc.add_heading('Charts', level=2)
            ch.paragraph_format.keep_with_next = True
            ch.runs[0].font.color.rgb = RGBColor(0, 0, 0)
            for chart in charts:
                if chart.get('data'):
                    # Caption stays with chart image
                    cap = doc.add_paragraph(chart.get('filename', ''))
                    cap.paragraph_format.keep_with_next = True
                    img_bytes = base64.b64decode(chart['data'])
                    img_buf = io.BytesIO(img_bytes)
                    doc.add_picture(img_buf, width=Inches(6))
                    doc.add_paragraph('')  # spacing

        # Save to bytes
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return base64.b64encode(buf.read()).decode('ascii')
    except Exception as e:
        print(f'DOCX generation error: {e}')
        import traceback
        traceback.print_exc()
        return ''


def _add_formatted_text(paragraph, text):
    """Add text to a paragraph, handling **bold** and *italic* markdown."""
    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith('*') and part.endswith('*'):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        else:
            paragraph.add_run(part)


# ============================================
# CATALYST SYNTHESIS
# ============================================

CATALYST_LENGTH_PRESETS = {
    'quick': {'label': 'Quick Take (1 paragraph)', 'words': '50-100', 'instruction': 'Write a single paragraph (3-5 sentences) with the headline conclusion and key investment implication.'},
    'summary': {'label': 'Summary (2-3 paragraphs)', 'words': '200-400', 'instruction': 'Write 2-3 paragraphs covering key findings and investment implications. Be concise but substantive.'},
    'standard': {'label': 'Standard (1-2 pages)', 'words': '600-1000', 'instruction': 'Write a 1-2 page report balancing positives and negatives with a clear investment conclusion. Include key data points.'},
    'deep': {'label': 'Deep Dive (3-4 pages)', 'words': '1500-2500', 'instruction': 'Write a detailed 3-4 page analysis with data tables, bull/bear cases, and a comprehensive investment conclusion.'},
    'comprehensive': {'label': 'Comprehensive (5+ pages)', 'words': '3000+', 'instruction': 'Write a comprehensive 5+ page deep-dive with full analysis, data appendices, competitive context, and detailed investment conclusion.'},
}

CATALYST_SYNTHESIS_PROMPT = """You are a senior equity research analyst writing a synthesis report for your portfolio manager.

## ATTRIBUTION RULES (CRITICAL — FOLLOW EXACTLY)
- NEVER reference any specific broker, sellside firm, or analyst by name (no Goldman Sachs, Morgan Stanley, JPMorgan, Bernstein, etc.)
- NEVER cite analyst counts, ratings, or consensus targets as arguments
- NEVER use sellside sentiment as thesis support (e.g., "the Street is constructive", "consensus targets suggest upside")
- DO synthesize the factual content, data points, and analysis from the sources
- You may reference "estimates" or "consensus" as factual data points in valuation sections only
- Write in FIRST PERSON perspective ("I think", "My view is", "I see the risk-reward as")

## TONE & STYLE
- Confident, direct analyst voice writing to their PM
- No AI filler phrases ("it's worth noting", "delve into", "poised to", "navigate the landscape")
- Lead with conclusions, support with evidence
- Be specific with numbers, data points, and dates

## LENGTH
{length_instruction}

## TOPIC
Ticker: {ticker}
Topic: {topic}
{custom_instructions}

## SOURCE DOCUMENTS
The following are source documents to synthesize:
{source_content}

Write the synthesis report now. Start with a clear title line, then the analysis."""


@app.route('/api/catalysts/folders', methods=['GET'])
def list_catalyst_folders():
    """List catalyst topic subfolders for a ticker (read by local agent)."""
    ticker = request.args.get('ticker', '').upper()
    if not ticker:
        return jsonify({'error': 'Ticker required'}), 400

    # Check if we have cached folder data from the local agent
    with get_db() as (_, cur):
        cur.execute('SELECT folders, updated_at FROM catalyst_folders WHERE ticker = %s', (ticker,))
        row = cur.fetchone()

    if row:
        folders = row['folders'] if isinstance(row['folders'], list) else json.loads(row['folders'] or '[]')
        return jsonify({'ticker': ticker, 'folders': folders, 'updatedAt': row['updated_at'].isoformat() if row['updated_at'] else None})
    return jsonify({'ticker': ticker, 'folders': []})


@app.route('/api/catalysts/folders/refresh', methods=['POST'])
def refresh_catalyst_folders():
    """Request the local agent to refresh folder listings."""
    ticker = request.args.get('ticker') or (request.get_json() or {}).get('ticker', '')
    ticker = ticker.upper() if ticker else ''
    # Create a lightweight job for the local agent to scan folders
    job_id = str(uuid.uuid4())
    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO research_pipeline_jobs (id, batch_id, ticker, job_type, status, progress, current_step, total_steps, steps_detail)
            VALUES (%s, %s, %s, 'scan_catalysts', 'queued', 0, 'Waiting for local agent', 1, %s)
        ''', (job_id, str(uuid.uuid4()), ticker or 'ALL', json.dumps({})))
    return jsonify({'jobId': job_id, 'message': 'Scan requested'})


@app.route('/api/catalysts/folders', methods=['POST'])
def update_catalyst_folders():
    """Local agent reports catalyst folder contents."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper()
    folders = data.get('folders', [])
    if not ticker:
        return jsonify({'error': 'Ticker required'}), 400
    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO catalyst_folders (ticker, folders, updated_at)
            VALUES (%s, %s, NOW())
            ON CONFLICT (ticker) DO UPDATE SET folders = EXCLUDED.folders, updated_at = NOW()
        ''', (ticker, json.dumps(folders)))
    return jsonify({'success': True})


@app.route('/api/catalysts/synthesize', methods=['POST'])
def start_catalyst_synthesis():
    """Create a catalyst synthesis job for the local agent."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper()
    topic = data.get('topic', '').strip()
    length = data.get('length', 'standard')
    custom_instructions = data.get('customInstructions', '')
    uploaded_files = data.get('uploadedFiles', [])  # [{name, data (base64), type}]

    if not ticker:
        return jsonify({'error': 'Ticker required'}), 400
    if not topic and not uploaded_files:
        return jsonify({'error': 'Topic folder or uploaded files required'}), 400
    if length not in CATALYST_LENGTH_PRESETS:
        length = 'standard'

    excluded_files = data.get('excludedFiles', [])

    job_id = str(uuid.uuid4())
    job_detail = {
        'topic': topic,
        'length': length,
        'customInstructions': custom_instructions,
        'uploadedFiles': uploaded_files,  # base64 file data included for upload path
        'excludedFiles': excluded_files,
    }

    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO research_pipeline_jobs (id, batch_id, ticker, job_type, status, progress, current_step, total_steps, steps_detail)
            VALUES (%s, %s, %s, 'synthesis', 'queued', 0, 'Waiting for processing', 4, %s)
        ''', (job_id, str(uuid.uuid4()), ticker, json.dumps(job_detail)))

    # If files were uploaded (no local agent needed), process on backend
    if uploaded_files and not topic:
        threading.Thread(target=_run_catalyst_synthesis_backend, args=(job_id, ticker, job_detail), daemon=True).start()

    return jsonify({'jobId': job_id, 'ticker': ticker})


@app.route('/api/catalysts/results/<job_id>', methods=['GET'])
def get_catalyst_result(job_id):
    """Get catalyst synthesis result."""
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_pipeline_jobs WHERE id = %s', (job_id,))
        job = cur.fetchone()
    if not job:
        return jsonify({'error': 'Not found'}), 404

    result = job['result'] if isinstance(job.get('result'), dict) else json.loads(job['result'] or '{}') if job.get('result') else {}
    steps = job['steps_detail'] if isinstance(job.get('steps_detail'), (dict, list)) else json.loads(job['steps_detail'] or '{}') if job.get('steps_detail') else {}

    return jsonify({
        'id': job['id'],
        'ticker': job['ticker'],
        'status': job['status'],
        'progress': job['progress'],
        'currentStep': job['current_step'],
        'result': result,
        'detail': steps if isinstance(steps, dict) else {},
        'createdAt': job['created_at'].isoformat() if job['created_at'] else None,
        'completedAt': job['completed_at'].isoformat() if job['completed_at'] else None,
    })


@app.route('/api/catalysts/history', methods=['GET'])
def catalyst_history():
    """List recent catalyst synthesis jobs."""
    ticker = request.args.get('ticker', '')
    limit = request.args.get('limit', 20, type=int)
    with get_db() as (_, cur):
        if ticker:
            cur.execute("SELECT id, ticker, status, progress, current_step, created_at, completed_at, steps_detail FROM research_pipeline_jobs WHERE job_type = 'synthesis' AND ticker = %s ORDER BY created_at DESC LIMIT %s", (ticker.upper(), limit))
        else:
            cur.execute("SELECT id, ticker, status, progress, current_step, created_at, completed_at, steps_detail FROM research_pipeline_jobs WHERE job_type = 'synthesis' ORDER BY created_at DESC LIMIT %s", (limit,))
        rows = cur.fetchall()
    return jsonify({'jobs': [{
        'id': r['id'], 'ticker': r['ticker'], 'status': r['status'], 'progress': r['progress'],
        'currentStep': r['current_step'],
        'topic': (json.loads(r['steps_detail']) if isinstance(r['steps_detail'], str) else r['steps_detail'] or {}).get('topic', ''),
        'length': (json.loads(r['steps_detail']) if isinstance(r['steps_detail'], str) else r['steps_detail'] or {}).get('length', ''),
        'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
        'completedAt': r['completed_at'].isoformat() if r['completed_at'] else None,
    } for r in rows]})


@app.route('/api/catalysts/result/<job_id>/docx', methods=['GET'])
def get_catalyst_docx(job_id):
    """Generate and return docx for a catalyst synthesis."""
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_pipeline_jobs WHERE id = %s', (job_id,))
        job = cur.fetchone()
    if not job or job['status'] != 'complete':
        return jsonify({'error': 'Job not found or not complete'}), 404

    result = job['result'] if isinstance(job.get('result'), dict) else json.loads(job['result'] or '{}')
    markdown = result.get('markdown', '')
    ticker = job['ticker']
    detail = job['steps_detail'] if isinstance(job.get('steps_detail'), (dict, list)) else json.loads(job['steps_detail'] or '{}')
    topic = detail.get('topic', 'Catalyst') if isinstance(detail, dict) else 'Catalyst'

    docx_b64 = _generate_note_docx(ticker, topic, markdown, [])
    if not docx_b64:
        return jsonify({'error': 'DOCX generation failed'}), 500

    return jsonify({'docx': docx_b64, 'filename': f'{ticker}_{topic.replace("/", "_").replace(" ", "_")}_Synthesis.docx'})


@app.route('/api/catalysts/result/<job_id>/pdf', methods=['GET'])
def get_catalyst_pdf(job_id):
    """Generate and return PDF for a catalyst synthesis."""
    from xhtml2pdf import pisa
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_pipeline_jobs WHERE id = %s', (job_id,))
        job = cur.fetchone()
    if not job or job['status'] != 'complete':
        return jsonify({'error': 'Job not found or not complete'}), 404

    result = job['result'] if isinstance(job.get('result'), dict) else json.loads(job['result'] or '{}')
    markdown = result.get('markdown', '')
    ticker = job['ticker']
    detail = job['steps_detail'] if isinstance(job.get('steps_detail'), (dict, list)) else json.loads(job['steps_detail'] or '{}')
    topic = detail.get('topic', 'Catalyst') if isinstance(detail, dict) else 'Catalyst'

    # Convert markdown to simple HTML for xhtml2pdf
    import re as _re
    lines = markdown.split('\n')
    html_parts = []
    in_list = False
    in_table = False
    is_header_row = True

    def _inline_fmt(text):
        text = _re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
        text = _re.sub(r'\*(.*?)\*', r'<i>\1</i>', text)
        return text

    for line in lines:
        line = line.rstrip()

        # Table detection: lines starting and ending with |
        if line.strip().startswith('|') and line.strip().endswith('|'):
            if in_list: html_parts.append('</ul>'); in_list = False
            # Skip separator rows (|---|---|)
            if _re.match(r'^\|[\s\-:|]+\|$', line.strip()):
                continue
            cells = [c.strip() for c in line.strip().strip('|').split('|')]
            if not in_table:
                in_table = True
                is_header_row = True
                table_col_count = len(cells)
                font_size = '7pt' if table_col_count >= 7 else '8pt' if table_col_count >= 5 else '9pt'
                col_width = f'{int(100 / max(table_col_count, 1))}%'
                html_parts.append(f'<table width="100%" cellpadding="3" cellspacing="0" style="border-collapse: collapse; margin: 8px 0; font-size: {font_size}; word-wrap: break-word;">')
            if is_header_row:
                html_parts.append('<tr>' + ''.join(f'<th width="{col_width}" style="border: 1px solid #999; background-color: #e8e8e8; padding: 3px 4px; font-weight: bold; text-align: left;">{_inline_fmt(c)}</th>' for c in cells) + '</tr>')
                is_header_row = False
            else:
                html_parts.append('<tr>' + ''.join(f'<td width="{col_width}" style="border: 1px solid #ccc; padding: 3px 4px;">{_inline_fmt(c)}</td>' for c in cells) + '</tr>')
            continue
        else:
            if in_table:
                html_parts.append('</table>')
                in_table = False

        if line.startswith('### '):
            if in_list: html_parts.append('</ul>'); in_list = False
            html_parts.append(f'<h3>{_inline_fmt(line[4:])}</h3>')
        elif line.startswith('## '):
            if in_list: html_parts.append('</ul>'); in_list = False
            html_parts.append(f'<h2>{_inline_fmt(line[3:])}</h2>')
        elif line.startswith('# '):
            if in_list: html_parts.append('</ul>'); in_list = False
            html_parts.append(f'<h1>{_inline_fmt(line[2:])}</h1>')
        elif line.startswith('- ') or line.startswith('* '):
            if not in_list: html_parts.append('<ul>'); in_list = True
            html_parts.append(f'<li>{_inline_fmt(line[2:])}</li>')
        elif line.strip() == '':
            if in_list: html_parts.append('</ul>'); in_list = False
            continue
        else:
            if in_list: html_parts.append('</ul>'); in_list = False
            html_parts.append(f'<p>{_inline_fmt(line)}</p>')
    if in_list: html_parts.append('</ul>')
    if in_table: html_parts.append('</table>')

    body_html = '\n'.join(html_parts)
    full_html = f"""<html><head><style>
        body {{ font-family: Calibri, Arial, sans-serif; font-size: 11pt; color: #000; margin: 40px; line-height: 1.5; }}
        h1 {{ font-size: 18pt; margin: 0 0 4px 0; }}
        h2 {{ font-size: 14pt; margin: 14px 0 4px 0; border-bottom: 1px solid #ccc; padding-bottom: 2px; }}
        h3 {{ font-size: 12pt; margin: 10px 0 2px 0; }}
        p {{ margin: 6px 0; }}
        ul {{ margin: 4px 0 4px 20px; padding: 0; }}
        li {{ margin: 2px 0; }}
    </style></head><body>
        <h1>{ticker} -- {topic}</h1>
        <p style="color: #666; font-size: 10pt;">Catalyst Synthesis | {datetime.utcnow().strftime('%B %d, %Y')}</p>
        <hr/>
        {body_html}
    </body></html>"""

    pdf_buf = io.BytesIO()
    pisa_status = pisa.CreatePDF(full_html, dest=pdf_buf)
    if pisa_status.err:
        return jsonify({'error': 'PDF generation failed'}), 500

    pdf_b64 = base64.b64encode(pdf_buf.getvalue()).decode('ascii')
    return jsonify({'pdf': pdf_b64, 'filename': f'{ticker}_{topic.replace("/", "_").replace(" ", "_")}_Synthesis.pdf'})


@app.route('/api/catalysts/result/<job_id>/save', methods=['POST'])
def save_catalyst_to_docs(job_id):
    """Save catalyst synthesis result to research_documents."""
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_pipeline_jobs WHERE id = %s', (job_id,))
        job = cur.fetchone()
    if not job or job['status'] != 'complete':
        return jsonify({'error': 'Job not found or not complete'}), 404

    result = job['result'] if isinstance(job.get('result'), dict) else json.loads(job['result'] or '{}')
    markdown = result.get('markdown', '')
    ticker = job['ticker']
    detail = job['steps_detail'] if isinstance(job.get('steps_detail'), (dict, list)) else json.loads(job['steps_detail'] or '{}')
    topic = detail.get('topic', 'Catalyst') if isinstance(detail, dict) else 'Catalyst'

    doc_id = f"catalyst-{job_id}"
    with get_db(commit=True) as (conn, cur):
        cat_id = 'cat-catalyst-synthesis'
        cur.execute("INSERT INTO research_categories (id, name, type) VALUES (%s, 'Catalyst Synthesis', 'topic') ON CONFLICT (id) DO NOTHING", (cat_id,))

        doc_name = f"{ticker} -- {topic}"
        cur.execute('''
            INSERT INTO research_documents (id, category_id, name, content, doc_type, created_at)
            VALUES (%s, %s, %s, %s, 'catalyst', NOW())
            ON CONFLICT (id) DO UPDATE SET content = EXCLUDED.content
        ''', (doc_id, cat_id, doc_name, markdown))

    cache.invalidate('research_categories')
    return jsonify({'success': True, 'docId': doc_id})


def _generate_source_provenance(source_names, synthesis_markdown, ticker, topic):
    """Generate a source provenance summary explaining which documents contributed what to the synthesis."""
    if not source_names:
        return None
    try:
        file_list = '\n'.join([f"- {name}" for name in source_names])
        prompt = f"""You just produced a synthesis report for {ticker} on the topic "{topic}". Below are the source documents you had access to and the synthesis you wrote.

SOURCE DOCUMENTS PROVIDED:
{file_list}

SYNTHESIS OUTPUT (first 3000 chars):
{synthesis_markdown[:3000]}

Now write a brief SOURCE PROVENANCE summary for the analyst's own reference. For each source document:
1. State the document name
2. Describe in 1-2 sentences what key information or data points from that document were used in the synthesis
3. If a document was not materially used (e.g., duplicate content, irrelevant), note that

Keep it concise and factual. Use markdown formatting with bullet points. This is an internal reference, not part of the analysis."""

        result = call_llm(
            messages=[{"role": "user", "content": prompt}],
            system="You are an analyst documenting your research process. Be specific about which data points came from which sources.",
            tier="fast",
            max_tokens=2048,
        )
        return result['text']
    except Exception as e:
        print(f"[source-provenance] Failed to generate: {e}")
        return f"Source provenance generation failed. {len(source_names)} files were provided: {', '.join(source_names)}"


def _run_catalyst_synthesis_backend(job_id, ticker, detail):
    """Backend-side catalyst synthesis when files are uploaded (no local agent needed)."""
    try:
        uploaded_files = detail.get('uploadedFiles', [])
        length = detail.get('length', 'standard')
        custom_instructions = detail.get('customInstructions', '')
        topic = detail.get('topic', 'Uploaded Documents')

        def update_job(step, progress, result=None, status='running'):
            with get_db(commit=True) as (conn, cur):
                if result:
                    cur.execute('UPDATE research_pipeline_jobs SET current_step=%s, progress=%s, result=%s, status=%s, completed_at=NOW(), updated_at=NOW() WHERE id=%s',
                               (step, progress, json.dumps(result), status, job_id))
                else:
                    cur.execute('UPDATE research_pipeline_jobs SET current_step=%s, progress=%s, status=%s, updated_at=NOW() WHERE id=%s',
                               (step, progress, status, job_id))

        update_job('Reading uploaded files...', 10)

        # Build source content from uploaded files
        source_parts = []
        for f in uploaded_files:
            name = f.get('name', 'unnamed')
            content = f.get('text', '')
            if not content and f.get('data'):
                content = f'[Binary file: {name} - content extracted on upload]'
            source_parts.append(f"### Source: {name}\n{content}\n")

        source_content = '\n---\n'.join(source_parts)
        if len(source_content) > 80000:
            source_content = source_content[:80000] + '\n\n[Content truncated for length]'

        update_job('Synthesizing report...', 40)

        length_config = CATALYST_LENGTH_PRESETS.get(length, CATALYST_LENGTH_PRESETS['standard'])
        custom_block = f"\nAdditional Instructions: {custom_instructions}" if custom_instructions else ""

        prompt = CATALYST_SYNTHESIS_PROMPT.format(
            length_instruction=length_config['instruction'],
            ticker=ticker,
            topic=topic,
            custom_instructions=custom_block,
            source_content=source_content,
        )

        result = call_llm(
            messages=[{"role": "user", "content": prompt}],
            system="You are a senior equity research analyst. Follow all instructions precisely.",
            tier="standard",
            max_tokens=8192,
        )

        markdown = result['text']

        # Generate source provenance (separate from synthesis)
        update_job('Analyzing source contributions...', 70)
        source_names = [f.get('name', 'unnamed') for f in uploaded_files]
        provenance = _generate_source_provenance(source_names, markdown, ticker, topic)

        update_job('Generating Word document...', 85)

        docx_b64 = _generate_note_docx(ticker, topic, markdown, [])

        update_job('Complete', 100, result={
            'markdown': markdown,
            'docx': docx_b64,
            'topic': topic,
            'length': length,
            'fileCount': len(uploaded_files),
            'sourceFiles': source_names,
            'sourceProvenance': provenance,
        }, status='complete')

        print(f"[catalyst-synthesis {job_id}] Complete: {ticker}/{topic}")

    except Exception as e:
        print(f"[catalyst-synthesis {job_id}] Failed: {e}")
        import traceback; traceback.print_exc()
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE research_pipeline_jobs SET status='error', current_step=%s, updated_at=NOW() WHERE id=%s",
                       (f"Failed: {str(e)[:200]}", job_id))


# ============================================
# RESEARCH ANALYSIS ENDPOINT
# ============================================

@app.route('/api/research-analyze', methods=['POST'])
def research_analyze():
    """
    Deep analysis of sell-side research using customizable prompt frameworks.
    Supports native PDF/image document blocks for full-fidelity analysis
    (charts, graphs, tables, images — not just extracted text).

    Request body:
    {
        "text": "Framework prompt text",
        "promptId": "executive-brief",
        "promptName": "Executive Brief",
        "apiKey": "sk-ant-...",
        "files": [{"data": "base64...", "type": "application/pdf", "name": "report.pdf"}, ...]  (optional)
    }
    """
    try:
        data = request.json
        text = data.get('text', '')
        prompt_id = data.get('promptId', '')
        prompt_name = data.get('promptName', '')
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '')
        files = data.get('files', [])

        if not text:
            return jsonify({'error': 'No text provided'}), 400

        if not api_key:
            return jsonify({'error': 'No API key provided. Please add your API key in Settings.'}), 400

        # Build content: if files provided, use native document/image blocks
        if files:
            content = []
            for f in files:
                f_type = f.get('type', '')
                f_data = f.get('data', '')
                f_name = f.get('name', '')
                if not f_data:
                    continue
                if f_type == 'application/pdf':
                    content.append({"type": "text", "text": f"--- {f_name} ---"})
                    content.append({"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": f_data}})
                elif f_type.startswith('image/'):
                    content.append({"type": "text", "text": f"--- {f_name} ---"})
                    content.append({"type": "image", "source": {"type": "base64", "media_type": f_type, "data": f_data}})
            # Append the prompt text at the end
            content.append({"type": "text", "text": text})
            user_msg = content
        else:
            user_msg = text

        result = call_llm(
            messages=[{"role": "user", "content": user_msg}],
            tier="standard",
            max_tokens=4096,
            timeout=180,
            anthropic_api_key=api_key,
        )

        return jsonify({
            'result': result["text"],
            'promptId': prompt_id,
            'promptName': prompt_name,
            'usage': result["usage"]
        })

    except LLMError as e:
        print(f"Research analysis LLM error: {e}")
        return jsonify({'error': str(e)}), 502
    except Exception as e:
        print(f"Research analysis error: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# RESEARCH HIERARCHICAL ENDPOINTS
# ============================================

# --- Categories ---
@app.route('/api/research-categories', methods=['GET'])
def get_research_categories():
    """Get all research categories"""
    try:
        cached = cache.get('research_categories')
        if cached is not None:
            return jsonify(cached)

        with get_db() as (_, cur):
            cur.execute('SELECT id, name, type, created_at FROM research_categories ORDER BY created_at DESC')
            rows = cur.fetchall()

        result = [{
            'id': row['id'],
            'name': row['name'],
            'type': row['type'],
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None
        } for row in rows]
        cache.set('research_categories', result, ttl=600)
        return jsonify(result)
    except Exception as e:
        print(f"Error getting research categories: {e}")
        return jsonify([])


@app.route('/api/save-research-category', methods=['POST'])
def save_research_category():
    """Save a research category"""
    try:
        data = request.json
        cat_id = data.get('id', '')
        
        if not cat_id:
            return jsonify({'error': 'Category ID is required'}), 400
        
        with get_db(commit=True) as (_, cur):
            cur.execute('''
                INSERT INTO research_categories (id, name, type)
                VALUES (%s, %s, %s)
                ON CONFLICT (id) DO UPDATE SET name = EXCLUDED.name, type = EXCLUDED.type
                RETURNING id
            ''', (cat_id, data.get('name', ''), data.get('type', 'ticker')))

        cache.invalidate('research_categories')
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error saving research category: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/delete-research-category', methods=['POST'])
def delete_research_category():
    """Delete a research category and all its documents/analyses"""
    try:
        data = request.json
        cat_id = data.get('id', '')

        if not cat_id:
            return jsonify({'error': 'Category ID is required'}), 400

        with get_db(commit=True) as (_, cur):
            cur.execute('DELETE FROM research_categories WHERE id = %s', (cat_id,))

        cache.invalidate('research_categories')
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting research category: {e}")
        return jsonify({'error': str(e)}), 500


# --- Documents ---
@app.route('/api/research-documents', methods=['GET'])
def get_research_documents():
    """Get all research documents"""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT id, category_id, name, content, file_names, smart_name, original_filename, published_date, doc_type, has_stored_files, created_at FROM research_documents ORDER BY created_at DESC')
            rows = cur.fetchall()

        return jsonify([{
            'id': row['id'],
            'categoryId': row['category_id'],
            'name': row['name'],
            'content': row['content'],
            'fileNames': row['file_names'] or [],
            'smartName': row['smart_name'],
            'originalFilename': row['original_filename'],
            'publishedDate': row['published_date'],
            'docType': row.get('doc_type') or 'other',
            'hasStoredFiles': row['has_stored_files'] or False,
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None
        } for row in rows])
    except Exception as e:
        print(f"Error getting research documents: {e}")
        return jsonify([])


@app.route('/api/save-research-document', methods=['POST'])
def save_research_document():
    """Save a research document"""
    try:
        data = request.json
        doc_id = data.get('id', '')
        
        print(f"📄 save_research_document: id={doc_id}, name={data.get('name', '')[:50]}")
        
        if not doc_id:
            return jsonify({'error': 'Document ID is required'}), 400
        
        with get_db(commit=True) as (_, cur):
            cur.execute('''
                INSERT INTO research_documents (id, category_id, name, content, file_names, smart_name, original_filename, published_date, doc_type, has_stored_files)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (id) DO UPDATE SET
                    name = EXCLUDED.name,
                    content = EXCLUDED.content,
                    file_names = EXCLUDED.file_names,
                    smart_name = EXCLUDED.smart_name,
                    original_filename = EXCLUDED.original_filename,
                    published_date = EXCLUDED.published_date,
                    doc_type = EXCLUDED.doc_type,
                    has_stored_files = EXCLUDED.has_stored_files
                RETURNING id
            ''', (
                doc_id,
                data.get('categoryId', ''),
                data.get('name', ''),
                data.get('content', ''),
                json.dumps(data.get('fileNames', [])),
                data.get('smartName'),
                data.get('originalFilename'),
                data.get('publishedDate'),
                data.get('docType', 'other'),
                data.get('hasStoredFiles', False)
            ))

        print(f"✅ Document saved: {doc_id}")
        return jsonify({'success': True})
    except Exception as e:
        print(f"❌ Error saving research document: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/save-research-file', methods=['POST'])
def save_research_file():
    """Save a file for a research document"""
    try:
        data = request.json
        if not data:
            print("❌ No JSON data received")
            return jsonify({'error': 'No JSON data received'}), 400

        document_id = data.get('documentId', '')
        filename = data.get('filename', '')
        file_type = data.get('fileType', '')
        file_data = data.get('fileData', '')
        file_size = data.get('fileSize', 0)

        print(f"📁 save_research_file: docId={document_id}, filename={filename}, fileType={file_type}, dataLen={len(file_data) if file_data else 0}, fileSize={file_size}")

        if not document_id or not filename:
            print(f"❌ Missing required fields: docId={document_id}, filename={filename}")
            return jsonify({'error': 'Document ID and filename are required'}), 400

        if not file_data:
            print(f"❌ No file data provided for {filename}")
            return jsonify({'error': 'No file data provided'}), 400

        with get_db(commit=True) as (conn, cur):
            # Check if document exists first
            cur.execute('SELECT id FROM research_documents WHERE id = %s', (document_id,))
            doc_exists = cur.fetchone()
            if not doc_exists:
                print(f"❌ Document {document_id} does not exist in research_documents table")
                return jsonify({'error': f'Document {document_id} not found - must save document first'}), 400

            print(f"✅ Document {document_id} exists, proceeding with file save")

            cur.execute('''
                INSERT INTO research_document_files (document_id, filename, file_type, file_data, file_size)
                VALUES (%s, %s, %s, %s, %s)
                RETURNING id
            ''', (document_id, filename, file_type, file_data, file_size))

            result = cur.fetchone()
            if result is None:
                print(f"❌ INSERT did not return an id")
                conn.rollback()
                return jsonify({'error': 'Insert failed - no id returned'}), 500

            inserted_id = result['id']

        print(f"✅ File saved successfully: id={inserted_id}, filename={filename}")
        return jsonify({'success': True, 'id': inserted_id})
    except Exception as e:
        print(f"❌ Error saving research file: {type(e).__name__}: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'{type(e).__name__}: {str(e)}'}), 500


@app.route('/api/research-document-files/<document_id>', methods=['GET'])
def get_research_document_files(document_id):
    """Get stored files for a research document"""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, filename, file_type, file_data, file_size, created_at
                FROM research_document_files
                WHERE document_id = %s
                ORDER BY created_at
            ''', (document_id,))
            rows = cur.fetchall()

        return jsonify([{
            'id': row['id'],
            'filename': row['filename'],
            'fileType': row['file_type'],
            'fileData': row['file_data'],
            'fileSize': row['file_size'],
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None
        } for row in rows])
    except Exception as e:
        print(f"Error getting research document files: {e}")
        return jsonify([])


@app.route('/api/delete-research-file/<int:file_id>', methods=['DELETE'])
def delete_research_file(file_id):
    """Delete a stored research file"""
    try:
        with get_db(commit=True) as (_, cur):
            cur.execute('DELETE FROM research_document_files WHERE id = %s', (file_id,))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting research file: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/delete-research-document', methods=['POST'])
def delete_research_document():
    """Delete a research document and all its analyses"""
    try:
        data = request.json
        doc_id = data.get('id', '')
        
        if not doc_id:
            return jsonify({'error': 'Document ID is required'}), 400
        
        with get_db(commit=True) as (_, cur):
            # CASCADE will delete analyses
            cur.execute('DELETE FROM research_documents WHERE id = %s', (doc_id,))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting research document: {e}")
        return jsonify({'error': str(e)}), 500


# --- Analyses ---
@app.route('/api/research-analyses', methods=['GET'])
def get_research_analyses():
    """Get all research analyses"""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT id, document_id, prompt_id, prompt_name, prompt_icon, result, usage, created_at FROM research_analyses ORDER BY created_at DESC')
            rows = cur.fetchall()

        return jsonify([{
            'id': row['id'],
            'documentId': row['document_id'],
            'promptId': row['prompt_id'],
            'promptName': row['prompt_name'],
            'promptIcon': row['prompt_icon'],
            'result': row['result'],
            'usage': row['usage'],
            'createdAt': row['created_at'].isoformat() if row['created_at'] else None
        } for row in rows])
    except Exception as e:
        print(f"Error getting research analyses: {e}")
        return jsonify([])


@app.route('/api/save-research-analysis', methods=['POST'])
def save_research_analysis():
    """Save a research analysis"""
    try:
        data = request.json
        analysis_id = data.get('id', '')
        
        if not analysis_id:
            return jsonify({'error': 'Analysis ID is required'}), 400
        
        with get_db(commit=True) as (_, cur):
            cur.execute('''
                INSERT INTO research_analyses (id, document_id, prompt_id, prompt_name, prompt_icon, result, usage)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (id) DO UPDATE SET
                    result = EXCLUDED.result,
                    usage = EXCLUDED.usage
                RETURNING id
            ''', (
                analysis_id,
                data.get('documentId', ''),
                data.get('promptId', ''),
                data.get('promptName', ''),
                data.get('promptIcon', ''),
                data.get('result', ''),
                json.dumps(data.get('usage', {}))
            ))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error saving research analysis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/delete-research-analysis', methods=['POST'])
def delete_research_analysis():
    """Delete a research analysis"""
    try:
        data = request.json
        analysis_id = data.get('id', '')
        
        if not analysis_id:
            return jsonify({'error': 'Analysis ID is required'}), 400
        
        with get_db(commit=True) as (_, cur):
            cur.execute('DELETE FROM research_analyses WHERE id = %s', (analysis_id,))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting research analysis: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/email-research', methods=['POST'])
def email_research():
    """Email a research result"""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    import markdown
    import re
    
    def preprocess_bullets(text):
        """Convert bullet characters to standard markdown format"""
        # Split into lines for processing
        lines = text.split('\n')
        processed_lines = []
        
        for line in lines:
            # Handle lines that start with bullet characters (with optional whitespace)
            line = re.sub(r'^(\s*)[•·▪▸►‣⁃]\s*', r'\1- ', line)
            
            # Handle inline bullets (mid-line) - add newline before them
            # This catches patterns like "text • more text"
            line = re.sub(r'\s+[•·▪▸►‣⁃]\s+', '\n- ', line)
            
            processed_lines.append(line)
        
        return '\n'.join(processed_lines)
    
    try:
        data = request.json
        recipient = data.get('email', '')
        subject = data.get('subject', 'Research Analysis')
        content = data.get('content', '')
        prompt_name = data.get('promptName', '')
        ticker = data.get('ticker', '')
        topic = data.get('topic', '')
        smtp_config = data.get('smtpConfig', {})
        minimal = data.get('minimal', False)
        charts = data.get('charts', [])

        if not recipient:
            return jsonify({'error': 'Recipient email is required'}), 400

        # Preprocess bullet characters before markdown conversion
        processed_content = preprocess_bullets(content)

        # Convert markdown to HTML with nl2br for line break preservation
        try:
            content_html = markdown.markdown(
                processed_content,
                extensions=['tables', 'fenced_code', 'nl2br']
            )
        except:
            content_html = f"<pre style='white-space: pre-wrap;'>{content}</pre>"

        # Inject inline styles on table elements for email client compatibility
        content_html = content_html.replace('<table>', '<table style="border-collapse: collapse; width: 100%; margin: 12px 0; font-size: 13px;">')
        content_html = content_html.replace('<th>', '<th style="border: 1px solid #999; background-color: #e8e8e8; padding: 6px 8px; font-weight: bold; text-align: left;">')
        content_html = content_html.replace('<td>', '<td style="border: 1px solid #ccc; padding: 6px 8px;">')

        # Build charts HTML if provided (using cid: references for email embedding)
        charts_html = ''
        chart_attachments = []  # list of (cid, base64_data)
        if charts:
            for idx, chart in enumerate(charts):
                if isinstance(chart, dict) and chart.get('data'):
                    chart_type = chart.get('type', 'chart').title()
                    label = f'{ticker} {chart_type} Breakdown' if ticker else f'{chart_type} Breakdown'
                    cid = f'chart_{idx}'
                    charts_html += f'<p style="margin:16px 0 4px 0;font-weight:bold;font-size:11pt;color:#1e293b;">{label}</p>'
                    charts_html += f'<img src="cid:{cid}" style="width:100%;max-width:500px;margin:0 0 16px 0;" />'
                    chart_attachments.append((cid, chart['data']))

        if minimal:
            # Clean email — just content with basic styling, no header/footer
            html_body = f"""
        <html>
        <head>
            <style>
                body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; line-height: 1.6; color: #333; max-width: 800px; margin: 0 auto; padding: 20px; }}
                h1, h2, h3 {{ color: #1a1a2e; margin-top: 1.5em; margin-bottom: 0.5em; }}
                ul {{ margin: 10px 0; padding-left: 25px; }}
                li {{ margin-bottom: 8px; line-height: 1.5; }}
                ul ul {{ margin-top: 8px; }}
                p {{ margin: 0.8em 0; }}
                strong {{ color: #1e293b; }}
                hr {{ border: none; border-top: 1px solid #e2e8f0; margin: 1.5em 0; }}
                img {{ max-width: 100%; height: auto; }}
                table {{ border-collapse: collapse; width: 100%; margin: 12px 0; font-size: 13px; }}
                th {{ border: 1px solid #999; background-color: #e8e8e8; padding: 6px 8px; font-weight: bold; text-align: left; }}
                td {{ border: 1px solid #ccc; padding: 6px 8px; }}
            </style>
        </head>
        <body>
            {content_html}
            {charts_html}
        </body>
        </html>
            """
        else:
            # Full decorated email with header/footer
            header_info = []
            if ticker:
                header_info.append(f"<strong>Ticker:</strong> {ticker}")
            if topic:
                header_info.append(f"<strong>Topic:</strong> {topic}")
            if prompt_name:
                header_info.append(f"<strong>Framework:</strong> {prompt_name}")
            header_html = " | ".join(header_info) if header_info else ""

            html_body = f"""
        <html>
        <head>
            <style>
                body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; line-height: 1.6; color: #333; max-width: 800px; margin: 0 auto; padding: 20px; }}
                h1, h2, h3 {{ color: #1a1a2e; margin-top: 1.5em; margin-bottom: 0.5em; }}
                .header {{ background: linear-gradient(135deg, #0f172a, #1e293b); color: white; padding: 20px; border-radius: 10px; margin-bottom: 20px; }}
                .header h1 {{ margin: 0 0 10px 0; color: white; }}
                .header-meta {{ font-size: 14px; opacity: 0.9; }}
                .content {{ background: #f8fafc; padding: 20px; border-radius: 10px; border: 1px solid #e2e8f0; }}
                ul {{ margin: 10px 0; padding-left: 25px; }}
                li {{ margin-bottom: 8px; line-height: 1.5; }}
                ul ul {{ margin-top: 8px; }}
                table {{ border-collapse: collapse; width: 100%; margin: 15px 0; }}
                th, td {{ border: 1px solid #e2e8f0; padding: 10px; text-align: left; }}
                th {{ background: #f1f5f9; }}
                code {{ background: #f1f5f9; padding: 2px 6px; border-radius: 4px; font-size: 14px; }}
                pre {{ background: #1e293b; color: #e2e8f0; padding: 15px; border-radius: 8px; overflow-x: auto; }}
                p {{ margin: 0.8em 0; }}
                strong {{ color: #1e293b; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>Research Analysis</h1>
                <div class="header-meta">{header_html}</div>
            </div>
            <div class="content">
                {content_html}
            </div>
        </body>
        </html>
            """
        
        # Send email
        if chart_attachments:
            from email.mime.image import MIMEImage
            msg = MIMEMultipart('related')
            msg['Subject'] = subject
            msg['To'] = recipient
            alt_part = MIMEMultipart('alternative')
            alt_part.attach(MIMEText(content, 'plain'))
            alt_part.attach(MIMEText(html_body, 'html'))
            msg.attach(alt_part)
            for cid, b64_data in chart_attachments:
                img = MIMEImage(base64.b64decode(b64_data), _subtype='png')
                img.add_header('Content-ID', f'<{cid}>')
                img.add_header('Content-Disposition', 'inline')
                msg.attach(img)
        else:
            msg = MIMEMultipart('alternative')
            msg['Subject'] = subject
            msg['To'] = recipient
            msg.attach(MIMEText(content, 'plain'))
            msg.attach(MIMEText(html_body, 'html'))
        
        if smtp_config.get('use_gmail') and smtp_config.get('gmail_user') and smtp_config.get('gmail_app_password'):
            msg['From'] = smtp_config['gmail_user']
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
                server.login(smtp_config['gmail_user'], smtp_config['gmail_app_password'])
                server.send_message(msg)
        else:
            return jsonify({'error': 'Email not configured. Please set up Gmail in Settings.'}), 400
        
        return jsonify({'success': True, 'message': 'Research email sent successfully'})
        
    except smtplib.SMTPAuthenticationError:
        return jsonify({'error': 'Gmail authentication failed. Check your email and app password.'}), 401
    except smtplib.SMTPException as e:
        return jsonify({'error': f'SMTP error: {str(e)}'}), 500
    except Exception as e:
        print(f"Error sending research email: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# MEETING PREP - HELPERS & PROMPTS
# ============================================

import re as _re

def parse_mp_json(text):
    """Parse JSON from AI response, handling markdown fencing and truncation."""
    text = text.strip()
    if text.startswith("```"):
        lines = text.split("\n")
        lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        text = "\n".join(lines).strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    for sc, ec in [("[", "]"), ("{", "}")]:
        s = text.find(sc)
        e = text.rfind(ec)
        if s != -1 and e != -1 and e > s:
            try:
                return json.loads(text[s:e + 1])
            except json.JSONDecodeError:
                continue
    # Try to repair truncated JSON by closing open brackets/braces
    for sc, ec in [("{", "}"), ("[", "]")]:
        s = text.find(sc)
        if s != -1:
            fragment = text[s:]
            # Count open vs close brackets
            for attempt in range(5):
                try:
                    return json.loads(fragment)
                except json.JSONDecodeError:
                    # Try closing with appropriate bracket
                    open_braces = fragment.count('{') - fragment.count('}')
                    open_brackets = fragment.count('[') - fragment.count(']')
                    # Remove trailing partial content after last comma or complete value
                    last_comma = fragment.rfind(',')
                    if last_comma > 0:
                        fragment = fragment[:last_comma]
                    fragment += '}' * max(0, open_braces) + ']' * max(0, open_brackets)
                    try:
                        return json.loads(fragment)
                    except json.JSONDecodeError:
                        break
    raise ValueError(f"Could not parse JSON from response: {text[:300]}")


MP_DOC_TYPE_PATTERNS = {
    "earnings_transcript": [
        r"earnings\s*(call|transcript)",
        r"q[1-4]\s*\d{4}\s*(call|transcript|results)",
        r"(quarterly|annual)\s*results\s*call",
    ],
    "conference_transcript": [
        r"conference\s*(transcript|presentation)",
        r"investor\s*(day|conference|presentation)",
        r"fireside\s*chat",
    ],
    "broker_report": [
        r"(initiat|maintain|reiterat|upgrad|downgrad|price\s*target)",
        r"(buy|sell|hold|overweight|underweight|neutral|outperform)\s*(rating)?",
        r"(equity\s*research|research\s*report|analyst\s*note)",
    ],
    "press_release": [
        r"press\s*release",
        r"(announces|reports)\s*(q[1-4]|quarterly|annual|full.year)",
    ],
    "filing": [
        r"(10-[kq]|8-k|def\s*14a|proxy|annual\s*report)",
        r"securities\s*and\s*exchange\s*commission",
    ],
}

def classify_mp_document(filename, text_sample=""):
    """Auto-classify document type from filename and first 2000 chars of content."""
    combined = f"{filename} {text_sample[:2000]}".lower()
    for doc_type, patterns in MP_DOC_TYPE_PATTERNS.items():
        for pattern in patterns:
            if _re.search(pattern, combined, _re.IGNORECASE):
                return doc_type
    return "other"


MP_ANALYSIS_PROMPT = """You are a senior equity research analyst assistant preparing for a management meeting.

Analyze the following {doc_type} for {ticker} ({company_name}).

Extract and organize the following into a structured JSON response:

1. **key_metrics**: Array of objects with {{metric, value, change, period}} — revenue, margins, EPS, guidance, segment data, KPIs. Be specific with numbers.
2. **management_claims**: Array of strings — specific commitments, promises, strategic statements management made. Quote where possible.
3. **guidance_changes**: Array of objects with {{metric, old_guidance, new_guidance, direction}} — any changes to forward guidance.
4. **risks_concerns**: Array of strings — risks flagged by management, analysts, or evident from data.
5. **catalysts**: Array of strings — upcoming events, product launches, regulatory decisions, etc. that could move the stock.
6. **contradictions**: Array of strings — anything that contradicts prior statements, guidance, or consensus.
7. **notable_quotes**: Array of objects with {{quote, speaker, context}} — important verbatim quotes worth referencing.
8. **key_numbers**: Object — the most important 5-10 data points someone should know before a meeting.

Return ONLY valid JSON, no markdown fencing."""

MP_SYNTHESIS_PROMPT = """You are preparing a senior equity research analyst for a management meeting with {ticker} ({company_name}, {sector} sector).

You have analyses of {doc_count} documents spanning the {timeframe} period. Your task is to synthesize these into a coherent picture that identifies what the analyst MUST explore in the meeting.

{analyses_text}

{past_questions_text}

Synthesize into a JSON object with:

1. **narrative_arc**: 3-5 sentences describing the story across these documents — what's the trajectory? What's changed?
2. **key_themes**: Array of objects with {{theme, description, supporting_evidence}} — the 4-7 major themes emerging.
3. **contradictions**: Array of objects with {{claim_1, source_1, claim_2, source_2, significance}} — where documents or statements conflict.
4. **information_gaps**: Array of strings — what's NOT being discussed that should be? What data is missing?
5. **tone_shifts**: Array of objects with {{topic, old_tone, new_tone, source}} — where management messaging has shifted.
6. **unresolved_from_prior**: Array of objects with {{question, original_date, why_still_relevant}} — past questions that need follow-up.
7. **consensus_vs_reality**: Array of strings — where consensus expectations seem misaligned with what documents reveal.

Return ONLY valid JSON, no markdown fencing."""

MP_QUESTION_PROMPT = """You are helping a senior equity research analyst prepare 25-30 sophisticated questions for a management meeting with {ticker} ({company_name}, {sector} sector).

Context from document synthesis:
{synthesis_text}

CRITICAL RULES FOR QUESTION TEXT:
- Questions must sound 100% proprietary — as if the analyst developed them entirely from their own deep, independent research on the company
- The ONLY acceptable references in question text are: the company's own filings, earnings calls, press releases, management's own statements/quotes, and publicly reported financial data
- NEVER reference or allude to sell-side in ANY form. This includes:
  - Direct: "Wells Fargo notes...", "according to Deutsche Bank...", "Jefferies estimates..."
  - Indirect: "analysts are cutting estimates", "the Street expects", "consensus is...", "the investment community", "market participants", "some observers note", "there's skepticism among investors"
- Instead, ground every question in the company's OWN words, data, and disclosures: "You guided to X but reported Y", "Your 10-K shows margin compression from Z to W", "On the Q3 call you said X, but Q4 results suggest otherwise"
- The analyst's edge comes from connecting dots across the company's own disclosures — not from citing what other analysts think
- The "source" and "context" fields ARE where you indicate which broker report or document the question was derived from — that metadata is for the analyst's private reference only, never surfaced in the question

Generate questions that are:
- **Proprietary-sounding**: Every question reads as if the analyst personally identified the issue through their own deep research
- **Specific**: Reference concrete data points, quotes, or metrics — but attribute them to the company's own disclosures, not to broker commentary
- **Probing**: Push management beyond their prepared talking points — ask about inconsistencies, gaps, and changes
- **Organized**: Group by dynamically determined topics relevant to THIS company/sector (NOT a generic template)
- **Prioritized**: Mark each as high (must-ask, 8-10 questions), medium (important, 10-12 questions), or low (if time permits, 5-8 questions)
- **Strategic**: Include questions about capital allocation, competitive dynamics, and forward catalysts

For each question, also provide:
- **context**: Why this question matters — what data point or observation prompted it. THIS is where you can reference the sell-side source for the analyst's private notes (1-2 sentences)
- **source**: Which document(s) the question draws from (e.g., "Wells Fargo report, p.3; Q4 Earnings Transcript, p.12"). This is private metadata for the analyst.
- **follow_up_angle**: What to ask if management gives an evasive or generic answer

{unresolved_text}

Return a JSON array of topic groups:
[
  {{
    "topic": "Revenue Growth Trajectory",
    "description": "Brief description of why this topic matters for this company",
    "questions": [
      {{
        "question": "Proprietary-sounding question with NO broker/analyst attribution",
        "context": "Private note: sourced from Wells Fargo report highlighting X — important because Y",
        "source": "Wells Fargo report p.3, Q4 Earnings Transcript p.12",
        "priority": "high",
        "follow_up_angle": "If they deflect, ask about..."
      }}
    ]
  }}
]

Return ONLY valid JSON, no markdown fencing."""


_local_file_manifest = {}
_pending_doc_upload_requests = {}  # ticker -> {'requested_at': timestamp, 'job_id': str}

# ============================================
# LOCAL AGENT API
# ============================================

@app.route('/api/agent/update-job', methods=['POST'])
def agent_update_job():
    """Allow local agent to update a pipeline job's progress."""
    data = request.get_json()
    job_id = data.get('jobId', '')
    status = data.get('status', '')
    current_step = data.get('currentStep', '')
    progress = data.get('progress')
    error = data.get('error')

    if not job_id:
        return jsonify({'error': 'No jobId provided'}), 400

    kwargs = {}
    if status:
        kwargs['status'] = status
    if current_step:
        kwargs['current_step'] = current_step
    if progress is not None:
        kwargs['progress'] = progress
    if error:
        kwargs['error'] = error
    result = data.get('result')
    if result:
        kwargs['result'] = json.dumps(result) if isinstance(result, dict) else result

    if kwargs:
        _update_pipeline_job(job_id, **kwargs)

    return jsonify({'success': True})


@app.route('/api/agent/save-note', methods=['POST'])
def agent_save_note():
    """Allow local agent to save a completed research note."""
    data = request.get_json()
    note_id = data.get('noteId', str(uuid.uuid4()))
    ticker = data.get('ticker', '').upper()
    version = data.get('version', '1.0')
    note_md = data.get('noteMarkdown', '')
    sources_md = data.get('sourcesMarkdown', '')
    changelog_md = data.get('changelogMarkdown', '')
    docx_b64 = data.get('noteDocx', '')
    charts = data.get('charts', [])
    charts_data = data.get('chartsData', [])  # includes base64 image data
    metadata = data.get('metadata', {})

    # Prefer chartsData (with image data) over charts (filenames only)
    charts_to_store = charts_data if charts_data else charts

    if not ticker:
        return jsonify({'error': 'No ticker provided'}), 400

    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO research_notes (id, ticker, version, note_markdown, sources_markdown, changelog_markdown, note_docx, charts, metadata)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (id) DO UPDATE SET
                    note_markdown = EXCLUDED.note_markdown,
                    sources_markdown = EXCLUDED.sources_markdown,
                    changelog_markdown = EXCLUDED.changelog_markdown,
                    note_docx = EXCLUDED.note_docx,
                    charts = EXCLUDED.charts,
                    metadata = EXCLUDED.metadata,
                    updated_at = NOW()
            ''', (note_id, ticker, version, note_md, sources_md, changelog_md, docx_b64,
                  json.dumps(charts_to_store), json.dumps(metadata)))

        return jsonify({'success': True, 'noteId': note_id})
    except Exception as e:
        print(f'Error saving agent note: {e}')
        return jsonify({'error': str(e)}), 500


_pending_local_syncs = []

@app.route('/api/agent/sync-to-local', methods=['GET'])
def agent_sync_to_local():
    """Return pending sync jobs for the local agent to write to iCloud."""
    syncs = list(_pending_local_syncs)
    _pending_local_syncs.clear()
    return jsonify({'syncs': syncs})


@app.route('/api/notes/<ticker>/sync', methods=['POST'])
def sync_note_to_local(ticker):
    """Queue a note for syncing back to iCloud via the local agent."""
    ticker = ticker.upper()
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_notes WHERE ticker = %s ORDER BY created_at DESC LIMIT 1', (ticker,))
        note = cur.fetchone()
    if not note:
        return jsonify({'error': 'No note found'}), 404

    _pending_local_syncs.append({
        'ticker': ticker,
        'noteMarkdown': note['note_markdown'],
        'sourcesMarkdown': note['sources_markdown'],
        'changelogMarkdown': note['changelog_markdown'],
        'version': note.get('version', '1.0'),
        'timestamp': datetime.utcnow().isoformat(),
    })
    return jsonify({'success': True, 'ticker': ticker})


@app.route('/api/agent/file-manifest', methods=['POST'])
def agent_file_manifest():
    """Receive file manifest from local agent."""
    data = request.get_json()
    manifest = data.get('manifest', {})
    global _local_file_manifest
    _local_file_manifest = {
        'manifest': manifest,
        'timestamp': data.get('timestamp', datetime.utcnow().isoformat()),
    }
    return jsonify({'success': True, 'tickers': len(manifest)})

@app.route('/api/agent/local-files/<ticker>', methods=['GET'])
def agent_local_files(ticker):
    """Get local iCloud files for a ticker from the cached manifest."""
    ticker = ticker.upper()
    global _local_file_manifest
    files = _local_file_manifest.get('manifest', {}).get(ticker, [])
    return jsonify({
        'ticker': ticker,
        'files': files,
        'lastUpdated': _local_file_manifest.get('timestamp'),
    })

@app.route('/api/agent/doc-requests', methods=['GET'])
def agent_doc_requests():
    """Return tickers that need urgent document upload from local agent."""
    global _pending_doc_upload_requests
    # Clean up stale requests (older than 2 minutes)
    now = datetime.utcnow()
    stale = [t for t, info in _pending_doc_upload_requests.items()
             if (now - info['requested_at']).total_seconds() > 120]
    for t in stale:
        _pending_doc_upload_requests.pop(t, None)

    requests_list = [
        {'ticker': t, 'jobId': info['job_id']}
        for t, info in _pending_doc_upload_requests.items()
    ]
    return jsonify({'requests': requests_list})


@app.route('/api/agent/doc-upload-complete', methods=['POST'])
def agent_doc_upload_complete():
    """Local agent confirms document upload is done for a ticker."""
    global _pending_doc_upload_requests
    data = request.get_json() or {}
    ticker = data.get('ticker', '').upper()
    count = data.get('count', 0)
    if ticker in _pending_doc_upload_requests:
        _pending_doc_upload_requests.pop(ticker, None)
        print(f'[agent] Document upload complete for {ticker}: {count} files')
    return jsonify({'success': True})


@app.route('/api/agent/health', methods=['GET'])
def agent_health():
    """Health check endpoint for the agent."""
    return jsonify({'status': 'ok', 'timestamp': datetime.utcnow().isoformat()})


# ============================================
# TRADING AGENTS (Multi-Agent Analysis)
# ============================================

TRADING_AGENT_PROVIDERS = {
    "anthropic": ["claude-haiku-4-5-20251001", "claude-sonnet-4-6", "claude-opus-4-6"],
    "openai": ["gpt-4.1-mini", "gpt-4.1", "o4-mini"],
    "google": ["gemini-2.0-flash", "gemini-2.5-pro"],
}

@app.route('/api/agents/providers', methods=['GET'])
def get_agent_providers():
    return jsonify(TRADING_AGENT_PROVIDERS)


@app.route('/api/agents/run', methods=['POST'])
def start_agent_run():
    data = request.get_json()
    ticker = data.get('ticker', '').upper().strip()
    date_str = data.get('date', '')
    provider = data.get('provider', 'anthropic')
    model = data.get('model', 'claude-haiku-4-5-20251001')

    if not ticker or not date_str:
        return jsonify({'error': 'ticker and date required'}), 400

    run_id = str(uuid.uuid4())
    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO agent_runs (id, ticker, analysis_date, llm_provider, llm_model, status)
            VALUES (%s, %s, %s, %s, %s, 'running')
        ''', (run_id, ticker, date_str, provider, model))

    threading.Thread(target=_run_trading_agent, args=(run_id, ticker, date_str, provider, model), daemon=True).start()

    return jsonify({'runId': run_id, 'ticker': ticker})


@app.route('/api/agents/status/<run_id>', methods=['GET'])
def get_agent_status(run_id):
    since = request.args.get('since', 0, type=int)
    with get_db() as (_, cur):
        cur.execute('SELECT status, logs, decision, report, created_at, completed_at FROM agent_runs WHERE id = %s', (run_id,))
        row = cur.fetchone()
    if not row:
        return jsonify({'error': 'not found'}), 404

    logs = row['logs'] if isinstance(row['logs'], list) else json.loads(row['logs'] or '[]')
    return jsonify({
        'status': row['status'],
        'newLogs': logs[since:],
        'totalLogs': len(logs),
        'decision': row['decision'],
        'report': row['report'],
        'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
        'completedAt': row['completed_at'].isoformat() if row['completed_at'] else None,
    })


@app.route('/api/agents/results', methods=['GET'])
def list_agent_results():
    limit = request.args.get('limit', 50, type=int)
    with get_db() as (_, cur):
        cur.execute('''
            SELECT id, ticker, analysis_date, llm_provider, llm_model, status, decision, created_at, completed_at
            FROM agent_runs ORDER BY created_at DESC LIMIT %s
        ''', (limit,))
        rows = cur.fetchall()
    return jsonify({
        'runs': [{
            'id': r['id'], 'ticker': r['ticker'],
            'analysisDate': r['analysis_date'].isoformat() if r['analysis_date'] else None,
            'provider': r['llm_provider'], 'model': r['llm_model'],
            'status': r['status'], 'decision': r['decision'],
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
            'completedAt': r['completed_at'].isoformat() if r['completed_at'] else None,
        } for r in rows]
    })


@app.route('/api/agents/save/<run_id>', methods=['POST'])
def save_agent_to_research(run_id):
    """Save a completed agent run as a research note."""
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM agent_runs WHERE id = %s', (run_id,))
        run = cur.fetchone()
    if not run:
        return jsonify({'error': 'not found'}), 404
    if run['status'] != 'complete':
        return jsonify({'error': 'run not complete'}), 400

    ticker = run['ticker']
    logs = run['logs'] if isinstance(run['logs'], list) else json.loads(run['logs'] or '[]')
    log_text = '\n'.join(f"[{l.get('ts','')[-8:]}] {l.get('agent','')}: {l.get('message','')}" for l in logs)

    note_md = f"""# TradingAgents Analysis: {ticker}

**Date Analyzed:** {run['analysis_date']}
**Model:** {run['llm_provider']} / {run['llm_model']}
**Run ID:** {run_id}

---

## Decision

{run['decision'] or 'No decision'}

---

## Agent Log
```
{log_text}
```
"""

    # Save to research_documents (Research tab)
    doc_id = f"agent-{run_id}"
    with get_db(commit=True) as (conn, cur):
        cur.execute('SELECT id FROM research_categories WHERE UPPER(name) = %s', (ticker,))
        cat_row = cur.fetchone()
        if cat_row:
            category_id = cat_row['id']
        else:
            category_id = f"cat-{ticker.lower()}"
            cur.execute('''
                INSERT INTO research_categories (id, name, type, created_at)
                VALUES (%s, %s, 'ticker', NOW())
                ON CONFLICT (id) DO NOTHING
            ''', (category_id, ticker))

        doc_name = f"{ticker} -- TradingAgents Analysis ({run['analysis_date']})"
        full_content = note_md + ("\n\n" + run['report'] if run.get('report') else "")
        cur.execute('''
            INSERT INTO research_documents (id, category_id, name, content, file_names, doc_type, published_date, created_at)
            VALUES (%s, %s, %s, %s, '[]', 'agent_analysis', %s, NOW())
            ON CONFLICT (id) DO UPDATE SET content = EXCLUDED.content, name = EXCLUDED.name
        ''', (doc_id, category_id, doc_name, full_content, str(run['analysis_date'])))

    return jsonify({'success': True, 'docId': doc_id})


@app.route('/api/agents/batch-run', methods=['POST'])
def start_agent_batch_run():
    """Run TradingAgents on multiple tickers."""
    data = request.get_json()
    tickers = data.get('tickers', [])
    date_str = data.get('date', '')
    provider = data.get('provider', 'anthropic')
    model = data.get('model', 'claude-haiku-4-5-20251001')

    if not tickers or not date_str:
        return jsonify({'error': 'tickers and date required'}), 400

    batch_id = str(uuid.uuid4())
    run_ids = []

    with get_db(commit=True) as (conn, cur):
        for ticker in tickers:
            run_id = str(uuid.uuid4())
            cur.execute('''
                INSERT INTO agent_runs (id, ticker, analysis_date, llm_provider, llm_model, status, report)
                VALUES (%s, %s, %s, %s, %s, 'queued', %s)
            ''', (run_id, ticker.upper().strip(), date_str, provider, model, json.dumps({'batchId': batch_id})))
            run_ids.append({'id': run_id, 'ticker': ticker.upper().strip()})

    # Process sequentially in background
    def _run_batch():
        for run_info in run_ids:
            try:
                _run_trading_agent(run_info['id'], run_info['ticker'], date_str, provider, model)
            except Exception as e:
                print(f"[agent-batch] Failed {run_info['ticker']}: {e}")

    threading.Thread(target=_run_batch, daemon=True).start()

    return jsonify({'batchId': batch_id, 'runs': run_ids, 'count': len(run_ids)})


@app.route('/api/agents/compare/<run_id>', methods=['GET'])
def compare_agent_vs_thesis(run_id):
    """Compare an agent run's decision against the existing thesis."""
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM agent_runs WHERE id = %s', (run_id,))
        run = cur.fetchone()
    if not run:
        return jsonify({'error': 'not found'}), 404

    ticker = run['ticker']

    # Get thesis
    with get_db() as (_, cur):
        cur.execute('SELECT analysis FROM portfolio_analyses WHERE ticker = %s', (ticker,))
        pa = cur.fetchone()

    thesis_data = {}
    if pa and pa['analysis']:
        analysis = pa['analysis'] if isinstance(pa['analysis'], dict) else json.loads(pa['analysis'])
        thesis = analysis.get('thesis', {})
        thesis_data = {
            'summary': thesis.get('summary', ''),
            'pillars': thesis.get('pillars', []),
            'conclusion': analysis.get('conclusion', ''),
        }

    # Extract agent signal
    decision_text = run['decision'] or ''
    signal = 'UNKNOWN'
    for s in ['BUY', 'SELL', 'HOLD']:
        if s in decision_text.upper().split('\n')[0]:
            signal = s
            break

    # Determine thesis direction
    conclusion = thesis_data.get('conclusion', '').lower()
    thesis_direction = 'UNKNOWN'
    if any(w in conclusion for w in ['own', 'buy', 'bullish', 'add', 'overweight']):
        thesis_direction = 'BUY'
    elif any(w in conclusion for w in ['sell', 'avoid', 'underweight', 'bearish']):
        thesis_direction = 'SELL'
    elif any(w in conclusion for w in ['hold', 'neutral', 'revisit', 'wait']):
        thesis_direction = 'HOLD'

    conflict = signal != thesis_direction and signal != 'UNKNOWN' and thesis_direction != 'UNKNOWN'

    return jsonify({
        'ticker': ticker,
        'agentSignal': signal,
        'agentDecision': decision_text,
        'thesisDirection': thesis_direction,
        'thesisSummary': thesis_data.get('summary', ''),
        'thesisPillars': thesis_data.get('pillars', []),
        'thesisConclusion': thesis_data.get('conclusion', ''),
        'conflict': conflict,
    })


@app.route('/api/agents/dashboard', methods=['GET'])
def agent_dashboard():
    """Return agent run analytics and tracking data."""
    with get_db() as (_, cur):
        # Summary stats
        cur.execute('''
            SELECT
                COUNT(*) as total_runs,
                COUNT(*) FILTER (WHERE status = 'complete') as completed,
                COUNT(*) FILTER (WHERE status = 'error') as failed,
                COUNT(*) FILTER (WHERE status = 'running' OR status = 'queued') as active,
                COUNT(DISTINCT ticker) as unique_tickers,
                SUM(estimated_cost) as total_cost,
                AVG(EXTRACT(EPOCH FROM (completed_at - created_at))) FILTER (WHERE status = 'complete') as avg_duration
            FROM agent_runs
        ''')
        stats = cur.fetchone()

        # Decision distribution
        cur.execute('''
            SELECT
                CASE
                    WHEN UPPER(SPLIT_PART(decision, E'\n', 1)) LIKE '%%BUY%%' THEN 'BUY'
                    WHEN UPPER(SPLIT_PART(decision, E'\n', 1)) LIKE '%%SELL%%' THEN 'SELL'
                    WHEN UPPER(SPLIT_PART(decision, E'\n', 1)) LIKE '%%HOLD%%' THEN 'HOLD'
                    ELSE 'OTHER'
                END as signal,
                COUNT(*) as count
            FROM agent_runs WHERE status = 'complete'
            GROUP BY signal
        ''')
        distribution = {r['signal']: r['count'] for r in cur.fetchall()}

        # Per-ticker history (latest decision per ticker)
        cur.execute('''
            SELECT DISTINCT ON (ticker) ticker, decision, analysis_date, llm_model, created_at
            FROM agent_runs WHERE status = 'complete'
            ORDER BY ticker, created_at DESC
        ''')
        latest_by_ticker = [{
            'ticker': r['ticker'],
            'decision': r['decision'][:100] if r['decision'] else '',
            'analysisDate': r['analysis_date'].isoformat() if r['analysis_date'] else None,
            'model': r['llm_model'],
            'runDate': r['created_at'].isoformat() if r['created_at'] else None,
        } for r in cur.fetchall()]

    return jsonify({
        'stats': {
            'totalRuns': stats['total_runs'],
            'completed': stats['completed'],
            'failed': stats['failed'],
            'active': stats['active'],
            'uniqueTickers': stats['unique_tickers'],
            'totalCost': float(stats['total_cost'] or 0),
            'avgDurationSeconds': round(stats['avg_duration'] or 0),
        },
        'distribution': distribution,
        'latestByTicker': latest_by_ticker,
    })


# ============================================
# DEEP RESEARCH (Multi-Agent Web Research)
# ============================================

def _tavily_search(query, max_results=5, search_depth='advanced'):
    """Search the web using Tavily API."""
    try:
        from tavily import TavilyClient
        api_key = os.environ.get('TAVILY_API_KEY', '')
        if not api_key:
            return []
        client = TavilyClient(api_key=api_key)
        response = client.search(query, max_results=max_results, search_depth=search_depth)
        return response.get('results', [])
    except Exception as e:
        print(f"Tavily search error: {e}")
        return []


RESEARCH_MODES = {
    'topic': {
        'agents': ['official', 'media', 'community', 'background'],
        'labels': ['Official Sources', 'Media Coverage', 'Community Reaction', 'Background & Competition'],
    },
    'company': {
        'agents': ['company_info', 'financial', 'reputation', 'industry'],
        'labels': ['Company Overview', 'Financial Analysis', 'Media & Reputation', 'Industry & Competition'],
    },
    'person': {
        'agents': ['career', 'works', 'media', 'influence'],
        'labels': ['Career & Background', 'Works & Publications', 'Media & Interviews', 'Reputation & Influence'],
    },
}


@app.route('/api/research/start', methods=['POST'])
def start_deep_research():
    data = request.get_json()
    query = data.get('query', '').strip()
    mode = data.get('mode', 'topic')

    if not query:
        return jsonify({'error': 'Query required'}), 400
    if mode not in RESEARCH_MODES:
        return jsonify({'error': f'Invalid mode: {mode}'}), 400

    run_id = str(uuid.uuid4())
    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO research_runs (id, query, mode, status, progress, current_step)
            VALUES (%s, %s, %s, 'running', 0, 'Starting research...')
        ''', (run_id, query, mode))

    threading.Thread(target=_run_deep_research, args=(run_id, query, mode), daemon=True).start()

    return jsonify({'runId': run_id, 'query': query, 'mode': mode})


@app.route('/api/research/status/<run_id>', methods=['GET'])
def get_research_status(run_id):
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_runs WHERE id = %s', (run_id,))
        row = cur.fetchone()
    if not row:
        return jsonify({'error': 'Not found'}), 404

    agents = row['agents'] if isinstance(row['agents'], dict) else json.loads(row['agents'] or '{}')
    return jsonify({
        'id': row['id'],
        'query': row['query'],
        'mode': row['mode'],
        'status': row['status'],
        'progress': row['progress'],
        'currentStep': row['current_step'],
        'agents': agents,
        'synthesis': row['synthesis'],
        'createdAt': row['created_at'].isoformat() if row['created_at'] else None,
        'completedAt': row['completed_at'].isoformat() if row['completed_at'] else None,
    })


@app.route('/api/research/history', methods=['GET'])
def list_research_runs():
    limit = request.args.get('limit', 20, type=int)
    with get_db() as (_, cur):
        cur.execute('SELECT id, query, mode, status, progress, created_at, completed_at FROM research_runs ORDER BY created_at DESC LIMIT %s', (limit,))
        rows = cur.fetchall()
    return jsonify({
        'runs': [{
            'id': r['id'], 'query': r['query'], 'mode': r['mode'],
            'status': r['status'], 'progress': r['progress'],
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
            'completedAt': r['completed_at'].isoformat() if r['completed_at'] else None,
        } for r in rows]
    })


@app.route('/api/research/save/<run_id>', methods=['POST'])
def save_research_to_docs(run_id):
    """Save completed research to research_documents."""
    with get_db() as (_, cur):
        cur.execute('SELECT * FROM research_runs WHERE id = %s', (run_id,))
        run = cur.fetchone()
    if not run or run['status'] != 'complete':
        return jsonify({'error': 'Run not found or not complete'}), 404

    agents = run['agents'] if isinstance(run['agents'], dict) else json.loads(run['agents'] or '{}')

    # Build full content
    content_parts = [f"# Deep Research: {run['query']}\n\nMode: {run['mode']}\nDate: {run['created_at']}\n\n---\n"]
    if run['synthesis']:
        content_parts.append(f"## Synthesis\n\n{run['synthesis']}\n\n---\n")
    for agent_key, agent_data in agents.items():
        content_parts.append(f"## {agent_data.get('label', agent_key)}\n\n{agent_data.get('content', '')}\n\n---\n")

    full_content = '\n'.join(content_parts)

    # Save to research_documents
    doc_id = f"research-{run_id}"
    with get_db(commit=True) as (conn, cur):
        # Ensure category
        cat_id = 'cat-deep-research'
        cur.execute("INSERT INTO research_categories (id, name, type) VALUES (%s, 'Deep Research', 'topic') ON CONFLICT (id) DO NOTHING", (cat_id,))

        doc_name = f"Research: {run['query'][:80]}"
        cur.execute('''
            INSERT INTO research_documents (id, category_id, name, content, doc_type, created_at)
            VALUES (%s, %s, %s, %s, 'deep_research', NOW())
            ON CONFLICT (id) DO UPDATE SET content = EXCLUDED.content
        ''', (doc_id, cat_id, doc_name, full_content))

    cache.invalidate('research_categories')
    return jsonify({'success': True, 'docId': doc_id})


def _run_deep_research(run_id, query, mode):
    """Run multi-agent deep research with web search."""
    try:
        config = RESEARCH_MODES[mode]
        agent_names = config['agents']
        agent_labels = config['labels']
        agents_data = {}

        def update_progress(step, progress):
            with get_db(commit=True) as (conn, cur):
                cur.execute('UPDATE research_runs SET current_step = %s, progress = %s, agents = %s WHERE id = %s',
                           (step, progress, json.dumps(agents_data), run_id))

        # Phase 1: Parallel web search for each agent
        update_progress('Searching the web...', 10)

        search_queries = _build_search_queries(query, mode, agent_names)
        search_results = {}

        def search_for_agent(agent_key, queries):
            results = []
            for q in queries[:3]:  # Max 3 searches per agent
                results.extend(_tavily_search(q, max_results=5))
            search_results[agent_key] = results

        threads = []
        for i, agent_key in enumerate(agent_names):
            t = threading.Thread(target=search_for_agent, args=(agent_key, search_queries[agent_key]))
            t.start()
            threads.append(t)
        for t in threads:
            t.join(timeout=30)

        update_progress('Analyzing search results...', 30)

        # Phase 2: Each agent analyzes its search results
        def run_agent(agent_key, label, results):
            sources_text = '\n\n'.join([
                f"**{r.get('title', 'Untitled')}** ({r.get('url', '')})\n{r.get('content', '')[:1000]}"
                for r in results[:10]
            ])

            prompt = f"""You are a research analyst investigating "{query}".
Your focus area: {label}

Analyze these search results and write a comprehensive section for a research report.
Include specific facts, data points, and quotes. Cite sources with URLs.

SEARCH RESULTS:
{sources_text}

Write a detailed, well-structured analysis (500-1000 words) with source citations."""

            result = call_llm(
                messages=[{"role": "user", "content": prompt}],
                system=f"You are an expert research analyst. Write factual, well-sourced analysis. Always cite URLs.",
                tier="fast",
                max_tokens=4096,
            )

            agents_data[agent_key] = {
                'label': label,
                'content': result['text'],
                'sources': [{'title': r.get('title', ''), 'url': r.get('url', '')} for r in results[:10]],
                'status': 'complete',
            }

        agent_threads = []
        for i, (agent_key, label) in enumerate(zip(agent_names, agent_labels)):
            results = search_results.get(agent_key, [])
            t = threading.Thread(target=run_agent, args=(agent_key, label, results))
            t.start()
            agent_threads.append(t)

        for i, t in enumerate(agent_threads):
            t.join(timeout=120)
            update_progress(f'Agent {i+1}/{len(agent_names)} complete...', 30 + (i+1) * 15)

        # Phase 3: Synthesis
        update_progress('Synthesizing findings...', 85)

        all_findings = '\n\n---\n\n'.join([
            f"## {data.get('label', key)}\n\n{data.get('content', '')}"
            for key, data in agents_data.items()
        ])

        synthesis_result = call_llm(
            messages=[{"role": "user", "content": f"""You have 4 research agents' findings about "{query}".
Synthesize them into a cohesive executive summary (300-500 words).
Highlight key findings, consensus views, and areas of disagreement.
Include the most important data points and cite sources.

AGENT FINDINGS:
{all_findings[:12000]}"""}],
            system="You are a senior research director synthesizing multiple analysts' work into a clear executive summary.",
            tier="standard",
            max_tokens=4096,
        )

        synthesis = synthesis_result['text']

        # Save final results
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                UPDATE research_runs SET status = 'complete', progress = 100, current_step = 'Complete',
                agents = %s, synthesis = %s, completed_at = NOW() WHERE id = %s
            ''', (json.dumps(agents_data), synthesis, run_id))

        print(f"[deep-research {run_id}] Complete: {query}")

    except Exception as e:
        print(f"[deep-research {run_id}] Failed: {e}")
        import traceback
        traceback.print_exc()
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE research_runs SET status = 'error', current_step = %s WHERE id = %s",
                       (f"Failed: {str(e)[:200]}", run_id))


def _build_search_queries(query, mode, agent_names):
    """Build search queries for each agent based on mode."""
    year = datetime.utcnow().year
    queries = {}

    if mode == 'topic':
        queries = {
            'official': [f'{query} official documentation {year}', f'{query} announcement blog post {year}', f'{query} technical architecture'],
            'media': [f'{query} news {year}', f'{query} TechCrunch OR Verge OR Wired {year}', f'{query} analysis review'],
            'community': [f'{query} reddit discussion {year}', f'{query} hacker news {year}', f'{query} twitter opinions'],
            'background': [f'{query} competitors comparison {year}', f'{query} market analysis industry', f'{query} research paper academic'],
        }
    elif mode == 'company':
        queries = {
            'company_info': [f'{query} company overview {year}', f'{query} products services', f'{query} CEO leadership team'],
            'financial': [f'{query} revenue earnings {year}', f'{query} funding valuation investment', f'{query} stock analysis financial'],
            'reputation': [f'{query} reviews reputation {year}', f'{query} glassdoor employee reviews', f'{query} customer reviews'],
            'industry': [f'{query} competitors market share {year}', f'{query} industry analysis SWOT', f'{query} market outlook forecast'],
        }
    elif mode == 'person':
        queries = {
            'career': [f'{query} biography career background', f'{query} education university', f'{query} LinkedIn profile'],
            'works': [f'{query} achievements publications', f'{query} papers Google Scholar', f'{query} patents awards'],
            'media': [f'{query} interview podcast {year}', f'{query} keynote conference talk', f'{query} quotes statements'],
            'influence': [f'{query} influence reputation industry', f'{query} opinions about Reddit Twitter', f'{query} advisory board investments'],
        }

    return queries


# ============================================
# RESEARCH EXPORT TO SLIDES/INFOGRAPHICS
# ============================================

RESEARCH_EXPORT_FORMATS = {
    'professional': {
        'name': 'Professional Presentation',
        'type': 'deck',
        'theme': 'general',
        'style': 'Clean corporate design. White/navy/accent color palette. Data tables and charts where appropriate. Sans-serif headers. Structured grid layout. Professional, polished, board-ready.',
    },
    'sketchnote': {
        'name': 'Sketchnote',
        'type': 'deck',
        'theme': 'sketchnote',
        'style': 'Hand-drawn visual style on warm beige/cream paper. Sketch-like icons and doodles. Warm earthy colors. Playful handwritten typography. Visual metaphors connecting ideas. Fun yet informative.',
    },
    'executive_brief': {
        'name': 'Executive Brief',
        'type': 'deck',
        'theme': 'general',
        'style': 'Dense, information-rich layout. Serif typography. Muted professional colors (dark navy, charcoal, gold accents). Key metrics in large callout boxes. Minimal decoration. Board-ready density.',
    },
    'visual_storytelling': {
        'name': 'Visual Storytelling',
        'type': 'deck',
        'theme': 'general',
        'style': 'Cinematic full-bleed imagery. One key insight per slide. Large bold typography. High contrast. Dramatic visual hierarchy. Minimal text, maximum impact. TED-talk style.',
    },
    'research_summary': {
        'name': 'Research Summary Infographic',
        'type': 'infographic',
        'theme': 'general',
        'style': 'Vertical infographic flow. Sections for key findings, statistics, competitive data, and sources. Professional color palette. Icons for each section. Single comprehensive page.',
    },
    'competitive_landscape': {
        'name': 'Competitive Landscape Map',
        'type': 'infographic',
        'theme': 'general',
        'style': 'Comparison matrix or quadrant layout. Side-by-side competitor columns. Company logos/icons. Strengths in green, weaknesses in red. Clean grid structure. Data-forward.',
    },
    'timeline': {
        'name': 'Timeline / Evolution',
        'type': 'infographic',
        'theme': 'general',
        'style': 'Horizontal or vertical timeline with date markers. Milestone icons. Progressive color gradient from left to right. Key events with brief descriptions. Chronological narrative.',
    },
}

SLIDE_COUNT_GUIDANCE = {
    'compact': {'min': 1, 'max': 3, 'instruction': 'Distill to only the most critical findings. Maximum density per slide. No filler.'},
    'standard': {'min': 4, 'max': 6, 'instruction': 'Cover key themes with supporting evidence. One major theme per slide.'},
    'detailed': {'min': 7, 'max': 10, 'instruction': 'Comprehensive coverage with data, quotes, and analysis. Include context and nuance.'},
    'comprehensive': {'min': 11, 'max': 20, 'instruction': 'Full deep-dive with appendix-level detail. Separate slides for each sub-topic, data deep-dives, and source analysis.'},
}


@app.route('/api/research/<run_id>/export', methods=['POST'])
def start_research_export(run_id):
    """Queue research export jobs (slides/infographics)."""
    try:
        # Validate research run exists and is complete
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM research_runs WHERE id = %s', (run_id,))
            run = cur.fetchone()
        if not run:
            return jsonify({'error': 'Research run not found'}), 404
        if run['status'] != 'complete':
            return jsonify({'error': 'Research run not complete yet'}), 400

        data = request.get_json()
        exports = data.get('exports', [])
        if not exports:
            return jsonify({'error': 'No exports specified'}), 400

        export_ids = []
        with get_db(commit=True) as (conn, cur):
            for exp in exports:
                fmt = exp.get('format', '')
                if fmt not in RESEARCH_EXPORT_FORMATS:
                    continue
                fmt_config = RESEARCH_EXPORT_FORMATS[fmt]
                exp_type = fmt_config['type']
                slide_count = exp.get('slideCount', 5) if exp_type == 'deck' else 1

                export_id = str(uuid.uuid4())
                cur.execute('''
                    INSERT INTO research_exports (id, research_run_id, format, type, slide_count, status, created_at)
                    VALUES (%s, %s, %s, %s, %s, 'queued', NOW())
                ''', (export_id, run_id, fmt, exp_type, slide_count))
                export_ids.append({'id': export_id, 'format': fmt, 'type': exp_type, 'status': 'queued'})

                # Spawn background thread for each export
                threading.Thread(
                    target=_run_research_export,
                    args=(export_id, run_id, fmt, exp_type, slide_count),
                    daemon=True
                ).start()

        return jsonify({'exports': export_ids})
    except Exception as e:
        print(f"Error starting research export: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/catalysts/<job_id>/export', methods=['POST'])
def start_catalyst_export(job_id):
    """Queue catalyst synthesis export jobs (slides/infographics) — reuses research export infrastructure."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM research_pipeline_jobs WHERE id = %s', (job_id,))
            job = cur.fetchone()
        if not job:
            return jsonify({'error': 'Catalyst job not found'}), 404
        if job['status'] != 'complete':
            return jsonify({'error': 'Catalyst job not complete yet'}), 400

        result = job['result'] if isinstance(job.get('result'), dict) else json.loads(job['result'] or '{}')
        markdown = result.get('markdown', '')
        detail = job['steps_detail'] if isinstance(job.get('steps_detail'), (dict, list)) else json.loads(job['steps_detail'] or '{}')
        topic = detail.get('topic', 'Catalyst') if isinstance(detail, dict) else 'Catalyst'
        query = f"{job['ticker']} -- {topic}"

        data = request.get_json()
        exports = data.get('exports', [])
        if not exports:
            return jsonify({'error': 'No exports specified'}), 400

        export_ids = []
        with get_db(commit=True) as (conn, cur):
            for exp in exports:
                fmt = exp.get('format', '')
                if fmt not in RESEARCH_EXPORT_FORMATS:
                    continue
                fmt_config = RESEARCH_EXPORT_FORMATS[fmt]
                exp_type = fmt_config['type']
                slide_count = exp.get('slideCount', 5) if exp_type == 'deck' else 1

                export_id = str(uuid.uuid4())
                cur.execute('''
                    INSERT INTO research_exports (id, research_run_id, format, type, slide_count, status, created_at)
                    VALUES (%s, %s, %s, %s, %s, 'queued', NOW())
                ''', (export_id, job_id, fmt, exp_type, slide_count))
                export_ids.append({'id': export_id, 'format': fmt, 'type': exp_type, 'status': 'queued'})

                threading.Thread(
                    target=_run_content_export,
                    args=(export_id, job_id, fmt, exp_type, slide_count, markdown, query),
                    daemon=True
                ).start()

        return jsonify({'exports': export_ids})
    except Exception as e:
        print(f"Error starting catalyst export: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/research/<run_id>/exports', methods=['GET'])
def get_research_exports(run_id):
    """Get all exports for a research run with status."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT re.*, sp.title as project_title, sp.status as project_status, sp.total_slides
                FROM research_exports re
                LEFT JOIN slide_projects sp ON re.slide_project_id = sp.id
                WHERE re.research_run_id = %s
                ORDER BY re.created_at
            ''', (run_id,))
            rows = cur.fetchall()

        exports = []
        for r in rows:
            fmt_config = RESEARCH_EXPORT_FORMATS.get(r['format'], {})
            # Get thumbnail from first slide if complete
            thumbnail = None
            if r['slide_project_id'] and r['status'] == 'complete':
                with get_db() as (_, cur2):
                    cur2.execute('SELECT LEFT(image_data, 200) as thumb_check FROM slide_items WHERE project_id = %s AND slide_number = 1 AND image_data IS NOT NULL', (r['slide_project_id'],))
                    thumb_row = cur2.fetchone()
                    if thumb_row and thumb_row['thumb_check']:
                        thumbnail = True  # Just indicate it exists; frontend fetches via slide image endpoint

            exports.append({
                'id': r['id'],
                'format': r['format'],
                'formatName': fmt_config.get('name', r['format']),
                'type': r['type'],
                'slideCount': r['slide_count'],
                'status': r['status'],
                'progress': r['progress'],
                'error': r['error_message'],
                'slideProjectId': r['slide_project_id'],
                'projectTitle': r.get('project_title'),
                'projectStatus': r.get('project_status'),
                'totalSlides': r.get('total_slides'),
                'hasThumbnail': thumbnail is not None,
                'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
                'completedAt': r['completed_at'].isoformat() if r['completed_at'] else None,
            })
        return jsonify({'exports': exports})
    except Exception as e:
        print(f"Error getting research exports: {e}")
        return jsonify({'exports': []})


def _run_research_export(export_id, run_id, fmt, exp_type, slide_count):
    """Background worker: generate slides/infographic from research."""
    try:
        # Read research data
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM research_runs WHERE id = %s', (run_id,))
            run = cur.fetchone()
        if not run:
            raise ValueError("Research run not found")

        agents = run['agents'] if isinstance(run['agents'], dict) else json.loads(run['agents'] or '{}')
        synthesis = run['synthesis'] or ''
        query = run['query']

        # Build research content for outline generation
        research_content = f"# Research: {query}\n\n## Executive Summary\n{synthesis}\n\n"
        for key, data in agents.items():
            research_content += f"## {data.get('label', key)}\n{data.get('content', '')}\n\n"

        _run_content_export(export_id, fmt, exp_type, slide_count, research_content, query)
    except Exception as e:
        print(f"[research-export {export_id}] Failed: {e}")
        try:
            with get_db(commit=True) as (conn, cur):
                cur.execute("UPDATE research_exports SET status = 'error', error_message = %s WHERE id = %s",
                           (str(e)[:500], export_id))
        except:
            pass


def _run_content_export(export_id, fmt, exp_type, slide_count, research_content, query):
    """Background worker: generate slides/infographic from content string. Used by both research and catalyst exports."""
    try:
        fmt_config = RESEARCH_EXPORT_FORMATS[fmt]

        def update_export(status=None, progress=None, error=None, project_id=None):
            updates = []
            params = []
            if status:
                updates.append("status = %s"); params.append(status)
            if progress is not None:
                updates.append("progress = %s"); params.append(progress)
            if error:
                updates.append("error_message = %s"); params.append(error)
            if project_id:
                updates.append("slide_project_id = %s"); params.append(project_id)
            if status == 'complete':
                updates.append("completed_at = NOW()")
            if updates:
                params.append(export_id)
                with get_db(commit=True) as (conn, cur):
                    cur.execute(f"UPDATE research_exports SET {', '.join(updates)} WHERE id = %s", params)

        update_export(status='generating', progress=5)

        # Step 1: Generate outline via LLM
        if exp_type == 'deck':
            count_instruction = f"Generate EXACTLY {slide_count} slides."
            if slide_count <= 3:
                count_instruction += " " + SLIDE_COUNT_GUIDANCE['compact']['instruction']
            elif slide_count <= 6:
                count_instruction += " " + SLIDE_COUNT_GUIDANCE['standard']['instruction']
            elif slide_count <= 10:
                count_instruction += " " + SLIDE_COUNT_GUIDANCE['detailed']['instruction']
            else:
                count_instruction += " " + SLIDE_COUNT_GUIDANCE['comprehensive']['instruction']

            outline_prompt = f"""Based on this research, create a slide deck outline.

FORMAT STYLE: {fmt_config['name']}
{count_instruction}

For each slide, provide:
- title: slide title
- content: the text content for the slide (key points, data, quotes)
- illustration_hints: list of 2-3 visual elements to illustrate

RESEARCH CONTENT:
{research_content[:15000]}

Return ONLY valid JSON in this format:
{{
  "slides": [
    {{"slide_number": 1, "title": "...", "type": "title", "content": "...", "illustration_hints": ["...", "..."]}},
    {{"slide_number": 2, "title": "...", "type": "content", "content": "...", "illustration_hints": ["...", "..."]}}
  ]
}}"""
        else:
            # Infographic - single slide
            infographic_type = {
                'research_summary': 'a comprehensive research summary infographic showing key findings, statistics, and insights',
                'competitive_landscape': 'a competitive landscape map showing key players, their strengths/weaknesses, and market positioning',
                'timeline': 'a timeline/evolution infographic showing key events, milestones, and developments chronologically',
            }.get(fmt, 'a research summary infographic')

            outline_prompt = f"""Based on this research, create {infographic_type}.

Generate exactly 1 slide that contains ALL the key information in a dense, single-page infographic format.

Provide:
- title: infographic title
- content: comprehensive content covering all key data points, findings, comparisons
- illustration_hints: list of 3-5 visual elements

RESEARCH CONTENT:
{research_content[:15000]}

Return ONLY valid JSON:
{{
  "slides": [
    {{"slide_number": 1, "title": "...", "type": "infographic", "content": "...", "illustration_hints": ["...", "..."]}}
  ]
}}"""

        update_export(progress=15)

        outline_result = call_llm(
            messages=[{"role": "user", "content": outline_prompt}],
            system="You are a presentation designer. Generate structured slide outlines from research content. Return ONLY valid JSON, no markdown fences.",
            tier="standard",
            max_tokens=8192,
        )

        # Parse outline JSON
        outline_text = outline_result['text'].strip()
        # Strip markdown code fences if present
        if outline_text.startswith('```'):
            outline_text = outline_text.split('\n', 1)[1] if '\n' in outline_text else outline_text[3:]
        if outline_text.endswith('```'):
            outline_text = outline_text[:-3]
        outline_text = outline_text.strip()
        if outline_text.startswith('json'):
            outline_text = outline_text[4:].strip()

        outline = json.loads(outline_text)
        slides_data = outline.get('slides', [])

        if not slides_data:
            raise ValueError("LLM returned empty outline")

        update_export(progress=25)

        # Step 2: Create slide project
        theme = fmt_config.get('theme', 'general')
        format_label = fmt_config['name']
        project_title = f"{query[:60]} -- {format_label} ({len(slides_data)} slides)"

        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO slide_projects (ticker, title, theme, status, total_slides)
                VALUES (NULL, %s, %s, 'generating', %s)
                RETURNING id
            ''', (project_title, theme, len(slides_data)))
            project_id = cur.fetchone()['id']

        update_export(progress=30, project_id=project_id)

        # Step 3: Populate slides
        with get_db(commit=True) as (conn, cur):
            for s in slides_data:
                content_hash = _compute_content_hash(s)
                cur.execute('''
                    INSERT INTO slide_items (project_id, slide_number, title, type, content, illustration_hints, no_header, content_hash, status)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, 'new')
                ''', (
                    project_id, s['slide_number'], s.get('title', ''), s.get('type', 'content'),
                    s.get('content', ''), json.dumps(s.get('illustration_hints', [])),
                    s.get('no_header', False), content_hash,
                ))

        # Step 4: Generate images via Gemini
        gemini_key = os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY', '')
        if not gemini_key:
            # No Gemini key - leave slides as text-only but mark complete
            with get_db(commit=True) as (conn, cur):
                cur.execute('UPDATE slide_projects SET status = %s WHERE id = %s', ('ready', project_id))
            update_export(status='complete', progress=100)
            print(f"[research-export {export_id}] Complete (text-only, no Gemini key): {fmt}")
            return

        total_slides = len(slides_data)
        for idx, slide in enumerate(slides_data):
            slide_data = {
                'slide_number': slide['slide_number'],
                'title': slide.get('title', ''),
                'type': slide.get('type', 'content'),
                'content': slide.get('content', ''),
                'illustration_hints': slide.get('illustration_hints', []),
                'no_header': slide.get('no_header', False),
            }

            # Build format-specific prompt
            style_override = fmt_config['style']
            base_prompt = _build_slide_prompt(slide_data, theme, project_title, total_slides)
            full_prompt = f"STYLE OVERRIDE: {style_override}\n\n{base_prompt}"

            image_b64 = _generate_slide_image(full_prompt, api_key=gemini_key)
            if image_b64:
                content_hash = _compute_content_hash(slide_data)
                with get_db(commit=True) as (conn, cur):
                    cur.execute('''
                        UPDATE slide_items SET image_data = %s, content_hash = %s, status = 'generated', updated_at = %s
                        WHERE project_id = %s AND slide_number = %s
                    ''', (image_b64, content_hash, datetime.utcnow(), project_id, slide['slide_number']))

            progress = 30 + int((idx + 1) / total_slides * 65)
            update_export(progress=progress)
            import time; time.sleep(2)  # Rate limit

        # Mark complete
        with get_db(commit=True) as (conn, cur):
            cur.execute('UPDATE slide_projects SET status = %s, updated_at = %s WHERE id = %s',
                        ('ready', datetime.utcnow(), project_id))
        update_export(status='complete', progress=100)
        print(f"[research-export {export_id}] Complete: {fmt} ({total_slides} slides) for '{query[:40]}'")

    except Exception as e:
        print(f"[research-export {export_id}] Failed: {e}")
        import traceback
        traceback.print_exc()
        try:
            with get_db(commit=True) as (conn, cur):
                cur.execute("UPDATE research_exports SET status = 'error', error_message = %s WHERE id = %s",
                           (str(e)[:500], export_id))
        except:
            pass


# ============================================
# MULTI-AI VALIDATION
# ============================================

VALIDATION_PROMPT = """You are a senior equity research quality reviewer. Review this investment analysis for ALL of the following:

## 1. FACTUAL ACCURACY
- Mathematical errors: Do EPS bridges add up? Are YoY growth rates calculated correctly? Do segment revenues sum to total?
- Unsupported claims: Is every factual statement (revenue figures, margins, guidance, dates) verifiable from the context?
- Contradictions: Does the narrative contradict any data tables or numbers cited elsewhere?
- Stale data: Are any metrics clearly outdated or referencing the wrong time period?

## 2. AI-GENERATED LANGUAGE DETECTION
- Filler phrases: "it's worth noting", "in the realm of", "a testament to", "poised to", "delve into", "navigate the landscape", "robust growth trajectory"
- Excessive hedging: "it could potentially", "this may or may not", "there is a possibility that"
- Generic conclusions that could apply to any company
- Repetitive structure or overly symmetrical bullet points
- Missing analyst voice — should sound like a confident analyst writing to their PM ("I think", "My view is")
- NOTE: Some hedging like "may", "could" is appropriate when genuine uncertainty exists — only flag if excessive or unnecessary

## 3. ATTRIBUTION COMPLIANCE
- FLAG: Broker/firm names (Goldman Sachs, Morgan Stanley, RBC, Citi, TD Cowen, etc.)
- FLAG: Analyst names or specific individuals at research firms
- FLAG: Sellside opinion as thesis support ("the Street is constructive", "consensus targets suggest upside", "analysts have Buy ratings")
- FLAG: Rating references ("Buy-rated", "Outperform consensus")
- OK: "consensus estimates" as data-source label, analyst's own valuation math, generic market language

For EACH issue found, provide:
- CATEGORY: "factual" or "ai_language" or "attribution"
- SEVERITY: "HIGH" or "MEDIUM" or "LOW"
- LOCATION: The specific sentence or phrase
- ISSUE: What's wrong
- FIX: How to correct it

Return ONLY valid JSON (no markdown fences):
{"score": 0-100, "issues": [{"category": "...", "severity": "...", "location": "...", "issue": "...", "fix": "..."}], "summary": "one line overall assessment"}"""

VALIDATION_MODELS = [
    {'provider': 'anthropic', 'model': 'claude-sonnet-4-6', 'name': 'Claude'},
    {'provider': 'openai', 'model': 'gpt-4.1-mini', 'name': 'GPT'},
    {'provider': 'gemini', 'model': 'gemini-2.0-flash', 'name': 'Gemini'},
]

# Map validation provider names to internal adapter keys
_VALIDATION_PROVIDER_MAP = {
    'anthropic': 'anthropic',
    'openai': 'openai',
    'google': 'gemini',
    'gemini': 'gemini',
}


def _extract_json_robust(text):
    """Extract JSON object from LLM response text."""
    text = text.strip()
    if '```' in text:
        m = re.search(r'```(?:json)?\s*([\s\S]*?)```', text)
        if m:
            text = m.group(1).strip()
    if not text.startswith('{'):
        pos = text.find('{')
        if pos >= 0:
            text = text[pos:]
    brace_count = 0
    end_pos = 0
    for ci, ch in enumerate(text):
        if ch == '{': brace_count += 1
        elif ch == '}': brace_count -= 1
        if brace_count == 0 and ci > 0:
            end_pos = ci + 1
            break
    if end_pos > 0:
        text = text[:end_pos]
    return json.loads(text)


@app.route('/api/validate', methods=['POST'])
def validate_content():
    """Run multi-AI consensus validation — same prompt to 3 models, then cross-reference."""
    data = request.get_json()
    content = data.get('content', '')
    content_type = data.get('type', 'thesis')
    ticker = data.get('ticker', '')

    if not content:
        return jsonify({'error': 'No content to validate'}), 400

    api_keys = _get_api_keys()
    review_prompt = VALIDATION_PROMPT + f"\n\nCONTENT TO REVIEW ({content_type} for {ticker}):\n\n{content[:15000]}"

    # Run ALL 3 models with the SAME prompt in parallel
    model_results = {}

    def run_model(config):
        provider_key = _VALIDATION_PROVIDER_MAP.get(config['provider'], config['provider'])
        api_key = api_keys.get(provider_key, '')
        try:
            if not api_key:
                raise ValueError(f"No API key for provider {provider_key}")
            result = _LLM_ADAPTERS[provider_key](
                messages=[{"role": "user", "content": review_prompt}],
                system="You are a senior equity research quality reviewer. Return ONLY valid JSON.",
                model=config['model'],
                max_tokens=4096,
                timeout=120,
                api_key=api_key,
            )
            parsed = _extract_json_robust(result['text'])
            parsed['provider'] = provider_key
            parsed['model'] = config['model']
            parsed['reviewerName'] = config['name']
            model_results[config['name']] = parsed
        except json.JSONDecodeError:
            model_results[config['name']] = {
                'score': 50, 'issues': [],
                'summary': f"Response was not valid JSON",
                'provider': provider_key, 'model': config['model'], 'reviewerName': config['name'],
            }
        except Exception as e:
            model_results[config['name']] = {
                'score': 0, 'issues': [],
                'summary': f"Review failed: {str(e)[:200]}",
                'provider': provider_key, 'model': config['model'], 'reviewerName': config['name'],
            }

    threads = []
    for config in VALIDATION_MODELS:
        t = threading.Thread(target=run_model, args=(config,))
        t.start()
        threads.append(t)
    for t in threads:
        t.join(timeout=120)

    # Consensus scoring — cross-reference issues across models
    # Build a list of all issues with which models flagged them
    all_issues_raw = []
    for model_name, result in model_results.items():
        for issue in result.get('issues', []):
            all_issues_raw.append({**issue, '_model': model_name})

    # Group similar issues by location similarity
    consensus_issues = []
    used = set()
    for i, issue_a in enumerate(all_issues_raw):
        if i in used:
            continue
        loc_a = (issue_a.get('location') or '').lower()[:50]
        matching_models = [issue_a['_model']]
        best_fix = issue_a.get('fix', '')
        for j, issue_b in enumerate(all_issues_raw):
            if j <= i or j in used:
                continue
            loc_b = (issue_b.get('location') or '').lower()[:50]
            # Match if locations share significant overlap or same category+similar text
            if (loc_a and loc_b and (loc_a in loc_b or loc_b in loc_a)) or \
               (issue_a.get('category') == issue_b.get('category') and
                loc_a and loc_b and len(set(loc_a.split()) & set(loc_b.split())) >= 3):
                matching_models.append(issue_b['_model'])
                if not best_fix and issue_b.get('fix'):
                    best_fix = issue_b['fix']
                used.add(j)
        used.add(i)

        agreement = len(matching_models)
        confidence = 'HIGH' if agreement == 3 else 'MEDIUM' if agreement == 2 else 'LOW'

        consensus_issues.append({
            'category': issue_a.get('category', 'other'),
            'severity': issue_a.get('severity', 'MEDIUM'),
            'confidence': confidence,
            'agreement': f"{agreement}/3",
            'models': matching_models,
            'location': issue_a.get('location', ''),
            'issue': issue_a.get('issue', ''),
            'fix': best_fix,
        })

    # Sort: high confidence first, then high severity
    conf_order = {'HIGH': 0, 'MEDIUM': 1, 'LOW': 2}
    sev_order = {'HIGH': 0, 'MEDIUM': 1, 'LOW': 2}
    consensus_issues.sort(key=lambda x: (conf_order.get(x['confidence'], 3), sev_order.get(x['severity'], 3)))

    # Overall scores
    scores = [r.get('score', 0) for r in model_results.values() if isinstance(r.get('score'), (int, float))]
    overall_score = round(sum(scores) / len(scores)) if scores else 0
    high_conf_issues = sum(1 for i in consensus_issues if i['confidence'] == 'HIGH')

    # Save validation run to history
    val_id = str(uuid.uuid4())
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO validation_runs (id, ticker, content_type, overall_score, total_issues, high_confidence_issues, consensus_issues, reviewers)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
            ''', (val_id, ticker, content_type, overall_score, len(consensus_issues), high_conf_issues,
                  json.dumps(consensus_issues), json.dumps({k: {kk: vv for kk, vv in v.items() if kk != 'issues'} for k, v in model_results.items()})))
    except Exception as e:
        print(f'Failed to save validation run: {e}')

    return jsonify({
        'overallScore': overall_score,
        'totalIssues': len(consensus_issues),
        'highConfidenceIssues': high_conf_issues,
        'highIssues': sum(1 for i in consensus_issues if i['severity'] == 'HIGH'),
        'consensusIssues': consensus_issues,
        'reviewers': model_results,
        'ticker': ticker,
        'contentType': content_type,
    })


@app.route('/api/validate/auto-fix', methods=['POST'])
def validate_auto_fix():
    """Auto-fix flagged issues in content."""
    data = request.get_json()
    content = data.get('content', '')
    issues = data.get('issues', [])

    if not content or not issues:
        return jsonify({'error': 'Content and issues required'}), 400

    issues_text = '\n'.join(f"- [{i.get('severity')}] {i.get('location')}: {i.get('issue')} -> FIX: {i.get('fix')}" for i in issues)

    result = call_llm(
        messages=[{"role": "user", "content": f"""Fix the following issues in this investment analysis. Apply ONLY the specific fixes listed below. Do not change anything else. Preserve all formatting, structure, and content that isn't flagged.

ISSUES TO FIX:
{issues_text}

CONTENT:
{content}

Return the FULL corrected content with fixes applied. No commentary, just the corrected text."""}],
        system="You are an expert editor. Apply only the requested fixes. Change nothing else.",
        tier="standard",
        max_tokens=16384,
    )

    return jsonify({'fixedContent': result['text'], 'issuesFixed': len(issues)})


@app.route('/api/validate/history/<ticker>', methods=['GET'])
def validation_history(ticker):
    ticker = ticker.upper()
    with get_db() as (_, cur):
        cur.execute('''
            SELECT id, ticker, content_type, overall_score, total_issues, high_confidence_issues, consensus_issues, reviewers, created_at
            FROM validation_runs WHERE ticker = %s ORDER BY created_at DESC LIMIT 20
        ''', (ticker,))
        rows = cur.fetchall()
    return jsonify({
        'runs': [{
            'id': r['id'], 'ticker': r['ticker'], 'contentType': r['content_type'],
            'overallScore': r['overall_score'], 'totalIssues': r['total_issues'],
            'highConfidenceIssues': r['high_confidence_issues'],
            'consensusIssues': r['consensus_issues'] if isinstance(r['consensus_issues'], list) else json.loads(r['consensus_issues'] or '[]'),
            'reviewers': r['reviewers'] if isinstance(r['reviewers'], dict) else json.loads(r['reviewers'] or '{}'),
            'createdAt': r['created_at'].isoformat() if r['created_at'] else None,
        } for r in rows]
    })


@app.route('/api/chat/context', methods=['POST'])
def context_chat():
    """Send a message in a contextual chat about thesis/note content."""
    data = request.get_json()
    ticker = data.get('ticker', '').upper()
    content_type = data.get('contentType', 'thesis')
    content = data.get('content', '')
    message = data.get('message', '')
    chat_id = data.get('chatId', '')

    if not ticker or not message:
        return jsonify({'error': 'ticker and message required'}), 400

    # Load or create chat
    history = []
    if chat_id:
        with get_db() as (_, cur):
            cur.execute('SELECT messages FROM content_chats WHERE id = %s', (chat_id,))
            row = cur.fetchone()
            if row:
                history = row['messages'] if isinstance(row['messages'], list) else json.loads(row['messages'] or '[]')
    else:
        chat_id = str(uuid.uuid4())

    # Build system prompt with content context
    system = f"""You are a senior equity research analyst reviewing your own work. You generated the following {content_type} for {ticker}.

The user (your PM) is asking questions about this {content_type}, challenging assumptions, or requesting edits.

Answer precisely and honestly. If the user points out an error, acknowledge it. If they ask why you made a choice, explain your reasoning from the source documents. If they request an edit, provide the corrected text clearly marked.

When providing edits, format them as:
EDIT: [section name]
BEFORE: [original text]
AFTER: [corrected text]

CONTENT:
{content[:12000]}"""

    # Build messages for LLM
    llm_messages = []
    for h in history[-10:]:  # Keep last 10 messages for context
        llm_messages.append({"role": h['role'], "content": h['content']})
    llm_messages.append({"role": "user", "content": message})

    result = call_llm(
        messages=llm_messages,
        system=system,
        tier="standard",
        max_tokens=4096,
    )

    response = result['text']

    # Save to history
    history.append({"role": "user", "content": message, "ts": datetime.utcnow().isoformat()})
    history.append({"role": "assistant", "content": response, "ts": datetime.utcnow().isoformat()})

    with get_db(commit=True) as (conn, cur):
        cur.execute('''
            INSERT INTO content_chats (id, ticker, content_type, messages, updated_at)
            VALUES (%s, %s, %s, %s, NOW())
            ON CONFLICT (id) DO UPDATE SET messages = %s, updated_at = NOW()
        ''', (chat_id, ticker, content_type, json.dumps(history), json.dumps(history)))

    return jsonify({
        'chatId': chat_id,
        'response': response,
        'provider': result.get('provider', ''),
        'model': result.get('model', ''),
    })


@app.route('/api/chat/history/<ticker>', methods=['GET'])
def chat_history(ticker):
    """Get chat history for a ticker."""
    ticker = ticker.upper()
    content_type = request.args.get('type', '')
    with get_db() as (_, cur):
        if content_type:
            cur.execute('SELECT * FROM content_chats WHERE ticker = %s AND content_type = %s ORDER BY updated_at DESC LIMIT 10', (ticker, content_type))
        else:
            cur.execute('SELECT * FROM content_chats WHERE ticker = %s ORDER BY updated_at DESC LIMIT 10', (ticker,))
        rows = cur.fetchall()
    return jsonify({
        'chats': [{
            'id': r['id'], 'ticker': r['ticker'], 'contentType': r['content_type'],
            'messageCount': len(r['messages'] if isinstance(r['messages'], list) else json.loads(r['messages'] or '[]')),
            'updatedAt': r['updated_at'].isoformat() if r['updated_at'] else None,
        } for r in rows]
    })


# ============================================
# MEETING PREP - MEETING ENDPOINTS
# ============================================

@app.route('/api/mp/meetings', methods=['POST'])
def mp_create_meeting():
    """Create a new meeting prep session."""
    try:
        data = request.json
        ticker = data.get('ticker', '').upper().strip()
        company_name = data.get('companyName', '')
        sector = data.get('sector', '')
        meeting_date = data.get('meetingDate') or None
        meeting_type = data.get('meetingType', 'other')
        notes = data.get('notes', '')

        if not ticker:
            return jsonify({'error': 'Ticker is required'}), 400

        with get_db(commit=True) as (_, cur):
            # Upsert company
            cur.execute('''
                INSERT INTO mp_companies (ticker, name, sector)
                VALUES (%s, %s, %s)
                ON CONFLICT (ticker) DO UPDATE SET
                    name = COALESCE(NULLIF(EXCLUDED.name, ''), mp_companies.name),
                    sector = COALESCE(NULLIF(EXCLUDED.sector, ''), mp_companies.sector)
                RETURNING id, ticker, name, sector
            ''', (ticker, company_name, sector))
            company = dict(cur.fetchone())

            # Create meeting
            cur.execute('''
                INSERT INTO mp_meetings (company_id, meeting_date, meeting_type, notes)
                VALUES (%s, %s, %s, %s)
                RETURNING id, company_id, meeting_date, meeting_type, status, notes, created_at, updated_at
            ''', (company['id'], meeting_date, meeting_type, notes))
            meeting = dict(cur.fetchone())
            meeting['ticker'] = company['ticker']
            meeting['company_name'] = company['name']
            meeting['sector'] = company['sector']

        return jsonify(meeting)
    except Exception as e:
        print(f"Error creating meeting: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/meetings', methods=['GET'])
def mp_list_meetings():
    """List all meeting prep sessions."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT m.*, c.ticker, c.name as company_name, c.sector,
                       (SELECT COUNT(*) FROM mp_documents WHERE meeting_id = m.id) as doc_count,
                       (SELECT COUNT(*) FROM mp_question_sets WHERE meeting_id = m.id AND status = 'ready') as qs_count
                FROM mp_meetings m
                JOIN mp_companies c ON m.company_id = c.id
                ORDER BY m.created_at DESC
            ''')
            rows = cur.fetchall()

        result = []
        for r in rows:
            result.append({
                'id': r['id'],
                'company_id': r['company_id'],
                'ticker': r['ticker'],
                'company_name': r['company_name'],
                'sector': r['sector'],
                'meeting_date': str(r['meeting_date']) if r['meeting_date'] else None,
                'meeting_type': r['meeting_type'],
                'status': r['status'],
                'notes': r['notes'],
                'doc_count': r['doc_count'],
                'qs_count': r['qs_count'],
                'created_at': r['created_at'].isoformat() if r['created_at'] else None,
            })

        return jsonify(result)
    except Exception as e:
        print(f"Error listing meetings: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/meetings/<int:meeting_id>', methods=['GET'])
def mp_get_meeting(meeting_id):
    """Get a meeting with its documents and latest question set."""
    try:
        with get_db() as (_, cur):
            # Get meeting
            cur.execute('''
                SELECT m.*, c.ticker, c.name as company_name, c.sector
                FROM mp_meetings m
                JOIN mp_companies c ON m.company_id = c.id
                WHERE m.id = %s
            ''', (meeting_id,))
            row = cur.fetchone()
            if not row:
                return jsonify({'error': 'Meeting not found'}), 404

            meeting = dict(row)
            meeting['meeting_date'] = str(meeting['meeting_date']) if meeting['meeting_date'] else None
            meeting['created_at'] = meeting['created_at'].isoformat() if meeting['created_at'] else None
            meeting['updated_at'] = meeting['updated_at'].isoformat() if meeting['updated_at'] else None

            # Get documents (without file_data to keep response small)
            cur.execute('''
                SELECT id, meeting_id, filename, doc_type, doc_date, page_count, token_estimate,
                       upload_order, file_size, created_at
                FROM mp_documents WHERE meeting_id = %s ORDER BY upload_order
            ''', (meeting_id,))
            docs = []
            for d in cur.fetchall():
                dd = dict(d)
                dd['created_at'] = dd['created_at'].isoformat() if dd['created_at'] else None
                docs.append(dd)

            # Get latest question set
            cur.execute('''
                SELECT * FROM mp_question_sets
                WHERE meeting_id = %s ORDER BY version DESC LIMIT 1
            ''', (meeting_id,))
            qs_row = cur.fetchone()
            question_set = None
            if qs_row:
                question_set = dict(qs_row)
                if question_set['topics_json']:
                    question_set['topics'] = json.loads(question_set['topics_json'])
                else:
                    question_set['topics'] = []
                question_set['created_at'] = question_set['created_at'].isoformat() if question_set['created_at'] else None

        return jsonify({
            'meeting': meeting,
            'documents': docs,
            'questionSet': question_set,
        })
    except Exception as e:
        print(f"Error getting meeting: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/meetings/<int:meeting_id>', methods=['DELETE'])
def mp_delete_meeting(meeting_id):
    """Delete a meeting and all related data."""
    try:
        with get_db(commit=True) as (_, cur):
            cur.execute('DELETE FROM mp_meetings WHERE id = %s', (meeting_id,))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting meeting: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# MEETING PREP - DOCUMENT ENDPOINTS
# ============================================

@app.route('/api/mp/meetings/<int:meeting_id>/documents', methods=['POST'])
def mp_upload_documents(meeting_id):
    """Upload PDF documents for a meeting. Expects JSON with base64 file data."""
    try:
        from PyPDF2 import PdfReader
        import io

        data = request.json
        documents = data.get('documents', [])

        if not documents:
            return jsonify({'error': 'No documents provided'}), 400

        with get_db(commit=True) as (_, cur):
            # Verify meeting exists
            cur.execute('SELECT id FROM mp_meetings WHERE id = %s', (meeting_id,))
            if not cur.fetchone():
                return jsonify({'error': 'Meeting not found'}), 404

            # Get current max upload_order
            cur.execute('SELECT COALESCE(MAX(upload_order), 0) AS max_order FROM mp_documents WHERE meeting_id = %s', (meeting_id,))
            order = cur.fetchone()['max_order']

            results = []
            for doc in documents:
                order += 1
                filename = doc.get('filename', 'unknown.pdf')
                file_data = doc.get('fileData', '')
                extracted_text = doc.get('extractedText', '')
                page_count = doc.get('pageCount')

                # If no extracted text provided, try extracting from base64 PDF
                if not extracted_text and file_data:
                    try:
                        pdf_bytes = base64.b64decode(file_data)
                        reader = PdfReader(io.BytesIO(pdf_bytes))
                        pages = []
                        for page in reader.pages:
                            t = page.extract_text()
                            if t:
                                pages.append(t)
                        extracted_text = '\n\n'.join(pages)
                        if page_count is None:
                            page_count = len(reader.pages)
                    except Exception as ex:
                        print(f"PDF extraction error for {filename}: {ex}")

                # Classify and estimate tokens
                doc_type = classify_mp_document(filename, extracted_text)
                token_estimate = len(extracted_text) // 4 if extracted_text else 0
                file_size = len(file_data) * 3 // 4 if file_data else 0

                cur.execute('''
                    INSERT INTO mp_documents (meeting_id, filename, file_data, doc_type, doc_date,
                        page_count, token_estimate, extracted_text, upload_order, file_size)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id, filename, doc_type, doc_date, page_count, token_estimate, upload_order, file_size, created_at
                ''', (meeting_id, filename, file_data, doc_type, doc.get('docDate'),
                      page_count, token_estimate, extracted_text, order, file_size))
                row = dict(cur.fetchone())
                row['created_at'] = row['created_at'].isoformat() if row['created_at'] else None
                results.append(row)

        return jsonify(results)
    except Exception as e:
        print(f"Error uploading documents: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/meetings/<int:meeting_id>/documents/<int:doc_id>', methods=['DELETE'])
def mp_delete_document(meeting_id, doc_id):
    """Delete a document from a meeting."""
    try:
        with get_db(commit=True) as (_, cur):
            cur.execute('DELETE FROM mp_documents WHERE id = %s AND meeting_id = %s', (doc_id, meeting_id))

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting document: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# MEETING PREP - PIPELINE ENDPOINTS
# ============================================

@app.route('/api/mp/analyze-document', methods=['POST'])
def mp_analyze_document():
    """Step 1: Analyze a single document. Uses native PDF document blocks when available
    so Claude can see charts, graphs, tables, and images — not just extracted text."""
    try:
        data = request.json
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '')
        if not api_key:
            return jsonify({'error': 'No API key provided. Please add your API key in Settings.'}), 400

        ticker = data.get('ticker', '')
        company_name = data.get('companyName', ticker)
        doc_type = data.get('docType', 'document')
        filename = data.get('filename', '')
        doc_id = data.get('docId')
        extracted_text = data.get('extractedText', '')

        # Try to fetch native PDF binary from database for full-fidelity analysis
        file_data = None
        if doc_id:
            try:
                with get_db() as (_, cur):
                    cur.execute('SELECT file_data FROM mp_documents WHERE id = %s', (doc_id,))
                    row = cur.fetchone()
                    if row and row['file_data']:
                        file_data = row['file_data']
            except Exception as e:
                print(f"Could not fetch file_data for doc {doc_id}: {e}")

        prompt = MP_ANALYSIS_PROMPT.format(
            doc_type=doc_type, ticker=ticker, company_name=company_name
        )

        # Build content: prefer native document blocks for PDFs/images
        fname_lower = (filename or '').lower()
        if file_data and fname_lower.endswith('.pdf'):
            # Native PDF — Claude can see charts, graphs, tables, images
            user_content = [
                {"type": "text", "text": f"Document: {filename}"},
                {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": file_data}},
                {"type": "text", "text": "Analyze this document thoroughly, including all charts, graphs, tables, and visual data."}
            ]
        elif file_data and any(fname_lower.endswith(ext) for ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp')):
            ext = fname_lower.rsplit('.', 1)[-1]
            mime_map = {'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'gif': 'image/gif', 'webp': 'image/webp'}
            user_content = [
                {"type": "text", "text": f"Document: {filename}"},
                {"type": "image", "source": {"type": "base64", "media_type": mime_map.get(ext, 'image/png'), "data": file_data}},
                {"type": "text", "text": "Analyze this image thoroughly, including all charts, graphs, tables, and visual data."}
            ]
        else:
            # Fallback to extracted text
            if not extracted_text:
                return jsonify({'error': 'No document text or file data available'}), 400
            if len(extracted_text) > 400000:
                extracted_text = extracted_text[:400000] + "\n\n[... document truncated for length ...]"
            user_content = f"Document: {filename}\n\n{extracted_text}"

        def generate():
            try:
                llm_result = None
                for chunk in call_llm_stream(
                    messages=[{"role": "user", "content": user_content}],
                    system=prompt,
                    tier="standard",
                    max_tokens=16384,
                    anthropic_api_key=api_key,
                ):
                    if isinstance(chunk, dict):
                        llm_result = chunk
                    else:
                        yield chunk

                if llm_result is None:
                    raise Exception("No result from LLM")

                analysis = parse_mp_json(llm_result["text"])
                tokens_used = llm_result["usage"]["input_tokens"] + llm_result["usage"]["output_tokens"]
                yield "\n" + json.dumps({'analysis': analysis, 'tokensUsed': tokens_used, 'filename': filename})
            except LLMError as e:
                print(f"MP analyze LLM error: {e}")
                yield "\n" + json.dumps({'error': str(e)})
            except Exception as e:
                print(f"MP analyze stream error: {e}")
                yield "\n" + json.dumps({'error': str(e)})

        return app.response_class(generate(), mimetype='text/plain')

    except Exception as e:
        print(f"MP analyze error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/synthesize', methods=['POST'])
def mp_synthesize():
    """Step 2: Cross-reference all document analyses. Uses streaming to avoid timeout."""
    try:
        data = request.json
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '')
        if not api_key:
            return jsonify({'error': 'No API key provided.'}), 400

        ticker = data.get('ticker', '')
        company_name = data.get('companyName', ticker)
        sector = data.get('sector', 'unknown')
        analyses = data.get('analyses', [])
        past_questions = data.get('pastQuestions', [])
        timeframe = data.get('timeframe', 'recent')

        if not analyses:
            return jsonify({'error': 'No analyses provided'}), 400

        analyses_parts = []
        for i, a in enumerate(analyses):
            analyses_parts.append(f"### Document {i+1}: {a.get('_source_filename', 'unknown')}\n{json.dumps(a, indent=2)}")
        analyses_text = "\n\n".join(analyses_parts)

        past_q_text = ""
        if past_questions:
            pq_items = []
            for pq in past_questions[:30]:
                status_note = f" [STATUS: {pq.get('status', '')}]" if pq.get('status') != 'asked' else ""
                response = f" — Response: {pq.get('response_notes', '')}" if pq.get('response_notes') else ""
                pq_items.append(f"- [{pq.get('meeting_date', '?')}] {pq.get('question', '')}{status_note}{response}")
            past_q_text = "PAST QUESTIONS FROM PRIOR MEETINGS (reference these and flag unresolved items):\n" + "\n".join(pq_items)

        prompt = MP_SYNTHESIS_PROMPT.format(
            ticker=ticker, company_name=company_name, sector=sector,
            doc_count=len(analyses), timeframe=timeframe,
            analyses_text=analyses_text, past_questions_text=past_q_text,
        )

        def generate():
            try:
                llm_result = None
                for chunk in call_llm_stream(
                    messages=[{"role": "user", "content": "Synthesize the above document analyses."}],
                    system=prompt,
                    tier="standard",
                    max_tokens=16384,
                    anthropic_api_key=api_key,
                ):
                    if isinstance(chunk, dict):
                        llm_result = chunk
                    else:
                        yield chunk

                if llm_result is None:
                    raise Exception("No result from LLM")

                synthesis = parse_mp_json(llm_result["text"])
                tokens_used = llm_result["usage"]["input_tokens"] + llm_result["usage"]["output_tokens"]
                yield "\n" + json.dumps({'synthesis': synthesis, 'tokensUsed': tokens_used})
            except LLMError as e:
                print(f"MP synthesize LLM error: {e}")
                yield "\n" + json.dumps({'error': str(e)})
            except Exception as e:
                print(f"MP synthesize stream error: {e}")
                yield "\n" + json.dumps({'error': str(e)})

        return app.response_class(generate(), mimetype='text/plain')

    except Exception as e:
        print(f"MP synthesize error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/generate-questions', methods=['POST'])
def mp_generate_questions():
    """Step 3: Generate questions from synthesis. Uses streaming to avoid timeout."""
    try:
        data = request.json
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '')
        if not api_key:
            return jsonify({'error': 'No API key provided.'}), 400

        ticker = data.get('ticker', '')
        company_name = data.get('companyName', ticker)
        sector = data.get('sector', 'unknown')
        synthesis = data.get('synthesis', {})
        unresolved = data.get('unresolvedQuestions', [])

        if not synthesis:
            return jsonify({'error': 'No synthesis provided'}), 400

        synthesis_text = json.dumps(synthesis, indent=2)

        unresolved_text = ""
        if unresolved:
            items = [f"- {q.get('question', '')} (from {q.get('meeting_date', '?')})" for q in unresolved[:15]]
            unresolved_text = "UNRESOLVED QUESTIONS FROM PRIOR MEETINGS (include follow-ups for these):\n" + "\n".join(items)

        prompt = MP_QUESTION_PROMPT.format(
            ticker=ticker, company_name=company_name, sector=sector,
            synthesis_text=synthesis_text, unresolved_text=unresolved_text,
        )

        def generate():
            try:
                llm_result = None
                for chunk in call_llm_stream(
                    messages=[{"role": "user", "content": "Generate the meeting preparation questions."}],
                    system=prompt,
                    tier="standard",
                    max_tokens=16384,
                    anthropic_api_key=api_key,
                ):
                    if isinstance(chunk, dict):
                        llm_result = chunk
                    else:
                        yield chunk

                if llm_result is None:
                    raise Exception("No result from LLM")

                topics = parse_mp_json(llm_result["text"])
                tokens_used = llm_result["usage"]["input_tokens"] + llm_result["usage"]["output_tokens"]
                yield "\n" + json.dumps({'topics': topics, 'tokensUsed': tokens_used})
            except LLMError as e:
                print(f"MP generate questions LLM error: {e}")
                yield "\n" + json.dumps({'error': str(e)})
            except Exception as e:
                print(f"MP generate questions stream error: {e}")
                yield "\n" + json.dumps({'error': str(e)})

        return app.response_class(generate(), mimetype='text/plain')

    except Exception as e:
        print(f"MP generate questions error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/save-results', methods=['POST'])
def mp_save_results():
    """Save pipeline results: question set + past questions."""
    try:
        data = request.json
        meeting_id = data.get('meetingId')
        topics = data.get('topics', [])
        synthesis_json = data.get('synthesisJson')
        total_tokens = data.get('totalTokens', 0)
        model = data.get('model', 'claude-sonnet-4-5-20250929')

        if not meeting_id:
            return jsonify({'error': 'meetingId is required'}), 400

        with get_db(commit=True) as (_, cur):
            # Get next version
            cur.execute('SELECT COALESCE(MAX(version), 0) + 1 AS next_ver FROM mp_question_sets WHERE meeting_id = %s', (meeting_id,))
            version = cur.fetchone()['next_ver']

            # Insert question set
            cur.execute('''
                INSERT INTO mp_question_sets (meeting_id, version, status, topics_json, synthesis_json, generation_model, generation_tokens)
                VALUES (%s, %s, 'ready', %s, %s, %s, %s)
                RETURNING id, version
            ''', (meeting_id, version, json.dumps(topics), json.dumps(synthesis_json) if synthesis_json else None, model, total_tokens))
            qs = dict(cur.fetchone())

            # Update meeting status
            cur.execute("UPDATE mp_meetings SET status = 'ready', updated_at = CURRENT_TIMESTAMP WHERE id = %s", (meeting_id,))

            # Save questions to past_questions
            cur.execute('SELECT company_id FROM mp_meetings WHERE id = %s', (meeting_id,))
            company_row = cur.fetchone()
            if company_row:
                company_id = company_row['company_id']
                for topic in (topics if isinstance(topics, list) else []):
                    topic_name = topic.get('topic', '') if isinstance(topic, dict) else ''
                    questions = topic.get('questions', []) if isinstance(topic, dict) else []
                    for q in questions:
                        q_text = q.get('question', '') if isinstance(q, dict) else ''
                        if q_text:
                            cur.execute('''
                                INSERT INTO mp_past_questions (company_id, meeting_id, question, topic, status)
                                VALUES (%s, %s, %s, %s, 'asked')
                            ''', (company_id, meeting_id, q_text, topic_name))

        return jsonify({'questionSetId': qs['id'], 'version': qs['version']})
    except Exception as e:
        print(f"Error saving results: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# MEETING PREP - HISTORY ENDPOINTS
# ============================================

@app.route('/api/mp/companies/<ticker>/past-questions', methods=['GET'])
def mp_get_past_questions(ticker):
    """Get past questions for a company."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM mp_companies WHERE ticker = %s', (ticker.upper(),))
            company = cur.fetchone()
            if not company:
                return jsonify({'company': None, 'pastQuestions': []})

            cur.execute('''
                SELECT pq.*, m.meeting_date, m.meeting_type
                FROM mp_past_questions pq
                LEFT JOIN mp_meetings m ON pq.meeting_id = m.id
                WHERE pq.company_id = %s
                ORDER BY pq.created_at DESC
                LIMIT 100
            ''', (company['id'],))
            rows = cur.fetchall()

        pqs = []
        for r in rows:
            d = dict(r)
            d['created_at'] = d['created_at'].isoformat() if d['created_at'] else None
            d['meeting_date'] = str(d['meeting_date']) if d['meeting_date'] else None
            pqs.append(d)

        return jsonify({'company': dict(company), 'pastQuestions': pqs})
    except Exception as e:
        print(f"Error getting past questions: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/past-questions/<int:pq_id>/note', methods=['POST'])
def mp_update_past_question(pq_id):
    """Update notes/status on a past question."""
    try:
        data = request.json

        updates = []
        params = []
        if 'responseNotes' in data:
            updates.append('response_notes = %s')
            params.append(data['responseNotes'])
        if 'status' in data:
            updates.append('status = %s')
            params.append(data['status'])

        if updates:
            params.append(pq_id)
            with get_db(commit=True) as (_, cur):
                cur.execute(f"UPDATE mp_past_questions SET {', '.join(updates)} WHERE id = %s", params)

        return jsonify({'success': True})
    except Exception as e:
        print(f"Error updating past question: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/mp/meetings/<int:meeting_id>/documents/<int:doc_id>/text', methods=['GET'])
def mp_get_document_text(meeting_id, doc_id):
    """Get extracted text for a document (needed by frontend pipeline)."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT id, filename, doc_type, extracted_text,
                       (file_data IS NOT NULL AND file_data != '') as has_file_data
                FROM mp_documents WHERE id = %s AND meeting_id = %s
            ''', (doc_id, meeting_id))
            row = cur.fetchone()

        if not row:
            return jsonify({'error': 'Document not found'}), 404

        return jsonify({
            'id': row['id'],
            'filename': row['filename'],
            'docType': row['doc_type'],
            'extractedText': row['extracted_text'] or '',
            'hasFileData': bool(row['has_file_data']),
        })
    except Exception as e:
        print(f"Error getting document text: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# MEETING PREP — GOOGLE DRIVE INTEGRATION
# ============================================

@app.route('/api/mp/search-drive', methods=['POST'])
def mp_search_drive():
    """Search Google Drive 'Research Reports Char' folder for documents matching a ticker."""
    try:
        from datetime import datetime, timedelta
        import requests as http_requests

        data = request.json
        access_token = data.get('accessToken', '')
        ticker = data.get('ticker', '')
        time_range = data.get('timeRange', '3months')
        keyword = data.get('keyword', '').strip()

        if not access_token:
            return jsonify({'error': 'Google access token required'}), 400
        if not ticker:
            return jsonify({'error': 'Ticker required'}), 400

        # Sanitize inputs to prevent Drive API query injection
        import re
        def sanitize_drive_query(s):
            """Escape single quotes and strip non-printable chars for Drive API queries"""
            return re.sub(r"['\\\x00-\x1f]", '', s.strip())[:100]

        ticker = sanitize_drive_query(ticker)
        keyword = sanitize_drive_query(keyword)

        headers = {'Authorization': f'Bearer {access_token}'}
        drive_api = 'https://www.googleapis.com/drive/v3/files'

        # Calculate date cutoff
        ranges = {
            'day': 1, 'week': 7, 'month': 30, '3months': 90,
            '6months': 180, 'year': 365, '3years': 1095
        }
        days = ranges.get(time_range, 90)
        cutoff = (datetime.utcnow() - timedelta(days=days)).strftime('%Y-%m-%dT%H:%M:%S')

        # Find "Research Reports Char" folder
        folder_query = "name = 'Research Reports Char' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_resp = http_requests.get(drive_api, headers=headers, params={
            'q': folder_query, 'fields': 'files(id, name)', 'pageSize': 5
        }, timeout=15)
        if folder_resp.status_code == 401:
            return jsonify({'error': 'Google authentication expired. Please re-authenticate.'}), 401
        folder_resp.raise_for_status()
        folders = folder_resp.json().get('files', [])

        if not folders:
            return jsonify({'error': 'Folder "Research Reports Char" not found in your Google Drive'}), 404

        folder_id = folders[0]['id']

        # Find subfolders whose name matches the ticker (e.g. "RTX" folder inside root)
        search_term = keyword if keyword else ticker
        subfolder_resp = http_requests.get(drive_api, headers=headers, params={
            'q': f"'{folder_id}' in parents and mimeType = 'application/vnd.google-apps.folder' and trashed = false and name contains '{ticker}'",
            'fields': 'files(id, name)',
            'pageSize': 50
        }, timeout=15)
        subfolder_resp.raise_for_status()
        subfolders = subfolder_resp.json().get('files', [])

        # Search root folder + any matching subfolders
        folder_ids = [folder_id] + [sf['id'] for sf in subfolders]

        all_files = []
        for fid in folder_ids:
            file_query = (
                f"'{fid}' in parents and trashed = false "
                f"and modifiedTime > '{cutoff}' "
                f"and (name contains '{search_term}' or fullText contains '{search_term}')"
            )
            file_resp = http_requests.get(drive_api, headers=headers, params={
                'q': file_query,
                'fields': 'files(id, name, mimeType, modifiedTime, size)',
                'pageSize': 50,
                'orderBy': 'modifiedTime desc'
            }, timeout=15)
            file_resp.raise_for_status()
            all_files.extend(file_resp.json().get('files', []))

        # Deduplicate by file ID and sort by modifiedTime desc
        seen = set()
        files = []
        for f in sorted(all_files, key=lambda x: x.get('modifiedTime', ''), reverse=True):
            if f['id'] not in seen:
                seen.add(f['id'])
                files.append(f)

        return jsonify({'files': files, 'folderId': folder_id, 'folderName': folders[0]['name']})

    except Exception as e:
        error_msg = str(e)
        print(f"Drive search error: {error_msg}")
        if '401' in error_msg:
            return jsonify({'error': 'Google authentication expired. Please re-authenticate.'}), 401
        return jsonify({'error': error_msg}), 500


@app.route('/api/mp/preview-zip', methods=['POST'])
def mp_preview_zip():
    """Download a zip file from Google Drive and list the PDF files inside."""
    try:
        import requests as http_requests
        import io
        import zipfile

        data = request.json
        access_token = data.get('accessToken', '')
        file_id = data.get('fileId', '')

        if not access_token or not file_id:
            return jsonify({'error': 'accessToken and fileId are required'}), 400

        headers = {'Authorization': f'Bearer {access_token}'}
        drive_api = 'https://www.googleapis.com/drive/v3/files'

        dl_resp = http_requests.get(f'{drive_api}/{file_id}', headers=headers, params={'alt': 'media'}, timeout=60)
        if dl_resp.status_code == 401:
            return jsonify({'error': 'Google authentication expired. Please re-authenticate.'}), 401
        dl_resp.raise_for_status()

        try:
            with zipfile.ZipFile(io.BytesIO(dl_resp.content)) as zf:
                pdfs = []
                for info in zf.infolist():
                    if info.filename.lower().endswith('.pdf') and not info.filename.startswith('__MACOSX'):
                        name = info.filename.split('/')[-1] if '/' in info.filename else info.filename
                        pdfs.append({
                            'zipPath': info.filename,
                            'name': name,
                            'size': info.file_size
                        })
                return jsonify({'pdfs': pdfs})
        except zipfile.BadZipFile:
            return jsonify({'error': 'Invalid or corrupted zip file'}), 400

    except Exception as e:
        error_msg = str(e)
        print(f"Zip preview error: {error_msg}")
        if '401' in error_msg:
            return jsonify({'error': 'Google authentication expired. Please re-authenticate.'}), 401
        return jsonify({'error': error_msg}), 500


@app.route('/api/mp/import-drive-files', methods=['POST'])
def mp_import_drive_files():
    """Download files from Google Drive and import into a meeting. Handles zip files by extracting PDFs."""
    try:
        import requests as http_requests
        from PyPDF2 import PdfReader
        import io
        import zipfile

        data = request.json
        access_token = data.get('accessToken', '')
        meeting_id = data.get('meetingId')
        files_to_import = data.get('files', [])

        if not access_token or not meeting_id or not files_to_import:
            return jsonify({'error': 'accessToken, meetingId, and files are required'}), 400

        headers = {'Authorization': f'Bearer {access_token}'}
        drive_api = 'https://www.googleapis.com/drive/v3/files'

        with get_db(commit=True) as (_, cur):
            # Verify meeting exists
            cur.execute('SELECT id FROM mp_meetings WHERE id = %s', (meeting_id,))
            if not cur.fetchone():
                return jsonify({'error': 'Meeting not found'}), 404

            cur.execute('SELECT COALESCE(MAX(upload_order), 0) AS max_order FROM mp_documents WHERE meeting_id = %s', (meeting_id,))
            order = cur.fetchone()['max_order']

            def import_pdf(pdf_bytes, filename, doc_date, cur, meeting_id, order):
                """Extract text from PDF bytes and insert into mp_documents."""
                extracted_text = ''
                page_count = None
                try:
                    reader = PdfReader(io.BytesIO(pdf_bytes))
                    pages = []
                    for page in reader.pages:
                        t = page.extract_text()
                        if t:
                            pages.append(t)
                    extracted_text = '\n\n'.join(pages)
                    page_count = len(reader.pages)
                except Exception as ex:
                    print(f"PDF extraction error for {filename}: {ex}")

                doc_type = classify_mp_document(filename, extracted_text)
                token_estimate = len(extracted_text) // 4 if extracted_text else 0
                file_size = len(pdf_bytes)

                cur.execute('''
                    INSERT INTO mp_documents (meeting_id, filename, file_data, doc_type, doc_date,
                        page_count, token_estimate, extracted_text, upload_order, file_size)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING id, filename, doc_type, doc_date, page_count, token_estimate, upload_order, file_size, created_at
                ''', (meeting_id, filename, '', doc_type, doc_date,
                      page_count, token_estimate, extracted_text, order, file_size))
                row = dict(cur.fetchone())
                row['created_at'] = row['created_at'].isoformat() if row['created_at'] else None
                return row

            results = []
            for file_info in files_to_import:
                file_id = file_info.get('id')
                filename = file_info.get('name', 'unknown')
                mime_type = file_info.get('mimeType', '')
                doc_date = file_info.get('modifiedTime', '')[:10] if file_info.get('modifiedTime') else None

                try:
                    # Download file content via REST API
                    if mime_type in ('application/vnd.google-apps.document', 'application/vnd.google-apps.spreadsheet'):
                        dl_resp = http_requests.get(f'{drive_api}/{file_id}/export', headers=headers, params={'mimeType': 'application/pdf'}, timeout=60)
                    else:
                        dl_resp = http_requests.get(f'{drive_api}/{file_id}', headers=headers, params={'alt': 'media'}, timeout=60)
                    dl_resp.raise_for_status()
                    file_content = dl_resp.content

                    is_zip = (mime_type == 'application/zip' or
                              mime_type == 'application/x-zip-compressed' or
                              filename.lower().endswith('.zip'))

                    if is_zip:
                        # Extract selected (or all) PDFs from the zip file
                        selected_pdfs = set(file_info.get('selectedPdfs', []))
                        try:
                            with zipfile.ZipFile(io.BytesIO(file_content)) as zf:
                                pdf_names = [n for n in zf.namelist()
                                             if n.lower().endswith('.pdf') and not n.startswith('__MACOSX')]
                                if selected_pdfs:
                                    pdf_names = [n for n in pdf_names if n in selected_pdfs]
                                if not pdf_names:
                                    results.append({'filename': filename, 'error': 'No matching PDF files found inside zip'})
                                    continue
                                for pdf_name in pdf_names:
                                    order += 1
                                    pdf_bytes = zf.read(pdf_name)
                                    pdf_filename = pdf_name.split('/')[-1] if '/' in pdf_name else pdf_name
                                    row = import_pdf(pdf_bytes, pdf_filename, doc_date, cur, meeting_id, order)
                                    row['fromZip'] = filename
                                    results.append(row)
                        except zipfile.BadZipFile:
                            results.append({'filename': filename, 'error': 'Invalid or corrupted zip file'})
                    else:
                        # Regular PDF file
                        order += 1
                        row = import_pdf(file_content, filename, doc_date, cur, meeting_id, order)
                        results.append(row)

                except Exception as ex:
                    print(f"Error importing {filename}: {ex}")
                    results.append({'filename': filename, 'error': str(ex)})

        return jsonify(results)

    except Exception as e:
        error_msg = str(e)
        print(f"Drive import error: {error_msg}")
        if 'invalid_grant' in error_msg.lower() or '401' in error_msg:
            return jsonify({'error': 'Google authentication expired. Please re-authenticate.'}), 401
        return jsonify({'error': error_msg}), 500


# ============================================
# GENERIC GOOGLE DRIVE DOWNLOAD (for Research, Summary, Portfolio tabs)
# ============================================

@app.route('/api/drive/download-files', methods=['POST'])
def drive_download_files():
    """Download files from Google Drive and extract text. Returns extracted data without storing in DB."""
    try:
        import requests as http_requests
        from PyPDF2 import PdfReader
        import io
        import zipfile
        import base64

        data = request.json
        access_token = data.get('accessToken', '')
        files_to_import = data.get('files', [])

        if not access_token or not files_to_import:
            return jsonify({'error': 'accessToken and files are required'}), 400

        headers = {'Authorization': f'Bearer {access_token}'}
        drive_api = 'https://www.googleapis.com/drive/v3/files'

        def extract_pdf(pdf_bytes, filename):
            """Extract text from PDF bytes and return data dict."""
            extracted_text = ''
            page_count = None
            try:
                reader = PdfReader(io.BytesIO(pdf_bytes))
                pages = []
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        pages.append(t)
                extracted_text = '\n\n'.join(pages)
                page_count = len(reader.pages)
            except Exception as ex:
                print(f"PDF extraction error for {filename}: {ex}")

            return {
                'filename': filename,
                'extractedText': extracted_text,
                'pageCount': page_count,
                'fileSize': len(pdf_bytes),
                'fileData': base64.b64encode(pdf_bytes).decode('utf-8')
            }

        results = []
        for file_info in files_to_import:
            file_id = file_info.get('id')
            filename = file_info.get('name', 'unknown')
            mime_type = file_info.get('mimeType', '')

            try:
                if mime_type in ('application/vnd.google-apps.document', 'application/vnd.google-apps.spreadsheet'):
                    dl_resp = http_requests.get(f'{drive_api}/{file_id}/export', headers=headers, params={'mimeType': 'application/pdf'}, timeout=60)
                else:
                    dl_resp = http_requests.get(f'{drive_api}/{file_id}', headers=headers, params={'alt': 'media'}, timeout=60)
                dl_resp.raise_for_status()
                file_content = dl_resp.content

                is_zip = (mime_type == 'application/zip' or
                          mime_type == 'application/x-zip-compressed' or
                          filename.lower().endswith('.zip'))

                if is_zip:
                    selected_pdfs = set(file_info.get('selectedPdfs', []))
                    try:
                        with zipfile.ZipFile(io.BytesIO(file_content)) as zf:
                            pdf_names = [n for n in zf.namelist()
                                         if n.lower().endswith('.pdf') and not n.startswith('__MACOSX')]
                            if selected_pdfs:
                                pdf_names = [n for n in pdf_names if n in selected_pdfs]
                            if not pdf_names:
                                results.append({'filename': filename, 'error': 'No matching PDF files found inside zip'})
                                continue
                            for pdf_name in pdf_names:
                                pdf_bytes = zf.read(pdf_name)
                                pdf_filename = pdf_name.split('/')[-1] if '/' in pdf_name else pdf_name
                                row = extract_pdf(pdf_bytes, pdf_filename)
                                row['fromZip'] = filename
                                results.append(row)
                    except zipfile.BadZipFile:
                        results.append({'filename': filename, 'error': 'Invalid or corrupted zip file'})
                else:
                    row = extract_pdf(file_content, filename)
                    results.append(row)

            except Exception as ex:
                print(f"Error downloading {filename}: {ex}")
                results.append({'filename': filename, 'error': str(ex)})

        return jsonify({'files': results})

    except Exception as e:
        error_msg = str(e)
        print(f"Drive download error: {error_msg}")
        if 'invalid_grant' in error_msg.lower() or '401' in error_msg:
            return jsonify({'error': 'Google authentication expired. Please re-authenticate.'}), 401
        return jsonify({'error': error_msg}), 500


# ============================================
# SLIDE GENERATOR ENDPOINTS
# ============================================

# --- Theme constants ---
SLIDE_THEMES = {
    'sketchnote': {
        'name': 'Sketchnote',
        'style_prefix': 'Background: Warm beige/cream colored textured paper (like aged notebook paper)\nIllustrations: Cute hand-drawn cartoon style with colorful doodles\nTitle text: Large, colorful, hand-lettered style typography (rounded, playful, multicolor gradients)\nBody text: Clean, readable text in dark gray/brown\nDecorations: Small stars, sparkles, arrows, dots, swirls scattered around\nLayout: Professional yet friendly, like a designer\'s sketchnote\nAspect ratio: 16:9 widescreen slide\nResolution: High quality, crisp text\nALL text MUST be in English\nDO NOT include any watermarks or AI generation notices',
        'illustration_guidance': 'Use topic-appropriate icons and characters. DO NOT include robots unless topic is AI/robotics.',
        'negative_guidance': 'DO NOT include robots or bot characters unless the topic is specifically about AI/robotics.',
    },
    'healthcare': {
        'name': 'Healthcare',
        'style_prefix': 'Background: Warm beige/cream colored textured paper\nIllustrations: Cute hand-drawn cartoon style with medical-themed doodles\nTitle text: Large, colorful, hand-lettered typography (soft blues, greens, whites)\nBody text: Clean, readable text in dark gray/brown\nDecorations: Stars, sparkles, medical icons scattered around\nLayout: Professional yet friendly sketchnote style\nAspect ratio: 16:9 widescreen\nResolution: High quality, crisp text\nALL text MUST be in English\nDO NOT include any watermarks',
        'illustration_guidance': 'Use medical icons: stethoscopes, hearts, hospitals, pill capsules, doctors, nurses, medical charts, DNA helixes.',
        'negative_guidance': 'DO NOT include robots unless topic is health AI. DO NOT include violent or graphic medical imagery.',
    },
    'technology': {
        'name': 'Technology',
        'style_prefix': 'Background: Warm beige/cream colored textured paper\nIllustrations: Cute hand-drawn cartoon style with tech-themed doodles\nTitle text: Large, colorful, hand-lettered typography (electric blues, purples, neon greens)\nBody text: Clean, readable text in dark gray/brown\nDecorations: Stars, sparkles, circuit-like arrows, dots\nLayout: Professional yet friendly sketchnote style\nAspect ratio: 16:9 widescreen\nResolution: High quality, crisp text\nALL text MUST be in English\nDO NOT include any watermarks',
        'illustration_guidance': 'Use tech icons: laptops, smartphones, cloud symbols, servers, code brackets, AI brain icons. Robot/AI characters are appropriate.',
        'negative_guidance': 'DO NOT include medical or finance-specific imagery unless warranted.',
    },
    'finance': {
        'name': 'Finance',
        'style_prefix': 'Background: Warm beige/cream colored textured paper\nIllustrations: Cute hand-drawn cartoon style with finance-themed doodles\nTitle text: Large, colorful, hand-lettered typography (deep blues, golds, greens)\nBody text: Clean, readable text in dark gray/brown\nDecorations: Stars, sparkles, dollar signs scattered around\nLayout: Professional yet friendly sketchnote style\nAspect ratio: 16:9 widescreen\nResolution: High quality, crisp text\nALL text MUST be in English\nDO NOT include any watermarks',
        'illustration_guidance': 'Use finance icons: stock charts, bank buildings, money bags, bull/bear characters, trading terminals.',
        'negative_guidance': 'DO NOT include robots unless discussing fintech/algo trading.',
    },
    'general': {
        'name': 'General',
        'style_prefix': 'Background: Warm beige/cream colored textured paper\nIllustrations: Cute hand-drawn cartoon style with colorful doodles\nTitle text: Large, colorful, hand-lettered typography\nBody text: Clean, readable text in dark gray/brown\nDecorations: Stars, sparkles, arrows scattered around\nLayout: Professional yet friendly sketchnote style\nAspect ratio: 16:9 widescreen\nResolution: High quality, crisp text\nALL text MUST be in English\nDO NOT include any watermarks',
        'illustration_guidance': 'Use topic-appropriate icons and characters.',
        'negative_guidance': 'DO NOT include robots unless topic is about AI/robotics.',
    },
}

# --- Thesis Infographic Styles ---
THESIS_INFOGRAPHIC_STYLES = {
    'professional': {
        'name': 'Professional',
        'prompt': 'Create a PROFESSIONAL CORPORATE INFOGRAPHIC with the following visual style:\n- Background: Deep navy (#1e293b) with subtle geometric patterns\n- Layout: Structured grid with clear sections separated by thin lines\n- Colors: Navy, white, slate blue accents, gold highlights for key numbers\n- Typography: Clean sans-serif, large bold headers, readable body text\n- Icons: Simple line icons (not filled) in white/gold\n- Data visualization: Clean bar charts, progress indicators, status badges\n- Overall feel: Fortune 500 boardroom presentation\n',
        'colorSchemes': [
            {'id': 'default', 'name': 'Navy & Teal', 'promptSuffix': ''},
            {'id': 'midnight', 'name': 'Midnight', 'promptSuffix': '\nCOLOR OVERRIDE: Use black (#0a0a0a) background with electric blue (#3b82f6) accents instead of navy/teal. White text, blue highlights.'},
            {'id': 'charcoal_amber', 'name': 'Charcoal & Amber', 'promptSuffix': '\nCOLOR OVERRIDE: Use dark charcoal (#1c1917) background with warm amber (#f59e0b) accents instead of navy/teal. Cream text, amber highlights.'},
        ],
    },
    'whiteboard': {
        'name': 'Whiteboard',
        'prompt': 'Create a WHITEBOARD STYLE INFOGRAPHIC with the following visual style:\n- Background: Clean white/very light gray\n- Layout: Organized boxes and sections with connecting arrows, mind-map feel\n- Colors: Blue, red, green, black on white background\n- Typography: Hand-written marker style text, neat and readable\n- Icons: Simple hand-drawn sketches, basic shapes\n- Data visualization: Hand-drawn charts, circled numbers, underlined key metrics\n- Overall feel: Clean strategy brainstorm — organized and professional, not cluttered\n',
        'colorSchemes': [
            {'id': 'default', 'name': 'Classic', 'promptSuffix': ''},
            {'id': 'pastel', 'name': 'Pastel', 'promptSuffix': '\nCOLOR OVERRIDE: Use soft pastel colors — light blue, blush pink, mint green, lavender — instead of bold primary markers. Gentle, modern feel.'},
            {'id': 'bold', 'name': 'Bold', 'promptSuffix': '\nCOLOR OVERRIDE: Use thick high-contrast lines — black outlines, bright red/blue/green fills. Bold, punchy, heavy marker weight.'},
        ],
    },
    'sketchnote': {
        'name': 'Sketchnote',
        'prompt': 'Create a SKETCHNOTE STYLE INFOGRAPHIC with the following visual style:\n- Background: Warm beige/cream textured paper like a Moleskine notebook\n- Layout: Organic flow with hand-lettered headers, doodle borders\n- Colors: Warm palette - brown, orange, teal, mustard yellow, coral\n- Typography: Hand-lettered headers (large, playful), neat handwriting for body\n- Icons: Cute cartoon doodles, tiny illustrations, emoji-like drawings\n- Data visualization: Doodle charts, numbered lists with fun bullets\n- Decorations: Stars, sparkles, arrows, sticky notes, washi tape borders\n- Overall feel: Designer sketchnote from a TED talk\n',
        'colorSchemes': [
            {'id': 'default', 'name': 'Warm', 'promptSuffix': ''},
            {'id': 'cool', 'name': 'Cool', 'promptSuffix': '\nCOLOR OVERRIDE: Use cool palette — indigo, teal, mint green, slate blue — instead of warm browns/oranges. Keep the notebook texture.'},
            {'id': 'mono', 'name': 'Monochrome', 'promptSuffix': '\nCOLOR OVERRIDE: Use monochrome palette — black ink with burnt orange (#ea580c) as the only accent color. Like a Moleskine with one orange pen.'},
        ],
    },
    'blueprint': {
        'name': 'Blueprint',
        'prompt': 'Create a BLUEPRINT STYLE INFOGRAPHIC with the following visual style:\n- Background: Deep navy/dark blue (#0f172a) with subtle grid lines\n- Layout: Technical drawing style with measurement marks and grid\n- Colors: Cyan/electric blue lines on dark navy, white text, neon accents\n- Typography: Monospaced/technical font, all-caps headers, clean body text\n- Icons: Wireframe/outline style icons, technical diagrams, schematic symbols\n- Data visualization: Technical gauges, radar charts, status indicators with glow\n- Decorations: Grid dots, coordinate markers, dimension lines, crosshairs\n- Overall feel: Engineering blueprint or technical schematic\n',
        'colorSchemes': [
            {'id': 'default', 'name': 'Cyan', 'promptSuffix': ''},
            {'id': 'matrix', 'name': 'Matrix', 'promptSuffix': '\nCOLOR OVERRIDE: Use dark black (#0a0a0a) background with bright green (#22c55e) lines and text — Matrix/terminal aesthetic. Green on black.'},
            {'id': 'amber', 'name': 'Amber', 'promptSuffix': '\nCOLOR OVERRIDE: Use dark navy background with warm amber (#f59e0b) lines and text instead of cyan. Retro radar/instrument panel feel.'},
        ],
    },
    'dashboard': {
        'name': 'Dashboard',
        'prompt': 'Create a DATA DASHBOARD STYLE INFOGRAPHIC with the following visual style:\n- Background: White/very light gray with subtle card shadows\n- Layout: Multi-panel grid of cards/widgets, clean modern UI\n- Colors: White cards, dark text, green/yellow/red traffic lights, blue charts\n- Typography: Modern sans-serif, bold metric numbers, small labels\n- Icons: Material design style, filled, inside colored circles\n- Data visualization: Donut charts, sparklines, KPI cards with big numbers, traffic light dots\n- Decorations: Card shadows, thin borders, rounded corners, status pills\n- Overall feel: Bloomberg terminal meets modern SaaS dashboard\n',
        'colorSchemes': [
            {'id': 'default', 'name': 'Light', 'promptSuffix': ''},
            {'id': 'dark', 'name': 'Dark Mode', 'promptSuffix': '\nCOLOR OVERRIDE: Use dark mode — dark gray (#1e1e1e) background, dark cards (#2d2d2d), white text, blue/green/red accent colors. Modern dark dashboard.'},
            {'id': 'minimal', 'name': 'Minimal', 'promptSuffix': '\nCOLOR OVERRIDE: Use pure white background with single teal (#0d9488) accent color. Minimal, clean, no gradients. Teal for highlights and charts only.'},
        ],
    },
    'editorial': {
        'name': 'Editorial',
        'prompt': 'Create an EDITORIAL MAGAZINE STYLE INFOGRAPHIC with the following visual style:\n- Background: Off-white/cream paper with slight texture\n- Layout: Magazine editorial layout with columns, pull quotes, sidebar boxes\n- Colors: High-contrast black and white with ONE accent color (red #dc2626)\n- Typography: Large bold serif headers (like The Economist), clean sans-serif body\n- Icons: Minimal, line-art style, used sparingly\n- Data visualization: Clean minimal charts with red accent, large oversized numbers\n- Decorations: Thin black rules/dividers, drop caps, red accent bars\n- Overall feel: The Economist or Barron\'s magazine feature article\n',
        'colorSchemes': [
            {'id': 'default', 'name': 'Classic Red', 'promptSuffix': ''},
            {'id': 'blue', 'name': 'Blue Edition', 'promptSuffix': '\nCOLOR OVERRIDE: Use navy blue (#1e40af) as the single accent color instead of red. Blue accent bars, blue highlights, blue charts.'},
            {'id': 'gold', 'name': 'Gold Edition', 'promptSuffix': '\nCOLOR OVERRIDE: Use deep gold (#b45309) as the single accent color instead of red. Gold accent bars, gold highlights. Luxury financial magazine feel.'},
        ],
    },
    'precision': {
        'name': 'Precision',
        'prompt': '',
        'engine': 'pillow',
        'colorSchemes': [
            {'id': 'default', 'name': 'Navy & Teal', 'colors': {'bg': (30, 41, 59), 'accent': (20, 184, 166), 'dark': (15, 23, 42)}},
            {'id': 'midnight', 'name': 'Midnight', 'colors': {'bg': (15, 10, 40), 'accent': (139, 92, 246), 'dark': (5, 3, 20)}},
            {'id': 'arctic', 'name': 'Arctic', 'colors': {'bg': (15, 30, 50), 'accent': (56, 189, 248), 'dark': (8, 15, 30)}},
        ],
    },
    'analyst_brief': {
        'name': 'Analyst Brief',
        'prompt': '',
        'engine': 'pillow',
        'colorSchemes': [
            {'id': 'default', 'name': 'Dark Navy', 'colors': {
                'bg': (15, 23, 42), 'accent': (20, 184, 166), 'dark': (8, 15, 30),
                'card': (22, 33, 55), 'text': (255, 255, 255), 'muted': (148, 163, 184),
                'risk_accent': (239, 68, 68), 'border': (45, 58, 80),
            }},
            {'id': 'light', 'name': 'Light', 'colors': {
                'bg': (248, 250, 252), 'accent': (15, 118, 110), 'dark': (241, 245, 249),
                'card': (255, 255, 255), 'text': (15, 23, 42), 'muted': (100, 116, 139),
                'risk_accent': (220, 38, 38), 'border': (226, 232, 240),
            }},
        ],
    },
    'quad_grid': {
        'name': 'Quad Grid',
        'prompt': '',
        'engine': 'pillow',
        'colorSchemes': [
            {'id': 'default', 'name': 'Dark Navy', 'colors': {
                'bg': (15, 23, 42), 'accent': (20, 184, 166), 'dark': (8, 15, 30),
                'card': (22, 33, 55), 'text': (255, 255, 255), 'muted': (148, 163, 184),
                'risk_accent': (239, 68, 68), 'border': (45, 58, 80),
            }},
            {'id': 'light', 'name': 'Light', 'colors': {
                'bg': (248, 250, 252), 'accent': (15, 118, 110), 'dark': (241, 245, 249),
                'card': (255, 255, 255), 'text': (15, 23, 42), 'muted': (100, 116, 139),
                'risk_accent': (220, 38, 38), 'border': (226, 232, 240),
            }},
        ],
    },
}


# ============================================
# PRECISION INFOGRAPHIC RENDERER (Pillow)
# ============================================

def _img_to_base64(img):
    """Convert PIL Image to base64 PNG string."""
    import io
    buf = io.BytesIO()
    img.save(buf, format='PNG', optimize=True)
    return base64.b64encode(buf.getvalue()).decode('utf-8')

def _wrap_text(draw, text, font, max_width):
    """Word-wrap text to fit within max_width pixels."""
    words = (text or '').split()
    lines = []; current = ''
    for word in words:
        test = f"{current} {word}".strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current: lines.append(current)
            current = word
    if current: lines.append(current)
    return lines or ['']

def _load_fonts():
    """Load DejaVu Sans fonts with fallback."""
    from PIL import ImageFont
    fonts = {}
    try:
        base = "/usr/share/fonts/truetype/dejavu/"
        fonts['bold_lg'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 44)
        fonts['bold_md'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 26)
        fonts['bold_sm'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 18)
        fonts['regular'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 16)
        fonts['small'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 13)
        fonts['tiny'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 11)
    except Exception:
        fonts['bold_lg'] = ImageFont.load_default(size=44)
        fonts['bold_md'] = ImageFont.load_default(size=26)
        fonts['bold_sm'] = ImageFont.load_default(size=18)
        fonts['regular'] = ImageFont.load_default(size=16)
        fonts['small'] = ImageFont.load_default(size=13)
        fonts['tiny'] = ImageFont.load_default(size=11)
    return fonts

def _draw_status_dot(draw, x, y, status, size=14):
    """Draw a colored status circle."""
    colors = {'green': (22, 163, 74), 'yellow': (202, 138, 4), 'red': (220, 38, 38)}
    c = colors.get((status or '').lower(), (148, 163, 184))
    draw.ellipse([x, y, x + size, y + size], fill=c)

def _get_pillow_colors(style_key, color_scheme):
    """Load color dict from THESIS_INFOGRAPHIC_STYLES for a Pillow-rendered style."""
    sdef = THESIS_INFOGRAPHIC_STYLES.get(style_key, {})
    schemes = sdef.get('colorSchemes', [])
    # Default to first scheme
    colors = schemes[0].get('colors', {}) if schemes else {}
    if color_scheme and color_scheme != 'default':
        for cs in schemes:
            if cs['id'] == color_scheme and 'colors' in cs:
                colors = cs['colors']
                break
    return colors


def _generate_analyst_brief_infographic(d, scorecard_data, mode, detail='full', show_risk_detail=False, color_scheme=None, include_company=False):
    """Layout 1: Analyst Brief — dense card-based layout, light theme, content-fitted."""
    import html as _html
    from PIL import Image, ImageDraw, ImageFont
    sp_data = _build_signpost_data(d['signposts'], scorecard_data)
    rk_data = _build_risk_data(d['threats'], scorecard_data)
    g, y_count, r = _tally_statuses(scorecard_data)
    W, H = 1920, 1080

    # Font sizes — large enough to fill the slide with readable content
    fonts = {}
    try:
        base = "/usr/share/fonts/truetype/dejavu/"
        fonts['ticker']  = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 38)
        fonts['company'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 22)
        fonts['section'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 22)
        fonts['pillar_num'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 24)
        fonts['pillar_title'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 20)
        fonts['body']    = ImageFont.truetype(f"{base}DejaVuSans.ttf", 18)
        fonts['summary'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 22)
        fonts['bullet']  = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 19)
        fonts['detail']  = ImageFont.truetype(f"{base}DejaVuSans.ttf", 18)
        fonts['conclusion'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 16)
        fonts['footer']  = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 16)
        # For 3-slide mode
        fonts['bold_md'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 26)
        fonts['bold_sm'] = ImageFont.truetype(f"{base}DejaVuSans-Bold.ttf", 18)
        fonts['regular'] = ImageFont.truetype(f"{base}DejaVuSans.ttf", 16)
        fonts['small']   = ImageFont.truetype(f"{base}DejaVuSans.ttf", 14)
    except Exception:
        for name, sz in [('ticker',38),('company',22),('section',22),('pillar_num',24),
                         ('pillar_title',20),('body',18),('summary',22),('bullet',19),
                         ('detail',17),('conclusion',16),('footer',16),('bold_md',26),
                         ('bold_sm',18),('regular',16),('small',14)]:
            fonts[name] = ImageFont.load_default(size=sz)

    ticker = d['ticker']; company = d['company']
    thesis = d['thesis']; conclusion = d['conclusion']
    pillars = _build_pillar_data(thesis, scorecard_data)

    def ue(text):
        """Unescape HTML entities and strip quotation marks."""
        if not text:
            return text
        text = _html.unescape(text)
        for ch in ['\u201c', '\u201d', '\u2018', '\u2019', '"']:
            text = text.replace(ch, '')
        return text

    # -- Color palette --
    BG           = (242, 243, 247)
    HEADER_BG    = (20, 26, 40)
    HEADER_ACC   = (14, 180, 164)
    CARD_BG      = (255, 255, 255)
    CARD_BORDER  = (222, 224, 230)
    TEXT_PRI     = (28, 32, 48)
    TEXT_SEC     = (90, 100, 120)
    TEXT_MUTED   = (130, 138, 155)
    THESIS_ACC   = (82, 74, 160)
    THESIS_LIGHT = (242, 240, 255)
    SP_ACC       = (18, 155, 82)
    SP_LIGHT     = (236, 252, 242)
    SP_DOT       = (30, 172, 95)
    RISK_ACC     = (210, 60, 50)
    RISK_LIGHT   = (254, 241, 240)
    RISK_DOT     = (220, 72, 60)
    CONCL_BG     = (248, 249, 252)

    def rr(draw, box, fill, outline=None, rad=10):
        x0, y0, x1, y1 = box
        r = min(rad, (x1-x0)//2, (y1-y0)//2)
        draw.rectangle([x0+r, y0, x1-r, y1], fill=fill)
        draw.rectangle([x0, y0+r, x1, y1-r], fill=fill)
        draw.pieslice([x0, y0, x0+2*r, y0+2*r], 180, 270, fill=fill)
        draw.pieslice([x1-2*r, y0, x1, y0+2*r], 270, 360, fill=fill)
        draw.pieslice([x0, y1-2*r, x0+2*r, y1], 90, 180, fill=fill)
        draw.pieslice([x1-2*r, y1-2*r, x1, y1], 0, 90, fill=fill)
        if outline:
            draw.line([x0+r, y0, x1-r, y0], fill=outline)
            draw.line([x0+r, y1, x1-r, y1], fill=outline)
            draw.line([x0, y0+r, x0, y1-r], fill=outline)
            draw.line([x1, y0+r, x1, y1-r], fill=outline)
            draw.arc([x0, y0, x0+2*r, y0+2*r], 180, 270, fill=outline)
            draw.arc([x1-2*r, y0, x1, y0+2*r], 270, 360, fill=outline)
            draw.arc([x0, y1-2*r, x0+2*r, y1], 90, 180, fill=outline)
            draw.arc([x1-2*r, y1-2*r, x1, y1], 0, 90, fill=outline)

    def sec_hdr(draw, x, y, w, text, accent, tint):
        h = 34
        rr(draw, [x, y, x+w, y+h], fill=tint, rad=6)
        draw.rectangle([x, y+4, x+4, y+h-4], fill=accent)
        bbox = draw.textbbox((0,0), text, font=fonts['section'])
        tw = bbox[2]-bbox[0]
        draw.text((x + (w-tw)//2, y+6), text, font=fonts['section'], fill=accent)
        return y + h + 10

    import re as _re
    def fit_sentences(draw, text, font, max_w, max_n):
        """Select complete sentences that fit within max_n wrapped lines.

        Instead of truncating mid-thought, this keeps only whole sentences
        that fit. If even the first sentence is too long, it takes the
        first clause (up to a comma/semicolon) that fits.
        """
        if not text:
            return []
        # First check if the full text fits
        all_lines = _wrap_text(draw, text, font, max_w)
        if len(all_lines) <= max_n:
            return all_lines
        # Split into sentences (keep the delimiter attached)
        sentences = _re.split(r'(?<=[.!?])\s+', text.strip())
        if not sentences:
            return all_lines[:max_n]
        # Greedily add sentences until we exceed max_n lines
        chosen = ''
        for s in sentences:
            candidate = (chosen + ' ' + s).strip() if chosen else s
            wrapped = _wrap_text(draw, candidate, font, max_w)
            if len(wrapped) <= max_n:
                chosen = candidate
            else:
                break
        # If we got at least one complete sentence, use it
        if chosen:
            return _wrap_text(draw, chosen, font, max_w)
        # First sentence alone exceeds max_n lines — try first clause
        first = sentences[0]
        for delim in ['; ', ', ']:
            parts = first.split(delim)
            for k in range(len(parts), 0, -1):
                fragment = delim.join(parts[:k])
                if not fragment.endswith('.'):
                    fragment += '.'
                wrapped = _wrap_text(draw, fragment, font, max_w)
                if len(wrapped) <= max_n:
                    return wrapped
        # Last resort: take max_n lines of the first sentence
        return _wrap_text(draw, first, font, max_w)[:max_n]

    if mode == '1':
        M = 28
        img = Image.new('RGB', (W, H), BG)
        draw = ImageDraw.Draw(img)

        # Detail-level config: full / summary / simple
        # Full: all detail shown | Summary: less detail | Simple: titles only
        show_pillar_desc = detail != 'simple'
        show_sp_detail = detail != 'simple'  # all modes except simple show detail
        show_rk_detail_text = detail != 'simple'
        max_p = 3 if detail == 'simple' else min(len(pillars), 5)
        max_sp = 10 if detail != 'simple' else 6
        max_rk = 8 if detail != 'simple' else 5
        # Detail lines: full gets more, summary gets 1 line each
        sp_detail_lines = 1  # signpost detail is always 1 line (target info)
        rk_detail_lines = 2 if detail == 'full' else 1

        # --- Header ---
        HH = 56
        draw.rectangle([0, 0, W, HH], fill=HEADER_BG)
        draw.rectangle([0, HH, W, HH+3], fill=HEADER_ACC)
        draw.text((M, 10), ticker, font=fonts['ticker'], fill=(255,255,255))
        if company:
            tb = draw.textbbox((0,0), ticker, font=fonts['ticker'])
            draw.text((M + tb[2]-tb[0] + 18, 18), ue(company), font=fonts['company'], fill=(175,185,205))

        yc = HH + 8

        # --- Summary — always show full text, up to 5 lines ---
        summary = ue(thesis.get('summary', ''))
        sum_font = fonts['summary']  # 22pt
        max_summary = 5 if detail != 'simple' else 2
        if summary:
            sl = fit_sentences(draw, summary, sum_font, W - M*2 - 28, max_summary)
            LH_SUM = 30  # line height for 22pt
            sh = len(sl) * LH_SUM + 16
            rr(draw, [M, yc, W-M, yc+sh], fill=CARD_BG, outline=CARD_BORDER, rad=8)
            sy = yc + 8
            for line in sl:
                draw.text((M+14, sy), line, font=sum_font, fill=TEXT_SEC)
                sy += LH_SUM
            yc += sh + 8

        # --- Three columns ---
        GAP = 14
        COL_W = (W - M*2 - GAP*2) // 3
        C1X = M
        C2X = M + COL_W + GAP
        C3X = M + (COL_W + GAP)*2
        col_top = yc
        col_bottom = H - 8

        LH_TITLE = 26   # pillar title line height (20pt bold)
        LH_BODY = 24    # body/description line height (18pt)
        LH_BULLET = 26  # signpost/risk title line height (19pt bold)

        # ===== COLUMN 1: Investment Thesis (Pillars) =====
        yl = col_top
        yl = sec_hdr(draw, C1X, yl, COL_W, "Investment Thesis", THESIS_ACC, THESIS_LIGHT)

        avail_col = col_bottom - yl
        card_gap = 6
        card_h = (avail_col - card_gap * max(max_p - 1, 0)) // max(max_p, 1)

        for i, p in enumerate(pillars[:max_p], 1):
            ptitle = ue(p.get('pillar', p.get('title', '')))
            pdesc = ue(p.get('detail', p.get('description', '')))
            rr(draw, [C1X, yl, C1X+COL_W, yl+card_h], fill=CARD_BG, outline=CARD_BORDER, rad=8)
            draw.text((C1X+12, yl+10), str(i), font=fonts['pillar_num'], fill=THESIS_ACC)
            ty = yl + 12
            title_lines = _wrap_text(draw, ptitle, fonts['pillar_title'], COL_W - 56)[:2]
            for line in title_lines:
                draw.text((C1X+44, ty), line, font=fonts['pillar_title'], fill=TEXT_PRI)
                ty += LH_TITLE
            if pdesc and show_pillar_desc:
                ty += 4
                remaining_h = (yl + card_h - 10) - ty
                max_desc = max(1, remaining_h // LH_BODY)
                desc_lines = fit_sentences(draw, pdesc, fonts['body'], COL_W - 56, max_desc)
                for line in desc_lines:
                    if ty + LH_BODY > yl + card_h - 4:
                        break
                    draw.text((C1X+44, ty), line, font=fonts['body'], fill=TEXT_SEC)
                    ty += LH_BODY
            yl += card_h + card_gap

        # ===== COLUMN 2: Signposts / Catalysts =====
        yr = col_top
        yr = sec_hdr(draw, C2X, yr, COL_W, "Signposts / Catalysts", SP_ACC, SP_LIGHT)

        # Build signpost items — always include detail (except simple)
        sp_rich = []
        for s in sp_data:
            title = ue(s['metric'])
            if show_sp_detail:
                parts = []
                if s.get('latest'):
                    parts.append(ue(s['latest']))
                if s.get('ltGoal') and s.get('ltGoal') != s.get('latest'):
                    parts.append(f"Target: {ue(s['ltGoal'])}")
                detail_str = " | ".join(parts) if parts else ""
            else:
                detail_str = ""
            sp_rich.append((title, detail_str))

        sp_card_h = col_bottom - yr
        rr(draw, [C2X, yr, C2X+COL_W, yr+sp_card_h], fill=CARD_BG, outline=CARD_BORDER, rad=8)
        sy = yr + 14
        sp_n = min(len(sp_rich), max_sp)
        # Compute even spacing
        per_item = LH_BULLET + (LH_BODY * sp_detail_lines if show_sp_detail else 0)
        total_items_h = sp_n * per_item
        sp_gap = max(6, (sp_card_h - 28 - total_items_h) // max(sp_n, 1))
        for idx in range(sp_n):
            title, det = sp_rich[idx]
            if sy + LH_BULLET > yr + sp_card_h - 10:
                break
            status = _normalize_status(sp_data[idx].get('status', ''))
            dot_color = (30, 172, 95) if status == 'green' else (234, 179, 8) if status == 'yellow' else (220, 72, 60) if status == 'red' else SP_DOT
            draw.ellipse([C2X+14, sy+6, C2X+28, sy+20], fill=dot_color)
            tlines = _wrap_text(draw, title, fonts['bullet'], COL_W - 54)
            draw.text((C2X+36, sy), tlines[0], font=fonts['bullet'], fill=TEXT_PRI)
            sy += LH_BULLET
            if det and sy + LH_BODY < yr + sp_card_h - 10:
                dlines = fit_sentences(draw, det, fonts['body'], COL_W - 54, sp_detail_lines)
                for dl in dlines:
                    draw.text((C2X+36, sy), dl, font=fonts['body'], fill=TEXT_MUTED)
                    sy += LH_BODY
            sy += sp_gap

        # ===== COLUMN 3: Key Risks =====
        yrk = col_top
        yrk = sec_hdr(draw, C3X, yrk, COL_W, "Key Risks", RISK_ACC, RISK_LIGHT)

        rk_rich = []
        for rk in rk_data:
            title = ue(rk['threat'])
            if show_rk_detail_text:
                trigger = ue(rk.get('triggers', '')) or ''
                status_note = ue(rk.get('statusNote', '')) or ''
                detail_str = trigger if trigger else status_note
            else:
                detail_str = ""
            rk_rich.append((title, detail_str))

        rk_card_h = col_bottom - yrk
        rr(draw, [C3X, yrk, C3X+COL_W, yrk+rk_card_h], fill=CARD_BG, outline=CARD_BORDER, rad=8)
        ry = yrk + 14
        rk_n = min(len(rk_rich), max_rk)
        # Estimate height per risk: title(1-2 lines) + detail
        rk_per = LH_BULLET * 2 + LH_BODY * rk_detail_lines if show_rk_detail_text else LH_BULLET * 2
        rk_gap = max(6, (rk_card_h - 28 - rk_n * rk_per) // max(rk_n, 1))
        for idx in range(rk_n):
            title, det = rk_rich[idx]
            if ry + LH_BULLET > yrk + rk_card_h - 10:
                break
            status = _normalize_status(rk_data[idx].get('status', ''))
            dot_color = (30, 172, 95) if status == 'green' else (234, 179, 8) if status == 'yellow' else (220, 72, 60) if status == 'red' else RISK_DOT
            draw.ellipse([C3X+14, ry+6, C3X+28, ry+20], fill=dot_color)
            # Title — wrap up to 2 lines with clean truncation
            tlines = fit_sentences(draw, title, fonts['bullet'], COL_W - 54, 2)
            for tl in tlines:
                draw.text((C3X+36, ry), tl, font=fonts['bullet'], fill=TEXT_PRI)
                ry += LH_BULLET
            # Trigger/detail text
            if det and show_rk_detail_text:
                dlines = fit_sentences(draw, det, fonts['body'], COL_W - 54, rk_detail_lines)
                for dl in dlines:
                    if ry + LH_BODY > yrk + rk_card_h - 10:
                        break
                    draw.text((C3X+36, ry), dl, font=fonts['body'], fill=TEXT_MUTED)
                    ry += LH_BODY
            ry += rk_gap

        return [_img_to_base64(img)]
    else:
        # ====== 3-SLIDE MODE ======
        images = []
        HH = 62; M = 40

        def slide_hdr(draw, title):
            draw.rectangle([0,0,W,HH], fill=HEADER_BG)
            draw.rectangle([0,HH,W,HH+3], fill=HEADER_ACC)
            draw.text((M,14), title, font=fonts['bold_md'], fill=(255,255,255))

        # Slide 1: Thesis
        img1 = Image.new('RGB', (W, H), BG)
        d1 = ImageDraw.Draw(img1)
        slide_hdr(d1, f"{ticker} — Investment Thesis (1/3)")
        y = HH + 18
        summary = ue(thesis.get('summary', ''))
        if summary:
            sl = _wrap_text(d1, summary, fonts['summary'], W-M*2-36)[:4]
            sh = len(sl)*22+20
            rr(d1, [M,y,W-M,y+sh], fill=CARD_BG, outline=CARD_BORDER, rad=8)
            sy = y+10
            for line in sl:
                d1.text((M+18,sy), line, font=fonts['summary'], fill=TEXT_SEC); sy+=22
            y += sh + 16
        max_p = 3 if detail == 'simple' else min(len(pillars), 6)
        avail = H - y - M
        card_h = max(55, (avail - 8*max(max_p-1,0)) // max(max_p,1))
        card_h = min(card_h, 130)
        for i, p in enumerate(pillars[:max_p], 1):
            pt = ue(p.get('pillar', p.get('title', '')))
            pd = ue(p.get('detail', p.get('description', '')))
            rr(d1, [M,y,W-M,y+card_h], fill=CARD_BG, outline=CARD_BORDER, rad=8)
            d1.text((M+14,y+10), str(i), font=fonts['pillar_num'], fill=THESIS_ACC)
            ty = y+12
            for line in _wrap_text(d1, pt, fonts['pillar_title'], W-M*2-72)[:2]:
                d1.text((M+46,ty), line, font=fonts['pillar_title'], fill=TEXT_PRI); ty+=22
            if pd:
                ty += 4
                for line in _wrap_text(d1, pd, fonts['body'], W-M*2-72)[:max(1,(card_h-(ty-y)-10)//20)]:
                    d1.text((M+46,ty), line, font=fonts['body'], fill=TEXT_SEC); ty+=20
            y += card_h + 8
        images.append(_img_to_base64(img1))

        # Slide 2: Signposts
        img2 = Image.new('RGB', (W, H), BG)
        d2 = ImageDraw.Draw(img2)
        slide_hdr(d2, f"{ticker} — Signposts (2/3)")
        y = HH + 18
        y = sec_hdr(d2, M, y, W-M*2, "Signposts / Catalysts", SP_ACC, SP_LIGHT)
        for s in sp_data[:12]:
            d2.ellipse([M+18, y+7, M+28, y+17], fill=SP_DOT)
            metric = ue(s['metric'])
            if detail != 'simple' and s.get('latest'):
                metric += f" — {ue(s['latest'])}"
            lines = _wrap_text(d2, metric, fonts['bullet'], W-M*2-56)
            for line in lines[:2]:
                d2.text((M+38,y), line, font=fonts['bullet'], fill=TEXT_PRI); y+=28
            y += 6
        images.append(_img_to_base64(img2))

        # Slide 3: Risks
        img3 = Image.new('RGB', (W, H), BG)
        d3 = ImageDraw.Draw(img3)
        slide_hdr(d3, f"{ticker} — Risk Assessment (3/3)")
        y = HH + 18
        y = sec_hdr(d3, M, y, W-M*2, "Key Risks", RISK_ACC, RISK_LIGHT)
        for rk in rk_data[:10]:
            d3.ellipse([M+18, y+7, M+28, y+17], fill=RISK_DOT)
            lines = _wrap_text(d3, ue(rk['threat']), fonts['bullet'], W-M*2-56)
            for line in lines[:2]:
                d3.text((M+38,y), line, font=fonts['bullet'], fill=TEXT_PRI); y+=28
            y += 6
        images.append(_img_to_base64(img3))
        return images


def _generate_quad_grid_infographic(d, scorecard_data, mode, detail='full', show_risk_detail=False, color_scheme=None, include_company=False):
    """Layout 2: Quad Grid — pillars left, signposts top-right, risks bottom-right."""
    from PIL import Image, ImageDraw
    sp_data = _build_signpost_data(d['signposts'], scorecard_data)
    rk_data = _build_risk_data(d['threats'], scorecard_data)
    g, y_count, r = _tally_statuses(scorecard_data)
    fonts = _load_fonts()
    W, H = 1920, 1080
    c = _get_pillow_colors('quad_grid', color_scheme)
    BG = c.get('bg', (15, 23, 42)); ACCENT = c.get('accent', (20, 184, 166))
    DARK = c.get('dark', (8, 15, 30)); CARD = c.get('card', (22, 33, 55))
    TEXT = c.get('text', (255, 255, 255)); MUTED = c.get('muted', (148, 163, 184))
    RISK_ACCENT = c.get('risk_accent', (239, 68, 68)); BORDER = c.get('border', (45, 58, 80))

    ticker = d['ticker']; company = d['company']
    thesis = d['thesis']; conclusion = d['conclusion']
    header_text = f"{ticker} — {company}" if company and include_company else ticker
    pillars = _build_pillar_data(thesis, scorecard_data)

    def draw_section_pill(draw, x, y, text, color, width=280):
        """Draw a colored section header pill."""
        draw.rectangle([x, y, x + width, y + 32], fill=color)
        draw.text((x + 12, y + 6), text, font=fonts['bold_sm'], fill=(255, 255, 255))
        return y + 42

    if mode == '1':
        # Landscape 1920x1080 — content fills the full canvas
        M = 30  # narrow margin
        img = Image.new('RGB', (W, H), BG)
        draw = ImageDraw.Draw(img)

        # Header bar
        draw.rectangle([0, 0, W, 68], fill=DARK)
        draw.text((M, 16), header_text, font=fonts['bold_md'], fill=TEXT)
        draw.rectangle([0, 68, W, 72], fill=ACCENT)

        FOOTER_H = 50
        yc = 80

        # Thesis summary (full width)
        summary = thesis.get('summary', '')
        if summary:
            for line in _wrap_text(draw, summary, fonts['regular'], W - M * 2)[:2]:
                draw.text((M, yc), line, font=fonts['regular'], fill=MUTED)
                yc += 24
            yc += 8

        # Divider
        draw.rectangle([M, yc, W - M, yc + 2], fill=BORDER)
        yc += 12

        # Two-column layout: 55% left (pillars), 45% right (signposts+risks)
        COL_GAP = 28
        LEFT_W = int((W - M * 2 - COL_GAP) * 0.55)
        LEFT_X = M
        RIGHT_X = M + LEFT_W + COL_GAP
        RIGHT_W = W - RIGHT_X - M
        col_top = yc
        col_bottom = H - FOOTER_H - 10

        y_left = col_top
        y_right = col_top

        # -- Left column: Investment Thesis pill --
        y_left = draw_section_pill(draw, LEFT_X, y_left, "Investment Thesis", ACCENT, width=240)

        # Adaptive card sizing to fill the full column height
        max_p = 3 if detail == 'simple' else min(len(pillars), 6)
        avail_pillar_h = col_bottom - y_left - 10
        card_gap = 10
        card_h_full = max(80, (avail_pillar_h - card_gap * max(max_p - 1, 0)) // max(max_p, 1))
        card_h_simple = min(60, card_h_full // 2)

        for i, p in enumerate(pillars[:max_p], 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pdesc = p.get('detail', p.get('description', ''))
            ch = card_h_full if detail == 'full' and pdesc else card_h_simple
            draw.rectangle([LEFT_X, y_left, LEFT_X + LEFT_W, y_left + ch], fill=CARD, outline=BORDER)
            # Large number
            draw.text((LEFT_X + 12, y_left + 6), str(i), font=fonts['bold_lg'], fill=ACCENT)
            # Title
            title_x = LEFT_X + 70
            title_lines = _wrap_text(draw, ptitle, fonts['bold_sm'], LEFT_W - 90)
            ty = y_left + 10
            for line in title_lines[:2]:
                draw.text((title_x, ty), line, font=fonts['bold_sm'], fill=TEXT)
                ty += 24
            # Description — fill remaining card space
            if detail == 'full' and pdesc:
                max_desc_lines = max(1, (ch - (ty - y_left) - 8) // 18)
                desc_lines = _wrap_text(draw, pdesc, fonts['small'], LEFT_W - 90)
                for line in desc_lines[:max_desc_lines]:
                    draw.text((title_x, ty), line, font=fonts['small'], fill=MUTED)
                    ty += 18
            y_left += ch + card_gap

        # -- Right column: Signposts + Risks --
        right_avail = col_bottom - y_right
        sp_zone = int(right_avail * 0.55)

        # Signposts pill
        y_right = draw_section_pill(draw, RIGHT_X, y_right, "Signposts / Catalysts", ACCENT, width=280)
        sp_bottom = col_top + sp_zone - 10
        for s in sp_data[:10]:
            if y_right > sp_bottom:
                break
            draw.ellipse([RIGHT_X + 4, y_right + 6, RIGHT_X + 14, y_right + 16], fill=ACCENT)
            metric = s['metric']
            if detail != 'simple' and s.get('latest'):
                metric += f": {s['latest']}"
            lines = _wrap_text(draw, metric, fonts['regular'], RIGHT_W - 28)
            for line in lines[:2]:
                draw.text((RIGHT_X + 22, y_right), line, font=fonts['regular'], fill=TEXT)
                y_right += 22
            y_right += 8

        # Risks pill
        y_right = col_top + sp_zone + 4
        y_right = draw_section_pill(draw, RIGHT_X, y_right, "Key Risks", RISK_ACCENT, width=160)
        for rk in rk_data[:8]:
            if y_right > col_bottom:
                break
            draw.ellipse([RIGHT_X + 4, y_right + 6, RIGHT_X + 14, y_right + 16], fill=RISK_ACCENT)
            lines = _wrap_text(draw, rk['threat'], fonts['regular'], RIGHT_W - 28)
            for line in lines[:2]:
                draw.text((RIGHT_X + 22, y_right), line, font=fonts['regular'], fill=TEXT)
                y_right += 22
            y_right += 8

        # Conclusion (full width, just above footer)
        if conclusion:
            cy = col_bottom - 4
            draw.rectangle([M, cy, W - M, cy + 2], fill=BORDER)
            cy += 8
            for line in _wrap_text(draw, conclusion, fonts['small'], W - M * 2)[:2]:
                draw.text((M, cy), line, font=fonts['small'], fill=MUTED)
                cy += 18

        # Footer health bar
        total = g + y_count + r
        draw.rectangle([0, H - FOOTER_H, W, H], fill=DARK)
        if total > 0:
            pct = round(g / total * 100)
            draw.text((M, H - 38), f"Thesis Health: {pct}%", font=fonts['bold_sm'], fill=ACCENT)
            bx = W - 240
            for label, cnt, clr in [('G', g, (22, 163, 74)), ('Y', y_count, (202, 138, 4)), ('R', r, (220, 38, 38))]:
                draw.rectangle([bx, H - 40, bx + 46, H - 12], fill=clr)
                draw.text((bx + 16, H - 38), str(cnt), font=fonts['bold_sm'], fill=(255, 255, 255))
                bx += 64

        return [_img_to_base64(img)]
    else:
        # 3-slide mode
        images = []
        # Slide 1: Thesis + Pillars
        img1 = Image.new('RGB', (W, H), BG)
        d1 = ImageDraw.Draw(img1)
        d1.rectangle([0, 0, W, 72], fill=DARK)
        d1.text((40, 18), f"{header_text} — Investment Thesis (1/3)", font=fonts['bold_md'], fill=TEXT)
        d1.rectangle([0, 72, W, 76], fill=ACCENT)
        y = 96
        y = draw_section_pill(d1, 40, y, "Investment Thesis", ACCENT, width=240)
        summary = thesis.get('summary', '')
        if detail in ('summary', 'simple') and '.' in summary:
            summary = summary[:summary.index('.') + 1]
        for line in _wrap_text(d1, summary, fonts['regular'], W - 100)[:4]:
            d1.text((40, y), line, font=fonts['regular'], fill=TEXT)
            y += 24
        y += 20
        max_p = 3 if detail == 'simple' else min(len(pillars), 6)
        for i, p in enumerate(pillars[:max_p], 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pdesc = p.get('detail', p.get('description', ''))
            card_h = 90 if detail == 'full' and pdesc else 55
            d1.rectangle([40, y, W - 40, y + card_h], fill=CARD, outline=BORDER)
            d1.text((56, y + 8), str(i), font=fonts['bold_lg'], fill=ACCENT)
            d1.text((110, y + 14), ptitle, font=fonts['bold_sm'], fill=TEXT)
            if detail == 'full' and pdesc:
                for line in _wrap_text(d1, pdesc, fonts['small'], W - 180)[:2]:
                    d1.text((110, y + 40), line, font=fonts['small'], fill=MUTED)
                    break
            y += card_h + 10
        images.append(_img_to_base64(img1))

        # Slide 2: Signposts
        img2 = Image.new('RGB', (W, H), BG)
        d2 = ImageDraw.Draw(img2)
        d2.rectangle([0, 0, W, 72], fill=DARK)
        d2.text((40, 18), f"{header_text} — Signposts (2/3)", font=fonts['bold_md'], fill=TEXT)
        d2.rectangle([0, 72, W, 76], fill=ACCENT)
        y = 96
        y = draw_section_pill(d2, 40, y, "Signposts / Catalysts", ACCENT, width=280)
        for s in sp_data[:12]:
            d2.ellipse([40, y + 6, 48, y + 14], fill=ACCENT)
            metric = s['metric']
            if detail != 'simple' and s.get('latest'):
                metric += f" — {s['latest']}"
            if detail == 'full' and s.get('ltGoal'):
                metric += f" (Target: {s['ltGoal']})"
            lines = _wrap_text(d2, metric, fonts['regular'], W - 120)
            for line in lines[:2]:
                d2.text((58, y), line, font=fonts['regular'], fill=TEXT)
                y += 24
            y += 8
        images.append(_img_to_base64(img2))

        # Slide 3: Risks
        img3 = Image.new('RGB', (W, H), BG)
        d3 = ImageDraw.Draw(img3)
        d3.rectangle([0, 0, W, 72], fill=DARK)
        d3.text((40, 18), f"{header_text} — Risk Assessment (3/3)", font=fonts['bold_md'], fill=TEXT)
        d3.rectangle([0, 72, W, 76], fill=ACCENT)
        y = 96
        y = draw_section_pill(d3, 40, y, "Key Risks", RISK_ACCENT, width=160)
        for rk in rk_data[:10]:
            d3.ellipse([40, y + 6, 48, y + 14], fill=RISK_ACCENT)
            text = rk['threat']
            lines = _wrap_text(d3, text, fonts['regular'], W - 120)
            for line in lines[:2]:
                d3.text((58, y), line, font=fonts['regular'], fill=TEXT)
                y += 24
            y += 8
        total = g + y_count + r
        if total > 0:
            d3.rectangle([0, H - 50, W, H], fill=DARK)
            pct = round(g / total * 100)
            d3.text((40, H - 40), f"Thesis Health: {pct}%", font=fonts['bold_sm'], fill=ACCENT)
        images.append(_img_to_base64(img3))
        return images


def _generate_precision_infographic(d, scorecard_data, mode, detail='full', show_risk_detail=False, color_scheme=None, include_company=False):
    """Generate deterministic infographic images using Pillow."""
    from PIL import Image, ImageDraw
    sp_data = _build_signpost_data(d['signposts'], scorecard_data)
    rk_data = _build_risk_data(d['threats'], scorecard_data)
    g, y_count, r = _tally_statuses(scorecard_data)
    fonts = _load_fonts()
    W, H = 1920, 1080
    BG = (30, 41, 59); TEAL = (20, 184, 166); WHITE = (255, 255, 255)
    SLATE = (148, 163, 184); DARK = (15, 23, 42); LIGHT = (203, 213, 225)

    # Override colors from color scheme
    if color_scheme and color_scheme != 'default':
        prec_def = THESIS_INFOGRAPHIC_STYLES.get('precision', {})
        for cs in prec_def.get('colorSchemes', []):
            if cs['id'] == color_scheme and 'colors' in cs:
                BG = cs['colors']['bg']
                TEAL = cs['colors']['accent']
                DARK = cs['colors']['dark']
                break

    def draw_header(draw, img, title_text, subtitle=''):
        draw.rectangle([0, 0, W, 80], fill=DARK)
        draw.text((30, 20), title_text, font=fonts['bold_md'], fill=WHITE)
        if subtitle:
            bbox = draw.textbbox((0, 0), subtitle, font=fonts['small'])
            draw.text((W - 30 - (bbox[2] - bbox[0]), 30), subtitle, font=fonts['small'], fill=SLATE)
        draw.rectangle([0, 80, W, 84], fill=TEAL)

    def draw_section_bar(draw, y_pos, text):
        draw.rectangle([30, y_pos, W - 30, y_pos + 32], fill=TEAL)
        draw.text((42, y_pos + 6), text, font=fonts['bold_sm'], fill=WHITE)
        return y_pos + 40

    def draw_signpost_table(draw, x, y_start, width, data, compact=False):
        font = fonts['small'] if compact else fonts['regular']
        hdr_font = fonts['bold_sm'] if not compact else fonts['small']
        row_h = 28 if compact else 34
        if detail == 'simple':
            # Simple: 2 columns only — Metric, Status
            col_w = [int(width * 0.75), int(width * 0.25)]
            headers = ['Metric', 'Status']
        else:
            col_w = [int(width * 0.32), int(width * 0.25), int(width * 0.25), int(width * 0.18)]
            headers = ['Metric', 'LT Goal', 'Latest', 'Status']
        draw.rectangle([x, y_start, x + width, y_start + row_h], fill=tuple(list(TEAL) + [200]) if len(TEAL) == 3 else TEAL)
        cx = x
        for i, h in enumerate(headers):
            draw.text((cx + 8, y_start + 6), h, font=hdr_font, fill=WHITE)
            cx += col_w[i]
        y = y_start + row_h
        for s in data:
            draw.rectangle([x, y, x + width, y + row_h], fill=DARK)
            draw.line([x, y, x + width, y], fill=(255, 255, 255, 30))
            cx = x
            metric = (s.get('metric', '') or '')[:30]
            draw.text((cx + 8, y + 6), metric, font=font, fill=LIGHT)
            cx += col_w[0]
            if detail != 'simple':
                draw.text((cx + 8, y + 6), (s.get('ltGoal', '') or '')[:20], font=font, fill=SLATE)
                cx += col_w[1]
                draw.text((cx + 8, y + 6), (s.get('latest', '') or '—')[:20], font=font, fill=SLATE)
                cx += col_w[2]
            _draw_status_dot(draw, cx + col_w[-1]//2 - 7, y + row_h//2 - 7, s.get('status', ''))
            y += row_h
        return y

    is_condensed = d.get('_condensed', False)
    def draw_risk_table(draw, x, y_start, width, data, compact=False):
        font = fonts['small'] if compact else fonts['regular']
        hdr_font = fonts['bold_sm'] if not compact else fonts['small']
        row_h = 28 if compact else 34
        if detail == 'simple' or is_condensed:
            # Simple or condensed: always 2 columns — Risk Factor, Status (no L/I)
            col_w = [int(width * 0.75), int(width * 0.25)]
            headers = ['Risk Factor', 'Status']
        elif show_risk_detail:
            col_w = [int(width * 0.75), int(width * 0.25)]
            headers = ['Risk Factor', 'Status']
        draw.rectangle([x, y_start, x + width, y_start + row_h], fill=(220, 38, 38, 180))
        cx = x
        for i, h in enumerate(headers):
            draw.text((cx + 8, y_start + 6), h, font=hdr_font, fill=WHITE)
            cx += col_w[i]
        y = y_start + row_h
        for rk in data:
            draw.rectangle([x, y, x + width, y + row_h], fill=DARK)
            draw.line([x, y, x + width, y], fill=(255, 255, 255, 30))
            cx = x
            max_chars = 60
            draw.text((cx + 8, y + 6), (rk.get('threat', '') or '')[:max_chars], font=font, fill=LIGHT)
            cx += col_w[0]
            _draw_status_dot(draw, cx + col_w[-1]//2 - 7, y + row_h//2 - 7, rk.get('status', ''))
            y += row_h
        return y

    ticker = d['ticker']; company = d['company']
    thesis = d['thesis']; conclusion = d['conclusion']
    header_text = f"{ticker} — {company}" if company and include_company else ticker

    if mode == '1':
        img = Image.new('RGB', (W, H), BG)
        draw = ImageDraw.Draw(img)
        draw_header(draw, img, header_text, 'Investment Thesis')
        # Thesis summary
        y = 100
        summary = thesis.get('summary', '')
        if detail in ('summary', 'simple') and '.' in summary:
            summary = summary[:summary.index('.') + 1]
        lines = _wrap_text(draw, summary, fonts['regular'], W - 80)
        for line in lines[:4]:
            draw.text((40, y), line, font=fonts['regular'], fill=WHITE)
            y += 22
        # Pillars
        y += 8
        max_pillars = 3 if detail == 'simple' else 5
        for i, p in enumerate(_build_pillar_data(thesis, scorecard_data)[:max_pillars], 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pdesc = p.get('detail', p.get('description', ''))
            if detail in ('summary', 'simple'):
                text = f"{i}. {ptitle}"
            else:
                text = f"{i}. {ptitle}: {pdesc[:100]}"
            draw.text((40, y), text[:90], font=fonts['small'], fill=TEAL)
            y += 20
        # Two-column: signposts + risks
        y += 10
        col_w = (W - 80) // 2 - 10
        y = draw_section_bar(draw, y, "SIGNPOSTS")
        sp_end = draw_signpost_table(draw, 30, y, col_w, sp_data[:7], compact=True)
        # Risks on right
        draw.rectangle([30 + col_w + 20, y - 40, W - 30, y - 8], fill=TEAL)
        draw.text((42 + col_w + 20, y - 34), "RISKS", font=fonts['bold_sm'], fill=WHITE)
        draw_risk_table(draw, 30 + col_w + 20, y, col_w, rk_data[:5], compact=True)
        # Conviction bar at bottom
        draw.rectangle([0, H - 60, W, H], fill=DARK)
        total_items = g + y_count + r
        if total_items > 0:
            pct = round(g / total_items * 100)
            draw.text((40, H - 48), f"Health: {pct}%", font=fonts['bold_sm'], fill=TEAL)
        gyr_x = W - 300
        for label, cnt, color in [('G', g, (22, 163, 74)), ('Y', y_count, (202, 138, 4)), ('R', r, (220, 38, 38))]:
            draw.rectangle([gyr_x, H - 50, gyr_x + 60, H - 14], fill=color)
            draw.text((gyr_x + 20, H - 48), str(cnt), font=fonts['bold_sm'], fill=WHITE)
            gyr_x += 80
        return [_img_to_base64(img)]
    else:
        images = []
        # Slide 1: Thesis + Pillars
        img1 = Image.new('RGB', (W, H), BG)
        d1 = ImageDraw.Draw(img1)
        draw_header(d1, img1, header_text, 'Investment Thesis — Slide 1/3')
        y = 110
        y = draw_section_bar(d1, y, "WHY DO WE OWN IT?")
        summary = thesis.get('summary', '')
        if detail in ('summary', 'simple') and '.' in summary: summary = summary[:summary.index('.') + 1]
        for line in _wrap_text(d1, summary, fonts['regular'], W - 80)[:6]:
            d1.text((40, y), line, font=fonts['regular'], fill=WHITE); y += 24
        y += 15
        max_pillars_s1 = 3 if detail == 'simple' else 6
        for i, p in enumerate(_build_pillar_data(thesis, scorecard_data)[:max_pillars_s1], 1):
            ptitle = p.get('pillar', p.get('title', ''))
            pdesc = p.get('detail', p.get('description', ''))
            d1.text((40, y), f"{i}.", font=fonts['bold_md'], fill=TEAL)
            d1.text((80, y), ptitle, font=fonts['bold_sm'], fill=WHITE)
            y += 28
            if detail not in ('summary', 'simple') and pdesc:
                for line in _wrap_text(d1, pdesc, fonts['small'], W - 120)[:2]:
                    d1.text((80, y), line, font=fonts['small'], fill=SLATE); y += 18
            y += 10
        conclusion_text = conclusion if isinstance(conclusion, str) else ''
        if detail == 'simple':
            conclusion_text = ''  # Omit conclusion in simple mode
        elif detail == 'summary' and '.' in conclusion_text:
            conclusion_text = conclusion_text[:conclusion_text.index('.') + 1]
        if conclusion_text:
            d1.rectangle([30, H - 100, W - 30, H - 30], fill=DARK)
            for line in _wrap_text(d1, conclusion_text, fonts['small'], W - 100)[:3]:
                d1.text((50, H - 90), line, font=fonts['small'], fill=LIGHT); break
        images.append(_img_to_base64(img1))

        # Slide 2: Signposts
        img2 = Image.new('RGB', (W, H), BG)
        d2 = ImageDraw.Draw(img2)
        draw_header(d2, img2, header_text, 'Key Signposts — Slide 2/3')
        y = 110
        y = draw_section_bar(d2, y, "WHAT ARE WE WATCHING?")
        draw_signpost_table(d2, 30, y, W - 60, sp_data[:12])
        images.append(_img_to_base64(img2))

        # Slide 3: Risks + Conviction
        img3 = Image.new('RGB', (W, H), BG)
        d3 = ImageDraw.Draw(img3)
        draw_header(d3, img3, header_text, 'Risk Assessment — Slide 3/3')
        y = 110
        y = draw_section_bar(d3, y, "WHAT COULD GO WRONG?")
        y_end = draw_risk_table(d3, 30, y, W - 60, rk_data[:10])
        # Conviction bar
        y_end += 30
        d3.rectangle([30, y_end, W - 30, y_end + 80], fill=DARK)
        total_items = g + y_count + r
        boxes_x = 50
        for label, cnt, color in [('GREEN', g, (22, 163, 74)), ('YELLOW', y_count, (202, 138, 4)), ('RED', r, (220, 38, 38))]:
            d3.rectangle([boxes_x, y_end + 10, boxes_x + 140, y_end + 65], fill=color)
            d3.text((boxes_x + 50, y_end + 15), str(cnt), font=fonts['bold_md'], fill=WHITE)
            d3.text((boxes_x + 10, y_end + 45), label, font=fonts['tiny'], fill=WHITE)
            boxes_x += 170
        if total_items > 0:
            pct = round(g / total_items * 100)
            d3.text((boxes_x + 40, y_end + 20), f"Health Score: {pct}%", font=fonts['bold_md'], fill=TEAL)
        images.append(_img_to_base64(img3))
        return images

SECTOR_THEME_MAP = {
    'Healthcare': 'healthcare',
    'Health Care': 'healthcare',
    'Technology': 'technology',
    'Information Technology': 'technology',
    'Communication Services': 'technology',
    'Financials': 'finance',
    'Financial Services': 'finance',
    'Energy': 'general',
    'Consumer Discretionary': 'general',
    'Consumer Staples': 'general',
    'Industrials': 'general',
    'Materials': 'general',
    'Real Estate': 'finance',
    'Utilities': 'general',
}


def _build_slide_prompt(slide_data, theme_name, project_title, total_slides):
    """Build the complete Gemini prompt for a single slide."""
    theme = SLIDE_THEMES.get(theme_name, SLIDE_THEMES['sketchnote'])
    parts = ["You are generating a presentation slide image.\n"]
    parts.append("VISUAL STYLE (MUST follow exactly):")
    parts.append(theme['style_prefix'])
    if theme.get('illustration_guidance'):
        parts.append("\nILLUSTRATION GUIDANCE:")
        parts.append(theme['illustration_guidance'])
    if theme.get('negative_guidance'):
        parts.append("\nIMPORTANT RESTRICTIONS:")
        parts.append(theme['negative_guidance'])
    if not slide_data.get('no_header', False):
        parts.append(f'\nAt the top-left corner, show small header text: "{project_title}"')
        parts.append(f'At the top-right corner, show: "Slide {slide_data["slide_number"]} / {total_slides}"')
    else:
        parts.append("\nDo NOT include any header or slide number on this slide.")
    parts.append(f"\nSLIDE CONTENT:\n{slide_data['content'].strip()}")
    hints = slide_data.get('illustration_hints', [])
    if hints:
        parts.append("\nILLUSTRATIONS TO INCLUDE:")
        for h in hints:
            parts.append(f"- {h}")
    return "\n".join(parts)


def _generate_slide_image(prompt, api_key=None, reference_image=None):
    """Generate a slide image via Gemini API. Returns base64 PNG string or None.
    If reference_image (base64 PNG) is provided, it's sent as a style reference."""
    key = api_key or os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY', '')
    if not key:
        return None
    client = genai.Client(api_key=key)
    model = "gemini-3-pro-image-preview"
    # Build contents: either text-only or multimodal with reference image
    if reference_image:
        ref_instruction = (
            "Replicate the EXACT layout, typography, positioning, section arrangement, "
            "colors, and visual style of this reference image. Only replace the text "
            "content with the new data below. Keep the same visual structure.\n\n"
        )
        contents = [
            genai_types.Part.from_bytes(data=base64.b64decode(reference_image), mime_type='image/png'),
            ref_instruction + prompt
        ]
    else:
        contents = prompt
    for attempt in range(3):
        try:
            response = client.models.generate_content(
                model=model,
                contents=contents,
                config=genai_types.GenerateContentConfig(
                    response_modalities=["TEXT", "IMAGE"],
                    image_config=genai_types.ImageConfig(aspect_ratio="16:9"),
                ),
            )
            for part in response.candidates[0].content.parts:
                if hasattr(part, "inline_data") and part.inline_data is not None:
                    return base64.b64encode(part.inline_data.data).decode('utf-8')
            import time; time.sleep(2)
        except Exception as e:
            print(f"Slide generation attempt {attempt+1} failed: {e}")
            if attempt < 2:
                import time; time.sleep((attempt + 1) * 5)
    return None


def _compute_content_hash(slide_data):
    """Compute SHA-256 hash of slide content for change detection."""
    import hashlib
    hints = ','.join(slide_data.get('illustration_hints', []))
    payload = f"{slide_data.get('title','')}|{slide_data.get('type','')}|{slide_data.get('content','')}|{hints}|{slide_data.get('no_header', False)}"
    return hashlib.sha256(payload.encode()).hexdigest()


@app.route('/api/slides/themes', methods=['GET'])
def get_slide_themes():
    """List available slide themes."""
    result = []
    for key, theme in SLIDE_THEMES.items():
        result.append({'id': key, 'name': theme['name']})
    return jsonify(result)


@app.route('/api/slides/projects', methods=['GET'])
def get_slide_projects():
    """List all slide projects."""
    try:
        with get_db() as (conn, cur):
            cur.execute('''
                SELECT id, ticker, title, theme, status, total_slides, created_at, updated_at
                FROM slide_projects ORDER BY updated_at DESC
            ''')
            rows = cur.fetchall()
        result = []
        for row in rows:
            result.append({
                'id': row['id'],
                'ticker': row['ticker'],
                'title': row['title'],
                'theme': row['theme'],
                'status': row['status'],
                'total_slides': row['total_slides'],
                'created_at': row['created_at'].isoformat() if row['created_at'] else None,
                'updated_at': row['updated_at'].isoformat() if row['updated_at'] else None,
            })
        return jsonify(result)
    except Exception as e:
        print(f"Error getting slide projects: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects', methods=['POST'])
def create_slide_project():
    """Create a new slide project."""
    try:
        data = request.json
        title = data.get('title', '')
        ticker = data.get('ticker', '').upper() if data.get('ticker') else None
        theme = data.get('theme', 'sketchnote')
        if not title:
            return jsonify({'error': 'Title is required'}), 400
        if theme not in SLIDE_THEMES:
            theme = 'sketchnote'
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO slide_projects (ticker, title, theme, status, total_slides)
                VALUES (%s, %s, %s, 'draft', 0)
                RETURNING id
            ''', (ticker, title, theme))
            project_id = cur.fetchone()['id']
        return jsonify({'success': True, 'id': project_id})
    except Exception as e:
        print(f"Error creating slide project: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>', methods=['GET'])
def get_slide_project(project_id):
    """Get a slide project with all its slides."""
    try:
        with get_db() as (conn, cur):
            cur.execute('SELECT * FROM slide_projects WHERE id = %s', (project_id,))
            project = cur.fetchone()
            if not project:
                return jsonify({'error': 'Project not found'}), 404
            cur.execute('''
                SELECT id, slide_number, title, type, content, illustration_hints,
                       no_header, image_data, content_hash, status
                FROM slide_items WHERE project_id = %s ORDER BY slide_number
            ''', (project_id,))
            slides = cur.fetchall()
        result = {
            'id': project['id'],
            'ticker': project['ticker'],
            'title': project['title'],
            'theme': project['theme'],
            'status': project['status'],
            'total_slides': project['total_slides'],
            'created_at': project['created_at'].isoformat() if project['created_at'] else None,
            'updated_at': project['updated_at'].isoformat() if project['updated_at'] else None,
            'slides': [],
        }
        for s in slides:
            result['slides'].append({
                'id': s['id'],
                'slide_number': s['slide_number'],
                'title': s['title'],
                'type': s['type'],
                'content': s['content'],
                'illustration_hints': s['illustration_hints'] or [],
                'no_header': s['no_header'],
                'has_image': bool(s['image_data']),
                'status': s['status'],
            })
        return jsonify(result)
    except Exception as e:
        print(f"Error getting slide project: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>', methods=['DELETE'])
def delete_slide_project(project_id):
    """Delete a slide project and all its slides."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM slide_projects WHERE id = %s', (project_id,))
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting slide project: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/slides/<int:slide_num>', methods=['PUT'])
def update_slide(project_id, slide_num):
    """Update a slide's content."""
    try:
        data = request.json
        with get_db(commit=True) as (conn, cur):
            updates = []
            params = []
            for field in ['title', 'content', 'type', 'no_header']:
                if field in data:
                    updates.append(f"{field} = %s")
                    params.append(data[field])
            if 'illustration_hints' in data:
                updates.append("illustration_hints = %s")
                params.append(json.dumps(data['illustration_hints']))
            if not updates:
                return jsonify({'error': 'No fields to update'}), 400
            # Recompute hash and mark as edited
            new_hash = _compute_content_hash({
                'title': data.get('title', ''),
                'type': data.get('type', 'content'),
                'content': data.get('content', ''),
                'illustration_hints': data.get('illustration_hints', []),
                'no_header': data.get('no_header', False),
            })
            updates.append("content_hash = %s")
            params.append(new_hash)
            updates.append("status = 'edited'")
            updates.append("updated_at = %s")
            params.append(datetime.utcnow())
            params.extend([project_id, slide_num])
            cur.execute(f'''
                UPDATE slide_items SET {', '.join(updates)}
                WHERE project_id = %s AND slide_number = %s
            ''', params)
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error updating slide: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/slides', methods=['POST'])
def add_slide(project_id):
    """Add a new slide to a project."""
    try:
        data = request.json
        title = data.get('title', 'New Slide')
        after = data.get('after')  # insert after this slide number
        content = data.get('content', f'[Edit content for: {title}]')
        slide_type = data.get('type', 'content')
        hints = data.get('illustration_hints', [])
        no_header = data.get('no_header', False)
        with get_db(commit=True) as (conn, cur):
            if after is not None:
                new_num = after + 1
                # Shift existing slides up
                cur.execute('''
                    UPDATE slide_items SET slide_number = slide_number + 1, updated_at = %s
                    WHERE project_id = %s AND slide_number >= %s
                ''', (datetime.utcnow(), project_id, new_num))
            else:
                cur.execute('SELECT COALESCE(MAX(slide_number), 0) + 1 as next_num FROM slide_items WHERE project_id = %s', (project_id,))
                new_num = cur.fetchone()['next_num']
            content_hash = _compute_content_hash({
                'title': title, 'type': slide_type, 'content': content,
                'illustration_hints': hints, 'no_header': no_header,
            })
            cur.execute('''
                INSERT INTO slide_items (project_id, slide_number, title, type, content, illustration_hints, no_header, content_hash, status)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, 'new')
                RETURNING id
            ''', (project_id, new_num, title, slide_type, content, json.dumps(hints), no_header, content_hash))
            slide_id = cur.fetchone()['id']
            # Update project total
            cur.execute('SELECT COUNT(*) as cnt FROM slide_items WHERE project_id = %s', (project_id,))
            total = cur.fetchone()['cnt']
            cur.execute('UPDATE slide_projects SET total_slides = %s, updated_at = %s WHERE id = %s',
                        (total, datetime.utcnow(), project_id))
        return jsonify({'success': True, 'id': slide_id, 'slide_number': new_num})
    except Exception as e:
        print(f"Error adding slide: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/slides/<int:slide_num>', methods=['DELETE'])
def delete_slide_item(project_id, slide_num):
    """Delete a slide and renumber the rest."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM slide_items WHERE project_id = %s AND slide_number = %s', (project_id, slide_num))
            # Renumber slides after the deleted one
            cur.execute('''
                UPDATE slide_items SET slide_number = slide_number - 1, updated_at = %s
                WHERE project_id = %s AND slide_number > %s
            ''', (datetime.utcnow(), project_id, slide_num))
            # Update project total
            cur.execute('SELECT COUNT(*) as cnt FROM slide_items WHERE project_id = %s', (project_id,))
            total = cur.fetchone()['cnt']
            cur.execute('UPDATE slide_projects SET total_slides = %s, updated_at = %s WHERE id = %s',
                        (total, datetime.utcnow(), project_id))
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error deleting slide: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/generate/<int:slide_num>', methods=['POST'])
def generate_one_slide(project_id, slide_num):
    """Generate image for a single slide."""
    try:
        data = request.json or {}
        gemini_key = data.get('geminiApiKey', '') or os.environ.get('GEMINI_API_KEY', '') or os.environ.get('GOOGLE_API_KEY', '')
        if not gemini_key:
            return jsonify({'error': 'Gemini API key not configured. Add it in Settings.'}), 400
        with get_db() as (conn, cur):
            cur.execute('SELECT * FROM slide_projects WHERE id = %s', (project_id,))
            project = cur.fetchone()
            if not project:
                return jsonify({'error': 'Project not found'}), 404
            cur.execute('SELECT * FROM slide_items WHERE project_id = %s AND slide_number = %s', (project_id, slide_num))
            slide = cur.fetchone()
            if not slide:
                return jsonify({'error': 'Slide not found'}), 404
            total = project['total_slides']
        slide_data = {
            'slide_number': slide['slide_number'],
            'title': slide['title'],
            'type': slide['type'],
            'content': slide['content'] or '',
            'illustration_hints': slide['illustration_hints'] or [],
            'no_header': slide['no_header'],
        }
        prompt = _build_slide_prompt(slide_data, project['theme'], project['title'], total)
        image_b64 = _generate_slide_image(prompt, api_key=gemini_key)
        if not image_b64:
            return jsonify({'error': 'Failed to generate image'}), 500
        content_hash = _compute_content_hash(slide_data)
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                UPDATE slide_items SET image_data = %s, content_hash = %s, status = 'generated', updated_at = %s
                WHERE project_id = %s AND slide_number = %s
            ''', (image_b64, content_hash, datetime.utcnow(), project_id, slide_num))
        return jsonify({'success': True, 'has_image': True})
    except Exception as e:
        print(f"Error generating slide: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/generate', methods=['POST'])
def generate_slides(project_id):
    """Generate images for slides (all, changed-only, or specific slide_numbers)."""
    try:
        data = request.json or {}
        changed_only = data.get('changed_only', True)
        slide_numbers = data.get('slide_numbers', None)  # optional list of specific slide numbers
        gemini_key = data.get('geminiApiKey', '') or os.environ.get('GEMINI_API_KEY', '') or os.environ.get('GOOGLE_API_KEY', '')
        if not gemini_key:
            return jsonify({'error': 'Gemini API key not configured. Add it in Settings.'}), 400
        with get_db() as (conn, cur):
            cur.execute('SELECT * FROM slide_projects WHERE id = %s', (project_id,))
            project = cur.fetchone()
            if not project:
                return jsonify({'error': 'Project not found'}), 404
            cur.execute('UPDATE slide_projects SET status = %s, updated_at = %s WHERE id = %s',
                        ('generating', datetime.utcnow(), project_id))
            conn.commit()
        import threading
        def _generate_all(pid, theme, title, only_changed, api_key, target_slides=None):
            import time
            try:
                with get_db() as (conn, cur):
                    if target_slides:
                        placeholders = ','.join(['%s'] * len(target_slides))
                        cur.execute(f'SELECT * FROM slide_items WHERE project_id = %s AND slide_number IN ({placeholders}) ORDER BY slide_number',
                                    [pid] + list(target_slides))
                    else:
                        cur.execute('SELECT * FROM slide_items WHERE project_id = %s ORDER BY slide_number', (pid,))
                    slides = cur.fetchall()
                    # Get total slides in project for slide numbering
                    cur.execute('SELECT COUNT(*) as cnt FROM slide_items WHERE project_id = %s', (pid,))
                    total = cur.fetchone()['cnt']
                # Build list of slides that actually need generation
                to_generate = []
                for slide in slides:
                    slide_data = {
                        'slide_number': slide['slide_number'],
                        'title': slide['title'],
                        'type': slide['type'],
                        'content': slide['content'] or '',
                        'illustration_hints': slide['illustration_hints'] or [],
                        'no_header': slide['no_header'],
                    }
                    current_hash = _compute_content_hash(slide_data)
                    if target_slides:
                        # When specific slides requested, always generate them
                        to_generate.append((slide, slide_data, current_hash))
                    elif only_changed and slide['content_hash'] == current_hash and slide['image_data']:
                        continue
                    else:
                        to_generate.append((slide, slide_data, current_hash))
                generated_count = 0
                for idx, (slide, slide_data, current_hash) in enumerate(to_generate):
                    # Update progress: store in project status as "generating:current:total"
                    with get_db(commit=True) as (conn2, cur2):
                        progress_status = f"generating:{idx+1}:{len(to_generate)}"
                        cur2.execute('UPDATE slide_projects SET status = %s, updated_at = %s WHERE id = %s',
                                    (progress_status, datetime.utcnow(), pid))
                    prompt = _build_slide_prompt(slide_data, theme, title, total)
                    image_b64 = _generate_slide_image(prompt, api_key=api_key)
                    if image_b64:
                        with get_db(commit=True) as (conn2, cur2):
                            cur2.execute('''
                                UPDATE slide_items SET image_data = %s, content_hash = %s, status = 'generated', updated_at = %s
                                WHERE id = %s
                            ''', (image_b64, current_hash, datetime.utcnow(), slide['id']))
                        generated_count += 1
                    time.sleep(2)
                with get_db(commit=True) as (conn, cur):
                    cur.execute('UPDATE slide_projects SET status = %s, updated_at = %s WHERE id = %s',
                                ('ready', datetime.utcnow(), pid))
                print(f"Slide generation complete: {generated_count} slides for project {pid}")
            except Exception as e:
                print(f"Slide generation error: {e}")
                with get_db(commit=True) as (conn, cur):
                    cur.execute('UPDATE slide_projects SET status = %s, updated_at = %s WHERE id = %s',
                                ('error', datetime.utcnow(), pid))
        thread = threading.Thread(
            target=_generate_all,
            args=(project_id, project['theme'], project['title'], changed_only, gemini_key, slide_numbers),
            daemon=True,
        )
        thread.start()
        return jsonify({'success': True, 'status': 'generating'})
    except Exception as e:
        print(f"Error starting slide generation: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/slides/<int:slide_num>/image', methods=['GET'])
def get_slide_image(project_id, slide_num):
    """Get the generated image for a slide."""
    try:
        with get_db() as (conn, cur):
            cur.execute('SELECT image_data FROM slide_items WHERE project_id = %s AND slide_number = %s', (project_id, slide_num))
            row = cur.fetchone()
            if not row or not row['image_data']:
                return jsonify({'error': 'No image available'}), 404
        return jsonify({'image_data': row['image_data']})
    except Exception as e:
        print(f"Error getting slide image: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/export/pdf', methods=['POST'])
def export_slide_pdf(project_id):
    """Export all slides as a PDF. Returns base64-encoded PDF."""
    try:
        import img2pdf
        with get_db() as (conn, cur):
            cur.execute('SELECT title FROM slide_projects WHERE id = %s', (project_id,))
            project = cur.fetchone()
            if not project:
                return jsonify({'error': 'Project not found'}), 404
            cur.execute('''
                SELECT image_data FROM slide_items
                WHERE project_id = %s AND image_data IS NOT NULL
                ORDER BY slide_number
            ''', (project_id,))
            slides = cur.fetchall()
        if not slides:
            return jsonify({'error': 'No generated slides to export'}), 400
        img_bytes_list = []
        for s in slides:
            img_bytes_list.append(base64.b64decode(s['image_data']))
        pdf_bytes = img2pdf.convert(img_bytes_list)
        pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
        return jsonify({
            'success': True,
            'pdf_data': pdf_b64,
            'filename': f"{project['title'].replace(' ', '_')}.pdf",
            'slide_count': len(slides),
        })
    except Exception as e:
        print(f"Error exporting PDF: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/generate-outline', methods=['POST'])
def generate_slide_outline():
    """Generate a slide outline from various sources using LLM."""
    try:
        data = request.json
        source_type = data.get('source_type', 'ticker')  # ticker | research | summary | meetingprep | custom
        source_ids = data.get('source_ids', [])
        custom_text = data.get('custom_text', '')
        ticker = data.get('ticker', '').upper()

        context_parts = []
        company_name = ''
        theme = 'sketchnote'
        presentation_subject = ''

        if source_type == 'ticker':
            # Original behavior: pull from stock_overviews + portfolio_analyses
            if not ticker:
                return jsonify({'error': 'Ticker is required for ticker-based generation'}), 400
            overview_data = None
            analysis_data = None
            with get_db() as (conn, cur):
                cur.execute('SELECT * FROM stock_overviews WHERE ticker = %s', (ticker,))
                overview_row = cur.fetchone()
                if overview_row:
                    overview_data = {
                        'company_name': overview_row['company_name'],
                        'company_overview': overview_row['company_overview'],
                        'business_model': overview_row['business_model'],
                        'business_mix': overview_row.get('business_mix', ''),
                        'opportunities': overview_row['opportunities'],
                        'risks': overview_row['risks'],
                        'conclusion': overview_row['conclusion'],
                    }
                cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
                analysis_row = cur.fetchone()
                if analysis_row:
                    analysis_data = {
                        'company': analysis_row['company'],
                        'analysis': analysis_row['analysis'],
                    }
            if not overview_data and not analysis_data:
                return jsonify({'error': f'No existing data found for {ticker}. Create an overview or analysis first.'}), 404
            company_name = (overview_data or {}).get('company_name', '') or (analysis_data or {}).get('company', ticker)
            presentation_subject = f"an equity research presentation about {ticker} ({company_name})"
            # Auto-detect theme from sector
            sector = None
            if analysis_data and analysis_data.get('analysis'):
                a = analysis_data['analysis']
                if isinstance(a, str):
                    try:
                        a = json.loads(a)
                    except Exception:
                        a = {}
                sector = a.get('sector', '')
            theme = SECTOR_THEME_MAP.get(sector, 'sketchnote')
            context_parts = [f"Ticker: {ticker}", f"Company: {company_name}"]
            if overview_data:
                for field in ['company_overview', 'business_model', 'business_mix', 'opportunities', 'risks', 'conclusion']:
                    if overview_data.get(field):
                        label = field.replace('_', ' ').title()
                        context_parts.append(f"{label}:\n{overview_data[field]}")
            if analysis_data and analysis_data.get('analysis'):
                a = analysis_data['analysis']
                if isinstance(a, dict):
                    context_parts.append(f"Investment Analysis:\n{json.dumps(a, indent=2)[:3000]}")

        elif source_type == 'research':
            if not source_ids:
                return jsonify({'error': 'Select at least one research document'}), 400
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(source_ids))
                cur.execute(f'SELECT id, name, smart_name, content FROM research_documents WHERE id IN ({placeholders})', source_ids)
                rows = cur.fetchall()
            if not rows:
                return jsonify({'error': 'No research documents found for the selected IDs'}), 404
            presentation_subject = "a presentation based on the following research documents"
            for row in rows:
                doc_name = row['smart_name'] or row['name'] or f"Document {row['id']}"
                content = (row['content'] or '')[:5000]
                context_parts.append(f"Document: {doc_name}\n{content}")

        elif source_type == 'summary':
            if not source_ids:
                return jsonify({'error': 'Select at least one meeting summary'}), 400
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(source_ids))
                cur.execute(f'SELECT id, title, topic, raw_notes, summary FROM meeting_summaries WHERE id IN ({placeholders})', source_ids)
                rows = cur.fetchall()
            if not rows:
                return jsonify({'error': 'No meeting summaries found for the selected IDs'}), 404
            presentation_subject = "a presentation based on the following meeting notes and summaries"
            for row in rows:
                title = row['title'] or f"Summary {row['id']}"
                context_parts.append(f"Meeting: {title}")
                if row.get('topic'):
                    context_parts.append(f"Topic: {row['topic']}")
                if row.get('raw_notes'):
                    context_parts.append(f"Notes:\n{row['raw_notes'][:3000]}")
                if row.get('summary'):
                    context_parts.append(f"Summary:\n{row['summary'][:3000]}")

        elif source_type == 'meetingprep':
            if not source_ids:
                return jsonify({'error': 'Select at least one meeting prep session'}), 400
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(source_ids))
                cur.execute(f'''SELECT m.id, m.notes, c.ticker, c.name as company_name, c.sector
                               FROM mp_meetings m JOIN mp_companies c ON m.company_id = c.id
                               WHERE m.id IN ({placeholders})''', source_ids)
                meetings = cur.fetchall()
                if not meetings:
                    return jsonify({'error': 'No meeting prep sessions found for the selected IDs'}), 404
                # Fetch documents for these meetings
                cur.execute(f'SELECT meeting_id, filename, extracted_text FROM mp_documents WHERE meeting_id IN ({placeholders}) ORDER BY upload_order', source_ids)
                docs = cur.fetchall()
            presentation_subject = "a presentation based on meeting prep materials"
            for meeting in meetings:
                context_parts.append(f"Meeting Prep: {meeting['company_name']} ({meeting['ticker']})")
                if meeting.get('notes'):
                    context_parts.append(f"Meeting Notes:\n{meeting['notes'][:2000]}")
                if meeting.get('sector'):
                    theme = SECTOR_THEME_MAP.get(meeting['sector'], 'sketchnote')
            for doc in docs:
                if doc.get('extracted_text'):
                    context_parts.append(f"Document ({doc['filename']}):\n{doc['extracted_text'][:3000]}")

        elif source_type == 'custom':
            if not custom_text.strip():
                return jsonify({'error': 'Please provide text content'}), 400
            presentation_subject = "a presentation based on the following content"
            context_parts.append(custom_text[:10000])

        else:
            return jsonify({'error': f'Unknown source_type: {source_type}'}), 400

        llm_prompt = f"""You are creating slide content for {presentation_subject}.

Based on this source material:

{chr(10).join(context_parts)}

Generate a JSON array of slides for a comprehensive presentation. Each slide should have:
- slide_number (integer starting at 1)
- title (string)
- type (one of: title, toc, section_divider, content, closing)
- content (detailed text describing what should appear on the slide - include specific data, numbers, quotes from the source material)
- illustration_hints (array of strings like "company", "growth", "money", "risk", "data", "leader", "opportunity")
- no_header (boolean - true only for title, section_divider, and closing slides)

Structure: title slide, table of contents, then 5-7 sections covering the key themes from the source material. Each section should have a section divider slide followed by 2-4 content slides. End with a closing slide. Target 30-40 slides total.

IMPORTANT: Fill in actual data, numbers, and specific details from the source material. Do NOT use placeholder brackets like [Company Name] - use real names and real data from the content provided.

Return ONLY valid JSON array, no markdown fencing."""

        # Call Claude for outline generation
        api_key = os.environ.get('ANTHROPIC_API_KEY', '') or data.get('apiKey', '') or data.get('api_key', '')
        if not api_key:
            return jsonify({'error': 'Anthropic API key not configured. Add it in Settings.'}), 500
        import httpx as _httpx
        client_ai = anthropic.Anthropic(api_key=api_key, timeout=_httpx.Timeout(120.0, connect=15.0))
        response = client_ai.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=8192,
            messages=[{"role": "user", "content": llm_prompt}],
        )
        response_text = response.content[0].text.strip()
        # Parse JSON from response
        if response_text.startswith('```'):
            response_text = response_text.split('\n', 1)[1].rsplit('```', 1)[0].strip()
        slides_data = json.loads(response_text)
        return jsonify({
            'success': True,
            'ticker': ticker,
            'company_name': company_name,
            'theme': theme,
            'slides': slides_data,
        })
    except json.JSONDecodeError as e:
        print(f"Error parsing slide outline JSON: {e}")
        return jsonify({'error': 'Failed to parse LLM response as JSON'}), 500
    except Exception as e:
        print(f"Error generating slide outline: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/slides/projects/<int:project_id>/populate', methods=['POST'])
def populate_slides(project_id):
    """Populate a project with slides from generated outline data."""
    try:
        data = request.json
        slides_data = data.get('slides', [])
        if not slides_data:
            return jsonify({'error': 'No slides data provided'}), 400
        with get_db(commit=True) as (conn, cur):
            # Clear existing slides
            cur.execute('DELETE FROM slide_items WHERE project_id = %s', (project_id,))
            for s in slides_data:
                content_hash = _compute_content_hash(s)
                cur.execute('''
                    INSERT INTO slide_items (project_id, slide_number, title, type, content, illustration_hints, no_header, content_hash, status)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, 'new')
                ''', (
                    project_id, s['slide_number'], s['title'], s.get('type', 'content'),
                    s.get('content', ''), json.dumps(s.get('illustration_hints', [])),
                    s.get('no_header', False), content_hash,
                ))
            total = len(slides_data)
            cur.execute('UPDATE slide_projects SET total_slides = %s, updated_at = %s WHERE id = %s',
                        (total, datetime.utcnow(), project_id))
        return jsonify({'success': True, 'total_slides': total})
    except Exception as e:
        print(f"Error populating slides: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# SLIDE SOURCE LISTING ENDPOINTS
# ============================================

@app.route('/api/slides/sources/research', methods=['GET'])
def slide_sources_research():
    """List research documents for slide outline source selection."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT id, name, smart_name, category_id, doc_type, published_date, LEFT(content, 200) as snippet, created_at FROM research_documents ORDER BY created_at DESC')
            rows = cur.fetchall()
        return jsonify([{
            'id': row['id'],
            'name': row['smart_name'] or row['name'] or f"Document {row['id']}",
            'category': row['category_id'],
            'doc_type': row.get('doc_type') or 'other',
            'date': row['published_date'] or (row['created_at'].isoformat() if row['created_at'] else None),
            'snippet': (row['snippet'] or '')[:150] + ('...' if row.get('snippet') and len(row['snippet']) > 150 else ''),
        } for row in rows])
    except Exception as e:
        print(f"Error listing research sources: {e}")
        return jsonify([])


@app.route('/api/slides/sources/summaries', methods=['GET'])
def slide_sources_summaries():
    """List meeting summaries for slide outline source selection."""
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT id, title, topic, topic_type, LEFT(summary, 200) as snippet, created_at FROM meeting_summaries ORDER BY created_at DESC')
            rows = cur.fetchall()
        return jsonify([{
            'id': row['id'],
            'name': row['title'] or f"Summary {row['id']}",
            'topic': row.get('topic') or 'General',
            'topic_type': row.get('topic_type') or 'other',
            'date': row['created_at'].isoformat() if row['created_at'] else None,
            'snippet': (row['snippet'] or '')[:150] + ('...' if row.get('snippet') and len(row['snippet']) > 150 else ''),
        } for row in rows])
    except Exception as e:
        print(f"Error listing summary sources: {e}")
        return jsonify([])


@app.route('/api/slides/sources/meetingprep', methods=['GET'])
def slide_sources_meetingprep():
    """List meeting prep sessions for slide outline source selection."""
    try:
        with get_db() as (_, cur):
            cur.execute('''
                SELECT m.id, m.meeting_date, m.meeting_type, m.status, m.notes,
                       c.ticker, c.name as company_name, c.sector,
                       (SELECT COUNT(*) FROM mp_documents WHERE meeting_id = m.id) as doc_count
                FROM mp_meetings m
                JOIN mp_companies c ON m.company_id = c.id
                ORDER BY m.created_at DESC
            ''')
            rows = cur.fetchall()
        return jsonify([{
            'id': row['id'],
            'name': f"{row['company_name']} ({row['ticker']})" + (f" - {row['meeting_type']}" if row.get('meeting_type') else ''),
            'ticker': row['ticker'],
            'company_name': row['company_name'],
            'date': str(row['meeting_date']) if row['meeting_date'] else None,
            'doc_count': row['doc_count'],
            'snippet': (row['notes'] or '')[:150] + ('...' if row.get('notes') and len(row['notes']) > 150 else ''),
        } for row in rows])
    except Exception as e:
        print(f"Error listing meeting prep sources: {e}")
        return jsonify([])


# ============================================
# STUDIO ENDPOINTS
# ============================================

def _gather_source_content(source_config):
    """Gather source content for studio generation. Returns (context_text, metadata)."""
    source_type = source_config.get('source_type', 'custom')
    source_ids = source_config.get('source_ids', [])
    custom_text = source_config.get('custom_text', '')
    context_parts = []
    metadata = {'source_type': source_type}

    if source_type == 'research':
        if source_ids:
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(source_ids))
                cur.execute(f'SELECT id, name, smart_name, content FROM research_documents WHERE id IN ({placeholders})', source_ids)
                rows = cur.fetchall()
            for row in rows:
                doc_name = row['smart_name'] or row['name'] or f"Document {row['id']}"
                content = (row['content'] or '')[:5000]
                context_parts.append(f"Document: {doc_name}\n{content}")

    elif source_type == 'summary':
        if source_ids:
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(source_ids))
                cur.execute(f'SELECT id, title, topic, raw_notes, summary FROM meeting_summaries WHERE id IN ({placeholders})', source_ids)
                rows = cur.fetchall()
            for row in rows:
                title = row['title'] or f"Summary {row['id']}"
                context_parts.append(f"Meeting: {title}")
                if row.get('topic'):
                    context_parts.append(f"Topic: {row['topic']}")
                if row.get('raw_notes'):
                    context_parts.append(f"Notes:\n{row['raw_notes'][:3000]}")
                if row.get('summary'):
                    context_parts.append(f"Summary:\n{row['summary'][:3000]}")

    elif source_type == 'meetingprep':
        if source_ids:
            with get_db() as (conn, cur):
                placeholders = ','.join(['%s'] * len(source_ids))
                cur.execute(f'''SELECT m.id, m.notes, c.ticker, c.name as company_name, c.sector
                               FROM mp_meetings m JOIN mp_companies c ON m.company_id = c.id
                               WHERE m.id IN ({placeholders})''', source_ids)
                meetings = cur.fetchall()
                cur.execute(f'SELECT meeting_id, filename, extracted_text FROM mp_documents WHERE meeting_id IN ({placeholders}) ORDER BY upload_order', source_ids)
                docs = cur.fetchall()
            for meeting in meetings:
                context_parts.append(f"Meeting Prep: {meeting['company_name']} ({meeting['ticker']})")
                if meeting.get('notes'):
                    context_parts.append(f"Meeting Notes:\n{meeting['notes'][:2000]}")
            for doc in docs:
                if doc.get('extracted_text'):
                    context_parts.append(f"Document ({doc['filename']}):\n{doc['extracted_text'][:3000]}")

    elif source_type == 'custom':
        if custom_text.strip():
            context_parts.append(custom_text[:10000])

    elif source_type == 'upload':
        # Handle uploaded files — extract text from binary formats
        if custom_text.strip():
            context_parts.append(custom_text[:10000])
        uploaded_files = source_config.get('uploaded_files', [])
        for uf in uploaded_files:
            fname = uf.get('filename', '')
            ftype = uf.get('type', '')
            fdata = uf.get('data', '')
            if not fdata:
                continue
            try:
                raw_bytes = base64.b64decode(fdata)
                if ftype in ('docx', 'doc') or fname.lower().endswith(('.docx', '.doc')):
                    # Extract text from DOCX using zipfile + XML
                    import zipfile as _zf
                    from xml.etree import ElementTree as _ET
                    text_parts = []
                    buf = io.BytesIO(raw_bytes)
                    try:
                        with _zf.ZipFile(buf) as z:
                            if 'word/document.xml' in z.namelist():
                                tree = _ET.parse(z.open('word/document.xml'))
                                ns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
                                for p in tree.iter(f'{ns}p'):
                                    texts = [t.text for t in p.iter(f'{ns}t') if t.text]
                                    if texts:
                                        text_parts.append(''.join(texts))
                    except Exception as ze:
                        text_parts.append(f'[Could not extract DOCX text: {ze}]')
                    if text_parts:
                        # Clean up encoding artifacts
                        text_parts = [t.replace('\\', '').replace('\x00', '') for t in text_parts]
                        context_parts.append(f"Document: {fname}\n" + '\n'.join(text_parts))
                elif ftype in ('xlsx', 'xls') or fname.lower().endswith(('.xlsx', '.xls')):
                    # Extract text from Excel
                    try:
                        import openpyxl
                        buf = io.BytesIO(raw_bytes)
                        wb = openpyxl.load_workbook(buf, data_only=True, read_only=True)
                        text_parts = []
                        for sheet_name in wb.sheetnames[:5]:
                            ws = wb[sheet_name]
                            text_parts.append(f"\n=== Sheet: {sheet_name} ===")
                            for row in ws.iter_rows(values_only=True):
                                row_text = '\t'.join(str(c) if c is not None else '' for c in row)
                                if row_text.strip():
                                    text_parts.append(row_text)
                        wb.close()
                        context_parts.append(f"Spreadsheet: {fname}\n" + '\n'.join(text_parts[:200]))
                    except Exception:
                        context_parts.append(f'[Spreadsheet: {fname} — could not extract text]')
                elif ftype == 'pdf' or fname.lower().endswith('.pdf'):
                    # PDFs — just note it; Claude can read PDF natively if sent as document block
                    context_parts.append(f"[PDF uploaded: {fname}]")
                    metadata['has_pdf'] = True
                    if 'pdf_files' not in metadata:
                        metadata['pdf_files'] = []
                    metadata['pdf_files'].append({'filename': fname, 'data': fdata})
                elif ftype == 'image' or fname.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp')):
                    context_parts.append(f"[Image uploaded: {fname}]")
                    if 'image_files' not in metadata:
                        metadata['image_files'] = []
                    metadata['image_files'].append({'filename': fname, 'data': fdata, 'mimeType': uf.get('mimeType', 'image/png')})
                else:
                    context_parts.append(f"[File uploaded: {fname}]")
            except Exception as e:
                context_parts.append(f"[Error processing {fname}: {str(e)[:100]}]")

    elif source_type == 'ticker':
        ticker = source_config.get('ticker', '').upper()
        if ticker:
            with get_db() as (conn, cur):
                cur.execute('SELECT * FROM stock_overviews WHERE ticker = %s', (ticker,))
                overview_row = cur.fetchone()
                cur.execute('SELECT * FROM portfolio_analyses WHERE ticker = %s', (ticker,))
                analysis_row = cur.fetchone()
            if overview_row:
                for field in ['company_overview', 'business_model', 'business_mix', 'opportunities', 'risks', 'conclusion']:
                    if overview_row.get(field):
                        context_parts.append(f"{field.replace('_', ' ').title()}:\n{overview_row[field]}")
            if analysis_row and analysis_row.get('analysis'):
                a = analysis_row['analysis']
                if isinstance(a, str):
                    try: a = json.loads(a)
                    except: a = {}
                if isinstance(a, dict):
                    context_parts.append(f"Investment Analysis:\n{json.dumps(a, indent=2)[:3000]}")

    return '\n\n'.join(context_parts), metadata


@app.route('/api/studio/init', methods=['POST'])
def init_studio_tables():
    """Manually create studio tables if they don't exist."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                CREATE TABLE IF NOT EXISTS studio_design_themes (
                    id SERIAL PRIMARY KEY,
                    name VARCHAR(255) NOT NULL,
                    description TEXT,
                    style_prompt TEXT NOT NULL,
                    preview_image TEXT,
                    is_default BOOLEAN DEFAULT FALSE,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS studio_outputs (
                    id SERIAL PRIMARY KEY,
                    title VARCHAR(500) NOT NULL,
                    type VARCHAR(30) NOT NULL,
                    status VARCHAR(50) DEFAULT 'pending',
                    theme_id INTEGER REFERENCES studio_design_themes(id) ON DELETE SET NULL,
                    source_config JSONB DEFAULT '{}',
                    settings JSONB DEFAULT '{}',
                    content JSONB DEFAULT '{}',
                    image_data TEXT,
                    progress_current INTEGER DEFAULT 0,
                    progress_total INTEGER DEFAULT 0,
                    error_message TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('''
                CREATE TABLE IF NOT EXISTS studio_slide_images (
                    id SERIAL PRIMARY KEY,
                    output_id INTEGER REFERENCES studio_outputs(id) ON DELETE CASCADE,
                    slide_number INTEGER NOT NULL,
                    image_data TEXT,
                    content_hash VARCHAR(64),
                    status VARCHAR(20) DEFAULT 'new',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            cur.execute('CREATE INDEX IF NOT EXISTS idx_studio_slide_images_output ON studio_slide_images(output_id)')
            # Seed default themes
            cur.execute('SELECT COUNT(*) as cnt FROM studio_design_themes')
            if cur.fetchone()['cnt'] == 0:
                for key, theme in SLIDE_THEMES.items():
                    cur.execute('''
                        INSERT INTO studio_design_themes (name, description, style_prompt, is_default)
                        VALUES (%s, %s, %s, %s)
                    ''', (theme['name'], theme.get('illustration_guidance', ''), theme['style_prefix'], True))
        return jsonify({'success': True, 'message': 'Studio tables created'})
    except Exception as e:
        print(f"Studio init error: {e}")
        return jsonify({'error': str(e)}), 500


STUDIO_DESIGN_THEMES = [
    {'name': 'Sketchnote', 'category': 'classic', 'colors': '#F5E6D3,#FF6B6B,#4ECDC4', 'style_prompt': "Background: Warm beige/cream colored textured paper (like aged notebook paper)\nIllustrations: Cute hand-drawn cartoon style with colorful doodles\nTitle text: Large, colorful, hand-lettered style typography (rounded, playful, multicolor gradients)\nBody text: Clean, readable text in dark gray/brown\nDecorations: Small stars, sparkles, arrows, dots, swirls scattered around\nLayout: Professional yet friendly, like a designer's sketchnote\nAspect ratio: 16:9 widescreen slide\nResolution: High quality, crisp text\nALL text MUST be in English\nDO NOT include any watermarks or AI generation notices"},
    {'name': 'Frosted Glass', 'category': 'corporate', 'colors': '#0F172A,#38BDF8,#2DD4BF', 'style_prompt': "Background: Dark navy (#0F172A) with soft gradient orbs of blue and purple floating behind frosted-glass panels\nLayout: Content inside semi-transparent white glass cards with backdrop blur and subtle 1px white/10% borders, rounded corners 16px\nTitle text: Clean white sans-serif typography (like Inter or SF Pro), bold weight\nBody text: Light gray (#CBD5E1) sans-serif\nAccent colors: Cool blue (#38BDF8) and teal (#2DD4BF) for highlights and decorative elements\nDecorations: Subtle glassmorphism layers, thin divider lines at 20% white opacity\nAspect ratio: 16:9 widescreen\nALL text MUST be in English\nDO NOT include watermarks"},
    {'name': 'Bold Editorial', 'category': 'creative', 'colors': '#FFFFFF,#1A1A1A,#DC2626', 'style_prompt': "Background: Stark white or off-white\nTitle text: Dramatic oversized serif typography (like Playfair Display) at 80-120pt in solid black, occupying 50-60% of slide area\nBody text: Light-weight sans-serif in dark gray\nAccent color: Single bold vermillion red (#DC2626) used sparingly for underlines, pull quotes, thin rules\nLayout: Strong asymmetric grid, text at unexpected positions, generous whitespace\nDecorations: Thin horizontal rules, high-contrast black-and-white elements\nPhotography style: High-contrast with tight crops\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Aurora Gradient', 'category': 'creative', 'colors': '#7C3AED,#EC4899,#14B8A6', 'style_prompt': "Background: Rich flowing mesh gradients blending deep violet (#7C3AED), magenta (#EC4899), teal (#14B8A6), and warm coral -- reminiscent of northern lights, with soft organic blobs of color\nTitle text: Clean white sans-serif bold typography with subtle text shadow for legibility\nBody text: White or very light gray sans-serif\nContent cards: Semi-transparent dark overlays (rgba black 30-40%) with soft rounded corners\nDecorations: Thin white accent lines, minimal iconography, optional subtle grain texture\nAspect ratio: 16:9 widescreen\nALL text MUST be in English\nDO NOT include watermarks"},
    {'name': 'Zen Minimal', 'category': 'minimalist', 'colors': '#F5F0EB,#8B9467,#6B7280', 'style_prompt': "Background: Warm off-white or soft cream (#F5F0EB) with subtle paper or linen texture, 60-70% intentionally blank\nTitle text: Elegant thin-weight serif typography (like Cormorant Garamond) in dark charcoal, never pure black\nBody text: Humanist sans-serif in warm gray\nPalette: Muted earth tones -- warm stone gray, soft clay, sage green (#8B9467), pale sand\nDecorations: Single carefully placed visual element per slide -- delicate ink-wash illustration or minimal line drawing, thin hairline rules\nLayout: Generous margins, radical restraint, Japanese wabi-sabi inspired\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Cyberpunk Neon', 'category': 'dark', 'colors': '#0A0A0A,#00F0FF,#FF006E', 'style_prompt': "Background: True black (#0A0A0A) or very dark blue-black (#0D1117) with subtle grid/dot matrix pattern at 5-8% opacity\nTitle text: Geometric sans-serif or monospace typography (like Space Grotesk) in neon cyan (#00F0FF) or white, with glow effects\nBody text: White or light gray monospace\nAccent colors: Electric cyan (#00F0FF), hot magenta (#FF006E), acid green (#39FF14)\nDecorations: Glowing neon line effects (box shadows with color spread), angular sharp-cornered containers\nData viz: Neon colors against dark canvas\nAspect ratio: 16:9 widescreen\nALL text MUST be in English\nDO NOT include watermarks"},
    {'name': 'Marble Luxe', 'category': 'luxury', 'colors': '#1A1A1A,#C9A96E,#FFFFFF', 'style_prompt': "Background: Dark charcoal or black marble texture with sophisticated gray and gold veining\nTitle text: Elegant high-contrast serif typography (like Didot or Bodoni) in metallic gold (#C9A96E) or white\nBody text: Light serif in white or soft gold\nAccent: Metallic gold for borders, thin decorative lines, icons, and flourishes\nDecorations: Thin gold line dividers and frame borders, generous margins, centered layouts\nContent areas: Subtle drop shadows, premium luxury feel\nAspect ratio: 16:9 widescreen\nALL text MUST be in English\nDO NOT include watermarks"},
    {'name': 'Neobrutalist', 'category': 'playful', 'colors': '#FFE600,#FF5CBE,#3B82F6', 'style_prompt': "Background: Bright saturated colors -- primary yellow (#FFE600), hot pink, electric blue -- changing per slide section\nTitle text: Bold chunky sans-serif (like Space Grotesk Black) in black or white\nBody text: Clean sans-serif in black\nBorders: Thick 3-4px black borders around ALL elements\nShadows: Hard-edge solid drop shadows offset 4-6px in black behind content cards\nDecorations: Geometric shapes (circles, rectangles, zigzag lines), NO gradients, NO blur -- everything flat and intentional\nLayout: Raw, unpolished, high contrast, generous padding\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Botanical', 'category': 'organic', 'colors': '#B7C4A0,#FAF7F2,#C2622D', 'style_prompt': "Background: Warm linen white or very pale sage (#FAF7F2)\nTitle text: Elegant serif (like Lora or DM Serif Display) in dark forest green or charcoal\nBody text: Clean humanist sans-serif in warm gray\nPalette: Soft sage green (#B7C4A0), warm cream, muted forest tones, terracotta accents (#C2622D)\nDecorations: Delicate botanical line illustrations -- leaves, ferns, eucalyptus branches, wildflowers as frames and corner accents. Occasional watercolor-wash texture in pale green or blush\nLayout: Open and airy with organic, slightly asymmetric placement, soft rounded containers\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Data Dashboard', 'category': 'analytical', 'colors': '#F8F9FA,#2563EB,#F59E0B', 'style_prompt': "Background: Clean white or very light gray (#F8F9FA) with structured multi-panel grid layout\nTitle text: Technical sans-serif (like IBM Plex Sans or Inter) bold in dark gray\nBody text: Clean sans-serif, data labels in monospace\nPrimary accent: Confident blue (#2563EB) with secondary teal, amber (#F59E0B), and coral for data differentiation\nLayout: Slide divided into 2-4 content zones separated by thin gray lines or subtle card boundaries\nDecorations: Prominent donut charts, bar graphs, sparklines, key metric callouts with large numbers, subtle gray icons, cards with very light 2px shadows\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Retro Analog', 'category': 'retro', 'colors': '#CC6B2C,#D4A843,#6B7F3E', 'style_prompt': "Background: Slightly textured surface mimicking aged paper or cardboard, warm cream (#FFF3E0)\nTitle text: Retro slab-serif typography (like Rockwell or Cooper) with rounded, friendly letterforms in burnt orange or brown\nBody text: Warm sans-serif in dark brown\nPalette: Burnt orange (#CC6B2C), mustard yellow (#D4A843), avocado green (#6B7F3E), warm brown, cream\nDecorations: Halftone dot patterns, subtle noise overlays for vintage texture, rounded rectangles and circles, hand-drawn or screen-printed illustration style\nPhotography: Warm, slightly faded vintage color grading\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Tech Futurist', 'category': 'tech', 'colors': '#0F172A,#3B82F6,#E2E8F0', 'style_prompt': "Background: Deep space navy (#0F172A) with ultra-thin geometric line art: wireframe grids, concentric circles, angular connector lines, node-network patterns in low-opacity cyan or silver\nTitle text: Clean geometric sans-serif in white, tracking slightly wider than normal\nBody text: Light sans-serif in silver/light gray (#E2E8F0)\nAccent: Electric blue (#3B82F6) and white\nDecorations: Thin line borders, micro-detail decorative elements (small crosses, dots at intersections, coordinate markers), subtle blue glow on key elements\nLayout: Mathematically precise alignment, clean structured blocks\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Pastel Cloud', 'category': 'soft', 'colors': '#E8DEFF,#DBEAFE,#FDE8EF', 'style_prompt': "Background: Soft diffused gradient blending pastel lavender (#E8DEFF), baby blue (#DBEAFE), blush pink (#FDE8EF), mint (#D1FAE5) -- gentle and cloudlike\nTitle text: Clean rounded sans-serif (like Nunito or Quicksand) medium-weight in dark gray or soft navy\nBody text: Regular weight sans-serif in medium gray\nContent cards: White or frosted-white with 16-20px border radius and very soft shadows\nDecorations: Rounded pill-shaped buttons and tags, simple friendly line icons with rounded caps\nLayout: Soft, accessible, modern, generous spacing\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Swiss Mono', 'category': 'minimalist', 'colors': '#000000,#FFFFFF,#6B7280', 'style_prompt': "Background: Pure white\nTitle text: Bold sans-serif (like Helvetica Neue Bold or Archivo Black) in pure black, large and confident\nBody text: Light-weight sans-serif in dark gray\nPalette: STRICTLY black, white, and 3-4 shades of gray -- NO color whatsoever\nLayout: Strong Swiss/International typographic grid with precise 12-column alignment, heavy whitespace\nDecorations: Thick horizontal rules, bold section dividers, large slide numbers as typographic elements\nPhotography: True black and white, high contrast\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Warm Terracotta', 'category': 'warm', 'colors': '#C2622D,#FAF7F2,#8B4513', 'style_prompt': "Background: Soft white (#FAF7F2) or light sand with subtle linen or canvas texture\nTitle text: Rounded approachable serif (like Fraunces) in rich terracotta (#C2622D) or deep clay\nBody text: Modern geometric sans-serif in warm charcoal (#3D3024)\nPalette: Terracotta/burnt sienna, deep clay (#8B4513), warm charcoal, dusty rose, sage green\nDecorations: Organic soft shapes -- irregular blobs, arched frames, rounded containers. Thin terracotta accent lines and dots\nPhotography: Warm sun-kissed color grading\nLayout: Balanced but relaxed, generous breathing room\nAspect ratio: 16:9 widescreen\nALL text MUST be in English"},
    {'name': 'Holographic', 'category': 'futuristic', 'colors': '#E0E7FF,#F0ABFC,#67E8F9', 'style_prompt': "Background: Light silver-white or soft gray base with holographic/iridescent color shifts -- rainbow refractions of pink (#F0ABFC), blue, green, purple shimmering like holographic foil\nTitle text: Clean modern sans-serif in dark charcoal or black for contrast, bold weight\nBody text: Medium sans-serif in dark gray\nAccent: Metallic silver and chrome for borders and decorative lines\nContent cards: Iridescent gradient borders or holographic background fills at low opacity\nDecorations: Subtle light-leak or lens-flare effects, minimalist centered layout with ample whitespace\nAspect ratio: 16:9 widescreen\nALL text MUST be in English\nDO NOT include watermarks"},
]


@app.route('/api/studio/seed-themes', methods=['POST'])
def seed_studio_themes():
    """Seed the expanded studio design themes."""
    try:
        with get_db(commit=True) as (conn, cur):
            # Clear existing themes
            cur.execute('DELETE FROM studio_design_themes')
            for t in STUDIO_DESIGN_THEMES:
                cur.execute('''
                    INSERT INTO studio_design_themes (name, description, style_prompt, is_default)
                    VALUES (%s, %s, %s, TRUE)
                ''', (t['name'], t.get('category', ''), t['style_prompt']))
        return jsonify({'success': True, 'count': len(STUDIO_DESIGN_THEMES)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- Studio Theme CRUD ---

_THEME_COLORS = {t['name']: t['colors'] for t in STUDIO_DESIGN_THEMES}

@app.route('/api/studio/themes', methods=['GET'])
def get_studio_themes():
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT id, name, description, style_prompt, is_default, preview_image IS NOT NULL as has_preview, created_at FROM studio_design_themes ORDER BY is_default DESC, name ASC')
            rows = cur.fetchall()
        result = []
        for row in rows:
            d = dict(row)
            d['colors'] = _THEME_COLORS.get(row['name'], '')
            result.append(d)
        return jsonify(result)
    except Exception as e:
        print(f"Error listing studio themes: {e}")
        return jsonify([])


@app.route('/api/studio/themes', methods=['POST'])
def create_studio_theme():
    try:
        data = request.json
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO studio_design_themes (name, description, style_prompt)
                VALUES (%s, %s, %s) RETURNING id
            ''', (data['name'], data.get('description', ''), data['style_prompt']))
            theme_id = cur.fetchone()['id']
        return jsonify({'id': theme_id, 'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/themes/<int:theme_id>', methods=['PUT'])
def update_studio_theme(theme_id):
    try:
        data = request.json
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                UPDATE studio_design_themes SET name=%s, description=%s, style_prompt=%s, updated_at=%s
                WHERE id=%s
            ''', (data['name'], data.get('description', ''), data['style_prompt'], datetime.utcnow(), theme_id))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/themes/<int:theme_id>', methods=['DELETE'])
def delete_studio_theme(theme_id):
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('SELECT is_default FROM studio_design_themes WHERE id=%s', (theme_id,))
            row = cur.fetchone()
            if row and row['is_default']:
                return jsonify({'error': 'Cannot delete default themes'}), 400
            cur.execute('DELETE FROM studio_design_themes WHERE id=%s AND is_default=FALSE', (theme_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/themes/<int:theme_id>/preview', methods=['POST'])
def generate_studio_theme_preview(theme_id):
    try:
        data = request.json or {}
        gemini_key = data.get('geminiApiKey') or os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY', '')
        with get_db() as (_, cur):
            cur.execute('SELECT style_prompt, name FROM studio_design_themes WHERE id=%s', (theme_id,))
            theme = cur.fetchone()
        if not theme:
            return jsonify({'error': 'Theme not found'}), 404
        prompt = f"Generate a sample presentation slide preview showcasing this visual style:\n\n{theme['style_prompt']}\n\nCreate a sample slide with the title '{theme['name']} Theme Preview' and some placeholder content demonstrating the visual style."
        image_data = _generate_slide_image(prompt, api_key=gemini_key)
        if image_data:
            with get_db(commit=True) as (conn, cur):
                cur.execute('UPDATE studio_design_themes SET preview_image=%s, updated_at=%s WHERE id=%s', (image_data, datetime.utcnow(), theme_id))
            return jsonify({'success': True})
        return jsonify({'error': 'Failed to generate preview image'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- Studio Output CRUD ---

@app.route('/api/studio/outputs', methods=['GET'])
def get_studio_outputs():
    try:
        type_filter = request.args.get('type')
        with get_db() as (_, cur):
            query = 'SELECT id, title, type, status, theme_id, source_config, settings, progress_current, progress_total, error_message, created_at, updated_at FROM studio_outputs'
            params = []
            if type_filter:
                query += ' WHERE type = %s'
                params.append(type_filter)
            query += ' ORDER BY created_at DESC'
            cur.execute(query, params)
            rows = cur.fetchall()
        return jsonify([dict(row) for row in rows])
    except Exception as e:
        print(f"Error listing studio outputs: {e}")
        return jsonify([])


@app.route('/api/studio/outputs', methods=['POST'])
def create_studio_output():
    try:
        data = request.json
        with get_db(commit=True) as (conn, cur):
            cur.execute('''
                INSERT INTO studio_outputs (title, type, source_config, settings, theme_id)
                VALUES (%s, %s, %s, %s, %s) RETURNING id
            ''', (
                data['title'], data['type'],
                json.dumps(data.get('source_config', {})),
                json.dumps(data.get('settings', {})),
                data.get('theme_id')
            ))
            output_id = cur.fetchone()['id']
        return jsonify({'id': output_id, 'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>', methods=['GET'])
def get_studio_output(output_id):
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM studio_outputs WHERE id=%s', (output_id,))
            row = cur.fetchone()
            if not row:
                return jsonify({'error': 'Output not found'}), 404
            result = dict(row)
            # Don't send image_data in the main response (it's large)
            result['image_data'] = None
            result['has_image'] = row['image_data'] is not None
            # For slides, get slide images metadata
            if row['type'] == 'slides':
                cur.execute('SELECT id, slide_number, status, content_hash, image_data IS NOT NULL as has_image FROM studio_slide_images WHERE output_id=%s ORDER BY slide_number', (output_id,))
                result['slide_images'] = [dict(si) for si in cur.fetchall()]
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>/image', methods=['GET'])
def get_studio_output_image(output_id):
    try:
        slide_num = request.args.get('slide', type=int)
        with get_db() as (_, cur):
            if slide_num is not None:
                cur.execute('SELECT image_data FROM studio_slide_images WHERE output_id=%s AND slide_number=%s', (output_id, slide_num))
            else:
                cur.execute('SELECT image_data FROM studio_outputs WHERE id=%s', (output_id,))
            row = cur.fetchone()
        if row and row['image_data']:
            return jsonify({'image_data': row['image_data']})
        return jsonify({'image_data': None}), 404
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>', methods=['DELETE'])
def delete_studio_output(output_id):
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute('DELETE FROM studio_outputs WHERE id=%s', (output_id,))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- Studio Generation ---

def _generate_studio_slides(output_id, source_content, settings, api_keys):
    """Generate slides (outline + images) in background thread."""
    import time
    try:
        slide_count = settings.get('slide_count', 30)
        theme_name = settings.get('theme_name', 'sketchnote')

        # Look up studio design theme if one is set
        custom_theme_prompt = ''
        with get_db(commit=True) as (conn, cur):
            cur.execute("SELECT theme_id FROM studio_outputs WHERE id=%s", (output_id,))
            out_row = cur.fetchone()
            if out_row and out_row.get('theme_id'):
                cur.execute("SELECT name, style_prompt FROM studio_design_themes WHERE id=%s", (out_row['theme_id'],))
                theme_row = cur.fetchone()
                if theme_row:
                    theme_name = theme_row['name'].lower().replace(' ', '')
                    custom_theme_prompt = theme_row['style_prompt']
            cur.execute("UPDATE studio_outputs SET status='generating', progress_current=0, progress_total=%s, updated_at=%s WHERE id=%s",
                        (slide_count, datetime.utcnow(), output_id))

        # Phase 1: Generate outline
        llm_prompt = f"""You are creating slide content for a presentation.

Based on this source material:

{source_content}

Generate a JSON array of exactly {slide_count} slides for a comprehensive presentation. Each slide should have:
- slide_number (integer starting at 1)
- title (string)
- type (one of: title, toc, section_divider, content, closing)
- content (detailed text describing what should appear on the slide - include specific data, numbers, quotes from the source material)
- illustration_hints (array of strings like "company", "growth", "money", "risk", "data", "leader", "opportunity")
- no_header (boolean - true only for title, section_divider, and closing slides)

Structure: title slide, table of contents, then 5-7 sections covering the key themes. Each section should have a section divider slide followed by 2-4 content slides. End with a closing slide.

IMPORTANT: Fill in actual data and specific details from the source material. Do NOT use placeholder brackets.

Return ONLY valid JSON array, no markdown fencing."""

        result = call_llm(
            messages=[{"role": "user", "content": llm_prompt}],
            tier="standard", max_tokens=8192,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )
        response_text = result['text'].strip()
        if response_text.startswith('```'):
            response_text = response_text.split('\n', 1)[1].rsplit('```', 1)[0].strip()
        slides_data = json.loads(response_text)

        # Store outline in content
        style_instructions = settings.get('style_instructions', '')
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET content=%s, progress_total=%s, updated_at=%s WHERE id=%s",
                        (json.dumps({'slides': slides_data, 'theme_name': theme_name, 'style_instructions': style_instructions, 'custom_theme_prompt': custom_theme_prompt}), len(slides_data), datetime.utcnow(), output_id))
            # Create slide image records
            for s in slides_data:
                content_hash = _compute_content_hash(s)
                cur.execute('''
                    INSERT INTO studio_slide_images (output_id, slide_number, content_hash, status)
                    VALUES (%s, %s, %s, 'new')
                ''', (output_id, s['slide_number'], content_hash))

        # Phase 2: Generate images
        gemini_key = api_keys.get('gemini', '')
        for i, slide in enumerate(slides_data):
            if custom_theme_prompt:
                # Use the studio design theme directly instead of built-in SLIDE_THEMES
                prompt = f"You are generating a presentation slide image.\n\nVISUAL STYLE (MUST follow exactly):\n{custom_theme_prompt}\n\nSLIDE CONTENT:\n{slide.get('content', '').strip()}"
                hints = slide.get('illustration_hints', [])
                if hints:
                    prompt += "\n\nILLUSTRATIONS TO INCLUDE:\n" + "\n".join(f"- {h}" for h in hints)
            else:
                prompt = _build_slide_prompt(slide, theme_name, '', len(slides_data))
            if style_instructions:
                prompt += f"\n\nADDITIONAL STYLE INSTRUCTIONS (follow these closely):\n{style_instructions}"
            image_data = _generate_slide_image(prompt, api_key=gemini_key)
            with get_db(commit=True) as (conn, cur):
                if image_data:
                    cur.execute("UPDATE studio_slide_images SET image_data=%s, status='done' WHERE output_id=%s AND slide_number=%s",
                                (image_data, output_id, slide['slide_number']))
                else:
                    cur.execute("UPDATE studio_slide_images SET status='error' WHERE output_id=%s AND slide_number=%s",
                                (output_id, slide['slide_number']))
                cur.execute("UPDATE studio_outputs SET progress_current=%s, status=%s, updated_at=%s WHERE id=%s",
                            (i + 1, f"generating:{i+1}:{len(slides_data)}", datetime.utcnow(), output_id))
            time.sleep(2)

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='ready', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio slides generation error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


def _generate_html_infographic(output_id, source_content, settings, api_keys):
    """Generate an infographic using HTML rendering for precise layout.

    Uses the source content structure faithfully, renders as styled HTML,
    then converts to PDF via xhtml2pdf. Returns HTML + PDF as base64.
    """
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='generating', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

        density = settings.get('density', 'balanced')
        style_instructions = settings.get('style_instructions', '')
        title = settings.get('title', 'Infographic')

        density_css = {
            'clean': 'font-size: 16px; line-height: 2; max-width: 1000px;',
            'balanced': 'font-size: 13px; line-height: 1.6; max-width: 1200px;',
            'dense': 'font-size: 11px; line-height: 1.4; max-width: 1400px;',
        }.get(density, 'font-size: 13px; line-height: 1.6;')

        prompt = f"""Convert the following SOURCE CONTENT into a beautiful, well-styled HTML infographic.

ABSOLUTE RULES:
1. Use ONLY the data provided in the SOURCE CONTENT below. Do NOT add ANY information from your training data or general knowledge.
2. If something is not explicitly stated in the source content, do NOT include it.
3. PRESERVE the EXACT structure, order, and organization of the source content.
4. If the source is a chronological table/list, render it as a styled HTML table or timeline in the SAME order.
5. Every item in the source must appear in the output in the same position.
6. Do NOT reorganize, summarize, or create new groupings that are not in the source.
7. Do NOT make up or infer additional context, history, or background information.
8. Make it visually beautiful with CSS styling - colors, borders, icons (use Unicode), backgrounds.
9. All text must be perfectly spelled and readable.
10. Do NOT use CSS features unsupported by xhtml2pdf: no border-radius, no gradients, no flexbox, no calc(). Use table layouts with explicit width on every td/th.
11. Use only simple CSS: solid borders, background-color (hex), padding, margin, text-align, font-weight, font-size, color.

STYLE: {style_instructions or 'Professional, clean design with a dark header and well-organized sections'}
DENSITY: {density} layout

Return ONLY the complete HTML (no markdown fences, no explanation). The HTML should be a self-contained document with inline CSS.
Use this base structure:
<html>
<head><style>
body {{ margin: 0; padding: 40px; font-family: Arial, sans-serif; background: #ffffff; color: #1a1a2e; {density_css} }}
h1 {{ text-align: center; color: white; background-color: #1a1a2e; padding: 24px; margin: -40px -40px 30px -40px; font-size: 28px; }}
table {{ width: 100%; border-collapse: collapse; margin-bottom: 20px; }}
th {{ background-color: #16213e; color: white; padding: 10px; text-align: left; }}
td {{ padding: 8px; border-bottom: 1px solid #cccccc; }}
</style></head>
<body>
<!-- Faithful rendering of source content -->
</body>
</html>

SOURCE CONTENT:
{title}

{source_content[:12000]}"""

        result = call_llm(
            messages=[{"role": "user", "content": prompt}],
            system="You are an expert HTML/CSS designer who creates beautiful infographic-style documents. Return ONLY valid HTML. Do NOT use border-radius, gradients, flexbox, or calc() — only simple CSS compatible with xhtml2pdf.",
            tier="standard",
            max_tokens=16384,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )

        html_text = result['text'].strip()
        # Strip markdown fences if present
        if html_text.startswith('```'):
            html_text = html_text.split('\n', 1)[1].rsplit('```', 1)[0]
        if html_text.startswith('html'):
            html_text = html_text[4:].strip()

        # Convert HTML to PDF using xhtml2pdf
        from xhtml2pdf import pisa
        import io

        pdf_buf = io.BytesIO()
        pisa_status = pisa.CreatePDF(html_text, dest=pdf_buf)
        pdf_bytes = pdf_buf.getvalue()

        html_b64 = base64.b64encode(html_text.encode('utf-8')).decode('ascii')
        pdf_b64 = base64.b64encode(pdf_bytes).decode('ascii') if pdf_bytes else ''

        content_json = json.dumps({
            'render_mode': 'precise',
            'html': html_b64,
            'pdf': pdf_b64,
            'format': 'html',
            'density': density,
        })

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='ready', content=%s, updated_at=%s WHERE id=%s",
                        (content_json, datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio HTML infographic error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


def _generate_studio_infographic(output_id, source_content, settings, api_keys):
    """Generate an infographic image."""
    try:
        # Check render mode — dispatch to HTML renderer if 'precise'
        render_mode = settings.get('render_mode', 'ai')
        if render_mode == 'precise':
            return _generate_html_infographic(output_id, source_content, settings, api_keys)

        orientation = settings.get('orientation', 'landscape')
        aspect_map = {'landscape': '16:9', 'portrait': '9:16', 'square': '1:1'}
        aspect = aspect_map.get(orientation, '16:9')

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='generating', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

        # Get summary points from LLM
        result = call_llm(
            messages=[{"role": "user", "content": f"Summarize the following content into key data points, statistics, and insights suitable for an infographic. Format as a clear, visual-friendly summary:\n\n{source_content[:6000]}"}],
            tier="standard", max_tokens=2048,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )
        summary = result['text'].strip()

        # Get theme style
        theme_style = ''
        with get_db() as (_, cur):
            cur.execute('SELECT so.theme_id, sdt.style_prompt FROM studio_outputs so LEFT JOIN studio_design_themes sdt ON so.theme_id = sdt.id WHERE so.id=%s', (output_id,))
            row = cur.fetchone()
            if row and row.get('style_prompt'):
                theme_style = row['style_prompt']

        # Layout density instruction
        density = settings.get('density', 'balanced')
        density_instruction = {
            'clean': 'Use a MINIMAL, clean layout with lots of white space. Maximum 5-6 key data points. Large fonts, simple icons. Less is more.',
            'balanced': 'Use a balanced layout with moderate information density. Clear sections with readable text.',
            'dense': 'Pack as much information as possible into the infographic. Use small fonts, multiple columns, and detailed data tables. Maximize information density.',
        }.get(density, '')

        # Preserve structure instruction
        preserve_structure = settings.get('preserve_structure', False)
        if preserve_structure:
            structure_instruction = """ABSOLUTE REQUIREMENT - PRESERVE STRUCTURE:
You MUST preserve the EXACT structure, organization, and order of the source document. This is non-negotiable.

Rules:
1. If the source is a TABLE (Date | Announcement), render it as a VISUAL TABLE or TIMELINE in the SAME chronological order. Do NOT group by category.
2. If the source has sections (January, February, March), keep those EXACT section headers in the SAME order.
3. Every row/item in the source must appear in the output in the SAME position.
4. Do NOT summarize, combine, or reorganize items. Every item from the source must be individually visible.
5. Do NOT create new groupings (e.g. Product Categories) that do not exist in the source.
6. Think of this as: take the source document layout and make it beautiful, but change NOTHING about the order or grouping.

If the source is a chronological list/table, the output MUST be a chronological visual - timeline, dated list, or visual table - NOT a category-grouped summary."""
        else:
            structure_instruction = "Organize the content in the most visually effective way for an infographic. You may reorganize, group, and prioritize information for maximum visual impact."

        prompt = f"""Create a beautiful, professional infographic image.

{f'VISUAL STYLE: {theme_style}' if theme_style else 'VISUAL STYLE: Clean, modern infographic with vibrant colors, clear hierarchy, and professional typography.'}

LAYOUT DENSITY: {density_instruction}

CONTENT ORGANIZATION: {structure_instruction}

Aspect ratio: {aspect}

CONTENT TO VISUALIZE:
{summary}

Create a visually stunning infographic that presents this information clearly with:
- A compelling title at the top
- Key statistics highlighted in large, bold numbers
- Clear sections with icons and illustrations
- A logical flow from top to bottom
- Professional color scheme
ALL text MUST be in English.

IMPORTANT: Ensure all text is spelled correctly. Do not add backslashes, special characters, or encoding artifacts to any text. Double-check company names, product names, and titles for accuracy."""

        style_instructions = settings.get('style_instructions', '')
        if style_instructions:
            prompt += f"\n\nADDITIONAL STYLE INSTRUCTIONS (follow these closely):\n{style_instructions}"

        gemini_key = api_keys.get('gemini', '')
        key = gemini_key or os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY', '')
        client = genai.Client(api_key=key)
        image_data = None
        for attempt in range(3):
            try:
                response = client.models.generate_content(
                    model="gemini-3-pro-image-preview",
                    contents=prompt,
                    config=genai_types.GenerateContentConfig(
                        response_modalities=["TEXT", "IMAGE"],
                        image_config=genai_types.ImageConfig(aspect_ratio=aspect),
                    ),
                )
                for part in response.candidates[0].content.parts:
                    if hasattr(part, "inline_data") and part.inline_data is not None:
                        image_data = base64.b64encode(part.inline_data.data).decode('utf-8')
                        break
                if image_data:
                    break
            except Exception as e:
                print(f"Infographic attempt {attempt+1} failed: {e}")
                import time; time.sleep((attempt + 1) * 5)

        with get_db(commit=True) as (conn, cur):
            if image_data:
                cur.execute("UPDATE studio_outputs SET status='ready', image_data=%s, content=%s, updated_at=%s WHERE id=%s",
                            (image_data, json.dumps({'orientation': orientation, 'prompt_used': prompt[:500], 'description': summary[:1000]}), datetime.utcnow(), output_id))
            else:
                cur.execute("UPDATE studio_outputs SET status='error', error_message='Failed to generate infographic image', updated_at=%s WHERE id=%s",
                            (datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio infographic error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


def _generate_studio_mindmap(output_id, source_content, settings, api_keys):
    """Generate a mind map JSON structure."""
    try:
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='generating', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

        result = call_llm(
            messages=[{"role": "user", "content": f"""Analyze the following content and create a hierarchical mind map structure.

Return a JSON object with this exact structure:
{{
  "root": {{
    "label": "Main Topic",
    "children": [
      {{
        "label": "Branch 1",
        "children": [
          {{ "label": "Sub-topic 1.1", "children": [] }},
          {{ "label": "Sub-topic 1.2", "children": [] }}
        ]
      }},
      {{
        "label": "Branch 2",
        "children": [...]
      }}
    ]
  }}
}}

Create 4-7 main branches, each with 2-5 sub-topics. Sub-topics can have their own children (up to 3 levels deep).
Use concise, descriptive labels (3-8 words each).

SOURCE CONTENT:
{source_content[:6000]}

Return ONLY valid JSON, no markdown fencing."""}],
            tier="standard", max_tokens=4096,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )
        response_text = result['text'].strip()
        if response_text.startswith('```'):
            response_text = response_text.split('\n', 1)[1].rsplit('```', 1)[0].strip()
        mindmap_data = json.loads(response_text)

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='ready', content=%s, updated_at=%s WHERE id=%s",
                        (json.dumps(mindmap_data), datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio mindmap error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


def _generate_studio_report(output_id, source_content, settings, api_keys):
    """Generate a markdown report."""
    try:
        length = settings.get('length', 'standard')
        length_guidance = {
            'brief': 'Keep it concise, around 500-800 words. Focus on key takeaways.',
            'standard': 'Write a thorough report, around 1500-2500 words with detailed analysis.',
            'detailed': 'Write a comprehensive, in-depth report, around 3000-5000 words with exhaustive coverage.',
        }

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='generating', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

        result = call_llm(
            messages=[{"role": "user", "content": f"""Write a professional report in Markdown format based on the following source material.

{length_guidance.get(length, length_guidance['standard'])}

Structure the report with:
- A clear title (# heading)
- Executive summary
- Main sections with ## headings
- Sub-sections with ### headings where needed
- Key findings, data points, and analysis
- Conclusion with actionable insights

Use bullet points, bold text, and other markdown formatting for readability.
Include specific numbers, quotes, and data from the source material.

SOURCE MATERIAL:
{source_content[:8000]}"""}],
            tier="standard", max_tokens=8192,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='ready', content=%s, updated_at=%s WHERE id=%s",
                        (json.dumps({'markdown': result['text'].strip()}), datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio report error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


def _generate_studio_quiz(output_id, source_content, settings, api_keys):
    """Generate quiz questions."""
    try:
        question_count = settings.get('question_count', 10)
        difficulty = settings.get('difficulty', 'medium')

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='generating', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

        result = call_llm(
            messages=[{"role": "user", "content": f"""Create {question_count} multiple-choice quiz questions based on the following content.
Difficulty level: {difficulty}

Return a JSON object with this structure:
{{
  "questions": [
    {{
      "id": 1,
      "question": "What is...",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct": 0,
      "explanation": "The correct answer is A because..."
    }}
  ]
}}

Rules:
- Each question must have exactly 4 options
- "correct" is the 0-based index of the correct answer
- Include a brief explanation for each answer
- Questions should test understanding, not just recall
- Vary question types: factual, conceptual, analytical
- Make incorrect options plausible but clearly wrong

SOURCE CONTENT:
{source_content[:6000]}

Return ONLY valid JSON, no markdown fencing."""}],
            tier="standard", max_tokens=4096,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )
        response_text = result['text'].strip()
        if response_text.startswith('```'):
            response_text = response_text.split('\n', 1)[1].rsplit('```', 1)[0].strip()
        quiz_data = json.loads(response_text)

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='ready', content=%s, updated_at=%s WHERE id=%s",
                        (json.dumps(quiz_data), datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio quiz error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


def _generate_studio_flashcard(output_id, source_content, settings, api_keys):
    """Generate flashcards."""
    try:
        card_count = settings.get('card_count', 20)

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='generating', updated_at=%s WHERE id=%s", (datetime.utcnow(), output_id))

        result = call_llm(
            messages=[{"role": "user", "content": f"""Create {card_count} flashcards based on the following content.

Return a JSON object with this structure:
{{
  "cards": [
    {{
      "id": 1,
      "front": "Question or concept to remember",
      "back": "Answer or explanation",
      "category": "Topic Category"
    }}
  ]
}}

Rules:
- Front should be a clear question or key term/concept
- Back should be a concise but complete answer
- Category should group related cards together
- Cover the most important concepts from the material
- Mix factual recall, definitions, and conceptual questions
- Keep answers concise but informative

SOURCE CONTENT:
{source_content[:6000]}

Return ONLY valid JSON, no markdown fencing."""}],
            tier="standard", max_tokens=4096,
            anthropic_api_key=api_keys.get('anthropic', ''),
            gemini_api_key=api_keys.get('gemini', ''),
        )
        response_text = result['text'].strip()
        if response_text.startswith('```'):
            response_text = response_text.split('\n', 1)[1].rsplit('```', 1)[0].strip()
        flashcard_data = json.loads(response_text)

        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='ready', content=%s, updated_at=%s WHERE id=%s",
                        (json.dumps(flashcard_data), datetime.utcnow(), output_id))

    except Exception as e:
        print(f"Studio flashcard error: {e}")
        with get_db(commit=True) as (conn, cur):
            cur.execute("UPDATE studio_outputs SET status='error', error_message=%s, updated_at=%s WHERE id=%s",
                        (str(e), datetime.utcnow(), output_id))


_STUDIO_GENERATORS = {
    'slides': _generate_studio_slides,
    'infographic': _generate_studio_infographic,
    'mindmap': _generate_studio_mindmap,
    'report': _generate_studio_report,
    'quiz': _generate_studio_quiz,
    'flashcard': _generate_studio_flashcard,
}


@app.route('/api/studio/outputs/<int:output_id>/generate', methods=['POST'])
def generate_studio_output(output_id):
    """Universal generation endpoint. Dispatches to type-specific generator in background thread."""
    try:
        data = request.json or {}
        with get_db() as (_, cur):
            cur.execute('SELECT * FROM studio_outputs WHERE id=%s', (output_id,))
            output = cur.fetchone()
        if not output:
            return jsonify({'error': 'Output not found'}), 404

        output_type = output['type']
        generator = _STUDIO_GENERATORS.get(output_type)
        if not generator:
            return jsonify({'error': f'Unknown output type: {output_type}'}), 400

        source_config = output['source_config'] if isinstance(output['source_config'], dict) else json.loads(output['source_config'] or '{}')
        settings = output['settings'] if isinstance(output['settings'], dict) else json.loads(output['settings'] or '{}')
        source_content, _ = _gather_source_content(source_config)
        if not source_content.strip():
            # Auto-generate from title alone when no source content provided
            source_content = f"Generate content about: {output['title']}"

        api_keys = {
            'anthropic': data.get('apiKey', '') or os.environ.get('ANTHROPIC_API_KEY', ''),
            'gemini': data.get('geminiApiKey', '') or os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY', ''),
        }

        import threading
        thread = threading.Thread(target=generator, args=(output_id, source_content, settings, api_keys))
        thread.daemon = True
        thread.start()

        return jsonify({'success': True, 'message': f'Generation started for {output_type}'})
    except Exception as e:
        print(f"Error starting studio generation: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>/html', methods=['GET'])
def get_studio_html(output_id):
    """Get the HTML version of a studio output for precise rendering.

    If the output was generated with render_mode='precise', returns the stored
    HTML and PDF. Otherwise, generates a new HTML infographic on the fly from
    the source content.
    """
    try:
        with get_db() as (_, cur):
            cur.execute('SELECT settings, source_config, content, title FROM studio_outputs WHERE id = %s', (output_id,))
            row = cur.fetchone()
        if not row:
            return jsonify({'error': 'Not found'}), 404

        # If already generated as HTML, return stored data
        content = row['content'] if isinstance(row['content'], dict) else json.loads(row['content'] or '{}')
        if content.get('render_mode') == 'precise' and content.get('html'):
            return jsonify({
                'html': content['html'],
                'pdf': content.get('pdf', ''),
                'format': 'html',
            })

        # Otherwise, generate on the fly from source content
        settings = row['settings'] if isinstance(row['settings'], dict) else json.loads(row['settings'] or '{}')
        source_config = row['source_config'] if isinstance(row['source_config'], dict) else json.loads(row['source_config'] or '{}')

        context_text, _ = _gather_source_content(source_config)
        if not context_text.strip():
            context_text = f"Generate content about: {row['title'] or 'Infographic'}"
        title = settings.get('title', row['title'] or 'Infographic')

        density = settings.get('density', 'balanced')
        style_instructions = settings.get('style_instructions', '')

        density_css = {
            'clean': 'font-size: 16px; line-height: 2; max-width: 1000px;',
            'balanced': 'font-size: 13px; line-height: 1.6; max-width: 1200px;',
            'dense': 'font-size: 11px; line-height: 1.4; max-width: 1400px;',
        }.get(density, 'font-size: 13px; line-height: 1.6;')

        prompt = f"""Convert the following SOURCE CONTENT into a beautiful, well-styled HTML infographic.

ABSOLUTE RULES:
1. Use ONLY the data provided in the SOURCE CONTENT below. Do NOT add ANY information from your training data or general knowledge.
2. If something is not explicitly stated in the source content, do NOT include it.
3. PRESERVE the EXACT structure, order, and organization of the source content.
4. If the source is a chronological table/list, render it as a styled HTML table or timeline in the SAME order.
5. Every item in the source must appear in the output in the same position.
6. Do NOT reorganize, summarize, or create new groupings that are not in the source.
7. Do NOT make up or infer additional context, history, or background information.
8. Make it visually beautiful with CSS styling - colors, borders, icons (use Unicode), backgrounds.
9. All text must be perfectly spelled and readable.
10. Do NOT use CSS features unsupported by xhtml2pdf: no border-radius, no gradients, no flexbox, no calc(). Use table layouts with explicit width on every td/th.
11. Use only simple CSS: solid borders, background-color (hex), padding, margin, text-align, font-weight, font-size, color.

STYLE: {style_instructions or 'Professional, clean design with a dark header and well-organized sections'}
DENSITY: {density} layout

Return ONLY the complete HTML (no markdown fences, no explanation). The HTML should be a self-contained document with inline CSS.

SOURCE CONTENT:
{title}

{context_text[:12000]}"""

        result = call_llm(
            messages=[{"role": "user", "content": prompt}],
            system="You are an expert HTML/CSS designer who creates beautiful infographic-style documents. Return ONLY valid HTML. Do NOT use border-radius, gradients, flexbox, or calc().",
            tier="standard",
            max_tokens=16384,
        )

        html_text = result['text'].strip()
        if html_text.startswith('```'):
            html_text = html_text.split('\n', 1)[1].rsplit('```', 1)[0]
        if html_text.startswith('html'):
            html_text = html_text[4:].strip()

        from xhtml2pdf import pisa
        import io

        pdf_buf = io.BytesIO()
        pisa.CreatePDF(html_text, dest=pdf_buf)
        pdf_bytes = pdf_buf.getvalue()

        html_b64 = base64.b64encode(html_text.encode('utf-8')).decode('ascii')
        pdf_b64 = base64.b64encode(pdf_bytes).decode('ascii') if pdf_bytes else ''

        return jsonify({
            'html': html_b64,
            'pdf': pdf_b64,
            'format': 'html',
        })
    except Exception as e:
        print(f"Error generating HTML infographic: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>/update-slide/<int:slide_num>', methods=['PUT'])
def update_studio_slide(output_id, slide_num):
    """Update a single slide's content (title, content, illustration_hints)."""
    try:
        data = request.json or {}
        with get_db() as (_, cur):
            cur.execute('SELECT content FROM studio_outputs WHERE id=%s', (output_id,))
            output = cur.fetchone()
        if not output:
            return jsonify({'error': 'Output not found'}), 404
        content = output['content'] if isinstance(output['content'], dict) else json.loads(output['content'] or '{}')
        slides = content.get('slides', [])
        updated = False
        for s in slides:
            if s['slide_number'] == slide_num:
                if 'title' in data:
                    s['title'] = data['title']
                if 'content' in data:
                    s['content'] = data['content']
                if 'illustration_hints' in data:
                    s['illustration_hints'] = data['illustration_hints']
                updated = True
                break
        if not updated:
            return jsonify({'error': f'Slide {slide_num} not found'}), 404
        with get_db(commit=True) as (conn, cur):
            cur.execute('UPDATE studio_outputs SET content=%s, updated_at=%s WHERE id=%s',
                        (json.dumps(content), datetime.utcnow(), output_id))
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>/regenerate-slide/<int:slide_num>', methods=['POST'])
def regenerate_studio_slide(output_id, slide_num):
    """Regenerate a single slide image with optional edit prompt."""
    try:
        data = request.json or {}
        gemini_key = data.get('geminiApiKey', '') or os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY', '')
        edit_prompt = data.get('edit_prompt', '')

        with get_db() as (_, cur):
            cur.execute('SELECT content FROM studio_outputs WHERE id=%s', (output_id,))
            output = cur.fetchone()
        if not output:
            return jsonify({'error': 'Output not found'}), 404

        content = output['content'] if isinstance(output['content'], dict) else json.loads(output['content'] or '{}')
        slides = content.get('slides', [])
        theme_name = content.get('theme_name', 'sketchnote')
        custom_theme_prompt = content.get('custom_theme_prompt', '')
        slide_data = next((s for s in slides if s['slide_number'] == slide_num), None)
        if not slide_data:
            return jsonify({'error': f'Slide {slide_num} not found'}), 404

        if custom_theme_prompt:
            prompt = f"You are generating a presentation slide image.\n\nVISUAL STYLE (MUST follow exactly):\n{custom_theme_prompt}\n\nSLIDE CONTENT:\n{slide_data.get('content', '').strip()}"
            hints = slide_data.get('illustration_hints', [])
            if hints:
                prompt += "\n\nILLUSTRATIONS TO INCLUDE:\n" + "\n".join(f"- {h}" for h in hints)
        else:
            prompt = _build_slide_prompt(slide_data, theme_name, '', len(slides))
        style_instructions = content.get('style_instructions', '')
        if style_instructions:
            prompt += f"\n\nADDITIONAL STYLE INSTRUCTIONS (follow these closely):\n{style_instructions}"
        if edit_prompt:
            prompt += f"\n\nADDITIONAL INSTRUCTIONS: {edit_prompt}"

        image_data = _generate_slide_image(prompt, api_key=gemini_key)
        if image_data:
            with get_db(commit=True) as (conn, cur):
                cur.execute("UPDATE studio_slide_images SET image_data=%s, status='done', content_hash=%s WHERE output_id=%s AND slide_number=%s",
                            (_compute_content_hash(slide_data), image_data, output_id, slide_num))
            return jsonify({'success': True, 'image_data': image_data})
        return jsonify({'error': 'Failed to generate image'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/studio/outputs/<int:output_id>/export', methods=['POST'])
def export_studio_output(output_id):
    """Export studio output as PDF or downloadable file."""
    try:
        data = request.json or {}
        fmt = data.get('format', 'pdf')

        with get_db() as (_, cur):
            cur.execute('SELECT * FROM studio_outputs WHERE id=%s', (output_id,))
            output = cur.fetchone()
        if not output:
            return jsonify({'error': 'Output not found'}), 404

        content = output['content'] if isinstance(output['content'], dict) else json.loads(output['content'] or '{}')

        if output['type'] == 'slides' and fmt == 'pdf':
            with get_db() as (_, cur):
                cur.execute('SELECT image_data FROM studio_slide_images WHERE output_id=%s AND image_data IS NOT NULL ORDER BY slide_number', (output_id,))
                images = cur.fetchall()
            if not images:
                return jsonify({'error': 'No slide images available for export'}), 400
            import img2pdf
            from io import BytesIO
            pdf_images = []
            for img_row in images:
                pdf_images.append(base64.b64decode(img_row['image_data']))
            pdf_bytes = img2pdf.convert(pdf_images)
            return Response(
                pdf_bytes,
                mimetype='application/pdf',
                headers={'Content-Disposition': f'attachment; filename="{output["title"]}.pdf"'}
            )

        elif output['type'] == 'report':
            markdown_text = content.get('markdown', '')
            return Response(
                markdown_text,
                mimetype='text/markdown',
                headers={'Content-Disposition': f'attachment; filename="{output["title"]}.md"'}
            )

        elif output['type'] == 'infographic':
            if output['image_data']:
                img_bytes = base64.b64decode(output['image_data'])
                return Response(
                    img_bytes,
                    mimetype='image/png',
                    headers={'Content-Disposition': f'attachment; filename="{output["title"]}.png"'}
                )
            return jsonify({'error': 'No image data available'}), 400

        else:
            # For quiz, flashcard, mindmap — export as JSON
            return Response(
                json.dumps(content, indent=2),
                mimetype='application/json',
                headers={'Content-Disposition': f'attachment; filename="{output["title"]}.json"'}
            )

    except Exception as e:
        print(f"Export error: {e}")
        return jsonify({'error': str(e)}), 500


# ============================================
# HEALTH CHECK
# ============================================

@app.route('/api/prices', methods=['GET'])
def get_stock_prices():
    """Get current prices for all tickers in the universe."""
    try:
        import yfinance as yf

        with get_db() as (_, cur):
            cur.execute('SELECT ticker FROM portfolio_analyses')
            tickers = [r['ticker'] for r in cur.fetchall()]

        if not tickers:
            return jsonify({'prices': {}})

        # Fetch all at once for efficiency
        data = yf.download(tickers, period='2d', progress=False, threads=True)

        prices = {}
        for ticker in tickers:
            try:
                if len(tickers) == 1:
                    close = data['Close']
                else:
                    close = data['Close'][ticker]
                if len(close.dropna()) >= 2:
                    current = float(close.dropna().iloc[-1])
                    prev = float(close.dropna().iloc[-2])
                    change_pct = ((current - prev) / prev) * 100
                    prices[ticker] = {'price': round(current, 2), 'changePct': round(change_pct, 2)}
                elif len(close.dropna()) >= 1:
                    prices[ticker] = {'price': round(float(close.dropna().iloc[-1]), 2), 'changePct': 0}
            except Exception:
                pass

        return jsonify({'prices': prices})
    except Exception as e:
        print(f'Price fetch error: {e}')
        return jsonify({'prices': {}, 'error': str(e)})


@app.route('/api/earnings-calendar', methods=['GET'])
def earnings_calendar():
    """Get upcoming earnings dates for stocks in the universe."""
    try:
        import yfinance as yf

        with get_db() as (_, cur):
            cur.execute('SELECT ticker FROM portfolio_analyses')
            tickers = [r['ticker'] for r in cur.fetchall()]

        earnings = []
        for ticker in tickers:
            try:
                stock = yf.Ticker(ticker)
                cal = stock.calendar
                if cal is not None and not cal.empty:
                    earnings_date = cal.iloc[0].get('Earnings Date', None)
                    if earnings_date:
                        earnings.append({
                            'ticker': ticker,
                            'earningsDate': str(earnings_date)[:10],
                        })
            except Exception:
                pass

        earnings.sort(key=lambda x: x.get('earningsDate', ''))
        return jsonify({'earnings': earnings})
    except Exception as e:
        return jsonify({'earnings': [], 'error': str(e)})


@app.route('/health', methods=['GET'])
def health():
    """Health check endpoint"""
    return jsonify({'status': 'ok', 'database': 'postgresql'})


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    print("=" * 50)
    print("TDL Equity Analyzer - Backend Server")
    print("With PostgreSQL Database Support")
    print("=" * 50)
    print(f"Starting server on http://0.0.0.0:{port}")
    print("=" * 50)
    app.run(host='0.0.0.0', port=port, debug=False)
